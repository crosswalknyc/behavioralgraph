"""
Brand Partnership IQ - Analysis Deck Builder
============================================

Generates a high-design analysis deck (.pptx) from a single BPIQ result
payload. Style direction:

  - Layout / typography / color: matched to the 2026 LISA Audience
    Profile deck (16:9, dark forest + cream alternating sections, big
    editorial typography, eyebrow + headline + stat-tile rhythm).
  - Narrative voice / data framing: matched to the MIDG / CAA Eiza
    Gonzalez deck - declarative dollar-impact headlines, pre vs post
    stat comparisons, "Read As:" plain-English translations, and a
    closing brand-by-brand wrap.

Usage:

    from migration.bpiq_deck_builder import build_deck
    pptx_bytes = build_deck(data, image_url=..., category=...)

The single public entry point is `build_deck(...)`. All slide builders
are private helpers and assume a fully-formed BPIQ result dict
(matching the shape served by /api/brand-partnership-iq/results/...).
"""
from __future__ import annotations

import io
import os
import re
import tempfile
import textwrap
import urllib.parse
from dataclasses import dataclass
from typing import Any, Optional

import requests
from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


# ─────────────────────────────────────────────────────────────────────
#  Palette - lifted from the LISA deck.
#
#  We keep this list small on purpose: every slide picks one BG_* and
#  the typography palette inherits from it. Accent colors are reserved
#  for stat numerals and the +X% lift callouts.
# ─────────────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0C, 0x16, 0x18)   # primary - dark forest
BG_CREAM  = RGBColor(0xE9, 0xE8, 0xE1)   # secondary - warm cream
BG_PINK   = RGBColor(0xE6, 0x82, 0xFF)   # accent insight slide
BG_LAV    = RGBColor(0xB7, 0xB3, 0xD8)   # final-insight slide
BG_WARM   = RGBColor(0xEA, 0xE8, 0xE2)   # variant cream

FG_LIGHT  = RGBColor(0xF4, 0xF1, 0xEA)   # on dark backgrounds
FG_DARK   = RGBColor(0x0C, 0x16, 0x18)   # on cream backgrounds
ACCENT    = RGBColor(0xC8, 0xE6, 0x00)   # Crosswalk green - hero callouts
LIFT_UP   = RGBColor(0x10, 0xB9, 0x81)   # positive lift
LIFT_DOWN = RGBColor(0xEF, 0x44, 0x44)   # negative lift
MUTED_LT  = RGBColor(0x6B, 0x72, 0x80)   # muted text on cream
MUTED_DK  = RGBColor(0x8F, 0x94, 0x8C)   # muted text on dark

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Font stack - python-pptx writes font names directly into the XML; if
# the viewer doesn't have these, PowerPoint substitutes automatically.
FONT_DISPLAY = "Helvetica Neue"
FONT_BODY    = "Helvetica Neue"
FONT_NUMERIC = "Helvetica Neue"


# ─────────────────────────────────────────────────────────────────────
#  Number / money formatting helpers
# ─────────────────────────────────────────────────────────────────────
def fmt_money(n: Any) -> str:
    try:
        n = float(n or 0)
    except (TypeError, ValueError):
        return "$0"
    if abs(n) >= 1_000_000_000:
        return f"${n/1_000_000_000:.2f}B"
    if abs(n) >= 1_000_000:
        return f"${n/1_000_000:.2f}M"
    if abs(n) >= 1_000:
        return f"${n/1_000:.0f}K"
    return f"${n:,.0f}"


def fmt_num(n: Any) -> str:
    try:
        n = float(n or 0)
    except (TypeError, ValueError):
        return "0"
    if abs(n) >= 1_000_000_000:
        return f"{n/1_000_000_000:.2f}B"
    if abs(n) >= 1_000_000:
        return f"{n/1_000_000:.1f}M"
    if abs(n) >= 1_000:
        return f"{n/1_000:.0f}K"
    return f"{int(n):,}"


def fmt_pct(n: Any, decimals: int = 1) -> str:
    try:
        n = float(n or 0)
    except (TypeError, ValueError):
        return "0%"
    return f"{n:.{decimals}f}%"


def fmt_signed_pct(n: Any) -> str:
    try:
        n = float(n or 0)
    except (TypeError, ValueError):
        return "0%"
    sign = "+" if n > 0 else ("" if n == 0 else "")
    return f"{sign}{n:.1f}%"


def safe_filename(s: str) -> str:
    return re.sub(r"[^A-Za-z0-9._-]+", "_", (s or "BPIQ_Deck"))[:120]


# ─────────────────────────────────────────────────────────────────────
#  Layout primitives - thin wrappers around python-pptx so the slide
#  builders read closer to a layout-DSL than raw XML wrangling.
# ─────────────────────────────────────────────────────────────────────
def _add_bg(slide, color: RGBColor):
    """Solid full-bleed background rectangle."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def _add_text(
    slide,
    left, top, width, height,
    text: str,
    *,
    size: float = 18,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = FG_DARK,
    font: str = FONT_BODY,
    align: str = "left",
    anchor: str = "top",
    spacing: float = 1.05,
    letter_spacing: float = 0.0,
):
    """Add a text box and return the shape. `text` may contain \\n
    for multi-paragraph blocks; each paragraph inherits the same run
    style."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = {
        "top":    MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(anchor, MSO_ANCHOR.TOP)
    align_enum = {
        "left":   PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right":  PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)
    paragraphs = str(text).split("\n")
    for i, para_text in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align_enum
        p.line_spacing = spacing
        r = p.add_run()
        r.text = para_text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return box


def _add_rect(slide, left, top, width, height, *,
              fill: Optional[RGBColor] = None,
              line: Optional[RGBColor] = None,
              line_w: Optional[float] = None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is None:
        rect.fill.background()
    else:
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        if line_w is not None:
            rect.line.width = Pt(line_w)
    rect.shadow.inherit = False
    return rect


def _eyebrow(slide, left, top, text: str, *, on_dark: bool = True, width=Inches(10)):
    """Tiny uppercase tracked eyebrow label - LISA's signature
    'THE LISA AUDIENCE · FRAGRANCE' bar."""
    return _add_text(
        slide, left, top, width, Inches(0.3),
        text.upper(),
        size=10, bold=True,
        color=(MUTED_DK if on_dark else MUTED_LT),
        letter_spacing=0.18,
    )


#: Small prepositions / articles / conjunctions that should NOT be
#: capitalized inside a display name (e.g., "World of Coca-Cola",
#: not "World Of Coca-Cola"). Fixes C12.
_TITLE_LOWERCASE_TOKENS = {
    "a", "an", "and", "as", "at", "but", "by", "for", "if", "in", "of",
    "on", "or", "the", "to", "vs", "vs.", "via", "with",
}


def _touchpoint_display_name(name: str) -> str:
    """Render a brand-property name for the touchpoint table. If the
    source string already contains any uppercase letters, treat it as
    author-cased and preserve it verbatim (payloads now provide
    exact-case names). Otherwise apply a title-case that lowercases
    small prepositions."""
    if not name:
        return ""
    s = str(name).strip()
    if any(c.isupper() for c in s):
        return s
    words = s.split()
    out = []
    for i, w in enumerate(words):
        wl = w.lower()
        if i > 0 and wl in _TITLE_LOWERCASE_TOKENS:
            out.append(wl)
        else:
            out.append(w[:1].upper() + w[1:].lower())
    return " ".join(out)


def _normalize_project_label(label: str) -> str:
    """Normalize a project label for footer / display use. Replaces
    em-dashes and en-dashes with commas so the footer never carries a
    stray dash (Liz item 20)."""
    if not label:
        return ""
    return (
        str(label)
        .replace(" — ", ", ")
        .replace(" – ", ", ")
        .replace("—", ",")
        .replace("–", ",")
    )


def _page_footer(slide, idx: int, total: int, project_label: str, *, on_dark: bool = True):
    """Bottom bar: project tag (left), CROSSWALK BehaviorGraph (center
    bottom-left), and `idx / total` (right). Echoes LISA's footer."""
    color = MUTED_DK if on_dark else MUTED_LT
    _add_text(
        slide, Inches(0.5), Inches(7.10), Inches(6), Inches(0.3),
        _normalize_project_label(project_label).upper(),
        size=8, color=color, letter_spacing=0.14, bold=True,
    )
    _add_text(
        slide, Inches(11.0), Inches(7.10), Inches(1.85), Inches(0.3),
        f"{idx:02d} / {total:02d}",
        size=8, color=color, align="right", bold=True, letter_spacing=0.14,
    )


def _crosswalk_mark(slide, left, top, *, on_dark: bool = True, size: float = 0.18):
    """Inline 'CROSSWALK · BehaviorGraph' wordmark drawn as text. Uses
    Pillow only when we need a bitmap stamp (cover slide); for inline
    footers a text shape is faster and crisper at scale."""
    color = FG_LIGHT if on_dark else FG_DARK
    _add_text(
        slide, left, top, Inches(4), Inches(0.3),
        "CROSSWALK · BehaviorGraph",
        size=size * 60,  # ~10-12pt depending on size arg
        bold=True, color=color, letter_spacing=0.18,
        font=FONT_DISPLAY,
    )


# ─────────────────────────────────────────────────────────────────────
#  Image helpers - download + crop to slide-ready PNGs.
# ─────────────────────────────────────────────────────────────────────
def _download_image(url: str, max_dim: int = 2000,
                    crop_aspect: Optional[float] = None) -> Optional[str]:
    """Fetch a remote image to a tmp JPEG. Optional `crop_aspect` (w/h)
    center-crops the image to fit a target slot - lets the cover slide
    place a fixed-aspect hero without distortion. Returns path or None
    on any failure (deck-builder must NEVER crash because of a bad
    photo URL)."""
    if not url or not isinstance(url, str):
        return None
    try:
        r = requests.get(url, timeout=8, stream=True, headers={
            "User-Agent": "BehaviorGraph-DeckBuilder/1.0",
        })
        if r.status_code != 200:
            return None
        raw = io.BytesIO(r.content)
        im = Image.open(raw).convert("RGB")
        # Center-crop to target aspect first so any subsequent resize
        # doesn't deform the subject.
        if crop_aspect and crop_aspect > 0:
            src_ar = im.width / im.height
            if src_ar > crop_aspect:
                # Source is wider than target → trim left/right
                new_w = int(im.height * crop_aspect)
                off = (im.width - new_w) // 2
                im = im.crop((off, 0, off + new_w, im.height))
            elif src_ar < crop_aspect:
                # Source is taller than target → trim top/bottom
                new_h = int(im.width / crop_aspect)
                # Bias the crop slightly above center so faces don't
                # lose their forehead - this is the classic editorial
                # crop heuristic.
                off = max(0, (im.height - new_h) // 3)
                im = im.crop((0, off, im.width, off + new_h))
        # Cap the long edge so we don't bloat the pptx unnecessarily.
        if max(im.size) > max_dim:
            ratio = max_dim / max(im.size)
            im = im.resize((int(im.width * ratio), int(im.height * ratio)),
                           Image.LANCZOS)
        out = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False)
        im.save(out.name, "JPEG", quality=88, optimize=True)
        return out.name
    except Exception:
        return None


def _generate_typographic_hero(text: str, palette: str = "dark") -> str:
    """Fallback hero image - bold typographic block with the brand
    name. Used when the partnership has no admin-uploaded photo so the
    cover slide still has visual weight (mirrors LISA's typographic
    insight slides)."""
    w, h = 1600, 1000
    if palette == "cream":
        bg = (233, 232, 225)
        fg = (12, 22, 24)
    elif palette == "pink":
        bg = (230, 130, 255)
        fg = (12, 22, 24)
    else:
        bg = (12, 22, 24)
        fg = (244, 241, 234)
    im = Image.new("RGB", (w, h), bg)
    d = ImageDraw.Draw(im)
    # Try a few candidate fonts; fall back to Pillow's default if none
    # are present (the test environment ships a couple of system fonts).
    font_path_candidates = [
        "/System/Library/Fonts/HelveticaNeue.ttc",
        "/System/Library/Fonts/Supplemental/Arial Black.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    ]
    font = None
    for p in font_path_candidates:
        if os.path.exists(p):
            try:
                font = ImageFont.truetype(p, 220)
                break
            except Exception:
                pass
    if font is None:
        font = ImageFont.load_default()
    txt = (text or "").upper()
    bbox = d.textbbox((0, 0), txt, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    d.text(((w - tw) / 2, (h - th) / 2 - 40), txt, fill=fg, font=font)
    # Soft vignette so the typographic plate doesn't read as flat.
    overlay = Image.new("RGB", (w, h), bg)
    mask = Image.new("L", (w, h), 0)
    md = ImageDraw.Draw(mask)
    md.ellipse((-200, -200, w + 200, h + 200), fill=0)
    md.ellipse((100, 100, w - 100, h - 100), fill=255)
    mask = mask.filter(ImageFilter.GaussianBlur(120))
    im = Image.composite(im, overlay, mask)
    out = tempfile.NamedTemporaryFile(suffix=".jpg", delete=False)
    im.save(out.name, "JPEG", quality=88, optimize=True)
    return out.name


def _render_crosswalk_logo_png(on_dark: bool = True) -> str:
    """Pillow-rendered Crosswalk wordmark for the cover slide. No SVG
    dependency required. Returns a path to a transparent PNG."""
    w, h = 760, 100
    im = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    d = ImageDraw.Draw(im)
    fg = (244, 241, 234, 255) if on_dark else (12, 22, 24, 255)
    accent = (200, 230, 0, 255)
    font_candidates = [
        "/System/Library/Fonts/Supplemental/Arial Black.ttf",
        "/System/Library/Fonts/HelveticaNeue.ttc",
        "/System/Library/Fonts/Helvetica.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    ]
    font = None
    for p in font_candidates:
        if os.path.exists(p):
            try:
                font = ImageFont.truetype(p, 64)
                break
            except Exception:
                pass
    if font is None:
        font = ImageFont.load_default()
    d.text((0, 8), "CROSSWALK", fill=fg, font=font)
    # Lozenge accent + product subtitle, BehaviorGraph-style
    try:
        small = ImageFont.truetype(font.path, 22) if hasattr(font, "path") else font
    except Exception:
        small = font
    d.rectangle((420, 60, 460, 84), fill=accent)
    d.text((470, 60), "BehaviorGraph", fill=fg, font=small)
    out = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    im.save(out.name, "PNG", optimize=True)
    return out.name


# ─────────────────────────────────────────────────────────────────────
#  Narrative helpers - derive headlines & "read as" sentences from the
#  data. Pure functions, easy to extend / tune in isolation.
# ─────────────────────────────────────────────────────────────────────
@dataclass
class DeckCtx:
    """All the derived numbers the slide builders need, computed once
    so individual slides stay readable."""
    project:        str
    brand:          str
    category_label: str
    is_auto:        bool
    cover_image:    Optional[str]
    logo:           str
    total_value:    float
    emv:            float
    bev:            float
    blv:            float
    cv:             float
    incr_users:     int
    target_size:    int
    audience_proj:  int
    pre_users:      int
    post_users:     int
    pre_users_proj: int
    post_users_proj:int
    pre_pen:        float
    post_pen:       float
    delta_pp:       float
    delta_rel:      float
    incr_lift_pp:   float
    incr_lift_rel:  float
    campaign_window:str
    attribution_d:  int
    panel_size:     int
    sample_pre:     int
    sample_post:    int
    has_conversions:bool
    has_yearly:     bool
    sources_line:   str
    data:           dict


def _derive_context(data: dict, image_url: Optional[str],
                    category: Optional[str]) -> DeckCtx:
    totals = data.get("totals") or {}
    val    = data.get("valuation") or {}
    rates  = val.get("rates") or {}
    cg     = data.get("control_group") or {}
    is_auto = (str(category or "").strip().lower() == "automotive")
    conv = data.get("conversions") or {}
    has_conv = (
        not is_auto
        and bool(conv.get("enabled"))
        and not bool(conv.get("low_signal"))
        and (conv.get("post_users_projected") or 0) > 0
    )
    yearly = data.get("yearly_breakdown") or []
    has_yearly = isinstance(yearly, list) and len(yearly) >= 2

    pre_pen  = float(totals.get("audience_pen_pre_pct")  or 0)
    post_pen = float(totals.get("audience_pen_post_pct") or 0)
    delta_pp = round(post_pen - pre_pen, 2)
    delta_rel = ((post_pen - pre_pen) / pre_pen * 100) if pre_pen else 0.0

    # Cover image is placed in a 7.6"×7.5" slot - center-crop to the
    # matching aspect (≈1.013) so faces stay framed and the slot fills
    # cleanly with no letterboxing.
    cover = _download_image(image_url, crop_aspect=7.6 / 7.5) if image_url else None
    if cover is None:
        cover = _generate_typographic_hero(data.get("brand_partner")
                                           or data.get("project_name") or "BPIQ")

    logo = _render_crosswalk_logo_png(on_dark=True)

    project = data.get("project_name") or "Brand Partnership Valuation"
    brand   = data.get("brand_partner") or "Brand"
    cat_lbl = (category or "").strip().title() or "Brand Partnership"
    start   = data.get("start_date") or ""
    end     = data.get("end_date")   or ""
    win     = f"{start} → {end}" if start and end else ""

    src = (
        f"Source: Crosswalk BehaviorGraph · 30M+ opted-in U.S. consumer panel · "
        f"observed sample n={int(data.get('audience_size') or 0):,} · "
        f"projected universe n={int(data.get('projected_audience_size') or 0):,} · "
        f"window {win} · "
        f"attribution {int(data.get('attribution_window_days') or 0)}d"
    )

    return DeckCtx(
        project=project,
        brand=brand,
        category_label=cat_lbl,
        is_auto=is_auto,
        cover_image=cover,
        logo=logo,
        total_value=float(val.get("total_brand_value") or 0),
        emv=float(val.get("earned_media_value") or 0),
        bev=float(val.get("brand_engagement_value") or 0),
        blv=float(val.get("brand_lift_value") or 0),
        cv=float(val.get("conversion_value") or 0) if has_conv else 0.0,
        incr_users=int(val.get("incremental_users") or 0),
        target_size=int(data.get("audience_size") or 0),
        audience_proj=int(data.get("projected_audience_size") or 0),
        pre_users=int(totals.get("pre_users") or 0),
        post_users=int(totals.get("post_users") or 0),
        pre_users_proj=int(totals.get("pre_users_projected") or 0),
        post_users_proj=int(totals.get("post_users_projected") or 0),
        pre_pen=pre_pen,
        post_pen=post_pen,
        delta_pp=delta_pp,
        delta_rel=delta_rel,
        incr_lift_pp=float(cg.get("incremental_lift_pp") or 0),
        incr_lift_rel=float(cg.get("incremental_lift_rel_pct") or 0),
        campaign_window=win,
        attribution_d=int(data.get("attribution_window_days") or 0),
        panel_size=int(data.get("audience_size") or 0),
        sample_pre=int(totals.get("pre_users") or 0),
        sample_post=int(totals.get("post_users") or 0),
        has_conversions=has_conv,
        has_yearly=has_yearly,
        sources_line=src,
        data=data,
    )


_MONTHS_SHORT = [
    "Jan", "Feb", "Mar", "Apr", "May", "Jun",
    "Jul", "Aug", "Sep", "Oct", "Nov", "Dec",
]


def _fmt_date(iso: str) -> str:
    """YYYY-MM-DD → 'DD Mon YYYY' for slide copy. Returns the input
    unchanged if parsing fails."""
    if not iso or not isinstance(iso, str):
        return ""
    parts = iso.split("-")
    if len(parts) != 3:
        return iso
    try:
        yyyy, mm, dd = int(parts[0]), int(parts[1]), int(parts[2])
        return f"{dd:02d} {_MONTHS_SHORT[mm-1]} {yyyy}"
    except Exception:
        return iso


def _period_label(period: dict, fallback: str = "") -> str:
    """Human-readable window label from a {start, end, days} block.
    'from DD Mon YYYY to DD Mon YYYY (N days)'. Used on the
    engagement-lift slide to avoid the C2 mislabel (calling a 15-day
    window 'the year before campaign launch')."""
    if not isinstance(period, dict):
        return fallback
    start = _fmt_date(str(period.get("start") or ""))
    end   = _fmt_date(str(period.get("end")   or ""))
    days  = period.get("days")
    if start and end:
        base = f"from {start} to {end}"
        if isinstance(days, (int, float)) and days > 0:
            base += f" ({int(days)}-day window)"
        return base
    return fallback


def _engagement_headline(ctx: DeckCtx) -> str:
    """MIDG-style declarative headline that names the lift, the
    consumer count, and the brand.

    Uses adaptive precision on the pp delta (C9): sub-percent deltas
    display with 2 decimals so a +0.15pp lift never renders as
    "0.1pt" while the pre/post percentages show 0.2pp apart. Two-line
    sentence with terminal punctuation on both lines."""
    if ctx.delta_pp <= 0:
        return (
            f"The campaign held {ctx.brand}'s baseline.\n"
            f"{fmt_num(ctx.post_users_proj)} consumers engaged with the brand "
            f"in the post-window."
        )
    # Adaptive precision: below 1pp shows 2 decimals; at/above 1pp
    # shows 1 decimal. Guarantees the headline delta never contradicts
    # the pre/post display precision (fixes C9).
    if abs(ctx.delta_pp) < 1.0:
        delta_str = f"{ctx.delta_pp:.2f}pt"
    else:
        delta_str = f"{ctx.delta_pp:.1f}pt"
    moved = max(0, ctx.post_users_proj - ctx.pre_users_proj)
    return (
        f"A {delta_str} lift in brand engagement,\n"
        f"moving {fmt_num(moved)} more consumers toward {ctx.brand}."
    )


def _cover_subtitle(ctx: DeckCtx) -> str:
    """Purpose-written cover subtitle. Two lines: category + a
    declarative summary describing what this deck measures. Always
    ends with terminal punctuation so it never renders as a truncated
    fragment (fixes C7)."""
    diag = ctx.data.get("diagnostics") or {}
    is_counterfactual = "counterfactual" in (ctx.project or "").lower()
    label = ctx.category_label or "Brand Partnership"
    if is_counterfactual:
        line1 = f"{label} · Counterfactual read."
        line2 = (
            f"Isolates the sponsorship signal by measuring {ctx.brand} "
            f"against the same audience and window as the treated brand."
        )
    else:
        # Standard partnership read - lead with the total value we
        # observed so the cover has a hook, without leaning on the
        # slide-4 headline which uses a comma-ended fragment.
        line1 = f"{label} · Total observed value: {fmt_money(ctx.total_value)}."
        line2 = (
            f"An analysis of {ctx.brand}'s consumer engagement across "
            f"the campaign window."
        )
    return f"{line1}\n{line2}"


def _final_insight_lines(ctx: DeckCtx) -> tuple[str, str]:
    """Two-line editorial close - mirrors LISA's 'LISA wears it, it
    sells out.' format. We synthesize from the data so every partnership
    gets a tuned message.

    Counterfactual brands (project name contains 'counterfactual')
    always route to the control-brand copy regardless of adjusted-
    lift sign. Small positive drift on a control brand is category
    noise, not a sponsorship effect (Liz's July 27 v2 inversion
    callout: Pepsi F1 was mislabeled 'consistent with a genuine
    but small sponsorship effect' when Pepsi was not sponsored).

    Guardrail: never claim compounding lift over time on a single-
    window run."""
    if ctx.total_value <= 0:
        return ("The partnership ran.", "The audience didn't move.")
    lift_rel = ctx.incr_lift_rel if ctx.incr_lift_rel is not None else ctx.delta_rel
    sig      = (ctx.data.get("diagnostics") or {}).get("significance") or {}
    is_sig   = bool(sig.get("significant"))
    is_counterfactual = "counterfactual" in (ctx.project or "").lower()
    # Counterfactual brands always read as controls, regardless of
    # sign or significance of the tiny residual drift.
    if is_counterfactual:
        return (
            "The campaign held baseline.",
            f"{ctx.brand} was not in the show, and moved with control "
            "drift. Any residual is category-wide noise, not "
            "sponsorship-attributable.",
        )
    if lift_rel <= 0:
        return (
            "The campaign held baseline.",
            f"{ctx.brand}'s adjusted lift falls at or below control "
            "drift.",
        )
    if not is_sig:
        return (
            f"A directional {lift_rel:.1f}% adjusted lift,",
            "not statistically distinguishable from zero at this sample "
            "size. Signal is consistent with a genuine but small "
            "sponsorship effect.",
        )
    if lift_rel >= 200:
        line1 = f"When {ctx.brand} shows up with"
        line2 = f"{ctx.project.split('x')[0].strip() or 'the talent'}, the audience triples."
    elif lift_rel >= 100:
        line1 = f"{ctx.brand} more than doubles"
        line2 = "its share of voice across the audience."
    elif lift_rel >= 25:
        line1 = f"{ctx.brand} won meaningful"
        line2 = f"new ground: {lift_rel:.0f}% adjusted lift, post-campaign."
    else:
        if ctx.has_yearly:
            line1 = "Steady adjusted lift,"
            line2 = "compounding reach year over year."
        else:
            line1 = f"A {lift_rel:.1f}% adjusted lift,"
            line2 = "significant against control drift."
    return (line1, line2)


# ─────────────────────────────────────────────────────────────────────
#  Slide builders. Each takes (prs, ctx, idx, total) and returns the
#  new slide. Indexed manually because we need explicit `idx / total`
#  in every footer.
# ─────────────────────────────────────────────────────────────────────
def _blank(prs):
    """python-pptx ships several layouts; the 'blank' one (typically
    index 6) is the cleanest base because there are no placeholder
    text boxes to fight."""
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def _slide_cover(prs, ctx: DeckCtx, idx: int, total: int):
    s = _blank(prs)
    _add_bg(s, BG_DARK)
    # Right-hand hero image fills a 7.6×7.5 slot (cropped to that
    # aspect during download so the photo isn't deformed here).
    if ctx.cover_image:
        try:
            s.shapes.add_picture(
                ctx.cover_image,
                Inches(5.8), Inches(0),
                width=Inches(7.6), height=Inches(7.5),
            )
        except Exception:
            pass
    # Left-hand title stack. Title font size is scaled to length so
    # long partnerships (e.g. "Penélope Cruz x CHANEL Endorsement")
    # still fit inside the available block without crashing into the
    # subtitle below.
    _add_text(s, Inches(0.6), Inches(0.6), Inches(5.2), Inches(0.4),
              "A BRAND PARTNERSHIP VALUATION · 2026",
              size=10, bold=True, color=MUTED_DK, letter_spacing=0.2)
    title_text = _normalize_project_label(ctx.project) + "."
    title_len  = len(title_text)
    if   title_len <= 20: title_size = 64
    elif title_len <= 30: title_size = 52
    elif title_len <= 42: title_size = 42
    else:                 title_size = 34
    _add_text(s, Inches(0.6), Inches(1.7), Inches(5.0), Inches(3.6),
              title_text,
              size=title_size, bold=True, color=FG_LIGHT,
              spacing=0.95, font=FONT_DISPLAY)
    # Subtitle - purpose-written, not derived from slide 4's split
    # headline (fixes C7: previously the subtitle pulled the first
    # line of the two-line engagement headline which ended on a
    # comma, so the cover subtitle rendered as a truncated fragment).
    # Two-line format: category + declarative one-line summary; the
    # summary carries terminal punctuation so it reads as a complete
    # thought at any font size.
    subtitle = _cover_subtitle(ctx)
    _add_text(s, Inches(0.6), Inches(5.4), Inches(5.0), Inches(1.2),
              subtitle,
              size=15, color=FG_LIGHT, spacing=1.25, font=FONT_BODY,
              anchor="bottom")
    # DRAFT date stamp - explicitly names the analysis build date so
    # a reader can tell whether they're holding an old draft or the
    # July 27 v2 rebuild (C7 mentions the DRAFT date line went
    # missing from the rerun cover; restored here).
    created = str(ctx.data.get("created_at") or "")[:10]
    if created:
        _add_text(s, Inches(0.6), Inches(6.5), Inches(5.0), Inches(0.25),
                  f"DRAFT · {created}",
                  size=9, bold=True, color=MUTED_DK, letter_spacing=0.2)
    # CONFIDENTIAL stamp top right
    _add_text(s, Inches(11.4), Inches(0.6), Inches(1.4), Inches(0.3),
              "CONFIDENTIAL",
              size=9, bold=True, color=MUTED_DK, align="right",
              letter_spacing=0.2)
    # Bottom-left logo (rendered PNG) + brand mark
    try:
        s.shapes.add_picture(ctx.logo, Inches(0.6), Inches(6.7),
                             height=Inches(0.45))
    except Exception:
        _add_text(s, Inches(0.6), Inches(6.95), Inches(5), Inches(0.3),
                  "CROSSWALK · BehaviorGraph",
                  size=11, bold=True, color=FG_LIGHT, letter_spacing=0.18)
    return s


def _section_label(idx: int, title: str) -> str:
    """Consistent 'NN · Section Title' eyebrow, driven off idx so we
    never desynchronize the section counter from the actual slide
    position (fixes the classic 10→12 jump when an optional slide is
    skipped)."""
    return f"{idx:02d} · {title}"


def _slide_methodology(prs, ctx: DeckCtx, idx: int, total: int):
    """MIDG 'Direct Observational Data' slide - credibility / what
    makes BehaviorGraph trustworthy."""
    s = _blank(prs)
    _add_bg(s, BG_CREAM)
    _eyebrow(s, Inches(0.6), Inches(0.6),
             _section_label(idx, "The Data Foundation"), on_dark=False)
    _add_text(s, Inches(0.6), Inches(1.0), Inches(11), Inches(1.4),
              "Direct observational data.\nNo surveys, no self-report.",
              size=44, bold=True, color=FG_DARK,
              spacing=0.95, font=FONT_DISPLAY)
    # Dashes in the body copy get normalized to commas so they don't
    # regress the item-20 dash cleanup (C11).
    body = (
        "Crosswalk's BehaviorGraph draws on a 30M+ opted-in U.S. consumer "
        "panel. For this partnership we observed a sample of "
        f"{fmt_num(ctx.panel_size)} panelists, the Wheel of Fortune "
        "Next Day Air viewer cohort, and projected their behavior to a "
        f"{fmt_num(ctx.audience_proj)}-consumer U.S. viewing audience "
        "universe. Every observation reflects real behavior: which sites "
        "they visited, which platforms they spent time on, and which "
        "brand touchpoints appeared in their sessions before, during, "
        "and after the partnership ran."
    )
    _add_text(s, Inches(0.6), Inches(3.0), Inches(7), Inches(2.2),
              body, size=13, color=FG_DARK, spacing=1.32)
    # Study-context clarifier line - C3 fix. When the payload
    # carries a diagnostics.study_context field (rerun analyses,
    # counterfactuals, etc.), print it as an italic footnote so no
    # reader can misread this as the first-flight deck.
    study_ctx = ((ctx.data.get("diagnostics") or {})
                 .get("study_context"))
    if study_ctx:
        _add_text(s, Inches(0.6), Inches(5.35), Inches(7), Inches(1.35),
                  str(study_ctx), size=10, italic=True,
                  color=MUTED_LT, spacing=1.3)
    # Right column - 3 stat tiles. Window tile uses the actual pre /
    # event / post dates so slide 2 never claims a window it doesn't
    # cover (C2 + C3 defensive coverage).
    pre_p  = ctx.data.get("pre_period")   or {}
    post_p = ctx.data.get("post_period")  or {}
    window_val = (
        f"{_fmt_date(pre_p.get('start') or '')} to "
        f"{_fmt_date(post_p.get('end')   or '')}"
    ) if pre_p and post_p else (ctx.campaign_window or "0")
    stats = [
        ("Panel sample", f"{fmt_num(ctx.panel_size)}", "Observed viewer cohort"),
        ("Window",       window_val,                    "Pre, event, post windows"),
        ("Attribution",  f"{ctx.attribution_d}d",       "Post-event window"),
    ]
    col_x = Inches(8.2)
    col_w = Inches(4.6)
    for i, (label, value, sub) in enumerate(stats):
        y = Inches(3.2 + i * 1.15)
        _add_rect(s, col_x, y, col_w, Inches(1.0),
                  fill=None, line=FG_DARK, line_w=0.75)
        _add_text(s, col_x + Inches(0.25), y + Inches(0.1),
                  Inches(2.0), Inches(0.3),
                  label.upper(), size=9, bold=True, color=MUTED_LT,
                  letter_spacing=0.16)
        _add_text(s, col_x + Inches(0.25), y + Inches(0.32),
                  Inches(col_w.inches - 0.5), Inches(0.6),
                  value, size=22, bold=True, color=FG_DARK,
                  font=FONT_NUMERIC)
        _add_text(s, col_x + Inches(0.25), y + Inches(0.72),
                  Inches(col_w.inches - 0.5), Inches(0.3),
                  sub, size=9, color=MUTED_LT, italic=True)
    _page_footer(s, idx, total, ctx.project, on_dark=False)
    return s


def _slide_audience_scale(prs, ctx: DeckCtx, idx: int, total: int):
    """LISA-style 'THE LISA AUDIENCE · SCALE' big stat slide."""
    s = _blank(prs)
    _add_bg(s, BG_DARK)
    _eyebrow(s, Inches(0.6), Inches(0.6),
             _section_label(idx, f"The {ctx.brand} Audience · Scale"),
             on_dark=True)
    _add_text(s, Inches(0.6), Inches(1.0), Inches(11), Inches(1.3),
              f"Reach and resonance,\nside by side.",
              size=44, bold=True, color=FG_LIGHT,
              spacing=0.95, font=FONT_DISPLAY)
    # Three giant stat tiles
    tiles = [
        (fmt_num(ctx.audience_proj),
         "U.S. CONSUMER AUDIENCE",
         f"Projected from {fmt_num(ctx.target_size)} observed panelists "
         f"(30M+ Crosswalk panel)."),
        (fmt_num(ctx.post_users_proj),
         f"POST-CAMPAIGN {ctx.brand.upper()} CONSUMERS",
         f"{fmt_pct(ctx.post_pen)} of the observed cohort had a brand "
         f"touchpoint in the post-window."),
        (fmt_money(ctx.total_value),
         "TOTAL BRAND VALUE DELIVERED",
         f"Sum of Earned Media, Brand Engagement, "
         f"{'Brand Lift, and Conversion' if ctx.has_conversions else 'and Brand Lift'} value."),
    ]
    tile_w = Inches(4.0)
    gutter = Inches(0.15)
    base_x = Inches(0.6)
    for i, (num, label, sub) in enumerate(tiles):
        x = base_x + (tile_w + gutter) * i
        y = Inches(3.4)
        _add_rect(s, x, y, tile_w, Inches(3.1),
                  fill=None, line=MUTED_DK, line_w=0.5)
        _add_text(s, x + Inches(0.25), y + Inches(0.25),
                  tile_w - Inches(0.5), Inches(1.4),
                  num, size=54, bold=True, color=ACCENT,
                  font=FONT_NUMERIC, spacing=0.95)
        _add_text(s, x + Inches(0.25), y + Inches(1.7),
                  tile_w - Inches(0.5), Inches(0.4),
                  label, size=10, bold=True, color=FG_LIGHT,
                  letter_spacing=0.14)
        _add_text(s, x + Inches(0.25), y + Inches(2.1),
                  tile_w - Inches(0.5), Inches(0.9),
                  sub, size=10, color=MUTED_DK, spacing=1.3)
    _page_footer(s, idx, total, ctx.project, on_dark=True)
    return s


def _slide_engagement_lift(prs, ctx: DeckCtx, idx: int, total: int):
    """MIDG centerpiece - 'A 3.2pt Lift That Moved $19M' headline +
    Pre vs Post stat blocks + Read As callout."""
    s = _blank(prs)
    _add_bg(s, BG_CREAM)
    _eyebrow(s, Inches(0.6), Inches(0.6),
             _section_label(idx, f"{ctx.brand} · Engagement Lift"),
             on_dark=False)
    headline_lines = _engagement_headline(ctx).split("\n")
    _add_text(s, Inches(0.6), Inches(1.0), Inches(11), Inches(1.6),
              headline_lines[0],
              size=42, bold=True, color=FG_DARK,
              spacing=0.95, font=FONT_DISPLAY)
    if len(headline_lines) > 1:
        _add_text(s, Inches(0.6), Inches(2.05), Inches(11), Inches(0.7),
                  headline_lines[1],
                  size=22, color=FG_DARK, font=FONT_DISPLAY)
    # Pre / Post visual comparison - two big stat plates side by side
    plate_y = Inches(3.4)
    plate_w = Inches(5.2)
    plate_h = Inches(2.8)
    # Pre plate
    _add_rect(s, Inches(0.6), plate_y, plate_w, plate_h,
              fill=FG_DARK, line=None)
    _add_text(s, Inches(0.9), plate_y + Inches(0.25), Inches(2), Inches(0.3),
              "PRE-CAMPAIGN", size=10, bold=True, color=ACCENT,
              letter_spacing=0.2)
    _add_text(s, Inches(0.9), plate_y + Inches(0.6), plate_w - Inches(0.5),
              Inches(1.3),
              fmt_pct(ctx.pre_pen), size=70, bold=True, color=FG_LIGHT,
              font=FONT_NUMERIC, spacing=0.95)
    # Pre-plate sub-label = the actual pre-window from the payload
    # (fixes C2: prior copy said "in the year before campaign launch"
    # which mislabels the trailing pre-period as a 12-month window).
    #
    # Additionally shows the RAW INCIDENCE line ("2.4M of 10.0M
    # panelists") directly under the projected-consumer count so
    # the numerator and denominator behind the pre/post % are
    # visible on the slide, not reverse-engineered. Answers Liz's
    # July 27 v2 C1 disclosure ask and the standing "supply the
    # raw engager counts" request.
    pre_p  = ctx.data.get("pre_period")  or {}
    post_p = ctx.data.get("post_period") or {}
    pre_label = _period_label(pre_p, fallback="in the pre-window")
    post_label = _period_label(post_p, fallback=f"in the {ctx.attribution_d}-day post-window")
    pre_incidence  = f"{fmt_num(ctx.pre_users)} of {fmt_num(ctx.panel_size)} panelists (incidence rate)"
    post_incidence = f"{fmt_num(ctx.post_users)} of {fmt_num(ctx.panel_size)} panelists (incidence rate)"
    _add_text(s, Inches(0.9), plate_y + Inches(1.85), plate_w - Inches(0.5),
              Inches(0.35),
              f"{fmt_num(ctx.pre_users_proj)} U.S. consumers projected",
              size=13, color=FG_LIGHT)
    _add_text(s, Inches(0.9), plate_y + Inches(2.15), plate_w - Inches(0.5),
              Inches(0.30),
              pre_incidence,
              size=10, italic=True, color=MUTED_DK)
    _add_text(s, Inches(0.9), plate_y + Inches(2.42), plate_w - Inches(0.5),
              Inches(0.35),
              f"engaged with the brand {pre_label}.",
              size=10, italic=True, color=MUTED_DK)
    # Post plate
    post_x = Inches(7.2)
    _add_rect(s, post_x, plate_y, plate_w, plate_h,
              fill=ACCENT, line=None)
    _add_text(s, post_x + Inches(0.3), plate_y + Inches(0.25),
              Inches(2), Inches(0.3),
              "POST-CAMPAIGN", size=10, bold=True, color=FG_DARK,
              letter_spacing=0.2)
    _add_text(s, post_x + Inches(0.3), plate_y + Inches(0.6),
              plate_w - Inches(0.5), Inches(1.3),
              fmt_pct(ctx.post_pen), size=70, bold=True, color=FG_DARK,
              font=FONT_NUMERIC, spacing=0.95)
    _add_text(s, post_x + Inches(0.3), plate_y + Inches(1.85),
              plate_w - Inches(0.5), Inches(0.35),
              f"{fmt_num(ctx.post_users_proj)} U.S. consumers projected",
              size=13, color=FG_DARK)
    _add_text(s, post_x + Inches(0.3), plate_y + Inches(2.15),
              plate_w - Inches(0.5), Inches(0.30),
              post_incidence,
              size=10, italic=True, color=FG_DARK)
    _add_text(s, post_x + Inches(0.3), plate_y + Inches(2.42),
              plate_w - Inches(0.5), Inches(0.35),
              f"engaged with the brand {post_label}.",
              size=10, italic=True, color=FG_DARK)
    # Read As callout + significance line, MIDG-style.
    moved = max(0, ctx.post_users_proj - ctx.pre_users_proj)
    read_as = (
        f"Read As: {fmt_num(moved)} more U.S. consumers were observed "
        f"engaging with {ctx.brand} after the campaign, a "
        f"+{ctx.delta_rel:.1f}% relative lift over baseline."
    )
    _add_text(s, Inches(0.6), Inches(6.25), Inches(12.1), Inches(0.35),
              read_as, size=12, italic=True, color=MUTED_LT, align="center")
    # Statistical significance line - pooled two-sample z on paired
    # marginals (same-panel pre/post design). We display n, delta,
    # z, and a 95% CI on the pp delta. See methodology slide.
    sig = _mcnemar_summary(ctx)
    if sig:
        _add_text(s, Inches(0.6), Inches(6.62), Inches(12.1), Inches(0.35),
                  sig, size=10, italic=True, color=MUTED_LT, align="center")
    _page_footer(s, idx, total, ctx.project, on_dark=False)
    return s


def _mcnemar_summary(ctx: DeckCtx) -> str:
    """Compose the on-slide significance strip from the payload's
    diagnostics.significance block.

    Per July 27 QC memo (C6): at large panel n the p-value carries
    no information (any Δ above ~0.04pp clears p=0.05), so on
    n>=1e6 panels we drop the p-value entirely and report point
    estimate, z, and the 95% CI. On smaller panels we keep p for
    readers who want it.

    Recomputes z and CI from panel size + observed pre/post counts
    when possible, so a stale-denominator payload still renders the
    right numbers on the corrected panel base."""
    from math import erf, sqrt

    diag  = (ctx.data.get("diagnostics") or {}).get("significance") or {}
    n_pre  = ctx.sample_pre
    n_post = ctx.sample_post
    n_panel = ctx.panel_size

    # Recompute on the panel base actually declared in the payload,
    # rather than trusting a possibly-stale precomputed z/CI written
    # against a different denominator.
    z_calc, se_pp_calc, delta_pp_calc = None, None, None
    if n_pre and n_post and n_panel:
        p_pool = (n_pre + n_post) / (2 * n_panel)
        se = sqrt(2 * p_pool * (1 - p_pool) / n_panel)
        delta = (n_post - n_pre) / n_panel
        z_calc = delta / se if se else 0
        se_pp_calc = se * 100
        delta_pp_calc = delta * 100

    # Prefer recomputed z/CI/Δ. Fall back to payload numbers only if
    # the recompute fails.
    if delta_pp_calc is not None:
        delta_pp = delta_pp_calc
        z        = z_calc
        ci       = [delta_pp - 1.96 * se_pp_calc, delta_pp + 1.96 * se_pp_calc]
    else:
        delta_pp = diag.get("delta_pp_point")
        z        = diag.get("primary_test_z")
        ci       = diag.get("delta_ci_95_pp") or diag.get("ci_95_pp")

    # Significance flag: recomputed z if available, else the payload
    # flag (which is only correct on the payload's declared base).
    if z is not None:
        sig = abs(z) >= 1.96
    else:
        sig = bool(diag.get("significant"))

    parts = [f"Panel n={fmt_num(n_panel)}"]
    if delta_pp is not None:
        parts.append(f"Δ={float(delta_pp):+.2f}pp")
    if z is not None:
        parts.append(f"z={float(z):+.2f}")
    # Drop p when n is large enough that p becomes uninformative
    # (item C6: detection floor at n=1e7 is ~0.04pp; p is not a
    # useful column on a client-facing slide at this scale).
    if n_panel < 1_000_000:
        p = diag.get("primary_test_p_value")
        if p is None and z is not None:
            p = 2 * (1 - 0.5 * (1 + erf(abs(z) / sqrt(2))))
        if p is not None:
            parts.append(f"p={float(p):.3f}")
    if ci:
        try:
            lo, hi = float(ci[0]), float(ci[1])
            parts.append(f"95% CI on Δ: [{lo:+.2f}pp, {hi:+.2f}pp]")
        except Exception:
            pass
    parts.append("SIGNIFICANT" if sig else "not distinguishable from zero")
    return " · ".join(parts)


def _slide_total_value_hero(prs, ctx: DeckCtx, idx: int, total: int):
    """Hero $ slide - LISA's 'ONE INSIGHT' moment, repurposed for the
    total brand value. Solid pink so it explodes off the page after the
    cream lift slide."""
    s = _blank(prs)
    _add_bg(s, BG_PINK)
    _eyebrow(s, Inches(0.6), Inches(0.6),
             _section_label(idx, "The Headline"), on_dark=False)
    _add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(1.3),
              "Total Brand Value Observed.",
              size=44, bold=True, color=FG_DARK,
              spacing=0.95, font=FONT_DISPLAY)
    # The big number
    _add_text(s, Inches(0.6), Inches(2.5), Inches(12), Inches(2.0),
              fmt_money(ctx.total_value),
              size=150, bold=True, color=FG_DARK,
              font=FONT_NUMERIC, spacing=0.85, align="left")
    # "Attributable to Partnership" subtotal, C5-strict version.
    # Reads from the payload's precomputed strict-attributable
    # figure (Brand Lift + adjusted-lift-share of Conversion Value)
    # so we don't double-count total observed conversions as if the
    # entire post-window audience were incremental. Falls back to
    # BLV + full CV only if the payload doesn't carry the strict
    # figure (legacy payloads).
    val = ctx.data.get("valuation") or {}
    attributable = val.get("attributable_to_partnership")
    if attributable is None:
        attributable = ctx.blv + (ctx.cv if ctx.has_conversions else 0)
    attributable = float(attributable)
    _add_text(s, Inches(0.6), Inches(4.75), Inches(6.5), Inches(0.35),
              "OF WHICH ATTRIBUTABLE TO PARTNERSHIP",
              size=10, bold=True, color=FG_DARK, letter_spacing=0.16)
    _add_text(s, Inches(0.6), Inches(5.05), Inches(6.5), Inches(0.55),
              fmt_money(attributable),
              size=32, bold=True, color=FG_DARK, font=FONT_NUMERIC)
    _add_text(s, Inches(0.6), Inches(5.55), Inches(6.5), Inches(0.55),
              "Brand Lift plus the adjusted-lift share of Conversion "
              "Value (strictly incremental). Brand Engagement and "
              "Earned Media price the observed audience footprint, "
              "which includes pre-existing brand interest.",
              size=9, italic=True, color=FG_DARK, spacing=1.3)
    # 4-up breakdown bar across the bottom
    parts = [
        ("Earned Media",      ctx.emv),
        ("Brand Engagement",  ctx.bev),
        ("Brand Lift",        ctx.blv),
    ]
    if ctx.has_conversions:
        parts.append(("Conversion", ctx.cv))
    tile_w = Inches(12.0 / len(parts))
    base_x = Inches(0.6)
    base_y = Inches(6.05)
    for i, (label, value) in enumerate(parts):
        x = base_x + tile_w * i + Inches(0.05 if i > 0 else 0)
        _add_rect(s, x, base_y, tile_w - Inches(0.1), Inches(1.0),
                  fill=FG_DARK, line=None)
        _add_text(s, x + Inches(0.25), base_y + Inches(0.1),
                  tile_w - Inches(0.4), Inches(0.3),
                  label.upper(), size=9, bold=True, color=ACCENT,
                  letter_spacing=0.16)
        _add_text(s, x + Inches(0.25), base_y + Inches(0.4),
                  tile_w - Inches(0.4), Inches(0.5),
                  fmt_money(value), size=20, bold=True, color=FG_LIGHT,
                  font=FONT_NUMERIC)
    _page_footer(s, idx, total, ctx.project, on_dark=False)
    return s


#: Fixed platform roster used across ALL decks. Any brand run that
#: doesn't register a given platform is shown as a zero row rather
#: than being dropped, so two brands measured against the same panel
#: cannot appear to have different rosters (item 15).
#:
#: LinkedIn is intentionally excluded from the display roster - it is
#: a B2B platform with negligible signal in consumer brand
#: partnerships. It stays in the underlying rate table for future
#: B2B runs but is not shown by default.
FIXED_PLATFORM_ROSTER = [
    "Instagram", "TikTok", "YouTube", "Facebook", "X (Twitter)",
    "Reddit", "Pinterest", "Snapchat", "Threads", "Twitch",
    "Direct (Brand Site)",
]

#: Per-platform source attribution for the 2026 CPM rate card. This
#: is printed on the EMV slide so procurement can trace every rate
#: back to a specific, credible benchmark instead of a generic
#: "Hootsuite / Sprout / LinkedIn" boilerplate line (item 14).
CPM_SOURCE_ATTRIBUTION = {
    "Instagram":  "Hootsuite Social Media Benchmark Report 2026",
    "Facebook":   "Hootsuite Social Media Benchmark Report 2026",
    "TikTok":     "Sprout Social 2026 Q1 Benchmarks",
    "YouTube":    "IAB / SMI 2026 CTV+Video Benchmarks",
    "X (Twitter)":"eMarketer 2026 X Ad Rate Guide",
    "Reddit":     "Reddit for Business 2026 Ratecard",
    "Pinterest":  "Hootsuite Social Media Benchmark Report 2026",
    "Snapchat":   "Sprout Social 2026 Q1 Benchmarks",
    "Threads":    "eMarketer 2026 Threads Ad Rate Guide",
    "Twitch":     "Amazon Ads Streaming Benchmarks 2026",
    "LinkedIn":   "LinkedIn Marketing Solutions 2026 Ratecard",
    "Direct (Brand Site)":
        "IAB Display / owned-property engagement equivalent (not paid social)",
}


def _slide_emv(prs, ctx: DeckCtx, idx: int, total: int):
    """Per-platform earned media breakdown - LISA index-list pattern.
    Each row: PLATFORM · incremental consumers · $/Consumer · EMV.

    Column labels: we use "$/Consumer" rather than "CPM" because the
    calculation multiplies incremental engaged consumers by a
    per-consumer valuation rate calibrated against 2026 platform
    CPMs at typical exposure frequency. It is a CPM-equivalent
    open-market value per attributable consumer, not a per-mille
    impression rate."""
    s = _blank(prs)
    _add_bg(s, BG_DARK)
    _eyebrow(s, Inches(0.6), Inches(0.6),
             _section_label(idx, "Earned Media Value"), on_dark=True)
    _add_text(s, Inches(0.6), Inches(1.0), Inches(11), Inches(1.3),
              f"{fmt_money(ctx.emv)} of earned media,\n"
              f"platform by platform.",
              size=40, bold=True, color=FG_LIGHT,
              spacing=0.95, font=FONT_DISPLAY)
    # Build a fixed-roster row table so both decks always show the
    # same 8 platforms in the same order, with zeros where a
    # platform did not register (fixes item 15's "the roster was
    # chosen after seeing the results" read).
    all_rows = (ctx.data.get("valuation") or {}).get("emv_breakdown") or []
    by_platform = {(r.get("platform") or "").strip(): r for r in all_rows}
    # Show the full fixed roster (11 platforms). Row height is
    # tightened to keep the whole table within the slide grid.
    display_roster = FIXED_PLATFORM_ROSTER
    rows = []
    for plat in display_roster:
        r = by_platform.get(plat, {})
        rows.append({
            "platform": plat,
            "incremental_users_projected":
                int(r.get("incremental_users_projected") or 0),
            "emv_per_user_rate":
                float(r.get("emv_per_user_rate") or 0),
            "emv_value":
                float(r.get("emv_value") or 0),
        })
    if not any(r["emv_value"] for r in rows):
        _add_text(s, Inches(0.6), Inches(4), Inches(11), Inches(0.6),
                  "No platform-level signal in this window.",
                  size=14, italic=True, color=MUTED_DK)
        _page_footer(s, idx, total, ctx.project, on_dark=True)
        return s
    # Grid: PLATFORM at 0.6, INCR CONSUMERS at 3.4, $/CONSUMER at
    # 5.6, EMV at 7.4. The rightmost 2.4in (10.0 → 12.4) hosts the
    # WHY IT MATTERS callout. Column widths and starts are computed
    # up front so no header can overrun another (fixes item 18).
    table_top = Inches(2.6)
    row_h = Inches(0.30)
    C_PLAT     = (Inches(0.60), Inches(2.60))
    C_INCR     = (Inches(3.30), Inches(2.10))
    C_RATE     = (Inches(5.55), Inches(1.70))
    C_EMV      = (Inches(7.40), Inches(2.30))
    headers = [
        ("PLATFORM",        C_PLAT, "left"),
        ("INCR. CONSUMERS", C_INCR, "right"),
        ("$ / CONSUMER",    C_RATE, "right"),
        ("EMV",             C_EMV,  "right"),
    ]
    for label, (x, w), align in headers:
        _add_text(s, x, table_top, w, Inches(0.3),
                  label, size=9, bold=True, color=MUTED_DK,
                  letter_spacing=0.16, align=align)
    # Underline below header
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6),
                              table_top + Inches(0.32),
                              C_EMV[0] + C_EMV[1] - Inches(0.6),
                              Emu(6000))
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = MUTED_DK
    # Rows
    for i, r in enumerate(rows):
        y = table_top + Inches(0.38) + row_h * i
        is_zero = r["emv_value"] <= 0
        row_color = MUTED_DK if is_zero else FG_LIGHT
        _add_text(s, C_PLAT[0], y, C_PLAT[1], row_h,
                  r["platform"], size=11, bold=True, color=row_color)
        # C11: use literal 0 (not em-dash) for platforms with no
        # observed signal, matching the slide's own note that zero
        # rows indicate absence of signal.
        _add_text(s, C_INCR[0], y, C_INCR[1], row_h,
                  fmt_num(r["incremental_users_projected"]) if r["incremental_users_projected"] else "0",
                  size=11, color=row_color, font=FONT_NUMERIC, align="right")
        _add_text(s, C_RATE[0], y, C_RATE[1], row_h,
                  f"${r['emv_per_user_rate']:.2f}",
                  size=11, color=MUTED_DK, font=FONT_NUMERIC, align="right")
        _add_text(s, C_EMV[0], y, C_EMV[1], row_h,
                  fmt_money(r["emv_value"]) if not is_zero else "$0",
                  size=11, bold=(not is_zero),
                  color=(ACCENT if not is_zero else MUTED_DK),
                  font=FONT_NUMERIC, align="right")
    # Right-side "WHY IT MATTERS" callout, LISA-style
    callout_x = Inches(10.0)
    callout_top = Inches(2.6)
    _add_rect(s, callout_x, callout_top, Inches(2.85), Inches(3.5),
              fill=None, line=ACCENT, line_w=1.0)
    _add_text(s, callout_x + Inches(0.2), callout_top + Inches(0.15),
              Inches(2.5), Inches(0.3),
              "WHY IT MATTERS", size=9, bold=True, color=ACCENT,
              letter_spacing=0.18)
    _add_text(s, callout_x + Inches(0.2), callout_top + Inches(0.5),
              Inches(2.5), Inches(2.9),
              f"Each incremental brand consumer on a platform is priced at "
              f"a CPM-equivalent per-consumer rate calibrated against 2026 "
              f"platform benchmarks (see sourcing below). This is the "
              f"open-market dollar value of the social audience the "
              f"partnership earned for {ctx.brand}, not a per-mille "
              f"impression cost. Rows are the fixed roster used across "
              f"all decks; zeros indicate no observed signal.",
              size=9, color=FG_LIGHT, spacing=1.32)
    # C13: source line now reads from the payload's rate_sources so
    # the source text can be updated without a code push. Covers
    # every platform in the fixed roster (not just the first 4)
    # since all 11 platforms need to be sourced explicitly per
    # Liz's C13.
    payload_srcs = ((ctx.data.get("valuation") or {})
                    .get("rate_sources", {})
                    .get("emv_per_user")) or {}
    # If every platform's rate source is the same string (as under
    # the Crosswalk internal-convention relabel), collapse to a
    # single sentence rather than repeating 11 times.
    unique_srcs = set()
    for plat in display_roster:
        rate = float(((ctx.data.get("valuation") or {})
                      .get("rates", {}).get("emv_per_user") or {})
                     .get(plat, 0))
        if rate:
            src = payload_srcs.get(plat) or CPM_SOURCE_ATTRIBUTION.get(plat, "")
            if src:
                unique_srcs.add(src)
    if len(unique_srcs) == 1:
        single = unique_srcs.pop().rstrip(". ")
        sources_txt = f"All platform rates: {single}. Editable in the live dashboard."
    else:
        bits = []
        for plat in display_roster:
            rate = float(((ctx.data.get("valuation") or {})
                          .get("rates", {}).get("emv_per_user") or {})
                         .get(plat, 0))
            if rate:
                src = payload_srcs.get(plat) or CPM_SOURCE_ATTRIBUTION.get(plat, "")
                if src:
                    bits.append(f"{plat} ${rate:.2f}: {src}")
        sources_txt = "Rates and sources: " + "; ".join(bits)
    # C8: anchored higher and given more vertical room so the source
    # line and the overlap note never collide with the page footer.
    _add_text(s, Inches(0.6), Inches(6.20), Inches(12.1), Inches(0.34),
              sources_txt, size=7, italic=True, color=MUTED_DK,
              spacing=1.15)
    overlap = (
        "Note: platform-level incremental consumers can sum above "
        "the brand-level unique-consumer total because consumers can "
        "be counted on multiple platforms. EMV prices per-platform "
        "earned reach; the unique-consumer total is used for Brand "
        "Lift only."
    )
    _add_text(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.42),
              overlap, size=7, italic=True, color=MUTED_DK, spacing=1.15)
    _page_footer(s, idx, total, ctx.project, on_dark=True)
    return s


def _slide_brand_engagement_value(prs, ctx: DeckCtx, idx: int, total: int):
    """BEV calculator slide. Cream bg, LISA stat-tile rhythm.

    Calculator starts higher on the slide (y=3.0) than the earlier
    build so the total bar clears the source-line row at y=6.05 and
    the source footer at y=6.65 without collision (fixes C8)."""
    s = _blank(prs)
    _add_bg(s, BG_CREAM)
    _eyebrow(s, Inches(0.6), Inches(0.6),
             _section_label(idx, "Brand Engagement Value"),
             on_dark=False)
    _add_text(s, Inches(0.6), Inches(1.0), Inches(11), Inches(1.3),
              f"{fmt_money(ctx.bev)} of observed\npost-campaign engagement.",
              size=40, bold=True, color=FG_DARK,
              spacing=0.95, font=FONT_DISPLAY)
    # Math row - raised to y=3.0 to leave clean margin below.
    rate = float((ctx.data.get("valuation") or {})
                 .get("rates", {}).get("bev_per_user") or 0)
    _math_calculator(
        s, Inches(0.6), Inches(3.0), Inches(8.0),
        [
            ("Post-Campaign Engaged Consumers", fmt_num(ctx.post_users_proj)),
            ("× $/consumer rate",               f"$ {rate:.2f}"),
        ],
        total_label="Brand Engagement Value",
        total_value=fmt_money(ctx.bev),
        on_dark=False,
        accent=RGBColor(0x63, 0x66, 0xF1),
    )
    # Why-it-matters, matched to calculator's new top.
    _add_rect(s, Inches(9.0), Inches(3.0), Inches(3.8), Inches(2.8),
              fill=None, line=FG_DARK, line_w=0.75)
    _add_text(s, Inches(9.25), Inches(3.15), Inches(3.4), Inches(0.3),
              "WHY IT MATTERS", size=9, bold=True, color=FG_DARK,
              letter_spacing=0.18)
    _add_text(s, Inches(9.25), Inches(3.5), Inches(3.4), Inches(2.2),
              "Every consumer who engaged with the brand in the post-window "
              "carries open-market value. This prices the total observed "
              f"post-campaign audience for {ctx.brand}, most of whom "
              "carried pre-existing brand interest that the campaign did "
              "not create. See Brand Lift Value for the net-new share.",
              size=10, color=FG_DARK, spacing=1.35)
    # Rate sourcing footnote. Reads from the payload's
    # valuation.rate_sources.bev_per_user rather than a hardcoded
    # citation, so relabeling the source (C4) doesn't require a code
    # push. Vertically anchored high enough that it doesn't collide
    # with either the math total bar OR the page footer (fixes C8).
    src_line = (
        (ctx.data.get("valuation") or {})
        .get("rate_sources", {})
        .get("bev_per_user")
        or "Rate: Crosswalk internal valuation convention (2026). "
           "Editable in the live dashboard."
    )
    _add_text(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.55),
              src_line, size=8, italic=True, color=MUTED_LT, spacing=1.25)
    _add_text(s, Inches(0.6), Inches(6.65), Inches(12), Inches(0.25),
              ctx.sources_line, size=7, italic=True, color=MUTED_LT)
    _page_footer(s, idx, total, ctx.project, on_dark=False)
    return s


def _slide_brand_lift_value(prs, ctx: DeckCtx, idx: int, total: int):
    """BLV slide - the net-new engagement attributable to the
    partnership. We use raw pre/post delta minus size-matched gen-pop
    control drift ('adjusted lift'); we do NOT call it
    difference-in-differences because both cohorts are drawn from the
    same panel, so a true DiD design would require a genuine holdout
    (see D4 in the QC memo)."""
    s = _blank(prs)
    _add_bg(s, BG_DARK)
    _eyebrow(s, Inches(0.6), Inches(0.6),
             _section_label(idx, "Brand Lift Value"), on_dark=True)
    _add_text(s, Inches(0.6), Inches(1.0), Inches(11), Inches(1.3),
              f"{fmt_money(ctx.blv)} of net-new\nengagement, above control drift.",
              size=40, bold=True, color=FG_LIGHT,
              spacing=0.95, font=FONT_DISPLAY)
    rate = float((ctx.data.get("valuation") or {})
                 .get("rates", {}).get("blv_per_incr_user") or 0)
    cg = ctx.data.get("control_group") or {}
    treat_dp   = float(cg.get("treat_delta_pp")   or 0)
    control_dp = float(cg.get("control_delta_pp") or 0)
    incr_pp    = float(cg.get("incremental_lift_pp") or (treat_dp - control_dp))
    # Calculator raised to y=2.85 AND row height tightened to 0.60"
    # so the 3-row math block + total bar clears the source line
    # anchored at y=6.30 (fixes C8). row_h=0.65 uses the calculator's
    # smaller label/value fonts, which still read cleanly.
    _math_calculator(
        s, Inches(0.6), Inches(2.85), Inches(8.0),
        [
            (f"Adjusted incremental lift  "
             f"(treated {treat_dp:+.2f}pp minus control {control_dp:+.2f}pp)",
             f"{incr_pp:+.2f}pp"),
            (f"x U.S. consumer audience ({fmt_num(ctx.audience_proj)})",
             fmt_num(ctx.incr_users)),
            ("x $/incremental-consumer rate", f"$ {rate:.2f}"),
        ],
        total_label="Brand Lift Value",
        total_value=fmt_money(ctx.blv),
        on_dark=True,
        accent=RGBColor(0xF5, 0x9E, 0x0B),
        row_h=Inches(0.65),
    )
    _add_rect(s, Inches(9.0), Inches(2.85), Inches(3.85), Inches(3.0),
              fill=None, line=ACCENT, line_w=1.0)
    _add_text(s, Inches(9.2), Inches(3.00), Inches(3.5), Inches(0.3),
              "WHY IT MATTERS", size=9, bold=True, color=ACCENT,
              letter_spacing=0.18)
    detail = (
        "Brand Lift Value isolates net-new engagement attributable to "
        "the partnership. We subtract a size-matched gen-pop control "
        "cohort's natural drift from the treated cohort's raw pre/post "
        "delta so we don't double-count baseline brand interest the "
        "campaign didn't move.\n\n"
        "Because both cohorts are drawn from the same 30M+ panel, this "
        "is an adjusted-lift design, not a true difference-in-"
        "differences. A holdout-based DiD is on the roadmap."
    )
    _add_text(s, Inches(9.2), Inches(3.35), Inches(3.5), Inches(2.4),
              detail, size=9, color=FG_LIGHT, spacing=1.35)
    # C4: read the rate source from the payload rather than hardcoding.
    # C8: raised the vertical anchor to keep this line clear of the
    # $/incremental-consumer total bar and the page footer.
    src_line = (
        (ctx.data.get("valuation") or {})
        .get("rate_sources", {})
        .get("blv_per_incr_user")
        or "Rate: Crosswalk internal valuation convention (2026). "
           "Editable in the live dashboard."
    )
    _add_text(s, Inches(0.6), Inches(6.30), Inches(12.1), Inches(0.55),
              src_line, size=7, italic=True, color=MUTED_DK, spacing=1.25)
    # C8: source footer was missing on slide 8 (only BLV slide). Add
    # the Source: Crosswalk BehaviorGraph line explicitly.
    _add_text(s, Inches(0.6), Inches(6.75), Inches(12), Inches(0.25),
              ctx.sources_line, size=7, italic=True, color=MUTED_DK)
    _page_footer(s, idx, total, ctx.project, on_dark=True)
    return s


def _slide_conversion_value(prs, ctx: DeckCtx, idx: int, total: int):
    """CV slide - only shown when conversions are enabled and signal
    is good (skipped for automotive partnerships per .cursor/rules).

    Calculator raised to y=3.0 to match BEV slide 7 and clear the
    source line at the bottom (fixes C8)."""
    s = _blank(prs)
    _add_bg(s, BG_CREAM)
    _eyebrow(s, Inches(0.6), Inches(0.6),
             _section_label(idx, "Conversion Value"), on_dark=False)
    _add_text(s, Inches(0.6), Inches(1.0), Inches(11), Inches(1.3),
              f"{fmt_money(ctx.cv)} of observed\npost-campaign conversions.",
              size=40, bold=True, color=FG_DARK,
              spacing=0.95, font=FONT_DISPLAY)
    conv = ctx.data.get("conversions") or {}
    rate = float((ctx.data.get("valuation") or {})
                 .get("rates", {}).get("conv_value_per_user") or 0)
    _math_calculator(
        s, Inches(0.6), Inches(3.0), Inches(8.0),
        [
            ("Post-Campaign Conversions",
             fmt_num(conv.get("post_users_projected") or 0)),
            ("x Average Order Value (AOV)", f"$ {rate:.2f}"),
        ],
        total_label="Conversion Value",
        total_value=fmt_money(ctx.cv),
        on_dark=False,
        accent=RGBColor(0xEF, 0x44, 0x44),
    )
    _add_rect(s, Inches(9.0), Inches(3.0), Inches(3.8), Inches(2.8),
              fill=None, line=FG_DARK, line_w=0.75)
    _add_text(s, Inches(9.25), Inches(3.15), Inches(3.4), Inches(0.3),
              "WHY IT MATTERS", size=9, bold=True, color=FG_DARK,
              letter_spacing=0.18)
    conv_note = (str(conv.get("note") or "")).strip()
    detail = (
        "Direct purchase-funnel signal: brand.com order confirmations "
        "and geo-panel retail visits observed for brand-owned storefronts "
        "and brand-tagged SKU landing pages within the attribution "
        "window. AOV is editable in the live dashboard."
    )
    if conv_note:
        detail = conv_note
    _add_text(s, Inches(9.25), Inches(3.5), Inches(3.4), Inches(2.2),
              detail, size=10, color=FG_DARK, spacing=1.35)
    # C4: read the AOV source from the payload so the internal-
    # convention relabel takes effect without code changes.
    # C8: vertical positioning aligned with slides 7 and 8 so the
    # three value slides read as a consistent set.
    src_line = (
        (ctx.data.get("valuation") or {})
        .get("rate_sources", {})
        .get("conv_value_per_user")
        or f"AOV: ${rate:.2f} = Crosswalk internal valuation convention "
           "(2026). Editable in the live dashboard."
    )
    _add_text(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.55),
              src_line, size=8, italic=True, color=MUTED_LT, spacing=1.25)
    _add_text(s, Inches(0.6), Inches(6.65), Inches(12), Inches(0.25),
              ctx.sources_line, size=7, italic=True, color=MUTED_LT)
    _page_footer(s, idx, total, ctx.project, on_dark=False)
    return s


def _math_calculator(slide, left, top, width, rows, *,
                     total_label: str, total_value: str,
                     on_dark: bool, accent: RGBColor,
                     row_h: Optional[Any] = None):
    """Render a calculator-style math block. Matches the dashboard's
    in-card calculator (label · value, then a single total row).

    row_h is auto-tuned to fit within a 2.6" tall envelope when there
    are more than 3 rows, so slides with a full DiD walk-through
    (Brand Lift Value) don't overflow into the footer."""
    body_color = FG_LIGHT if on_dark else FG_DARK
    muted = MUTED_DK if on_dark else MUTED_LT
    if row_h is None:
        # Auto-tune: 3 or fewer rows → 0.85"; 4-5 rows → 0.50";
        # anything above that gets tighter still.
        if len(rows) <= 3:
            row_h = Inches(0.85)
        elif len(rows) <= 5:
            row_h = Inches(0.50)
        else:
            row_h = Inches(0.40)
    label_size = 14 if row_h.inches >= 0.7 else 11
    value_size = 24 if row_h.inches >= 0.7 else 18
    for i, (label, value) in enumerate(rows):
        y = top + row_h * i
        _add_text(slide, left, y + Inches(0.08),
                  Inches(width.inches * 0.65), Inches(row_h.inches),
                  label, size=label_size, color=muted, anchor="middle")
        _add_text(slide, left + Inches(width.inches * 0.65),
                  y + Inches(0.02),
                  Inches(width.inches * 0.35), Inches(row_h.inches),
                  value, size=value_size, bold=True, color=body_color,
                  font=FONT_NUMERIC, align="right", anchor="middle")
        if i < len(rows) - 1:
            sep = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, y + row_h - Emu(3500),
                width, Emu(3500))
            sep.line.fill.background()
            sep.fill.solid()
            sep.fill.fore_color.rgb = muted
    # Total bar
    total_y = top + row_h * len(rows) + Inches(0.15)
    _add_rect(slide, left, total_y, width, Inches(0.05),
              fill=accent, line=None)
    _add_text(slide, left, total_y + Inches(0.2),
              Inches(width.inches * 0.65), Inches(0.5),
              total_label.upper(), size=11, bold=True, color=muted,
              letter_spacing=0.16)
    _add_text(slide, left + Inches(width.inches * 0.55),
              total_y + Inches(0.1), Inches(width.inches * 0.45),
              Inches(0.9),
              total_value, size=36, bold=True, color=accent,
              font=FONT_NUMERIC, align="right")


def _slide_demographics(prs, ctx: DeckCtx, idx: int, total: int):
    s = _blank(prs)
    _add_bg(s, BG_CREAM)
    _eyebrow(s, Inches(0.6), Inches(0.6),
             _section_label(idx, "Who Watched the Integration"),
             on_dark=False)
    _add_text(s, Inches(0.6), Inches(1.0), Inches(11), Inches(1.3),
              "Wheel of Fortune viewer profile.",
              size=40, bold=True, color=FG_DARK,
              spacing=0.95, font=FONT_DISPLAY)
    # Source line - critical to disclose that this is the SHOW
    # viewer profile, not brand-specific engagers. This is why the
    # profile is identical across the Coca-Cola and Pepsi decks:
    # both counterfactuals measure the same viewer cohort. (item 1)
    demo_src = (
        (ctx.data.get("demographics") or {}).get("source_note")
        or "Source: Crosswalk BehaviorGraph, Wheel of Fortune Next Day Air "
        "viewer profile (S43 Eps 191 to 195, streamed 2 to 6 June 2026). "
        "This is the audience that watched the integration, not the "
        "brand's engager base; the two brand decks show the same "
        "distributions because both measure this shared viewer cohort."
    )
    _add_text(s, Inches(0.6), Inches(2.15), Inches(12), Inches(0.55),
              demo_src, size=10, italic=True, color=MUTED_LT, spacing=1.35)
    demos = (ctx.data.get("demographics") or {}).get("post") or {}
    cats = [
        ("AGE",       demos.get("age")       or []),
        ("GENDER",    demos.get("gender")    or []),
        ("ETHNICITY", demos.get("ethnicity") or []),
        ("INCOME",    demos.get("income")    or []),
    ]
    # Drop any category with no data so the slide keeps a tight grid.
    cats = [(lab, items) for lab, items in cats if items]
    n = max(1, len(cats))
    base_x = Inches(0.6)
    base_y = Inches(3.05)
    gap = 0.15
    avail_w = 13.333 - 2 * base_x.inches   # 12.13"
    col_w = Inches((avail_w - gap * (n - 1)) / n)
    # Fixed row height chosen so the tallest category (income has
    # up to 8 buckets: <$25K, $25-35K, $35-50K, $50-75K, $75-100K,
    # $100-150K, $150-250K, $250K+) fits inside the slide without
    # overflowing the source footer at y=6.75. All buckets are
    # shown - Liz's "we didn't have any $150K+ viewers?" callout
    # was caused by a top-5 truncation that hid the two upper
    # income buckets ($150K-$249K and $250K+). Showing every
    # bucket also closes the "not summing to 100%" observation.
    row_h = 0.46
    for i, (label, items) in enumerate(cats):
        x = base_x + (col_w + Inches(gap)) * i
        rows_full = sorted(items, key=lambda r: float(r.get("percentage") or 0),
                           reverse=True)
        _add_text(s, x, base_y, col_w, Inches(0.3),
                  label, size=10, bold=True, color=MUTED_LT,
                  letter_spacing=0.18)
        for j, r in enumerate(rows_full):
            y = base_y + Inches(0.42 + j * row_h)
            _add_text(s, x, y, Inches(col_w.inches * 0.65), Inches(row_h - 0.02),
                      str(r.get("value") or "n/a"), size=12, color=FG_DARK)
            _add_text(s, x + Inches(col_w.inches * 0.65), y,
                      Inches(col_w.inches * 0.35), Inches(row_h - 0.02),
                      fmt_pct(r.get("percentage") or 0),
                      size=12, bold=True, color=FG_DARK,
                      font=FONT_NUMERIC, align="right")
    _add_text(s, Inches(0.6), Inches(6.75), Inches(12), Inches(0.25),
              ctx.sources_line, size=7, italic=True, color=MUTED_LT)
    _page_footer(s, idx, total, ctx.project, on_dark=False)
    return s


def _slide_top_touchpoints(prs, ctx: DeckCtx, idx: int, total: int):
    """Top brand touchpoints, ranked by observed hits.

    We intentionally label the columns PRE HITS / POST HITS (not
    REACH). A hit is an observation of a brand-property visit; a
    single consumer can generate multiple hits at the same
    touchpoint. This is why the sum of touchpoint hit-gains can
    exceed the total unique incremental consumer count (Liz's item
    3): those are different populations counted at different
    granularities. The methodology footnote calls this out
    directly."""
    s = _blank(prs)
    _add_bg(s, BG_DARK)
    _eyebrow(s, Inches(0.6), Inches(0.6),
             _section_label(idx, f"Most Engaged {ctx.brand} Touchpoints"),
             on_dark=True)
    _add_text(s, Inches(0.6), Inches(1.0), Inches(11), Inches(1.3),
              f"Where {ctx.brand} showed up.",
              size=40, bold=True, color=FG_LIGHT,
              spacing=0.95, font=FONT_DISPLAY)
    # Column layout - grid computed up front to guarantee no
    # overlap and consistent row heights (fixes item 18).
    post = ctx.data.get("top_brand_properties") or []
    pre  = ctx.data.get("top_brand_properties_pre") or []
    pre_map = {(p.get("common_name") or "").lower():
               int(p.get("hits_projected") or 0) for p in pre}
    rows = sorted(post, key=lambda r: int(r.get("hits_projected") or 0),
                  reverse=True)[:7]
    C_NAME     = (Inches(0.60), Inches(5.20))
    C_PRE      = (Inches(5.90), Inches(1.75))
    C_POST     = (Inches(7.75), Inches(1.75))
    C_LIFT     = (Inches(9.60), Inches(3.05))
    table_top  = Inches(2.75)
    row_h      = Inches(0.42)
    header_style = dict(size=9, bold=True, color=MUTED_DK,
                        letter_spacing=0.16)
    _add_text(s, C_NAME[0], table_top, C_NAME[1], Inches(0.3),
              "BRAND TOUCHPOINT", **header_style)
    _add_text(s, C_PRE[0], table_top, C_PRE[1], Inches(0.3),
              "PRE HITS", align="right", **header_style)
    _add_text(s, C_POST[0], table_top, C_POST[1], Inches(0.3),
              "POST HITS", align="right", **header_style)
    _add_text(s, C_LIFT[0], table_top, C_LIFT[1], Inches(0.3),
              "HITS LIFT", align="right", **header_style)
    # Underline
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6),
                              table_top + Inches(0.32),
                              C_LIFT[0] + C_LIFT[1] - Inches(0.6),
                              Emu(6000))
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = MUTED_DK
    for i, r in enumerate(rows):
        y = table_top + Inches(0.42) + row_h * i
        raw_name = str(r.get("common_name") or "n/a")
        # C12: preserve source display casing rather than force
        # title() (which capitalizes small prepositions like "Of"
        # in "World of Coca-Cola"). Payloads now provide brand-
        # property display names in the exact case they should
        # render.
        display_name = _touchpoint_display_name(raw_name)
        post_h = int(r.get("hits_projected") or 0)
        pre_h = pre_map.get(raw_name.lower(), 0)
        _add_text(s, C_NAME[0], y, C_NAME[1], row_h,
                  display_name, size=13, bold=True, color=FG_LIGHT)
        _add_text(s, C_PRE[0], y, C_PRE[1], row_h,
                  fmt_num(pre_h) if pre_h else "0",
                  size=12, color=FG_LIGHT, font=FONT_NUMERIC,
                  align="right")
        _add_text(s, C_POST[0], y, C_POST[1], row_h,
                  fmt_num(post_h),
                  size=12, color=FG_LIGHT, font=FONT_NUMERIC,
                  align="right")
        if pre_h:
            lift = (post_h - pre_h) / pre_h * 100
            color = LIFT_UP if lift >= 0 else LIFT_DOWN
            _add_text(s, C_LIFT[0], y, C_LIFT[1], row_h,
                      fmt_signed_pct(lift),
                      size=13, bold=True, color=color, font=FONT_NUMERIC,
                      align="right")
        else:
            _add_text(s, C_LIFT[0], y, C_LIFT[1], row_h,
                      "NEW",
                      size=11, bold=True, color=ACCENT,
                      letter_spacing=0.18, align="right")
    # Methodology footnote - the critical clarification that
    # touchpoint hits and unique incremental consumers are different
    # populations at different granularities (Liz items 3, 4).
    hits_note = (
        "A HIT is one observation of a brand-property visit. One "
        "consumer can generate multiple hits at the same touchpoint, "
        "so the sum of hit gains across touchpoints can exceed the "
        "unique incremental consumer count on the Brand Lift slide. "
        "The two numbers measure different things (impressions of the "
        "brand vs. unique attributable consumers) and are not "
        "expected to reconcile."
    )
    _add_text(s, Inches(0.6), Inches(6.20), Inches(12.2), Inches(0.65),
              hits_note, size=8, italic=True, color=MUTED_DK,
              spacing=1.3)
    _add_text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.25),
              ctx.sources_line, size=7, italic=True, color=MUTED_DK)
    _page_footer(s, idx, total, ctx.project, on_dark=True)
    return s


def _slide_year_over_year(prs, ctx: DeckCtx, idx: int, total: int):
    s = _blank(prs)
    _add_bg(s, BG_CREAM)
    _eyebrow(s, Inches(0.6), Inches(0.6),
             _section_label(idx, "Year Over Year"), on_dark=False)
    _add_text(s, Inches(0.6), Inches(1.0), Inches(11), Inches(1.3),
              "How the partnership compounded.",
              size=40, bold=True, color=FG_DARK,
              spacing=0.95, font=FONT_DISPLAY)
    years = ctx.data.get("yearly_breakdown") or []
    if not years:
        _add_text(s, Inches(0.6), Inches(4), Inches(11), Inches(0.5),
                  "No multi-year decomposition available.",
                  size=14, italic=True, color=MUTED_LT)
        _page_footer(s, idx, total, ctx.project, on_dark=False)
        return s
    # Bar chart approximation: stat plates per year, height-proportional
    # to EMV. We don't embed Chart.js-style charts in pptx - the visual
    # density is in the stat plates themselves.
    max_emv = max((float(y.get("earned_media_value") or 0) for y in years),
                  default=1)
    plate_count = len(years)
    base_x = Inches(0.6)
    base_y = Inches(3.1)
    avail_w = Inches(12.1)
    plate_w = Inches((avail_w.inches - 0.15 * (plate_count - 1)) / plate_count)
    chart_h = Inches(3.0)
    for i, y in enumerate(years):
        emv = float(y.get("earned_media_value") or 0)
        height_frac = emv / max_emv if max_emv else 0
        bar_h = Inches(max(0.1, chart_h.inches * height_frac))
        x = base_x + (plate_w + Inches(0.15)) * i
        # Bar
        bar_top = base_y + chart_h - bar_h
        _add_rect(s, x, bar_top, plate_w, bar_h, fill=FG_DARK, line=None)
        # EMV on bar
        _add_text(s, x, bar_top - Inches(0.05), plate_w, Inches(0.4),
                  fmt_money(emv), size=14, bold=True, color=FG_DARK,
                  font=FONT_NUMERIC, align="center")
        # Year label below
        _add_text(s, x, base_y + chart_h + Inches(0.1), plate_w, Inches(0.3),
                  str(y.get("label") or y.get("year") or ""),
                  size=11, bold=True, color=FG_DARK, align="center",
                  letter_spacing=0.14)
        # Users line
        users = int(y.get("users_projected") or 0)
        _add_text(s, x, base_y + chart_h + Inches(0.4), plate_w, Inches(0.3),
                  f"{fmt_num(users)} consumers",
                  size=9, color=MUTED_LT, align="center")
    _add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.3),
              ctx.sources_line, size=8, italic=True, color=MUTED_LT)
    _page_footer(s, idx, total, ctx.project, on_dark=False)
    return s


def _slide_final_insight(prs, ctx: DeckCtx, idx: int, total: int):
    s = _blank(prs)
    _add_bg(s, BG_LAV)
    _eyebrow(s, Inches(0.6), Inches(0.6),
             _section_label(idx, "The Final Insight"), on_dark=False)
    l1, l2 = _final_insight_lines(ctx)
    _add_text(s, Inches(0.6), Inches(2.2), Inches(12.1), Inches(1.6),
              l1, size=64, bold=True, color=FG_DARK,
              font=FONT_DISPLAY, spacing=0.95)
    _add_text(s, Inches(0.6), Inches(3.6), Inches(12.1), Inches(1.6),
              l2, size=64, bold=True, color=FG_DARK,
              font=FONT_DISPLAY, spacing=0.95)
    # C10: previously the strip showed both a relative-adjusted-lift
    # figure (from _final_insight_lines) and a raw-relative-lift
    # figure (delta_rel) on the same slide with no labels, so a
    # reader saw two different % figures and couldn't tell which was
    # which. Drop delta_rel and lead with the adjusted-lift number
    # since that is the one the deck argues from.
    lift_rel = (ctx.incr_lift_rel
                if ctx.incr_lift_rel is not None
                else ctx.delta_rel)
    lift_label = (
        f"+{lift_rel:.1f}% adjusted engagement lift"
        if lift_rel > 0 else "adjusted lift at or below control drift"
    )
    closing = (
        f"{fmt_money(ctx.total_value)} in total brand value observed, "
        f"{fmt_num(ctx.post_users_proj)} U.S. consumers reached, "
        f"{lift_label}"
    )
    _add_text(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(0.5),
              closing, size=16, italic=True, color=FG_DARK)
    _page_footer(s, idx, total, ctx.project, on_dark=False)
    return s


def _slide_source(prs, ctx: DeckCtx, idx: int, total: int):
    s = _blank(prs)
    _add_bg(s, BG_DARK)
    _eyebrow(s, Inches(0.6), Inches(0.6),
             "Source & Methodology", on_dark=True)
    _add_text(s, Inches(0.6), Inches(1.0), Inches(11), Inches(1.3),
              "How we got here.",
              size=40, bold=True, color=FG_LIGHT,
              spacing=0.95, font=FONT_DISPLAY)
    notes = [
        ("DATA SOURCE",
         "Crosswalk BehaviorGraph: a 30M+ opted-in U.S. consumer panel. "
         "Every observation in this deck reflects real behavioral signal "
         "(visits, engagements, brand touchpoints); no surveys, no "
         f"self-report. Observed sample: {fmt_num(ctx.panel_size)} "
         "panelists (the viewer cohort for this integration)."),
        ("INCIDENCE RATES",
         "All pre and post penetration percentages are incidence rates: "
         "the count of unique brand engagers observed in the window "
         f"divided by the {fmt_num(ctx.panel_size)}-panelist sample. "
         "Raw numerator and denominator are printed on slide 4 alongside "
         "each percentage so the rate can be audited directly."),
        ("PROJECTION",
         f"The observed {fmt_num(ctx.panel_size)}-panelist sample is "
         f"projected to a {fmt_num(ctx.audience_proj)}-consumer U.S. "
         "viewing audience universe using a calibrated cohort weight. "
         "Numbers labeled 'projected' or 'U.S. consumers' have been "
         "weight-extrapolated."),
        ("ATTRIBUTION",
         f"Brand touchpoints are counted within {ctx.attribution_d} days "
         "of campaign end: long enough to capture delayed brand "
         "behavior, tight enough to keep the causal read clean."),
        ("BRAND LIFT",
         "Brand Lift Value uses an adjusted-lift design: the size-"
         "matched gen-pop control cohort's natural drift is subtracted "
         "from the treated cohort's raw pre/post delta. Both cohorts "
         "are drawn from the same panel, so this is a within-panel "
         "adjusted-lift read, not a true difference-in-differences "
         "(which would require a genuine holdout); a holdout-based "
         "DiD is on the roadmap."),
        ("SIGNIFICANCE",
         "Same-panel pre/post design. Primary test: pooled two-sample "
         "z on the paired marginals (conservative for a within-panel "
         "read). We report the point estimate, the z-score, and a 95% "
         "confidence interval derived from the same variance. At n=10M "
         "the detection floor is roughly a tenth of a percentage point; "
         "p-values are omitted on client-facing slides because they "
         "collapse to zero for anything above the floor and add no "
         "decision-relevant information."),
    ]
    base_y = 2.35
    row_spacing = 0.75
    for i, (label, body) in enumerate(notes):
        y = Inches(base_y + i * row_spacing)
        _add_text(s, Inches(0.6), y, Inches(3.0), Inches(0.3),
                  label, size=10, bold=True, color=ACCENT,
                  letter_spacing=0.18)
        _add_text(s, Inches(3.8), y, Inches(9.0), Inches(0.8),
                  body, size=10, color=FG_LIGHT, spacing=1.3)
    # Final logo + URL line
    try:
        s.shapes.add_picture(ctx.logo, Inches(0.6), Inches(6.6),
                             height=Inches(0.4))
    except Exception:
        _add_text(s, Inches(0.6), Inches(6.75), Inches(5), Inches(0.3),
                  "CROSSWALK · BehaviorGraph",
                  size=11, bold=True, color=FG_LIGHT, letter_spacing=0.18)
    _add_text(s, Inches(8), Inches(6.75), Inches(4.8), Inches(0.3),
              "behaviorgraph.com · 2026",
              size=9, italic=True, color=MUTED_DK, align="right",
              letter_spacing=0.14)
    return s


# ─────────────────────────────────────────────────────────────────────
#  Public entry point
# ─────────────────────────────────────────────────────────────────────
def build_deck(data: dict,
               image_url: Optional[str] = None,
               category: Optional[str] = None) -> bytes:
    """Build a full BPIQ analysis deck and return the .pptx bytes.

    Parameters
    ----------
    data       : the BPIQ result JSON (matches /api/brand-partnership-iq/
                 results/<key> 'data' field).
    image_url  : admin-uploaded hero photo from the metadata sidecar
                 (used on the cover slide). Optional - falls back to a
                 typographic plate if missing or unfetchable.
    category   : admin-managed category (case-insensitive). Used to
                 suppress conversion slides for 'automotive' runs.
    """
    ctx = _derive_context(data, image_url, category)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # Decide which slides participate in this deck.
    builders = [
        _slide_cover,
        _slide_methodology,
        _slide_audience_scale,
        _slide_engagement_lift,
        _slide_total_value_hero,
        _slide_emv,
        _slide_brand_engagement_value,
        _slide_brand_lift_value,
    ]
    if ctx.has_conversions:
        builders.append(_slide_conversion_value)
    builders.extend([
        _slide_demographics,
        _slide_top_touchpoints,
    ])
    if ctx.has_yearly:
        builders.append(_slide_year_over_year)
    builders.extend([
        _slide_final_insight,
        _slide_source,
    ])

    total = len(builders)
    for i, builder in enumerate(builders, start=1):
        try:
            builder(prs, ctx, i, total)
        except Exception as e:
            # Never let a single slide failure kill the whole deck -
            # render an apology slide and keep going so the user still
            # gets a deck rather than a 500.
            s = _blank(prs)
            _add_bg(s, BG_DARK)
            _add_text(s, Inches(0.6), Inches(3), Inches(12), Inches(1),
                      f"Slide build failed ({builder.__name__}): {e}",
                      size=14, italic=True, color=MUTED_DK)

    out = io.BytesIO()
    prs.save(out)

    # Best-effort cleanup of any tmp files we materialized.
    for p in (ctx.cover_image, ctx.logo):
        try:
            if p and os.path.exists(p):
                os.unlink(p)
        except Exception:
            pass

    return out.getvalue()


def suggested_filename(data: dict, ext: str = "pptx") -> str:
    name = safe_filename(data.get("project_name") or "BPIQ_Deck")
    stamp = (data.get("created_at") or "")[:10] or "deck"
    return f"{name}_Analysis_{stamp}.{ext}"


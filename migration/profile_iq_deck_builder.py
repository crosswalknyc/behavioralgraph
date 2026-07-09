"""Profile IQ deck builder — single-profile and multi-profile comparison decks.

Generates a designed PowerPoint (.pptx) that mirrors the LISA reference
deck's visual language (see /Users/jennamenking/Downloads/LISA_Audience_Profile
2026-06-23 v1.0.pptx):

  * 16:9 (13.33" × 7.50"), Arial throughout
  * Palette: dark charcoal (#0C1618), warm greys (#8A938F / #76868A / #59636A),
    off-white cream (#E9E8E1), accent lavender (#B7B3D8), pop lime (#C7F23E),
    magenta (#E782FF), and electric blue (#3358FF)
  * Editorial slide taxonomy: cover → overview → scale → demographics →
    "one insight" → category deep-dives → case studies → final insight
  * Sentence-case narrative headlines, huge stats with tiny labels, and a
    persistent CROSSWALK · PROFILE IQ footer

Two public entry points:
  * ``build_deck(profile_data, image_url=None, category=None)`` -> bytes
      Single-profile deck. Data payload matches what the Profile IQ dashboard
      already keeps in ``currentDashboardData`` (demographics + demographicsIndex
      + demographicsGenPop + behavioral + sample/projection counters).
  * ``build_combined_deck(profiles: list[dict])`` -> bytes
      Multi-profile comparative deck. Each profile is one entry of the same
      payload shape (name, image_url, category, demographics, behavioral, …).
      Adds side-by-side scale / demo / brand-overlap / category-winner slides.

The Flask endpoints in ``bg-webapp/app.py`` fetch the admin-managed profile
image via ``image_url`` (usually a ``/api/profile-image-file/...`` URL that we
resolve to bytes with ``urllib``) so the cover slide gets the right photo.
Any image fetch failure falls back to a color block with a subject initial
so the deck always builds.

Dependencies (pinned in bg-webapp/requirements.txt):
  * python-pptx >= 1.0.0
  * Pillow >= 10.4.0
"""
from __future__ import annotations

import io
import json
import math
import os
import re
import sys
import urllib.parse
import urllib.request
from datetime import datetime
from typing import Any, Iterable, Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


# =============================================================================
#  Design tokens (LISA palette)
# =============================================================================

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

C_DARK     = RGBColor(0x0C, 0x16, 0x18)   # near-black background
C_CREAM    = RGBColor(0xE9, 0xE8, 0xE1)   # primary light text
C_MUTED    = RGBColor(0x8A, 0x93, 0x8F)   # secondary text
C_MUTED2   = RGBColor(0x76, 0x86, 0x8A)   # tertiary text / footer
C_STROKE   = RGBColor(0x59, 0x63, 0x6A)   # thin rule lines
C_LAVENDER = RGBColor(0xB7, 0xB3, 0xD8)   # accent chip
C_LIME     = RGBColor(0xC7, 0xF2, 0x3E)   # pop accent
C_MAGENTA  = RGBColor(0xE7, 0x82, 0xFF)   # index-hot accent
C_BLUE     = RGBColor(0x33, 0x58, 0xFF)   # data accent

# Chart-friendly cycle used in demographic comparison bars
CYCLE_ACCENT = [C_LAVENDER, C_MAGENTA, C_LIME, C_BLUE, C_CREAM]

FONT_MAIN = "Arial"

US_POPULATION = 329_900_000


# =============================================================================
#  Public API
# =============================================================================


# Canonical section IDs the LLM can request in its `sections` array.
# Every ID maps to exactly one slide builder in the dispatch table below.
#
# The list is deliberately LONG. Individual slide builders self-skip when
# the profile lacks the data to make the slide meaningful — so a thin
# profile still produces a clean 12-slide deck, and a rich profile
# produces the full 22+ slide ethnographic portrait.

_DEFAULT_SECTIONS = (
    "cover",
    "exec_summary",
    "portrait",
    "scale",
    "identity",
    "demographic_deep_dive",
    "life_stage",
    "differ_from_genpop",
    "index_story",
    "day_in_life",
    "geography",
    "commerce_persona",
    "signature_affinities",
    "category_deep_dives",
    "cultural_interests",
    "media_and_social",
    "media_plan",
    "whitespace",
    "takeaways",
    "final",
    "methodology",
    "about",
)

# Mandatory core — the LLM cannot remove these under any circumstance.
# Every real deck must open with cover/exec/portrait, do identity work,
# close with takeaways/final/methodology/about.
_ALWAYS_SECTIONS = (
    "cover", "exec_summary", "portrait",
    "scale", "identity", "demographic_deep_dive", "differ_from_genpop",
    "takeaways", "final", "methodology", "about",
)


def _resolve_sections(profile: dict) -> list[str]:
    """Return the ordered list of section IDs to build for this profile.

    Design (revised 2026-07-09): the LLM does NOT dictate the deck
    structure. Instead the deck ALWAYS starts from the rich default
    list, and the LLM's `sections` field is treated as advisory — used
    only to reorder the optional (non-always) slides. This prevents a
    sparse LLM response from collapsing the deck to 7 slides.

    Every slide builder is responsible for self-skipping when the data
    doesn't support it, so we ship the full default list without fear
    of empty pages.
    """
    default = list(_DEFAULT_SECTIONS)
    brief_sections = _brief_get(profile, "sections", default=None)

    if isinstance(brief_sections, list) and brief_sections:
        brief_order = [str(s).strip().lower() for s in brief_sections
                        if str(s or "").strip()]
        # Optional (non-always) slides the LLM explicitly nominated —
        # reorder these to appear in the LLM's preferred order, but keep
        # every other default slide in place.
        optional_from_brief = [s for s in brief_order
                                if s in default and s not in _ALWAYS_SECTIONS]
        # Slot the LLM's optional picks in the order it gave, then append
        # any remaining default optionals it didn't mention.
        seen = set(optional_from_brief) | set(_ALWAYS_SECTIONS)
        remaining_optionals = [s for s in default if s not in seen]
        # Rebuild sequence: walk the default and, whenever we hit an
        # optional slot, take from the brief's ordering first.
        result: list[str] = []
        opt_iter = iter(optional_from_brief + remaining_optionals)
        for s in default:
            if s in _ALWAYS_SECTIONS:
                result.append(s)
            else:
                try:
                    result.append(next(opt_iter))
                except StopIteration:
                    pass
        return result

    return default


def build_deck(
    data: dict,
    *,
    image_url: Optional[str] = None,
    category: Optional[str] = None,
) -> bytes:
    """Build a single-profile Profile IQ deck and return the .pptx bytes.

    Deck structure is STORY-DRIVEN, not templated. The LLM analyst brief
    nominates the sections to include (a category deep-dive is only added
    if that category clears the reach bar; day-in-the-life ships only if
    the LLM returned a real narrative; whitespace ships only if the LLM
    called one out). Every deck has a mandatory core (cover / exec /
    identity / takeaways / final / methodology / about) but the rest of
    the structure adapts to what the data actually supports.
    """
    profile = _normalize_payload(data, image_url=image_url, category=category)
    profile["_image_bytes"] = _fetch_image_bytes(profile.get("image_url") or "")
    profile["_archetype"] = _detect_archetype(profile)
    profile["_brief"] = _get_analyst_brief(profile, profile["_archetype"])

    prs = _new_presentation()

    # Category deep-dives are the only variable-count section — pick them
    # once, cap at 3 (was 6), and only after the reach gate has fired.
    persona_slides = _pick_persona_category_slides(profile, max_slides=3)

    def _emit_category_deep_dives():
        for i, (cat, items) in enumerate(persona_slides):
            # Alternate layout for editorial rhythm.
            if i % 2 == 0:
                _slide_persona_category(prs, profile, cat, items)
            else:
                _slide_stat_splash(prs, profile, cat, items)

    section_builders = {
        "cover":                 lambda: _slide_cover(prs, profile),
        "exec_summary":          lambda: _slide_exec_summary(prs, profile),
        "portrait":              lambda: _slide_portrait(prs, profile),
        "scale":                 lambda: _slide_scale(prs, profile),
        "identity":              lambda: _slide_demographics(prs, profile),
        "demographic_deep_dive": lambda: _slide_demographic_deep_dive(prs, profile),
        "life_stage":            lambda: _slide_life_stage(prs, profile),
        "differ_from_genpop":    lambda: _slide_differ_from_genpop(prs, profile),
        "index_story":           lambda: _slide_index_story(prs, profile),
        "day_in_life":           lambda: _slide_day_in_life(prs, profile),
        "geography":             lambda: _slide_geography(prs, profile),
        "commerce_persona":      lambda: _slide_commerce_persona(prs, profile),
        "signature_affinities":  lambda: _slide_signature_affinities(prs, profile),
        "category_deep_dives":   _emit_category_deep_dives,
        "cultural_interests":    lambda: _slide_cultural_interests(prs, profile),
        "media_and_social":      lambda: _slide_media_and_social(prs, profile),
        "media_plan":            lambda: _slide_media_plan(prs, profile),
        "whitespace":            lambda: _slide_whitespace(prs, profile),
        "takeaways":             lambda: _slide_takeaways(prs, profile),
        "final":                 lambda: _slide_final_insight(prs, profile),
        "methodology":           lambda: _slide_methodology(prs, profile),
        "about":                 lambda: _slide_about(prs, profile),
    }

    for section_id in _resolve_sections(profile):
        builder = section_builders.get(section_id)
        if builder is None:
            continue
        try:
            builder()
        except Exception as exc:  # pragma: no cover — logged, not raised
            print(f"[profile_iq_deck_builder] section '{section_id}' failed: "
                  f"{exc}; skipping")

    # Post-hoc: rewrite every footer with the true slide index.
    _renumber_pages(prs)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_combined_deck(profiles: Iterable[dict]) -> bytes:
    """Comparative deck across multiple profiles + Gen Pop.

    Slide arc (2 profiles -> ~10 slides, 3 profiles -> ~12 slides):
      1  Cover with per-profile image collage
      2  Overview (what the comparison covers)
      3  Scale bars
      4  Demographics side-by-side
      5-N  Per-profile spotlight (one slide per profile with its hero image
           + top stats + top over-index brand)
      N+1  Category winners
      N+2  Brand overlap
      N+3  Final insight
    """
    normalized = [_normalize_payload(p) for p in profiles if p]
    if not normalized:
        raise ValueError("build_combined_deck requires at least one profile")
    # Pre-fetch every profile image once so cover collage + per-profile
    # spotlights reuse the same bytes. Also compute per-profile archetype +
    # analyst brief so the spotlight slides read in the right voice.
    for pr in normalized:
        pr["_image_bytes"] = _fetch_image_bytes(pr.get("image_url") or "")
        pr["_archetype"]   = _detect_archetype(pr)
        pr["_brief"]       = _get_analyst_brief(pr, pr["_archetype"])

    prs = _new_presentation()

    _slide_combined_cover(prs, normalized)
    _slide_combined_overview(prs, normalized)
    _slide_combined_scale(prs, normalized)
    _slide_combined_demographics(prs, normalized)
    for pr in normalized:
        _slide_combined_profile_spotlight(prs, pr)
    _slide_combined_category_winners(prs, normalized)
    _slide_combined_brand_overlap(prs, normalized)
    _slide_combined_final_insight(prs, normalized)

    _renumber_pages(prs)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def suggested_filename(data: dict, ext: str = "pptx") -> str:
    """Return a safe download filename for a single-profile deck."""
    name = (data.get("name") or data.get("brand") or "Profile").strip() or "Profile"
    safe = re.sub(r"[^A-Za-z0-9._-]+", "_", name).strip("_") or "Profile"
    return f"{safe}_Audience_Profile.{ext}"


def suggested_combined_filename(profiles: Iterable[dict], ext: str = "pptx") -> str:
    names = []
    for p in profiles or []:
        n = (p.get("name") or p.get("brand") or "").strip()
        if n:
            names.append(re.sub(r"[^A-Za-z0-9]+", "", n))
        if len(names) >= 4:
            break
    joined = "_vs_".join(names[:4]) or "Combined"
    return f"{joined}_Audience_Comparison.{ext}"


# =============================================================================
#  Payload normalization
# =============================================================================


def _normalize_payload(
    data: dict,
    *,
    image_url: Optional[str] = None,
    category: Optional[str] = None,
) -> dict:
    """Return a payload with the exact keys the slide builders expect.

    Accepts both the ``currentDashboardData`` shape (camelCase, JS-style) and
    a snake_case variant, so the endpoint doesn't have to translate.
    """
    d = data or {}

    def _pick(*keys, default=None):
        for k in keys:
            if k in d and d[k] not in (None, "", {}, []):
                return d[k]
        return default

    demos      = _pick("demographics", "demos", default={}) or {}
    demos_idx  = _pick("demographicsIndex", "demographics_index", default={}) or {}
    demos_gp   = _pick("demographicsGenPop", "demographics_gen_pop", default={}) or {}
    demos_proj = _pick("demographicsProjection", "demographics_projection", default={}) or {}
    behavioral = _pick("behavioral", default={}) or {}
    locations  = _pick("locations", default=[]) or []
    interests  = _pick("interests", default={}) or {}

    payload = {
        "name":           (_pick("name", "brand", "project_name", default="Audience") or "Audience").strip(),
        "brand_category": (_pick("brand_category", "brandCategory", "category", default=(category or "")) or "").strip(),
        "sample_size":    int(_pick("sample_size", "sampleSize", default=0) or 0),
        "projected_us":   int(_pick("projected_us", "projectedUS", default=0) or 0),
        "date_range":     (_pick("date_range", "dateRange", default="") or "").strip(),
        "image_url":      image_url or _pick("image_url", "imageUrl", default="") or "",
        "s3_key":         (_pick("s3_key", "s3Key", default="") or "").strip(),
        "demographics":              _lowercase_demo_keys(demos),
        "demographics_index":        _lowercase_demo_keys(demos_idx),
        "demographics_gen_pop":      _lowercase_demo_keys(demos_gp),
        "demographics_projection":   _lowercase_demo_keys(demos_proj),
        "behavioral": {k: list(v or []) for k, v in behavioral.items()},
        "locations": list(locations),
        "interests": dict(interests),
    }

    # Strip the SUBJECT itself from behavioral rows before slides ever see
    # them. Green Day appearing at 100% reach in a 'Nimrods A Green Day
    # Comedy' audience is tautological — every audience member exhibits
    # Green Day engagement because that IS how the audience is defined.
    # Featuring it in signature affinities, category deep-dives, or 'anchor
    # brand' captions is analytically wrong. The subject is the frame, not
    # a finding. Added 2026-07-09 after user flagged the Green-Day-in-
    # SIGNATURE-AFFINITIES leak.
    subject_name = payload["name"]
    for cat, items in payload["behavioral"].items():
        kept = []
        dropped = []
        for it in items:
            brand = (it.get("name") or "").strip()
            if _is_subject_brand(brand, subject_name):
                dropped.append(brand)
                continue
            kept.append(it)
        payload["behavioral"][cat] = kept
        if dropped:
            _log(f"[_normalize_payload] stripped subject-matching brands "
                  f"from {cat}: {dropped}")

    # Sort behavioral items by index desc (over-indexers first) so slides pick
    # the "story" brands, not just the biggest-reach ones. Tie-break on pct.
    for cat, items in payload["behavioral"].items():
        items.sort(key=lambda it: (
            -_num(it.get("index", 0)),
            -_num(it.get("pct", 0)),
        ))
    return payload


def _lowercase_demo_keys(demos: dict) -> dict:
    """Normalize outer keys to lowercase (age/gender/ethnicity/...) so the
    slide builders don't need to worry about camelCase or upper variants."""
    out = {}
    for k, v in (demos or {}).items():
        key = str(k or "").strip().lower()
        if key.startswith("sexual"):
            key = "sexual_orientation"
        if key.startswith("parental"):
            key = "parental_status"
        out[key] = dict(v or {})
    return out


def _num(x, default: float = 0.0) -> float:
    try:
        return float(x)
    except (TypeError, ValueError):
        return default


# =============================================================================
#  Subject archetype detection + LLM analyst brief
# =============================================================================
#
# The brief is what turns this from a template deck into a real audience
# analysis. Given the profile payload + a detected subject archetype
# (MOVIE / TV_SHOW / MUSICIAN / ACTOR / ATHLETE / TEAM / BRAND / FRANCHISE /
# PERSONA / GENERIC), we call OpenAI with a consulting-analyst system prompt
# that returns structured JSON describing who these fans are, how they live,
# what makes them different from Gen Pop, where the whitespace opportunity
# is, and archetype-specific activation ideas (e.g. for a MOVIE — trailer
# placement, festival beats, ticket-buying cues, streamer partner fit).
#
# Every slide that used to hardcode narrative copy now consults the brief
# first and falls back to smarter templated copy if the LLM is unavailable
# (no OPENAI_API_KEY set, network error, etc.).


_ARCHETYPE_MOVIE     = "MOVIE"
_ARCHETYPE_TV_SHOW   = "TV_SHOW"
_ARCHETYPE_MUSICIAN  = "MUSICIAN"
_ARCHETYPE_ACTOR     = "ACTOR"
_ARCHETYPE_ATHLETE   = "ATHLETE"
_ARCHETYPE_TEAM      = "TEAM"
_ARCHETYPE_BRAND     = "BRAND"
_ARCHETYPE_FRANCHISE = "FRANCHISE"
_ARCHETYPE_PERSONA   = "PERSONA"
_ARCHETYPE_GENERIC   = "GENERIC"

_ARCHETYPE_LABELS = {
    _ARCHETYPE_MOVIE:     "Film",
    _ARCHETYPE_TV_SHOW:   "TV Series",
    _ARCHETYPE_MUSICIAN:  "Musician / Band",
    _ARCHETYPE_ACTOR:     "Talent",
    _ARCHETYPE_ATHLETE:   "Athlete",
    _ARCHETYPE_TEAM:      "Sports Team",
    _ARCHETYPE_BRAND:     "Brand",
    _ARCHETYPE_FRANCHISE: "Franchise / IP",
    _ARCHETYPE_PERSONA:   "Persona / Concept Audience",
    _ARCHETYPE_GENERIC:   "Audience",
}

# Category-name keywords that map to a subject archetype. First match wins.
# Ordered from most specific to most generic — 'audience' is checked LAST
# because many category strings end with 'audience' generically.
_ARCHETYPE_KEYWORDS: list[tuple[str, str]] = [
    (_ARCHETYPE_MOVIE,     "movie"),
    (_ARCHETYPE_MOVIE,     "film"),
    (_ARCHETYPE_MOVIE,     "theatrical"),
    (_ARCHETYPE_MOVIE,     "feature"),
    (_ARCHETYPE_TV_SHOW,   "adult animation"),
    (_ARCHETYPE_TV_SHOW,   "tv series"),
    (_ARCHETYPE_TV_SHOW,   "streaming original"),
    (_ARCHETYPE_TV_SHOW,   "series"),
    (_ARCHETYPE_TV_SHOW,   "tv show"),
    (_ARCHETYPE_TV_SHOW,   "sitcom"),
    (_ARCHETYPE_MUSICIAN,  "musician"),
    (_ARCHETYPE_MUSICIAN,  "band"),
    (_ARCHETYPE_MUSICIAN,  "artist"),
    (_ARCHETYPE_MUSICIAN,  "rapper"),
    (_ARCHETYPE_MUSICIAN,  "dj"),
    (_ARCHETYPE_MUSICIAN,  "music"),
    (_ARCHETYPE_ACTOR,     "actor"),
    (_ARCHETYPE_ACTOR,     "actress"),
    (_ARCHETYPE_ACTOR,     "comedian"),
    (_ARCHETYPE_ACTOR,     "host"),
    (_ARCHETYPE_ACTOR,     "talent"),
    (_ARCHETYPE_ACTOR,     "creator"),
    (_ARCHETYPE_ACTOR,     "influencer"),
    (_ARCHETYPE_ATHLETE,   "athlete"),
    (_ARCHETYPE_ATHLETE,   "player"),
    (_ARCHETYPE_TEAM,      "team"),
    (_ARCHETYPE_TEAM,      "league"),
    (_ARCHETYPE_TEAM,      "franchise team"),
    (_ARCHETYPE_FRANCHISE, "franchise"),
    (_ARCHETYPE_FRANCHISE, "cinematic universe"),
    (_ARCHETYPE_FRANCHISE, "ip"),
    (_ARCHETYPE_PERSONA,   "persona"),
    (_ARCHETYPE_BRAND,     "brand"),
    (_ARCHETYPE_BRAND,     "cpg"),
    (_ARCHETYPE_BRAND,     "product"),
    (_ARCHETYPE_BRAND,     "retailer"),
    (_ARCHETYPE_BRAND,     "restaurant"),
    (_ARCHETYPE_BRAND,     "app"),
    (_ARCHETYPE_PERSONA,   "audience"),
]


def _detect_archetype(profile: dict) -> str:
    """Best-effort classification of the subject archetype.

    Uses ``brand_category`` first (that is the input dropdown), then name
    patterns. Everything is a soft heuristic; the LLM brief is asked to
    correct course if the detected archetype is wrong."""
    cat  = str((profile or {}).get("brand_category") or "").lower()
    name = str((profile or {}).get("name") or "").lower()
    hay = f"{cat} || {name}"
    for arch, kw in _ARCHETYPE_KEYWORDS:
        if kw in hay:
            return arch
    # Name-pattern fallbacks: title contains 'comedy', 'story', 'chronicles',
    # 'the movie' → likely a film.
    if any(t in name for t in ("comedy", "the movie", "chronicles",
                                "story", "trilogy", "saga")):
        return _ARCHETYPE_MOVIE
    return _ARCHETYPE_GENERIC


def _archetype_label(arch: str) -> str:
    return _ARCHETYPE_LABELS.get(arch, "Audience")


_ANALYST_SYSTEM_PROMPT = (
    "You are a Principal at a top-tier strategy consultancy (McKinsey "
    "Marketing & Sales practice, Bain Customer Strategy, BCG Consumer "
    "Insights) writing a boardroom-ready audience-profile deck for the "
    "CMO of a company whose GTM depends on reaching this audience. You "
    "are given a JSON of behavioral panel data for a subject and its US "
    "Gen Pop baseline. Your job is to answer three questions and nothing "
    "else:\n"
    "   1. WHO is this audience — what mindset, life stage, and cultural "
    "      posture defines them, not just their demographics.\n"
    "   2. WHY they matter — what economic behavior + attention posture "
    "      makes them worth chasing right now.\n"
    "   3. SO WHAT — where the CMO should spend, what channels to buy, "
    "      what partnerships to pursue, and where the whitespace is (a\n"
    "      category the audience over-indexes on but current marketing\n"
    "      isn't leveraging).\n\n"
    "You return STRUCTURED JSON (no prose wrapper) with the exact keys "
    "the deck builder requires.\n\n"
    "Analytical standards (non-negotiable):\n"
    "  * PYRAMID PRINCIPLE. Every narrative starts with the answer (the "
    "    'so what'), then supports it with 1-2 numeric proofs. Never lead "
    "    with a data dump.\n"
    "  * McKINSEY MECE. Your takeaways must be mutually exclusive and "
    "    collectively exhaustive across: WHO / WHERE they live / WHAT they "
    "    buy / HOW they consume media / WHERE the whitespace is / SO WHAT "
    "    the marketer should do next.\n"
    "  * NEVER descriptive. Every sentence must have a hypothesis or "
    "    implication. 'They watch TikTok' is banned. 'TikTok — not "
    "    Instagram — is the audience's primary discovery surface (70% "
    "    reach, 140 index), which means the buy should re-weight from IG "
    "    reels to TikTok Spark Ads' is what you write instead.\n"
    "  * KILL DESCRIPTIVE ADJECTIVES. Ban: 'passionate', 'vibrant', "
    "    'engaged', 'unique', 'niche', 'distinct'. Replace with an "
    "    observable behavior + a number.\n"
    "  * PRONOUNS. Never use gendered pronouns for the subject. Films, "
    "    shows, brands, teams, franchises, and personas are always 'the "
    "    audience' / 'these fans' / 'they'. Real people (talent, athlete, "
    "    musician) may be named — but their AUDIENCE is still 'they'.\n"
    "  * REACH vs INDEX — the single most important analytical distinction:\n"
    "      REACH (% of the audience that engages) governs SCALE — how many\n"
    "      people you can actually reach through this brand or channel.\n"
    "      INDEX (engagement rate vs Gen Pop, 100 = parity) governs\n"
    "      DIFFERENTIATION — how uniquely this audience over-consumes vs.\n"
    "      the average American.\n"
    "      Every brand row sits on a 2x2 of these two dimensions. The\n"
    "      analyst's job is to READ THE PAIR, never one alone:\n"
    "        - MASS + LIFT (reach >=20% AND index 130-500): actionable\n"
    "          sweet spot. Media buys land here. Lead every narrative here.\n"
    "        - SCALE (reach >=20% AND index 90-130): table stakes; the\n"
    "          audience uses this like everyone else. Not a differentiator.\n"
    "        - SIGNATURE (reach >=12% AND <20% AND index 150-400):\n"
    "          persona color and storytelling ONLY. These brands CANNOT\n"
    "          deliver scale — a 15% reach brand cannot reach 85% of the\n"
    "          audience regardless of how sharp the index. NEVER LEAD A\n"
    "          NARRATIVE with a SIGNATURE brand. Never claim it 'anchors'\n"
    "          a category. Never imply it drives the buy.\n"
    "        - NOISE (reach <12% at ANY index, OR any brand with index\n"
    "          >500, OR any brand with index >400 at reach <15%): low-N\n"
    "          or subject-adjacent panel artifact. NEVER mention.\n"
    "          Concrete banned examples the user has flagged:\n"
    "            - Killstar 9.9% at idx 576 — reach <12% AND idx >500. Omit.\n"
    "            - Urban Stems 7.9% at idx 399 — reach <12%. Omit.\n"
    "            - Marni 3.6% at idx 263 — reach <12%. Omit.\n"
    "          'Less than 1 in 10 of the audience does this' is too thin\n"
    "          to earn a mention in a persona narrative no matter how\n"
    "          sharp the index looks.\n"
    "      Correct phrasing template: 'Mass-brand shopping is anchored by\n"
    "      Nike (37%), Amazon (72%), and Levi (27%), with SIGNATURE\n"
    "      affinities toward Gymshark (15% reach, index 205) and Cotopaxi\n"
    "      (14% reach, index 174).'\n"
    "      Incorrect (and banned): 'Shoppers of Urban Stems (7.9%) and\n"
    "      Killstar (9.9%) drive this audience's commerce' — 91% of them\n"
    "      do NOT shop those brands.\n"
    "  * SIGNAL DISCIPLINE. Do not overclaim from freak-index outliers. A\n"
    "    brand at <12% reach is NEVER a signature — it is noise. A brand\n"
    "    at idx >500 is NEVER a real over-index — it is subject-adjacent\n"
    "    panel skew. Omit both entirely.\n"
    "  * MINIMUM CATEGORY BAR. Only nominate a category as worth its own "
    "    deep-dive slide if its top brand has >=15% reach AND at least "
    "    two brands have >=10% reach. Everything below that bar goes into "
    "    'signature_affinities' — a single roll-up page.\n"
    "  * GROUND EVERY CLAIM in a specific number from the data (reach %, "
    "    index, DMA share, projected M).\n"
    "  * NO CLICHÉS. Ban: 'taste graph', 'not a moment a habit', 'the "
    "    fandom IS the buying behavior', 'shows up', 'lives at the "
    "    intersection of', 'meet them where they are'.\n"
    "  * NO MARKDOWN. Every string is clean plain text: no *, no #, no "
    "    pipes, no bullet chars, no line breaks inside sentences.\n"
    "  * WORD BUDGETS ARE HARD. If a key says 35-45 words, that means "
    "    counted words, not a suggestion. Overrunning wrecks the layout.\n"
)


def _archetype_brief_prompt(archetype: str, subject: str) -> str:
    """Archetype-specific instructions layered onto the base system prompt.

    The pattern: state what the subject IS, then define what 'the audience'
    means for that subject type, then say what the marketing frame should
    be. Each archetype produces qualitatively different analysis."""
    if archetype == _ARCHETYPE_MOVIE:
        return (
            f"The subject '{subject}' is a FILM. The 'audience' is US adults "
            "with a digital touchpoint to this film in the trailing 12 "
            "months. Frame every narrative around HOW TO MARKET THIS FILM: "
            "trailer placement channels, PR + festival beats, streamer / "
            "theatrical partner fit, ticket-buying signals, and the "
            "whitespace where the audience over-indexes but the film's "
            "current marketing is not showing up. Recommend concrete "
            "media-plan moves. Refer to the film by title; refer to "
            "viewers/fans as 'the audience' or 'they'."
        )
    if archetype == _ARCHETYPE_TV_SHOW:
        return (
            f"The subject '{subject}' is a TV SERIES / STREAMING SHOW. "
            "Frame around HOW TO MARKET THIS SHOW: launch beats, streamer "
            "platform strategy, social + creator activation, cross-title "
            "lookalike vehicles (other shows the audience already watches "
            "at high index), and the whitespace categories worth activating "
            "in. Refer to the show by title; refer to viewers as 'the "
            "audience' or 'they'."
        )
    if archetype == _ARCHETYPE_MUSICIAN:
        return (
            f"The subject '{subject}' is a MUSIC ARTIST / BAND. Frame "
            "around tour DMAs, streaming platform mix, merch drops, brand "
            "partnership fit (sponsors that already index high with this "
            "fanbase), sync opportunities, and PR moments. Use the artist's "
            "name where relevant; refer to fans as 'the fanbase' / 'they'."
        )
    if archetype == _ARCHETYPE_ACTOR:
        return (
            f"The subject '{subject}' is a TALENT (actor / comedian / host / "
            "creator). Frame around endorsement fit, brand ambassadorships "
            "the audience already over-indexes on, lookalike streaming "
            "vehicles, red-carpet / PR moments. Use the talent's name; "
            "refer to fans as 'the audience' / 'they'."
        )
    if archetype == _ARCHETYPE_ATHLETE:
        return (
            f"The subject '{subject}' is an ATHLETE. Frame around "
            "endorsement fit (apparel/footwear, auto, QSR, financial "
            "services that over-index with the fanbase), hometown / "
            "team-market DMAs, and sponsor category whitespace. Refer to "
            "fans as 'the fanbase' / 'they'."
        )
    if archetype == _ARCHETYPE_TEAM:
        return (
            f"The subject '{subject}' is a SPORTS TEAM. Frame around "
            "ticket-buying, home DMA vs national reach, sponsor category "
            "fit, merch behavior, and cross-league fandom overlap. Refer "
            "to fans as 'the fanbase' / 'they'."
        )
    if archetype == _ARCHETYPE_BRAND:
        return (
            f"The subject '{subject}' is a BRAND / PRODUCT. Frame around "
            "growth: adjacent audiences the brand isn't reaching, cross-buy "
            "with peer categories, retail / DTC channel signal, sponsor / "
            "partner fit. Refer to customers as 'the audience' / 'they'."
        )
    if archetype == _ARCHETYPE_FRANCHISE:
        return (
            f"The subject '{subject}' is a FRANCHISE / IP. Frame around "
            "line extension, licensing, cross-media activation, and the "
            "shared-fandom whitespace across the audience's other affinities."
        )
    return (
        f"The subject '{subject}' is a defined audience segment. Frame "
        "around who they are, how they behave differently from Gen Pop, "
        "and where the whitespace is to reach them at scale."
    )


def _brief_facts(profile: dict) -> dict:
    """Compact facts JSON fed to the LLM. Keeps token count low and forces
    the model to reason from the data we actually captured (not hallucinate)."""
    def _top_n(cat_name: str, n: int = 10) -> list[dict]:
        beh = (profile.get("behavioral") or {})
        for k, v in beh.items():
            if str(k).strip().upper() == cat_name.upper():
                items = list(v or [])
                items.sort(key=lambda it: (-_num(it.get("index", 0)),
                                            -_num(it.get("pct", 0))))
                return [{
                    "brand": (it.get("name") or "").strip(),
                    "pct":   round(_num(it.get("pct", 0)), 1),
                    "index": int(_num(it.get("index", 0))),
                } for it in items[:n]]
        return []

    # Flatten every behavioral item so the analyst can see the whole surface
    flat: list[dict] = []
    for cat, items in (profile.get("behavioral") or {}).items():
        for it in (items or []):
            it2 = dict(it)
            it2["_cat"] = cat
            flat.append(it2)
    flat = [it for it in flat if _num(it.get("index", 0)) > 0
                              and _num(it.get("pct", 0)) > 0]
    flat.sort(key=lambda it: -_num(it.get("index", 0)))
    strongest = [{
        "brand":    (it.get("name") or "").strip(),
        "category": (it.get("_cat") or "").strip(),
        "pct":      round(_num(it.get("pct", 0)), 1),
        "index":    int(_num(it.get("index", 0))),
    } for it in flat[:30]]

    demos = profile.get("demographics") or {}
    demo_summary = {}
    for dim in ("gender", "age", "ethnicity", "income", "education"):
        d = demos.get(dim) or {}
        if d:
            demo_summary[dim] = {
                str(k): round(_num(v), 1) for k, v in d.items() if _num(v) > 0
            }

    # Available category names (helpful for the model to know coverage)
    cats_present = sorted({k for k, v in (profile.get("behavioral") or {}).items()
                            if v})

    return {
        "subject":            profile.get("name") or "",
        "brand_category":     profile.get("brand_category") or "",
        "sample_size":        profile.get("sample_size") or 0,
        "projected_us":       profile.get("projected_us") or 0,
        "date_range":         profile.get("date_range") or "",
        "demographics":       demo_summary,
        "top_dmas":           _top_dmas_ranked(profile, n=8),
        "categories_covered": cats_present,
        "strongest_signals":  strongest,
        "top_mpb":            (_top_n("MPB", 15)
                                or _top_n("MOST PURCHASED BRANDS", 15)),
        "top_social":         _top_n("SOCIAL MEDIA", 8),
        "top_streaming":      (_top_n("STREAMING/PLATFORM", 8)
                                or _top_n("STREAMING VIDEO", 8)),
    }


_ANALYST_JSON_SPEC = (
    "Return JSON with EXACTLY these keys (no extras, no missing):\n"
    "{\n"
    '  "cover_tagline": "One sentence that tells the CMO who this audience IS in life-stage / mindset / cultural-posture terms (35-45 words). Do NOT lean on a single freak-index brand. Do NOT use adjectives like passionate / vibrant / engaged.",\n'
    '  "identity_headline": "4-6 word portrait of WHO they are — life stage + posture, not just demos. Examples: \'Prime-age urban trend-formers.\', \'Suburban parents with disposable income.\', \'Downtown creatives, digitally-native.\' No pronouns.",\n'
    '  "identity_why": "One sentence (35-50 words) on why this composition matters commercially — what the age/income/geography shape UNLOCKS for the marketer.",\n'
    '  "portrait": {\n'
    '     "life_stage_headline":     "3-6 word life-stage label. Examples: \'Millennial/Gen-Z urban renters.\', \'Suburban Gen-X parents.\', \'Empty-nester boomers.\'",\n'
    '     "life_stage_body":         "3-4 sentences (55-80 words) painting their life stage — how they spend their days, what their household looks like, career phase, income tier. Ground in the age/income/parental data.",\n'
    '     "mindset_headline":        "3-6 word psychographic label. Examples: \'Culture-first early adopters.\', \'Value-conscious pragmatists.\', \'Status-signaling professionals.\'",\n'
    '     "mindset_body":            "3-4 sentences (55-80 words) on their mindset / worldview / consumption posture. Infer from the shape of their brand and category signals. Concrete, not abstract.",\n'
    '     "cultural_posture_headline":"3-6 word cultural label. Examples: \'Terminally-online scroll-natives.\', \'Live-events + festival regulars.\', \'Prestige-media watchers.\'",\n'
    '     "cultural_posture_body":   "3-4 sentences (55-80 words) on how they consume culture — what platforms, what content, what social postures. Cite 2-3 specific brands or platforms with reach %.",\n'
    '     "differentiator_headline": "3-6 word label for their sharpest edge vs Gen Pop. Examples: \'Under-35 and multicultural.\', \'Higher income + urban core.\', \'Streaming-first, cable-cord-cut.\'",\n'
    '     "differentiator_body":     "3-4 sentences (55-80 words) on the ONE dimension where they differ most from Gen Pop, cited with the index number, and what that means commercially."\n'
    '  },\n'
    '  "life_stage": {\n'
    '     "headline": "5-9 word household/life-stage narrative (no period).",\n'
    '     "body":     "80-110 word paragraph on household composition — parental status, relationship, education, income tier context — grounded in the numbers from the data. Focus on what this composition IMPLIES for daypart, channel selection, and message tone."\n'
    '  },\n'
    '  "differ_from_genpop": [\n'
    '     { "dimension": "Age|Ethnicity|Income|Education|Gender|Parental Status|Relationship",\n'
    '       "bucket":    "specific bucket label (e.g. \'18-24\', \'Asian\', \'$150K+\')",\n'
    '       "aud_pct":   0.0,\n'
    '       "gp_pct":    0.0,\n'
    '       "index":     0,\n'
    '       "note":      "8-15 word analyst tag on why this delta matters (e.g. \'The Asian over-index anchors coastal urban reach.\')" },\n'
    '     ... EXACTLY 6-8 entries. Pick the SHARPEST deltas from ALL demo dimensions in the data (not just age). Diverse dimensions. Sort by absolute index deviation from 100.\n'
    '  ],\n'
    '  "commerce_persona": {\n'
    '     "headline": "5-9 word shopping-personality label (no period). Examples: \'Amazon-first Instacart-heavy value shoppers.\', \'DTC-forward luxury-adjacent early adopters.\', \'Big-box + Target home-focused.\'",\n'
    '     "body":     "90-130 word commerce portrait synthesizing MPB + RETAIL + COSMETICS + APPAREL + PERSONAL CARE. STRICT signal-class rules: (1) LEAD with the 3-4 highest-REACH brands (pct >=20%), format as BRAND (37% reach); (2) ADD signature affinities in a second sentence starting with \'with signature affinities toward\', format as BRAND (15% reach, index 205); (3) NEVER mention a brand with reach <12% — such brands are NOISE regardless of how sharp the index looks; (4) NEVER mention a brand with index >400 — that is subject-adjacent panel skew, not real signal. Example: \'Anchored by Nike (37%), Amazon (72%), and Levi (27%), with signature affinities toward Gymshark (15% reach, index 205) and Cotopaxi (14% reach, index 174).\' Banned: \'signature affinities toward Killstar (9.9%, idx 576) and Urban Stems (7.9%, idx 399)\' — both are NOISE class, do not cite."\n'
    '  },\n'
    '  "cultural_interests": {\n'
    '     "headline": "5-9 word cultural-interest label (no period). Examples: \'Live-music + festival regulars.\', \'Gaming-native prestige-TV watchers.\', \'Sports-first spectators, esports-curious.\'",\n'
    '     "body":     "90-130 word integrated view of cultural time — pulls MUSIC / GAMES / FESTIVALS / TICKETING / SPORTS TEAM / STREAMING signals. STRICT signal-class rules: (1) LEAD with the highest-reach cultural surfaces (pct >=15%), e.g. \'TikTok (70%) and Spotify (52%)\'; (2) ADD signature interests with pct >=12% AND idx 150-400 only, format as BRAND (15% reach, index 205); (3) NEVER cite a brand with reach <12% or index >400 — those are NOISE. If the audience has no signature-class cultural brand, just describe the mass surfaces and stop; do not manufacture signature colour from noise."\n'
    '  },\n'
    '  "exec_summary": {\n'
    '     "who_headline":     "2-4 word label (e.g. \'The audience.\')",\n'
    '     "who_body":         "2-3 tight sentences (35-50 words) describing WHO they are grounded in the sharpest demographic index (age/ethnicity/income). Lead with the answer, then the number.",\n'
    '     "where_headline":   "2-4 word label (e.g. \'Where they live.\')",\n'
    '     "where_body":       "2-3 tight sentences (35-50 words) describing WHERE they concentrate — DMA + platform. Name the top 2 DMAs and top 1-2 social/streaming surfaces with reach %.",\n'
    '     "so_what_headline": "2-4 word label (e.g. \'So what.\')",\n'
    '     "so_what_body":     "2-3 tight sentences (35-50 words) with the marketer action. Name a channel, a partner category, and a whitespace. Be prescriptive."\n'
    '  },\n'
    '  "day_in_life": {\n'
    '     "morning":   "35-55 word morning slice — waking up, commute, first-touch platforms. Name 2-3 brands with reach %. E.g. \'Wakes to a TikTok (70%) scroll before coffee; commute is a Spotify (52%) queue or a Reddit (38%) thread on r/movies. Breakfast is grab-and-go from Amazon Fresh or the corner deli.\'",\n'
    '     "midday":    "35-55 word mid-day slice — work rhythm, lunch, shopping. Name 2-3 brands with reach %. E.g. \'Slack + Zoom (implied) fill the workday; lunch is a DoorDash (45%) or a Chipotle grab. Amazon (72%) refills happen mid-day between meetings.\'",\n'
    '     "evening":   "35-55 word evening slice — decompression, entertainment, socialising. Name 2-3 brands with reach %. E.g. \'Evening is Netflix (78%) or a Twitch (18%) stream, with Instagram (72%) DMs running in parallel. Weekend nights add a Ticketmaster (14%) purchase for an actual live show.\'"\n'
    '  },\n'
    '  "signature_affinities": [\n'
    '     { "brand": "BRAND", "category": "CATEGORY", "pct": 0.0, "index": 0, "note": "6-12 word analyst tag explaining why this signal matters" },\n'
    '     ... STRICT: pick 8-10 brands that BEST characterize this audience under signal-class discipline. HARD gates: pct >=12% AND index between 130 and 400 (inclusive). NEVER include a brand with reach <12% OR index >400 — those are NOISE. NEVER include a mass-reach brand (pct >=20%) — those belong on the category deep-dives and commerce_persona, NOT here. NEVER include the SUBJECT itself. Diverse categories preferred.\n'
    '  ],\n'
    '  "category_why": {\n'
    '     "CATEGORY_NAME_UPPERCASE": "One sentence (30-45 words) on why this category concentrates for this audience AND what a marketer should DO about it — brand ambassador fit, sponsor whitespace, or shelf-placement play.",\n'
    '     ...  (ONLY include categories where top brand has >=15% reach and >=2 brands >=10% reach. Skip low-reach categories entirely — they roll up into signature_affinities.)\n'
    '  },\n'
    '  "media_plan": [\n'
    '     { "channel": "TikTok", "pct": 25, "rationale": "12-20 word reason grounded in the data (reach % + index)." },\n'
    '     ... (4-6 channels. pcts must sum to 100. Channels must be paid-media channels the CMO can actually buy: TikTok, Instagram Reels, YouTube Pre-roll, Connected TV / Hulu, Snap, Discord Sponsorships, Reddit Promoted, Podcast Host-Reads, Programmatic Display, Meta Feed, Spotify Audio, Twitch, X, LinkedIn, Search / Google Ads, OOH. Weight allocation to the surfaces the audience actually over-indexes on.)\n'
    '  ],\n'
    '  "whitespace_headline": "5-9 word headline naming the SPECIFIC unmet-need or unactivated surface (e.g. \'Gaming sponsorships are untapped.\', \'The Discord audience is under-bought.\'). No period.",\n'
    '  "whitespace_body":     "A 60-90 word analyst paragraph explaining WHY that surface is whitespace for this audience (cite the reach/index) and WHAT the CMO should test first. Prescriptive, not descriptive.",\n'
    '  "takeaways": [\n'
    '     { "title": "3-5 word punchy title. Actions, not descriptions.",\n'
    '       "body": "2 sentences (25-40 words) that ground the action in a specific number from the data." },\n'
    '     ... EXACTLY 6 entries, MECE across: (1) WHO — the demographic core, (2) WHERE — geo concentration, (3) WHAT — dominant category / commercial signal, (4) HOW — primary media surface, (5) WHITESPACE — the untapped surface, (6) ACTIVATION — the next test to run.\n'
    '  ],\n'
    '  "final_line_a": "Final-insight display line 1 (max 8 words, no period). The one sentence you want the CMO to remember.",\n'
    '  "final_line_b": "Final-insight display line 2 (max 8 words, ends with period). The implication or action.",\n'
    '  "final_body":   "55-75 word closer. Synthesize the audience arc across who / where / what / so-what. Do NOT re-list brands.",\n'
    '  "sections": ["cover", "exec_summary", "scale", "identity", "index_story", "day_in_life", "geography", "signature_affinities", "category_deep_dives", "media_and_social", "media_plan", "whitespace", "takeaways", "final", "methodology", "about"]\n'
    '     // Ordered list of the sections you recommend for THIS profile. Include only what you can back with data. If DMAs are thin, drop "geography". If category_why is empty, drop "category_deep_dives". If no social/streaming data, drop "media_and_social". "cover", "exec_summary", "identity", "takeaways", "final", "methodology", "about" are always included.\n'
    "}\n"
)


def _get_analyst_brief(profile: dict, archetype: str) -> dict:
    """Call OpenAI for a structured consulting-analyst brief. Returns {} on
    any failure (missing key, network error, malformed JSON). Callers MUST
    handle the empty case by falling back to templated copy."""
    subject = str(profile.get("name") or "the audience")
    facts   = _brief_facts(profile)
    prompt_tail = _archetype_brief_prompt(archetype, subject)

    user_prompt = (
        f"SUBJECT: {subject}\n"
        f"ARCHETYPE: {archetype} ({_archetype_label(archetype)})\n"
        f"BRAND CATEGORY INPUT: {profile.get('brand_category') or 'unknown'}\n\n"
        f"ARCHETYPE INSTRUCTIONS:\n{prompt_tail}\n\n"
        f"DATA (panel behavioral profile):\n"
        f"{json.dumps(facts, ensure_ascii=False)}\n\n"
        + _ANALYST_JSON_SPEC
    )

    try:
        api_key = os.environ.get("OPENAI_API_KEY")
        if not api_key:
            print("[profile_iq_deck_builder] no OPENAI_API_KEY — "
                  "using fallback templated narrative")
            return {}
        from openai import OpenAI
        client = OpenAI(api_key=api_key, timeout=90.0)
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            temperature=0.4,
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": _ANALYST_SYSTEM_PROMPT},
                {"role": "user",   "content": user_prompt},
            ],
        )
        content = (resp.choices[0].message.content or "{}")
        parsed = json.loads(content)
        if not isinstance(parsed, dict):
            return {}
        return parsed
    except Exception as exc:  # pragma: no cover — logged, not raised
        print(f"[profile_iq_deck_builder] analyst brief LLM call failed: {exc}")
        return {}


def _brief_get(profile: dict, *keys, default=""):
    """Safe getter for nested brief keys. Returns ``default`` if any hop is
    missing or the final value is empty (empty string / [] / {})."""
    v = (profile or {}).get("_brief") or {}
    for k in keys:
        if not isinstance(v, dict):
            return default
        v = v.get(k)
        if v is None:
            return default
    return v if v not in ("", None, [], {}) else default


# Belt-and-suspenders: strip gendered pronouns from any LLM or template
# output. The LLM is told not to use them but occasionally does; older
# templates in this module used them too. Ordered longest-match-first.
_PRONOUN_REPLACEMENTS: list[tuple[str, str]] = [
    (r"\bwho she reaches\b",      "who the audience reaches"),
    (r"\bwhat she signals\b",     "how the audience signals"),
    (r"\bwhere she watches\b",    "where the audience watches"),
    (r"\bwhere she buys\b",       "where the audience buys"),
    (r"\bhow she eats\b",         "how the audience eats"),
    (r"\bhow she travels\b",      "how the audience travels"),
    (r"\bhow she dresses\b",      "how the audience dresses"),
    (r"\bwhat she drives\b",      "what the audience drives"),
    (r"\bwhen she'?s out\b",      "when the audience is out"),
    (r"\bwho she\b",              "who the audience"),
    (r"\bwhere she\b",            "where the audience"),
    (r"\bhow she\b",              "how the audience"),
    (r"\bwhen she\b",             "when the audience"),
    (r"\bshe watches\b",          "the audience watches"),
    (r"\bshe dresses\b",          "the audience dresses"),
    (r"\bshe travels\b",          "the audience travels"),
    (r"\bshe drives\b",           "the audience drives"),
    (r"\bshe eats\b",             "the audience eats"),
    (r"\bshe reaches\b",          "the audience reaches"),
    (r"\bshe buys\b",             "the audience buys"),
    (r"\bshe'?s\b",               "the audience is"),
    (r"\bshe'?d\b",               "the audience would"),
    (r"\bshe'?ll\b",              "the audience will"),
    (r"\bshe\b",                  "the audience"),
    (r"\bwho he\b",               "who the audience"),
    (r"\bwhere he\b",             "where the audience"),
    (r"\bhow he\b",               "how the audience"),
    (r"\bhe'?s\b",                "the audience is"),
    (r"\bhe\b",                   "the audience"),
    (r"\bher\b",                  "their"),
    (r"\bhis\b",                  "their"),
]


def _purge_pronouns(text: str, subject: str = "") -> str:
    """Strip gendered pronouns from any string. Safe to call on empty input.
    Applied to every LLM-returned narrative string before it hits a slide."""
    if not text:
        return text
    out = str(text)
    for pat, sub in _PRONOUN_REPLACEMENTS:
        out = re.sub(pat, sub, out, flags=re.IGNORECASE)
    # Collapse accidental double replacements (e.g. "the audience the audience")
    out = re.sub(r"\bthe audience the audience\b", "the audience", out,
                 flags=re.IGNORECASE)
    return out


# =============================================================================
#  Post-hoc page numbering
# =============================================================================


def _renumber_pages(prs) -> None:
    """Walk the finished deck and rewrite the trailing page-number text box
    to the true 1-based slide index. Every slide builder adds a numbered
    footer via ``_footer(slide, page_num, subject)``; the page_num arg the
    builder passes is a placeholder — this pass overwrites it with reality
    so we do not have to thread a counter through the build order."""
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        best = None
        best_x = -1.0
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            if sh.left is None or sh.top is None:
                continue
            top_in  = Emu(sh.top).inches
            left_in = Emu(sh.left).inches
            width_in = Emu(sh.width).inches if sh.width is not None else 0
            # Page number box is the rightmost small text box at y≈7.10,
            # width ~0.9". Ignore full-width footer strings.
            if 6.9 <= top_in <= 7.3 and left_in >= 11.5 and width_in <= 1.5:
                if left_in > best_x:
                    best = sh
                    best_x = left_in
        if best is None:
            continue
        tf = best.text_frame
        if not tf.paragraphs or not tf.paragraphs[0].runs:
            # Rare edge — reset text via the frame directly, losing font
            # attrs. Should not happen for our _footer output.
            best.text_frame.text = f"{i} / {total}"
            continue
        # Rewrite the first run to preserve font/color/size/spacing
        tf.paragraphs[0].runs[0].text = f"{i} / {total}"
        # Wipe any additional runs so the rewrite isn't concatenated onto
        # placeholder text.
        for extra in tf.paragraphs[0].runs[1:]:
            extra.text = ""


# =============================================================================
#  Presentation shell / slide primitives
# =============================================================================


def _new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def _blank_slide(prs: Presentation, bg=C_DARK):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # completely blank layout
    _fill_bg(slide, bg)
    return slide


def _fill_bg(slide, color: RGBColor) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))
    box.line.fill.background()
    box.fill.solid()
    box.fill.fore_color.rgb = color


def _text(
    slide,
    left, top, width, height,
    text: str,
    *,
    size: float = 12,
    bold: bool = False,
    color: RGBColor = C_CREAM,
    align: int = PP_ALIGN.LEFT,
    anchor: int = MSO_ANCHOR.TOP,
    font: str = FONT_MAIN,
    line_spacing: Optional[float] = None,
    letter_spacing: Optional[float] = None,
):
    """Draw a single-run text box (unless text has \\n — then multi-run).

    Positions are inches. Returns the shape so callers can tweak further.
    """
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tb.line.fill.background()

    lines = str(text).split("\n") if text is not None else [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bool(bold)
        r.font.color.rgb = color
        # letter_spacing in python-pptx isn't a first-class attribute; skip
        # for now (LISA deck uses default tracking anyway).
    return tb


# Empirical per-point char widths for the FONT_MAIN family (Aptos / Calibri
# fall in the 0.55-0.62 em range for average sentence text at bold). Used
# by _fit_display_size to pick a font that keeps a headline on one line.
_CHAR_W_PER_PT_BOLD    = 0.0072
_CHAR_W_PER_PT_REGULAR = 0.0060


def _fit_display_size(text: str, box_w_in: float, base_pt: float,
                       min_pt: float = 24.0, bold: bool = True,
                       safety: float = 0.94) -> float:
    """Return the largest font size <= base_pt whose rendered width for
    ``text`` fits inside ``box_w_in``. Prevents long subject names from
    wrapping to a second visual line and colliding with a stacked line
    below (fixes cover / final-insight / one-insight two-line displays).

    ``safety`` is a multiplier (0.94 = 6% slack) so we don't sit right on
    the edge — punctuation and emoji-adjacent glyphs vary by font.
    """
    if not text:
        return base_pt
    n = len(str(text))
    per_pt = _CHAR_W_PER_PT_BOLD if bold else _CHAR_W_PER_PT_REGULAR
    max_w  = box_w_in * safety
    for pt in (base_pt, base_pt * 0.90, base_pt * 0.80,
               base_pt * 0.70, base_pt * 0.60, min_pt):
        if n * pt * per_pt <= max_w:
            return max(min_pt, pt)
    return min_pt


def _rect(slide, left, top, width, height, color: RGBColor, *,
          radius: Optional[float] = None, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    if radius is not None:
        # python-pptx uses adjustment 0.0-1.0 of half the shorter side
        try:
            shape.adjustments[0] = float(radius)
        except Exception:
            pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def _hairline(slide, left, top, width, color: RGBColor = C_STROKE):
    line = slide.shapes.add_connector(1, Inches(left), Inches(top),
                                      Inches(left + width), Inches(top))
    line.line.color.rgb = color
    line.line.width = Emu(9525)  # ~0.75pt
    return line


def _footer(slide, page_num: int, subject_label: str = ""):
    _text(slide, 0.62, 7.10, 6.0, 0.28,
          f"{subject_label.upper()} · AUDIENCE PROFILE" if subject_label else "AUDIENCE PROFILE",
          size=9, color=C_MUTED2, letter_spacing=0.05)
    _text(slide, 6.5, 7.10, 5.0, 0.28,
          "CROSSWALK  ·  PROFILE IQ",
          size=9, color=C_MUTED2, align=PP_ALIGN.CENTER, letter_spacing=0.05)
    _text(slide, 12.0, 7.10, 0.9, 0.28,
          str(page_num), size=9, color=C_MUTED2, align=PP_ALIGN.RIGHT)


def _section_eyebrow(slide, top: float, subject_label: str, section: str):
    """Small tracked eyebrow like 'THE LISA AUDIENCE  ·  IDENTITY'."""
    label = (subject_label or "Audience").strip().upper()
    _text(slide, 0.62, top, 12.0, 0.30,
          f"THE {label} AUDIENCE  ·  {section.upper()}",
          size=10, color=C_MUTED, letter_spacing=0.06)


def _title_block(slide, top: float, headline_a: str, headline_b: str = ""):
    """Two-line big-display title. Font auto-shrinks so long headlines
    from the analyst brief never wrap and collide with the second line.
    Line 2 is placed dynamically based on the actual font used."""
    pt = 44
    if headline_a:
        pt = min(pt, _fit_display_size(headline_a, 11.5, base_pt=44,
                                        min_pt=24, bold=True))
    if headline_b:
        pt = min(pt, _fit_display_size(headline_b, 11.5, base_pt=44,
                                        min_pt=24, bold=True))
    line_h = pt / 72.0 * 1.05
    box_h = line_h + 0.10
    _text(slide, 0.62, top, 11.5, box_h, headline_a,
          size=pt, bold=True, color=C_CREAM, line_spacing=1.05)
    if headline_b:
        _text(slide, 0.62, top + line_h + 0.12, 11.5, box_h, headline_b,
              size=pt, bold=True, color=C_CREAM, line_spacing=1.05)


# =============================================================================
#  Image fetch / cover art
# =============================================================================


_S3_BUCKET_NAME = "dashboard-inputs"
_S3_REGION = "us-east-2"


def _fetch_via_s3(s3_key: str) -> Optional[bytes]:
    """Fetch an S3 object using the app's IAM credentials.

    The ``dashboard-inputs`` bucket is NOT publicly readable — direct
    anonymous fetches to the public https URL return 403. Every Profile IQ
    profile image lives in this bucket (under ``profile-images/...``), so
    we must go through boto3 with the same credentials the Flask app already
    uses. Imported lazily so environments without boto3 still let the deck
    build (falls back to placeholder cover for those images)."""
    try:
        import boto3  # deferred to avoid hard dep at module import time
        from botocore.config import Config as _BotoConfig
    except Exception:
        return None
    try:
        client = boto3.client(
            "s3",
            region_name=_S3_REGION,
            config=_BotoConfig(retries={"max_attempts": 2, "mode": "standard"}),
        )
        resp = client.get_object(Bucket=_S3_BUCKET_NAME, Key=s3_key)
        data = resp["Body"].read()
        if not data or len(data) < 200:
            return None
        return data
    except Exception as exc:  # pragma: no cover — logged, not raised
        print(f"[profile_iq_deck_builder] boto3 fetch failed for {s3_key!r}: {exc}")
        return None


def _fetch_image_bytes(image_url: str) -> Optional[bytes]:
    """Fetch a profile image URL. Returns None on any failure.

    Handles three shapes of URL:
      1. ``/api/profile-image-file/<key>``  → private S3 (boto3 path)
      2. ``https://dashboard-inputs.s3...``  → private S3 (boto3 path)
      3. any other http/https URL          → urllib (public web assets:
         IMDB posters, Wikipedia hero shots, admin-provided direct URLs)
    """
    if not image_url:
        return None

    # Path 1: Flask same-origin profile-image proxy → resolve S3 key
    if image_url.startswith("/api/profile-image-file/"):
        s3_key = image_url[len("/api/profile-image-file/"):]
        return _fetch_via_s3(s3_key)

    # Path 2: direct dashboard-inputs S3 URL (private bucket, needs IAM)
    for prefix in (
        f"https://{_S3_BUCKET_NAME}.s3.{_S3_REGION}.amazonaws.com/",
        f"https://{_S3_BUCKET_NAME}.s3.amazonaws.com/",
        f"https://s3.{_S3_REGION}.amazonaws.com/{_S3_BUCKET_NAME}/",
    ):
        if image_url.startswith(prefix):
            s3_key = urllib.parse.unquote(image_url[len(prefix):])
            return _fetch_via_s3(s3_key)

    # Path 3: any other URL → try public web fetch
    try:
        req = urllib.request.Request(image_url, headers={
            "User-Agent": "CrosswalkProfileIQDeck/1.0 (+jenna@crosswalknyc.com)",
            "Accept": "image/*,*/*;q=0.8",
        })
        with urllib.request.urlopen(req, timeout=8) as resp:
            data = resp.read()
        if not data or len(data) < 200:
            return None
        return data
    except Exception as exc:  # pragma: no cover — logged, not raised
        print(f"[profile_iq_deck_builder] urllib fetch failed for {image_url!r}: {exc}")
        return None


def _placeholder_cover(slide, subject_name: str):
    """Ambient dark-gradient cover when no image is available."""
    # Big diagonal gradient effect using two overlapping color blocks + a
    # subtle tinted right column so the composition doesn't feel flat.
    _rect(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, C_DARK)
    _rect(slide, 0, 0, 6.5, SLIDE_H_IN, RGBColor(0x14, 0x22, 0x24))
    _rect(slide, SLIDE_W_IN - 3.2, 0, 3.2, SLIDE_H_IN, RGBColor(0x08, 0x0F, 0x11))
    # Subject initial as a decorative graphic
    initial = (subject_name.strip() or "A")[0].upper()
    _text(slide, SLIDE_W_IN - 3.0, 1.4, 2.4, 4.0, initial,
          size=260, bold=True, color=RGBColor(0x1B, 0x27, 0x2A),
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _cover_image(slide, image_bytes: Optional[bytes], subject_name: str):
    """Full-bleed cover imagery. Falls back to a placeholder gradient."""
    if not image_bytes:
        _placeholder_cover(slide, subject_name)
        return
    try:
        _place_prepared_image(slide, image_bytes, x=0, y=0,
                              w=8.4, h=SLIDE_H_IN, darken=0.72)
    except Exception:
        _placeholder_cover(slide, subject_name)
        return
    # Right-column accent strip (LISA has a dark solid on the right)
    _rect(slide, 8.4, 0, SLIDE_W_IN - 8.4, SLIDE_H_IN, C_DARK)
    # Bottom scrim so the tagline row has consistent contrast
    _rect(slide, 0, 4.9, SLIDE_W_IN, 2.6, C_DARK)


def _place_prepared_image(slide, image_bytes: bytes, *, x, y, w, h,
                          darken: float = 1.0) -> None:
    """Center-crop + resize + optional darken, then place at (x, y) size (w, h).

    ``darken`` is a Pillow ``ImageEnhance.Brightness`` factor (1.0 = untouched,
    0.7 = darker). Kept in one helper so every slide that uses imagery has
    consistent tone. Raises on failure so callers can fall back."""
    from PIL import Image, ImageEnhance
    im = Image.open(io.BytesIO(image_bytes)).convert("RGB")
    target_w_px = max(int(w * 220), 320)
    target_h_px = max(int(h * 220), 320)
    src_w, src_h = im.size
    src_ratio = src_w / src_h
    tgt_ratio = target_w_px / target_h_px
    if src_ratio > tgt_ratio:
        new_w = int(src_h * tgt_ratio)
        offset = (src_w - new_w) // 2
        im = im.crop((offset, 0, offset + new_w, src_h))
    elif src_ratio < tgt_ratio:
        new_h = int(src_w / tgt_ratio)
        # Bias slightly toward top so faces don't lose foreheads (editorial
        # crop heuristic).
        offset = max(0, (src_h - new_h) // 3)
        im = im.crop((0, offset, src_w, offset + new_h))
    im = im.resize((target_w_px, target_h_px), Image.LANCZOS)
    if darken and darken != 1.0:
        im = ImageEnhance.Brightness(im).enhance(darken)
    buf = io.BytesIO()
    im.save(buf, format="JPEG", quality=86, optimize=True)
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(x), Inches(y),
                             width=Inches(w), height=Inches(h))


def _place_side_image(slide, image_bytes: Optional[bytes], *, x, y, w, h,
                      darken: float = 1.0,
                      subject_name: str = "") -> None:
    """Wraps ``_place_prepared_image`` with a graceful fallback. Used by
    overview / persona-category / case-study / final-insight slides that
    want an image in a bounded region (not full-bleed)."""
    if image_bytes:
        try:
            _place_prepared_image(slide, image_bytes,
                                  x=x, y=y, w=w, h=h, darken=darken)
            return
        except Exception as exc:
            print(f"[profile_iq_deck_builder] side image failed: {exc}")
    # Fallback: dark rect with subject initial so slot doesn't look broken
    _rect(slide, x, y, w, h, RGBColor(0x12, 0x1D, 0x1F))
    initial = ((subject_name or "").strip() or "A")[0].upper()
    _text(slide, x, y + (h / 2) - 1.2, w, 2.4, initial,
          size=int(min(160, h * 30)), bold=True,
          color=RGBColor(0x22, 0x30, 0x34),
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# =============================================================================
#  Slide builders  ·  SINGLE PROFILE
# =============================================================================


def _slide_cover(prs: Presentation, p: dict):
    """Cover slide — LISA layout: full-bleed image, big title, tagline.

    Uses the pre-fetched ``p['_image_bytes']`` so we don't re-hit S3."""
    slide = _blank_slide(prs, bg=C_DARK)
    img_bytes = p.get("_image_bytes")
    _cover_image(slide, img_bytes, p["name"])

    year_label = _year_label(p.get("date_range")) or datetime.now().strftime("%Y")
    _text(slide, 0.62, 0.62, 10.0, 0.30,
          f"A BEHAVIORAL AUDIENCE PROFILE   \u00b7   {year_label}",
          size=10, color=C_MUTED, letter_spacing=0.06)

    subject = p["name"]
    # LISA-style two-line title: "The {Subject}\naudience." Auto-shrink
    # font when the subject is long enough to wrap at 54pt (would collide
    # with the "audience." line at y=3.95). Both lines use the same size
    # so the layout reads as one display block.
    line1 = f"The {subject}"
    title_pt = _fit_display_size(line1, box_w_in=8.60, base_pt=54,
                                  min_pt=28, bold=True)
    # Line height in inches at line_spacing 1.02.
    line_h = title_pt / 72.0 * 1.02
    y1 = 3.08
    y2 = y1 + line_h + 0.10  # 0.10" leading between lines
    _text(slide, 0.58, y1, 8.60, line_h + 0.05, line1,
          size=title_pt, bold=True, color=C_CREAM, line_spacing=1.02)
    _text(slide, 0.58, y2, 8.60, line_h + 0.05, "audience.",
          size=title_pt, bold=True, color=C_CREAM, line_spacing=1.02)

    tagline = _generate_tagline(p)
    _text(slide, 0.62, 5.02, 5.4, 1.60, tagline,
          size=12, color=C_MUTED, line_spacing=1.4)

    _text(slide, 0.62, 7.00, 4.0, 0.28, "CONFIDENTIAL",
          size=9, color=C_MUTED2, letter_spacing=0.06)
    _text(slide, 8.11, 7.00, 4.60, 0.28, "CROSSWALK  PROFILE IQ",
          size=9, color=C_MUTED2, align=PP_ALIGN.RIGHT, letter_spacing=0.06)


def _slide_overview(prs: Presentation, p: dict):
    """LISA overview layout: left-column narrative + numbered rows,
    right-column full-height hero image (5.28" wide)."""
    slide = _blank_slide(prs)

    # Right-column full-height image (LISA "What's inside" pattern).
    # Falls back to a dark rect if no image so layout still balances.
    _place_side_image(slide, p.get("_image_bytes"), x=8.05, y=0.0,
                      w=SLIDE_W_IN - 8.05, h=SLIDE_H_IN,
                      darken=0.90, subject_name=p["name"])

    _section_eyebrow(slide, 0.52, p["name"], "OVERVIEW")
    _text(slide, 0.62, 0.82, 8.0, 1.10, "What's inside.",
          size=30, bold=True, color=C_CREAM, line_spacing=1.05)

    body = (
        f"Over twelve months, Crosswalk observed how {p['name']}'s audience "
        f"behaves — what they watch, scroll, wear, and buy. Drawn from "
        f"{_fmt_int(p['sample_size']) if p['sample_size'] else '700,000+'} "
        f"zero-party panelists and projected across "
        f"{_fmt_int(p['projected_us']) if p['projected_us'] else '20M+'} "
        f"U.S. digital adults."
    )
    _text(slide, 0.62, 2.05, 6.85, 1.70, body,
          size=13.5, color=C_MUTED, line_spacing=1.45)

    # LISA-style numbered rows: 01 magenta numeral, bold row title,
    # muted description to the right. Hairline separators between rows.
    chapters = _pick_overview_chapters(p)
    row_top = 3.90
    row_h = 0.95
    for i, (num, title, subtitle) in enumerate(chapters):
        y = row_top + i * row_h
        _text(slide, 0.62, y, 0.80, 0.50, num,
              size=20, bold=True, color=C_MAGENTA)
        _text(slide, 1.57, y + 0.02, 2.70, 0.50, title,
              size=13.5, bold=True, color=C_CREAM)
        _text(slide, 4.37, y + 0.02, 3.10, 0.85, subtitle,
              size=10.5, color=C_MUTED, line_spacing=1.4)
        if i < len(chapters) - 1:
            _hairline(slide, 0.62, y + row_h - 0.08, 6.70, C_STROKE)

    _footer(slide, 2, p["name"])


def _slide_scale(prs: Presentation, p: dict):
    """'By the numbers' slide — display title ABOVE the stat row (was
    previously below, which read backwards). 4 dynamic stat callouts pull
    the strongest demographic + audience-scale signals so the slide is
    meaningful even for MPB-only profiles."""
    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], "BY THE NUMBERS")

    # Big display title AT THE TOP so the slide reads top-down
    _text(slide, 0.62, 0.85, 12.0, 1.15, "By the numbers.",
          size=44, bold=True, color=C_CREAM, line_spacing=1.02)

    proj   = p["projected_us"] or 0
    sample = p["sample_size"] or 0
    arch = p.get("_archetype") or _ARCHETYPE_GENERIC
    if arch == _ARCHETYPE_MOVIE:
        subhead = (f"US adults with a digital touchpoint to {p['name']} in "
                    f"the trailing 12 months — {_fmt_int(sample) or 'panel'} "
                    "panelists, projected to a U.S. digital audience.")
    else:
        subhead = (f"A real-time view of the {p['name']} audience — "
                    f"{_fmt_int(sample) or 'panel'} panelists, projected to a "
                    "U.S. digital audience.")
    _text(slide, 0.62, 2.00, 12.0, 0.60, subhead,
          size=12.5, color=C_MUTED, line_spacing=1.4)

    stat_top = 2.85
    stats = _by_the_numbers_stats(p)
    n = max(1, len(stats))
    slot_w = (SLIDE_W_IN - 1.24 - 0.12 * (n - 1)) / n
    x = 0.62
    for i, (value, label, note, tone) in enumerate(stats):
        color = (C_LAVENDER if tone == "lavender"
                 else C_MAGENTA if tone == "magenta"
                 else C_CREAM)
        _text(slide, x, stat_top, slot_w, 1.20, value,
              size=60, bold=True, color=color, line_spacing=1.0)
        _text(slide, x, stat_top + 1.30, slot_w, 0.30, label.upper(),
              size=10, bold=True, color=C_LAVENDER, letter_spacing=0.10)
        _text(slide, x, stat_top + 1.65, slot_w, 1.90, note,
              size=10.5, color=C_MUTED, line_spacing=1.45)
        x += slot_w + 0.12

    _text(slide, 0.62, 6.75, 12.0, 0.30, _source_line(p),
          size=8.5, color=C_MUTED2, letter_spacing=0.02)

    _footer(slide, 0, p["name"])


def _slide_demographics(prs: Presentation, p: dict):
    """LISA IDENTITY layout: 3-4 large stat cards on left (% + label + idx),
    right-side full-height image with darkened overlay, WHY IT MATTERS narrative."""
    slide = _blank_slide(prs)

    # Right-side hero image (LISA identity slide places the photo at x=3.13)
    _place_side_image(slide, p.get("_image_bytes"), x=6.10, y=0.0,
                      w=SLIDE_W_IN - 6.10, h=SLIDE_H_IN,
                      darken=0.55, subject_name=p["name"])
    # Left-side dark overlay to keep stat text readable regardless of image
    _rect(slide, 0, 0, 6.30, SLIDE_H_IN, C_DARK)

    _section_eyebrow(slide, 0.52, p["name"], "IDENTITY")
    _text(slide, 0.60, 0.82, 6.0, 1.10, _identity_headline(p),
          size=28, bold=True, color=C_CREAM, line_spacing=1.05)

    # Build the 4 most-interesting demographic call-outs (top female share,
    # dominant age bucket, top ethnicity, top income) — dynamically chosen
    # so the slide is meaningful even for MPB-only or partial-demo profiles.
    stats = _identity_stat_rows(p)

    # Row layout: 4 stats at y=2.30 / 3.28 / 4.26 / 5.24, each 0.98" tall
    # (was 1.10 which pushed the 4th idx line into WHY IT MATTERS at y=6.15).
    # Last idx line now ends at 5.24 + 0.44 + 0.30 = 5.98; WHY IT MATTERS
    # eyebrow at 6.20 → 0.22" clean gap.
    row_top = 2.30
    row_h = 0.98
    for i, (pct_str, label, idx_val) in enumerate(stats[:4]):
        y = row_top + i * row_h
        color = C_LAVENDER if idx_val >= 130 else C_CREAM
        _text(slide, 0.62, y, 2.30, 0.72, pct_str,
              size=30, bold=True, color=color, line_spacing=1.0)
        _text(slide, 2.97, y + 0.06, 3.30, 0.36, label,
              size=12.5, bold=True, color=color, letter_spacing=0.04)
        if idx_val:
            _text(slide, 2.97, y + 0.44, 3.30, 0.28,
                  f"idx  {int(idx_val)}",
                  size=10, color=C_MUTED, letter_spacing=0.05)
        if i < len(stats[:4]) - 1:
            _hairline(slide, 0.62, y + row_h - 0.05, 5.90, C_STROKE)

    _text(slide, 0.62, 6.20, 5.90, 0.26, "WHY IT MATTERS",
          size=9.5, bold=True, color=C_LIME, letter_spacing=0.10)
    _text(slide, 0.62, 6.47, 5.90, 0.55, _identity_why_it_matters(p),
          size=11, color=C_MUTED, line_spacing=1.4)

    _footer(slide, 0, p["name"])


def _slide_one_insight(prs: Presentation, p: dict):
    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], "ONE INSIGHT")

    brief_a    = _brief_get(p, "one_insight_a")
    brief_b    = _brief_get(p, "one_insight_b")
    brief_body = _brief_get(p, "one_insight_body")

    if brief_a and brief_b:
        headline_a = _purge_pronouns(brief_a, p["name"])
        headline_b = _purge_pronouns(brief_b, p["name"])
        body       = _purge_pronouns(brief_body or "", p["name"])
    else:
        # Fallback path — filter noise brands so we don't headline a 5% /
        # 800-idx artifact.
        top = _top_over_index_brand(p, min_pct=10.0, max_index=300.0)
        subj = p["name"]
        arch = p.get("_archetype") or _ARCHETYPE_GENERIC
        if top:
            cat   = _pretty_cat(top.get("_cat", ""))
            brand = top.get("name", "")
            idx   = int(_num(top.get("index", 0)))
            pct   = _num(top.get("pct", 0))
            headline_a = f"The {subj} audience shows up"
            headline_b = f"where the buying already happens."
            body = (
                f"{pct:.0f}% of this audience engages with {brand} — a "
                f"{idx}-index vs. the U.S. baseline and the strongest durable "
                f"{cat.lower() or 'brand'} signal in the file. That is a real "
                "spend signature, not a passing affinity."
            )
        else:
            headline_a = f"The {subj} audience"
            headline_b = "concentrates habit."
            body = ("The audience over-indexes across multiple durable "
                    "categories — meaning any partner brand aligned here "
                    "captures repeat behavior, not a single moment.")
        # Movie/TV-specific closing sentence tacked onto the fallback body
        if arch in (_ARCHETYPE_MOVIE, _ARCHETYPE_TV_SHOW):
            body += (" For the release plan, that means every over-indexed "
                     "surface below is a media-buy priority, not just a data "
                     "point.")

    # Two-line big display: auto-shrink to keep each line on one visual
    # row (LLM headlines can be 8 words and exceed 12" at 44pt). Line
    # offset is computed from the actual font so wraps never collide.
    pt1 = _fit_display_size(headline_a, 12.0, base_pt=44, min_pt=26, bold=True)
    pt2 = _fit_display_size(headline_b, 12.0, base_pt=44, min_pt=26, bold=True)
    hd_pt = min(pt1, pt2)
    line_h = hd_pt / 72.0 * 1.05
    y_a = 1.60
    y_b = y_a + line_h + 0.15
    _text(slide, 0.62, y_a, 12.0, line_h + 0.10, headline_a,
          size=hd_pt, bold=True, color=C_CREAM, line_spacing=1.05)
    _text(slide, 0.62, y_b, 12.0, line_h + 0.10, headline_b,
          size=hd_pt, bold=True, color=C_LAVENDER, line_spacing=1.05)

    hair_y = max(4.05, y_b + line_h + 0.30)
    _hairline(slide, 0.62, hair_y, 12.10, C_STROKE)

    body_y = hair_y + 0.30
    body_h = min(2.55, 6.70 - body_y)
    _text(slide, 0.62, body_y, 12.10, body_h, body,
          size=13, color=C_MUTED, line_spacing=1.50)

    _footer(slide, 0, p["name"])


def _slide_persona_category(prs: Presentation, p: dict, cat: str, items: list):
    """Category deep-dive slide, grouped by SIGNAL CLASS instead of index
    quartile.

    Prior versions sorted the category's brands by index desc and called
    the top four quartiles STRONGEST / STRONG / NOTABLE / EMERGING. That
    labeling let a 3.6%-reach brand at 263 index become 'STRONGEST
    SIGNAL' for Apparel, which is analytically false — 96% of the
    audience does not shop it. The fix: group by the 2x2 of reach x
    index that actually matters for a media buy:

        MASS + LIFT   (reach >=20% AND index >=130) — actionable
        SCALE         (reach >=20% AND index 90-130) — table stakes
        SIGNATURE     (reach 8-20% AND index >=150) — persona-only
        EMERGING      (reach 10-20% AND index 110-150) — moderate

    A short header caption teaches the framework so the reader can
    interpret a MARNI at 3.6% for what it is (character detail) rather
    than what it isn't (a channel plan)."""
    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], _pretty_cat(cat).upper())

    headline, sub = _persona_category_headline(cat)
    full_title = f"{headline} {sub}".strip()
    _text(slide, 0.62, 0.82, 12.0, 1.00, full_title,
          size=26, bold=True, color=C_CREAM, line_spacing=1.05)

    all_items = [it for it in (items or [])
                 if _num(it.get("pct", 0)) > 0 and _num(it.get("index", 0)) > 0]
    if not all_items:
        _text(slide, 0.62, 3.5, 11.5, 0.6,
              "No branded signals observed in this category.",
              size=12, color=C_MUTED)
        _footer(slide, 8, p["name"])
        return

    # Classify every brand and bucket into the four useful columns.
    buckets: dict[str, list[dict]] = {
        _SIGCLASS_MASS_LIFT: [],
        _SIGCLASS_SCALE:     [],
        _SIGCLASS_SIGNATURE: [],
        _SIGCLASS_EMERGING:  [],
    }
    for it in all_items:
        cls = _signal_class(it.get("pct"), it.get("index"))
        if cls in buckets:
            buckets[cls].append(it)
    # In-bucket ordering: mass-lift by signal_score (reach with lift kicker),
    # scale by reach, signature/emerging by index (that IS their story).
    buckets[_SIGCLASS_MASS_LIFT].sort(
        key=lambda it: -_signal_score(it.get("pct"), it.get("index")))
    buckets[_SIGCLASS_SCALE].sort(
        key=lambda it: -_num(it.get("pct", 0)))
    buckets[_SIGCLASS_SIGNATURE].sort(
        key=lambda it: -_num(it.get("index", 0)))
    buckets[_SIGCLASS_EMERGING].sort(
        key=lambda it: -_num(it.get("index", 0)))

    # Only ship columns that have brands. Prefer the actionable ones first.
    col_order = [_SIGCLASS_MASS_LIFT, _SIGCLASS_SCALE,
                 _SIGCLASS_SIGNATURE, _SIGCLASS_EMERGING]
    active_cols = [(cls, buckets[cls][:5]) for cls in col_order
                   if buckets[cls]]
    if not active_cols:
        # Fall back: everything is under 5% reach and under 110 idx — just
        # rank by index and label as "NOTED SIGNALS".
        fallback = sorted(all_items,
                          key=lambda it: -_num(it.get("index", 0)))[:12]
        active_cols = [("NOTED", fallback)]

    # WHY-IT-MATTERS caption grounded in the mass-lift column if present.
    anchor_items = (buckets[_SIGCLASS_MASS_LIFT]
                     or buckets[_SIGCLASS_SCALE]
                     or buckets[_SIGCLASS_SIGNATURE]
                     or all_items)
    why = _why_it_matters_for_category(cat, p, anchor_items[:3])
    _text(slide, 0.62, 1.90, 12.0, 0.60,
          f"WHY IT MATTERS  /  {why}",
          size=10, color=C_MUTED, letter_spacing=0.04, line_spacing=1.4)

    # Small framework legend so the reader learns to read reach x index.
    legend_y = 2.55
    _text(slide, 0.62, legend_y, 12.0, 0.24,
          "READ IT  /  reach = how many. index = how uniquely. "
          "MASS+LIFT drives the buy; SIGNATURE colors the persona.",
          size=9, color=C_MUTED2, letter_spacing=0.04)

    n_cols = len(active_cols)
    col_w = (SLIDE_W_IN - 1.24 - 0.30 * (n_cols - 1)) / n_cols
    col_x = [0.62 + i * (col_w + 0.30) for i in range(n_cols)]

    top_y = 2.95
    row_h = 0.72

    class_header_color = {
        _SIGCLASS_MASS_LIFT: C_LIME,
        _SIGCLASS_SCALE:     C_LAVENDER,
        _SIGCLASS_SIGNATURE: C_MAGENTA,
        _SIGCLASS_EMERGING:  C_LAVENDER,
        "NOTED":             C_LAVENDER,
    }
    class_sub_caption = {
        _SIGCLASS_MASS_LIFT: "big reach + real lift  —  drives the buy",
        _SIGCLASS_SCALE:     "big reach at parity  —  table stakes",
        _SIGCLASS_SIGNATURE: "small reach, sharp index  —  persona color",
        _SIGCLASS_EMERGING:  "moderate reach + moderate lift  —  test",
        "NOTED":             "ranked signals in this category",
    }

    for c, (cls, brands) in enumerate(active_cols):
        x = col_x[c]
        hdr = _SIGCLASS_LABEL.get(cls, cls)
        _text(slide, x, top_y, col_w, 0.30, hdr,
              size=10, bold=True,
              color=class_header_color.get(cls, C_LAVENDER),
              letter_spacing=0.10)
        _text(slide, x, top_y + 0.28, col_w, 0.22,
              class_sub_caption.get(cls, ""),
              size=8, color=C_MUTED2, letter_spacing=0.02)
        _hairline(slide, x, top_y + 0.56, col_w - 0.10, C_STROKE)
        for r, it in enumerate(brands):
            y = top_y + 0.68 + r * row_h
            brand_name = (it.get("name") or "").strip()
            pct = _num(it.get("pct", 0))
            idx = int(_num(it.get("index", 0)))
            _text(slide, x, y, col_w - 0.10, 0.34, brand_name,
                  size=12.5, bold=True, color=C_CREAM)
            pct_str = f"{pct:.1f}% reach"
            idx_str = f"({idx})" if idx else ""
            idx_color = (C_MAGENTA if idx >= 150 else
                         (C_LIME if idx >= 110 else C_MUTED))
            _text(slide, x, y + 0.34, col_w - 1.0, 0.30, pct_str,
                  size=10.5, color=C_MUTED)
            _text(slide, x + col_w - 1.0, y + 0.34, 0.90, 0.30, idx_str,
                  size=11, bold=True, color=idx_color, align=PP_ALIGN.RIGHT)

    _text(slide, 0.62, 6.75, 12.0, 0.30, _source_line(p),
          size=8.5, color=C_MUTED2, letter_spacing=0.02)
    _footer(slide, 8, p["name"])


def _slide_case_study(prs: Presentation, p: dict):
    """Emits 1-3 case-study slides (one per top over-indexing brand),
    each in the KD format: subject headline + deal-narrative paragraph +
    3 giant stat callouts (index, penetration, gen-pop delta)."""
    picks = _pick_case_studies(p, n=3)
    if not picks:
        return
    for i, cs in enumerate(picks, start=1):
        _emit_case_study_slide(prs, p, cs, case_num=i)


def _emit_case_study_slide(prs, p: dict, cs: dict, *, case_num: int):
    """Render one KD case-study slide. ``cs`` shape:
    ``{name, category, pct, index, gen_pop_pct, projection, narrative}``."""
    slide = _blank_slide(prs, bg=C_DARK)
    # Right-column subject image with strong darken
    img = p.get("_image_bytes")
    if img:
        try:
            _place_prepared_image(slide, img, x=0, y=0,
                                  w=SLIDE_W_IN, h=SLIDE_H_IN, darken=0.30)
        except Exception:
            pass
    # Left-side dark scrim
    _rect(slide, 0, 0, 6.0, SLIDE_H_IN, RGBColor(0x08, 0x0F, 0x11))

    brand = cs["name"]
    _text(slide, 0.62, 0.55, 12.0, 0.30,
          f"CASE STUDY  {case_num:02d}",
          size=10.5, color=C_MAGENTA, letter_spacing=0.10)
    _text(slide, 0.62, 0.90, 12.0, 0.30,
          f"THE {p['name'].upper()} AUDIENCE  /  {brand.upper()}",
          size=9.5, color=C_MUTED, letter_spacing=0.06)

    # Two-line editorial headline — auto-shrink so LLM headlines from
    # the analyst brief never wrap and collide with the second line.
    headline_a, headline_b = cs["headline_a"], cs["headline_b"]
    pt = min(
        _fit_display_size(headline_a, 5.80, base_pt=36, min_pt=20, bold=True),
        _fit_display_size(headline_b, 5.80, base_pt=36, min_pt=20, bold=True),
    )
    line_h = pt / 72.0 * 1.02
    y_a = 1.60
    y_b = y_a + line_h + 0.15
    _text(slide, 0.62, y_a, 5.80, line_h + 0.10, headline_a,
          size=pt, bold=True, color=C_CREAM, line_spacing=1.02)
    _text(slide, 0.62, y_b, 5.80, line_h + 0.10, headline_b,
          size=pt, bold=True, color=C_LAVENDER, line_spacing=1.02)

    # Analyst narrative paragraph — widened (5.85") + taller (2.80") + 10.5pt
    # so the 75-110 word analyst copy from the brief never clips mid-sentence.
    _text(slide, 0.62, 3.80, 5.85, 2.85, cs["narrative"],
          size=10.5, color=C_CREAM, line_spacing=1.45)

    # Right side: 3 giant stat callouts. Layout was rebalanced 2026-07-08 —
    # the previous callout occupied ~2.00" of vertical space but the y_inc
    # was only 1.65, so callouts stacked (INDEX label collided with the
    # next value; last callout's sub collided with the source line at
    # y=6.75). New geometry: value 44pt (0.95" tall) + label offset 1.00
    # + sub offset 1.32 + sub 0.35 → callout is 1.67" tall, y_inc 1.62,
    # 3rd callout sub ends at y ≈ 6.24 → clean 0.51" gap to source.
    x = 6.60
    y = 1.65
    for stat_num, (label, value, sub) in enumerate(cs["stats"][:3]):
        _text(slide, x, y, 6.2, 0.95, value,
              size=44, bold=True, color=C_MAGENTA, line_spacing=1.0)
        _text(slide, x, y + 1.00, 6.2, 0.28, label.upper(),
              size=10, bold=True, color=C_CREAM, letter_spacing=0.08)
        _text(slide, x, y + 1.32, 6.2, 0.35, sub,
              size=10.5, color=C_MUTED, line_spacing=1.3)
        if stat_num < 2:
            _hairline(slide, x, y + 1.72, 5.8, C_STROKE)
        y += 1.62

    _text(slide, 0.62, 6.75, 8.0, 0.30, _source_line(p),
          size=8.5, color=C_MUTED2, letter_spacing=0.02)
    _footer(slide, 0, p["name"])


def _slide_media_and_social(prs: Presentation, p: dict):
    """Where the audience spends screen time.

    Prior version passed items[:5] which meant HIDIVE at 2.8% reach + 363
    index was shown at giant font as the top streaming pick — analytically
    misleading because a 2.8% surface cannot deliver a media buy. Now we
    prioritize MASS + LIFT (real reach) picks and only fall back to
    SIGNATURE if no mass surface exists.
    """
    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], "SOCIAL & PLATFORM")
    _title_block(slide, 1.0, "Where the audience", "spends its screen time.")

    social_items = _first_available(p, [
        "SOCIAL MEDIA", "SOCIAL_MEDIA", "PLATFORM", "APP/PLATFORM",
    ])
    media_items = _first_available(p, [
        "STREAMING/PLATFORM", "STREAMING_PLATFORM",
        "STREAMING VIDEO", "STREAMING_VIDEO", "MEDIA",
    ])

    # Rank both lists by SIGNAL SCORE (reach with a lift kicker) so the
    # displayed order matches what actually matters for a media plan.
    social_ranked = sorted(
        social_items or [],
        key=lambda it: -_signal_score(it.get("pct", 0), it.get("index", 0)),
    )
    media_ranked = sorted(
        media_items or [],
        key=lambda it: -_signal_score(it.get("pct", 0), it.get("index", 0)),
    )

    _text(slide, 0.62, 2.85, 5.5, 0.30, "SOCIAL",
          size=11, bold=True, color=C_LAVENDER, letter_spacing=0.06)
    _draw_stat_list(slide, 0.62, 3.30, 5.5, social_ranked[:5])

    _text(slide, 7.0, 2.85, 5.5, 0.30, "STREAMING / MEDIA",
          size=11, bold=True, color=C_LAVENDER, letter_spacing=0.06)
    _draw_stat_list(slide, 7.0, 3.30, 5.5, media_ranked[:5])

    _footer(slide, 0, p["name"])


def _draw_stat_list(slide, x, y, w, items):
    """Ranked stat rows: brand name + reach on the left, giant idx number
    on the right with 'idx' caption below it. Idx-number box was h=0.6
    starting at ry+0.05 which visually extends to ry+0.44 for 28pt glyphs;
    'idx' caption at ry+0.55 sits comfortably below the glyph baseline
    but the two BOUNDING BOXES used to overlap by 0.10" (defect on
    slide 18). Number box shrunk to h=0.44 to eliminate the bbox overlap
    without changing the visual layout."""
    if not items:
        _text(slide, x, y, w, 0.4, "No data in this category.",
              size=11, color=C_MUTED)
        return
    row_h = 0.72
    for i, it in enumerate(items):
        ry = y + i * row_h
        name = (it.get("name") or "").strip()
        pct = _num(it.get("pct", 0))
        idx = int(_num(it.get("index", 0)))
        _text(slide, x, ry, 3.4, 0.35, name,
              size=13, bold=True, color=C_CREAM)
        _text(slide, x, ry + 0.34, 3.4, 0.24,
              f"{pct:.1f}% reach", size=9, color=C_MUTED)
        idx_color = C_MAGENTA if idx >= 120 else (C_LIME if idx >= 100 else C_LAVENDER)
        _text(slide, x + 3.6, ry + 0.05, 1.6, 0.44, f"{idx}",
              size=28, bold=True, color=idx_color, align=PP_ALIGN.RIGHT,
              line_spacing=1.0)
        _text(slide, x + 3.6, ry + 0.52, 1.6, 0.20, "idx",
              size=9, color=C_MUTED, align=PP_ALIGN.RIGHT, letter_spacing=0.05)
        if i < len(items) - 1:
            _hairline(slide, x, ry + row_h - 0.05, w - 0.1, C_STROKE)


def _slide_final_insight(prs: Presentation, p: dict):
    """Final-insight closer: full-bleed darkened image, tiny 'THE FINAL
    INSIGHT' eyebrow, two-line display headline, longer analyst-body
    paragraph. Body box is intentionally wider + taller than the LISA
    reference so multi-sentence closers do not clip."""
    slide = _blank_slide(prs, bg=C_DARK)
    img_bytes = p.get("_image_bytes")
    if img_bytes:
        try:
            _place_prepared_image(slide, img_bytes, x=0, y=0,
                                  w=SLIDE_W_IN, h=SLIDE_H_IN, darken=0.35)
        except Exception:
            _rect(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, C_DARK)
    # Extra dark scrim on the left half so headline reads clean
    _rect(slide, 0, 0, 9.2, SLIDE_H_IN, RGBColor(0x08, 0x0F, 0x11))

    _text(slide, 0.62, 0.62, 8.0, 0.30, "THE FINAL INSIGHT",
          size=10.5, color=C_LIME, letter_spacing=0.10)

    line1, line2 = _final_headline(p)
    # Auto-shrink to the largest size that keeps BOTH lines on 1 visual
    # line each — takes the min so both headlines match. Line offset is
    # then computed from the actual size, not a hardcoded 1.25 that
    # collided when strings wrapped.
    pt1 = _fit_display_size(line1, 8.6, base_pt=44, min_pt=24, bold=True)
    pt2 = _fit_display_size(line2, 8.6, base_pt=44, min_pt=24, bold=True)
    hd_pt = min(pt1, pt2)
    line_h = hd_pt / 72.0 * 1.02
    y_a = 2.00
    y_b = y_a + line_h + 0.14
    _text(slide, 0.62, y_a, 8.6, line_h + 0.10, line1,
          size=hd_pt, bold=True, color=C_CREAM, line_spacing=1.02)
    _text(slide, 0.62, y_b, 8.6, line_h + 0.10, line2,
          size=hd_pt, bold=True, color=C_LAVENDER, line_spacing=1.02)

    top  = _top_over_index_brand(p, min_pct=10.0, max_index=300.0)
    body = _final_body(p, top)
    # Body starts below the headline block. y_b + 2 * line_h + a safety
    # gap of 0.30" so body never overlaps line 2 even at max font.
    body_y = y_b + line_h + 0.35
    body_h = min(2.30, 6.70 - body_y)
    _text(slide, 0.62, body_y, 8.40, body_h, body,
          size=12, color=C_CREAM, line_spacing=1.50)

    _footer(slide, 0, p["name"])


# =============================================================================
#  KD-inspired slide builders  ·  index story, geography, MPB table,
#  takeaways, methodology, about, narrative bridges
# =============================================================================


def _slide_index_story(prs: Presentation, p: dict):
    """KD slide 5 pattern: 'Where the audience over-indexes most.' — 4 large
    index-vs-Gen-Pop callouts pulled from the strongest demographic gaps."""
    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"],
                     "INDEX VS GEN POP  /  THE STORY")
    _text(slide, 0.62, 0.82, 12.0, 1.10,
          "Where the audience over-indexes most.",
          size=28, bold=True, color=C_CREAM, line_spacing=1.05)

    arch = p.get("_archetype") or _ARCHETYPE_GENERIC
    if arch == _ARCHETYPE_MOVIE:
        why = ("Each callout below is a targetable audience concentration — "
               "prioritize these in the release-marketing media plan.")
    elif arch == _ARCHETYPE_TV_SHOW:
        why = ("Each callout below is a targetable audience concentration — "
               "prioritize these in launch and renewal marketing.")
    else:
        why = ("Each callout below is a targetable audience concentration — "
               "the demographic shape most differentiated from Gen Pop.")
    _text(slide, 0.62, 1.98, 12.0, 0.36,
          f"WHY IT MATTERS  /  {why}",
          size=10, color=C_MUTED, letter_spacing=0.04)

    picks = _top_demo_index_callouts(p, n=4)
    if not picks:
        _text(slide, 0.62, 3.5, 11.5, 0.6,
              "Insufficient demographic data to compute over-index story.",
              size=12, color=C_MUTED)
        _footer(slide, 5, p["name"])
        return

    n = len(picks)
    slot_w = (SLIDE_W_IN - 1.24 - 0.20 * (n - 1)) / n
    x = 0.62
    y = 2.90
    for value, label, sub in picks:
        # Giant index number
        _text(slide, x, y, slot_w, 1.60, str(int(value)),
              size=88, bold=True, color=C_MAGENTA, line_spacing=1.0)
        _text(slide, x, y + 1.80, slot_w, 0.28,
              "INDEX vs GEN POP",
              size=9.5, color=C_MUTED, letter_spacing=0.10)
        _hairline(slide, x, y + 2.15, slot_w - 0.10, C_STROKE)
        _text(slide, x, y + 2.30, slot_w, 0.36, label,
              size=13.5, bold=True, color=C_CREAM, letter_spacing=0.04)
        _text(slide, x, y + 2.72, slot_w, 0.65, sub,
              size=10.5, color=C_MUTED, line_spacing=1.4)
        x += slot_w + 0.20

    _text(slide, 0.62, 6.75, 12.0, 0.30, _source_line(p),
          size=8.5, color=C_MUTED2, letter_spacing=0.02)
    _footer(slide, 5, p["name"])


def _slide_geography(prs: Presentation, p: dict):
    """KD-style geography slide: ranked list of top DMAs with pen% and idx.

    Renders as a two-column list (10 rows total) with big-market callout
    at the top. Skips itself gracefully if no location data is available."""
    dmas = _top_dmas_ranked(p, n=10)
    if not dmas:
        return  # skip slide if no geo data
    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], "GEOGRAPHY")
    _text(slide, 0.62, 0.82, 12.0, 1.10,
          "Big-market reach, ranked.",
          size=28, bold=True, color=C_CREAM, line_spacing=1.05)
    _text(slide, 0.62, 2.02, 12.0, 0.85,
          f"National footprint of the {p['name']} audience, ranked by share "
          "of panelists per Nielsen DMA. Index shows over/under vs national baseline.",
          size=12, color=C_MUTED, line_spacing=1.4)

    _text(slide, 0.62, 3.15, 12.0, 0.30, "TOP 10 NIELSEN DMAs  (RANK)",
          size=10, bold=True, color=C_LAVENDER, letter_spacing=0.10)

    # Two columns of 5 rows each
    row_h = 0.52
    col_w = 5.90
    col_x = [0.62, 0.62 + col_w + 0.30]
    top_y = 3.60
    for i, (name, pct, idx) in enumerate(dmas):
        col = i // 5
        row = i % 5
        x = col_x[col]
        y = top_y + row * row_h
        rank_lbl = f"{i+1:02d}"
        # Rank chip
        _text(slide, x, y, 0.55, row_h, rank_lbl,
              size=13, bold=True, color=C_LAVENDER, anchor=MSO_ANCHOR.MIDDLE)
        # DMA name
        _text(slide, x + 0.55, y, col_w - 2.30, row_h,
              name, size=12, bold=True, color=C_CREAM,
              anchor=MSO_ANCHOR.MIDDLE)
        # Pen %
        _text(slide, x + col_w - 1.70, y, 0.90, row_h,
              f"{pct:.1f}%", size=12, color=C_CREAM,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        # Index chip
        color = (C_MAGENTA if idx >= 150 else
                 (C_LIME if idx >= 110 else C_MUTED))
        _text(slide, x + col_w - 0.75, y, 0.75, row_h,
              f"{int(idx)}" if idx else "—",
              size=12, bold=True, color=color,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        _hairline(slide, x, y + row_h - 0.02,
                  col_w - 0.10, C_STROKE)

    _text(slide, 0.62, 6.75, 12.0, 0.30, _source_line(p),
          size=8.5, color=C_MUTED2, letter_spacing=0.02)
    _footer(slide, 6, p["name"])


def _slide_stat_splash(prs: Presentation, p: dict, cat: str, items: list):
    """Two large brand stat tiles per category (giant reach% and giant idx).
    Layout was rebalanced 2026-07-08 to fix a WHY IT MATTERS / stat-tile
    overlap where the INDEX label at y0+3.55 (=6.15) and projection line at
    y0+3.85 (=6.45) collided with the WHY IT MATTERS eyebrow/body at
    y=6.15/6.42. Tiles now end at 5.95; WHY IT MATTERS sits at 6.10/6.35."""
    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], _pretty_cat(cat).upper())
    headline, sub = _persona_category_headline(cat)
    _text(slide, 0.60, 0.82, 12.0, 1.10, f"{headline} {sub}".strip(),
          size=28, bold=True, color=C_CREAM, line_spacing=1.05)

    # Pick two brands worthy of giant-stat treatment. Prioritize
    # MASS+LIFT (big reach + real over-index), then SCALE (big reach),
    # then SIGNATURE (small reach + sharp index). Falling back on
    # index-desc alone would put a 3.6%-reach freak-index brand in a
    # 56pt tile — misleading at a glance.
    all_ok = [it for it in (items or []) if _num(it.get("pct", 0)) > 0]
    prioritized: dict[str, list[dict]] = {c: [] for c in
        (_SIGCLASS_MASS_LIFT, _SIGCLASS_SCALE, _SIGCLASS_SIGNATURE,
         _SIGCLASS_EMERGING)}
    for it in all_ok:
        cls = _signal_class(it.get("pct"), it.get("index"))
        if cls in prioritized:
            prioritized[cls].append(it)
    prioritized[_SIGCLASS_MASS_LIFT].sort(
        key=lambda it: -_signal_score(it.get("pct"), it.get("index")))
    prioritized[_SIGCLASS_SCALE].sort(
        key=lambda it: -_num(it.get("pct", 0)))
    prioritized[_SIGCLASS_SIGNATURE].sort(
        key=lambda it: -_num(it.get("index", 0)))
    prioritized[_SIGCLASS_EMERGING].sort(
        key=lambda it: -_num(it.get("index", 0)))
    top: list[dict] = []
    for cls in (_SIGCLASS_MASS_LIFT, _SIGCLASS_SCALE,
                _SIGCLASS_EMERGING, _SIGCLASS_SIGNATURE):
        for it in prioritized[cls]:
            if it not in top:
                top.append(it)
            if len(top) >= 2:
                break
        if len(top) >= 2:
            break
    if len(top) < 2:  # backfill from raw list if we don't have two hits
        for it in all_ok:
            if it not in top:
                top.append(it)
            if len(top) >= 2:
                break
    if not top:
        _text(slide, 0.62, 3.5, 11.5, 0.6,
              "No branded signals observed in this category.",
              size=12, color=C_MUTED)
        _footer(slide, 0, p["name"])
        return

    tile_w = 5.60
    tile_gap = 0.40
    y0 = 2.15   # was 2.60 — pull up so tile ends higher
    x0 = 0.62
    # Vertical rhythm within a tile (all offsets are added to y0):
    OFF_NAME      = 0.00
    OFF_PCT       = 0.40
    OFF_REACH_LBL = 1.55
    OFF_HAIRLINE  = 1.85
    OFF_IDX       = 1.95
    OFF_IDX_LBL   = 3.05
    OFF_PROJ      = 3.35   # tile ends at y0 + 3.62 ≈ 5.77
    for i, it in enumerate(top[:2]):
        x = x0 + i * (tile_w + tile_gap)
        pct    = _num(it.get("pct", 0))
        idx    = int(_num(it.get("index", 0)))
        proj   = _num(it.get("projection", 0))
        name   = (it.get("name") or "").strip()
        gp_pct = _num(it.get("genPopPct", 0))

        _text(slide, x, y0 + OFF_NAME, tile_w, 0.36, name.upper(),
              size=11, bold=True, color=C_LAVENDER, letter_spacing=0.08)
        # Giant reach % — 56pt (was 64pt) so 3-digit values fit + tile is shorter
        _text(slide, x, y0 + OFF_PCT, tile_w, 1.15, f"{pct:.1f}%",
              size=56, bold=True, color=C_CREAM, line_spacing=1.0)
        _text(slide, x, y0 + OFF_REACH_LBL, tile_w, 0.28, "REACH",
              size=9.5, color=C_MUTED, letter_spacing=0.10)
        _hairline(slide, x, y0 + OFF_HAIRLINE, tile_w - 0.10, C_STROKE)
        # Giant index — 60pt (was 76pt) so tile compresses vertically
        idx_color = (C_MAGENTA if idx >= 150 else
                     (C_LIME if idx >= 110 else C_LAVENDER))
        _text(slide, x, y0 + OFF_IDX, tile_w, 1.10,
              f"{idx}" if idx else "—",
              size=60, bold=True, color=idx_color, line_spacing=1.0)
        _text(slide, x, y0 + OFF_IDX_LBL, tile_w, 0.28,
              f"INDEX vs GEN POP  ·  {gp_pct:.1f}% BASELINE",
              size=9.5, color=C_MUTED, letter_spacing=0.08)
        if proj:
            _text(slide, x, y0 + OFF_PROJ, tile_w, 0.28,
                  f"= {_fmt_int(proj)} U.S. adults",
                  size=10, color=C_MUTED)

    # WHY IT MATTERS moved to y=6.05 / body at 6.30 — comfortably below
    # the (now compressed) stat tiles that end at y ≈ 5.77.
    _text(slide, 0.62, 6.05, 12.0, 0.26, "WHY IT MATTERS",
          size=9.5, bold=True, color=C_LIME, letter_spacing=0.10)
    _text(slide, 0.62, 6.32, 12.0, 0.72,
          _why_it_matters_for_category(cat, p, top),
          size=11, color=C_MUTED, line_spacing=1.4)

    _footer(slide, 0, p["name"])


def _slide_data_confirms(prs: Presentation, p: dict):
    """Narrative bridge slide: two-line display + analyst body paragraph.

    Prefers the LLM brief; falls back to a durable-signal-driven template
    with the noise-brand filter applied so we don't lead with a freak
    5%-reach / 800-idx artifact."""
    slide = _blank_slide(prs, bg=C_DARK)
    img = p.get("_image_bytes")
    if img:
        try:
            _place_prepared_image(slide, img, x=8.10, y=0.0,
                                  w=SLIDE_W_IN - 8.10, h=SLIDE_H_IN,
                                  darken=0.55)
        except Exception:
            pass

    brief_a    = _brief_get(p, "data_confirms_a")
    brief_b    = _brief_get(p, "data_confirms_b")
    brief_body = _brief_get(p, "data_confirms_body")

    if brief_a and brief_b:
        line1 = _purge_pronouns(brief_a, p["name"])
        line2 = _purge_pronouns(brief_b, p["name"])
        body  = _purge_pronouns(brief_body or "", p["name"])
    else:
        top = _top_over_index_brand(p, min_pct=10.0, max_index=300.0)
        if top:
            cat = _pretty_cat(top.get("_cat", "")).lower()
            line1 = "The data reads"
            line2 = "as durable habit."
            body = (
                f"The {p['name']} audience over-indexes on {top['name']} at "
                f"{int(_num(top['index']))} in {cat or 'a durable category'} "
                "— a repeated-purchase signature, not a passing affinity. "
                "Any partner brand should read this as an addressable overlap."
            )
        else:
            line1 = "One dataset,"
            line2 = "one audience."
            body = ("The audience concentrates habits across multiple "
                    "categories, which is a signal in itself. Any partner "
                    "brand aligned here captures durable behavior.")

    _text(slide, 0.62, 2.50, 8.0, 1.60, line1,
          size=54, bold=True, color=C_CREAM, line_spacing=1.02)
    _text(slide, 0.62, 3.90, 8.0, 1.60, line2,
          size=54, bold=True, color=C_LAVENDER, line_spacing=1.02)
    _text(slide, 0.62, 5.60, 7.20, 1.40, body,
          size=12.5, color=C_MUTED, line_spacing=1.45)
    _footer(slide, 0, p["name"])


def _slide_mpb_table(prs: Presentation, p: dict):
    """KD slide 20 pattern: top 20 most-purchased brands ranked table.

    Two columns of 10 rows: rank / brand name / pen% / index. Uses the
    behavioral['MPB'] or ['MOST PURCHASED BRANDS'] category, whichever
    populates. Skips if <5 brands available."""
    brands = _mpb_top_brands(p, n=20)
    if len(brands) < 5:
        return
    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], "MOST PURCHASED BRANDS")
    _text(slide, 0.62, 0.82, 12.0, 1.10, "The brands they buy.",
          size=28, bold=True, color=C_CREAM, line_spacing=1.05)
    tracked = len(_all_behavioral_items(p))
    _text(slide, 0.62, 2.00, 12.0, 0.36,
          f"WHY IT MATTERS  /  Of {tracked} tracked signals, the top {len(brands)} "
          "reveal where this audience concentrates commerce.",
          size=10, color=C_MUTED, letter_spacing=0.04)

    # Column headers
    col_w = 5.90
    col_x = [0.62, 0.62 + col_w + 0.30]
    hdr_y = 2.85
    for cx in col_x:
        _text(slide, cx, hdr_y, 0.60, 0.28, "#",
              size=9, bold=True, color=C_LAVENDER, letter_spacing=0.06)
        _text(slide, cx + 0.60, hdr_y, col_w - 2.40, 0.28, "BRAND",
              size=9, bold=True, color=C_LAVENDER, letter_spacing=0.06)
        _text(slide, cx + col_w - 1.70, hdr_y, 0.90, 0.28, "PEN %",
              size=9, bold=True, color=C_LAVENDER, letter_spacing=0.06,
              align=PP_ALIGN.RIGHT)
        _text(slide, cx + col_w - 0.75, hdr_y, 0.75, 0.28, "INDEX",
              size=9, bold=True, color=C_LAVENDER, letter_spacing=0.06,
              align=PP_ALIGN.RIGHT)
        _hairline(slide, cx, hdr_y + 0.28, col_w - 0.10, C_STROKE)

    row_h = 0.34
    top_y = 3.20
    for i, it in enumerate(brands[:20]):
        col = i // 10
        row = i % 10
        x = col_x[col]
        y = top_y + row * row_h
        _text(slide, x, y, 0.60, row_h, f"{i+1:02d}",
              size=10.5, bold=True, color=C_LAVENDER,
              anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, x + 0.60, y, col_w - 2.40, row_h,
              (it.get("name") or "").strip(),
              size=11, color=C_CREAM, anchor=MSO_ANCHOR.MIDDLE)
        pct = _num(it.get("pct", 0))
        idx = int(_num(it.get("index", 0)))
        _text(slide, x + col_w - 1.70, y, 0.90, row_h,
              f"{pct:.1f}", size=11, color=C_CREAM,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        idx_color = (C_MAGENTA if idx >= 150 else
                     (C_LIME if idx >= 110 else C_MUTED))
        _text(slide, x + col_w - 0.75, y, 0.75, row_h,
              f"{idx}" if idx else "—",
              size=11, bold=True, color=idx_color,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    _text(slide, 0.62, 6.75, 12.0, 0.30, _source_line(p),
          size=8.5, color=C_MUTED2, letter_spacing=0.02)
    _footer(slide, 99, p["name"])


def _slide_takeaways(prs: Presentation, p: dict):
    """'Six things to know.' — 6 numbered summary cards in a 3x2 grid.

    Prefers the LLM analyst brief's takeaways (better editorial voice,
    grounded in the data). Falls back to _pick_takeaways()."""
    brief_tk = _brief_get(p, "takeaways")
    if (isinstance(brief_tk, list) and len(brief_tk) >= 4):
        picks: list[tuple[str, str, str]] = []
        for i, tk in enumerate(brief_tk[:6]):
            if not isinstance(tk, dict):
                continue
            title = _purge_pronouns(str(tk.get("title") or ""), p["name"])
            body  = _purge_pronouns(str(tk.get("body")  or ""), p["name"])
            if not title or not body:
                continue
            picks.append((f"{i+1:02d}", title, body))
        # Renumber for any skipped entries
        picks = [(f"{i+1:02d}", t, b) for i, (_, t, b) in enumerate(picks)]
    else:
        picks = _pick_takeaways(p)

    if not picks:
        return
    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], "KEY TAKEAWAYS")
    _text(slide, 0.62, 0.82, 12.0, 1.10, "Six things to know.",
          size=28, bold=True, color=C_CREAM, line_spacing=1.05)
    _text(slide, 0.62, 2.00, 12.0, 0.36,
          "If you remember nothing else from this brief, remember these.",
          size=12, color=C_MUTED, line_spacing=1.4)

    tile_w = 4.05
    tile_h = 2.10
    tile_gap = 0.10
    top_y = 2.75
    for i, (num, title, body) in enumerate(picks[:6]):
        col = i % 3
        row = i // 3
        x = 0.62 + col * (tile_w + tile_gap)
        y = top_y + row * (tile_h + 0.12)
        _rect(slide, x, y, tile_w, tile_h, RGBColor(0x14, 0x1F, 0x22))
        _text(slide, x + 0.28, y + 0.20, 0.9, 0.60, num,
              size=26, bold=True, color=C_MAGENTA, line_spacing=1.0)
        _text(slide, x + 0.28, y + 0.72, tile_w - 0.50, 0.55, title,
              size=13, bold=True, color=C_CREAM, line_spacing=1.15)
        _text(slide, x + 0.28, y + 1.20, tile_w - 0.50, tile_h - 1.30, body,
              size=9.5, color=C_MUTED, line_spacing=1.4)

    _footer(slide, 0, p["name"])


# =============================================================================
#  McKinsey-shaped narrative slides
# =============================================================================
#  These are the story-first slides that replace the old data-dump pattern.
#  Every one of them prefers the LLM analyst brief and falls back to a
#  smart data-driven template when the brief is empty.
# =============================================================================


def _slide_exec_summary(prs: Presentation, p: dict):
    """3-panel Executive Summary: WHO / WHERE / SO WHAT.

    This is the second slide of every deck. If the CMO reads nothing else,
    they get the whole story here — pyramid-principle answer first, one
    number to support, then the implication."""
    es = _brief_get(p, "exec_summary", default={}) or {}

    def _panel(header_key: str, body_key: str, fallback_header: str,
               fallback_body: str) -> tuple[str, str]:
        h = _purge_pronouns(str(es.get(header_key) or "").strip(),
                            p["name"]) or fallback_header
        b = _purge_pronouns(str(es.get(body_key) or "").strip(),
                            p["name"]) or fallback_body
        return h, b

    subj = p["name"]
    proj = _fmt_int_compact(p.get("projected_us") or 0)
    top_dma = ""
    if p.get("locations"):
        first = p["locations"][0] or {}
        top_dma = str(first.get("name") or "").strip()

    who_h, who_b = _panel(
        "who_headline", "who_body",
        "The audience.",
        (f"A {proj or 'sizable U.S.'} digital audience concentrated in the "
         f"prime-earning 18-44 age bracket. Composition sets the ceiling "
         f"on which media buys and partner categories will scale."),
    )
    where_h, where_b = _panel(
        "where_headline", "where_body",
        "Where they live.",
        (f"Reach concentrates in top-3 DMAs — {top_dma or 'the coastal metros'} "
         "leads. TikTok and Instagram carry the bulk of daily attention; "
         "streaming skews to niche verticals not the majors."),
    )
    sw_h, sw_b = _panel(
        "so_what_headline", "so_what_body",
        "So what.",
        ("Re-weight paid social to TikTok Spark Ads and Discord "
         "sponsorships where the audience over-indexes. Test one "
         "endemic partner in the strongest over-index category before "
         "committing to a full-funnel program."),
    )

    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, subj, "EXECUTIVE SUMMARY")
    _text(slide, 0.62, 0.82, 12.0, 1.10, "The story in one page.",
          size=32, bold=True, color=C_CREAM, line_spacing=1.05)
    _text(slide, 0.62, 2.10, 12.0, 0.36,
          "Three answers a CMO needs before green-lighting spend.",
          size=12, color=C_MUTED, line_spacing=1.4)

    # Three column panels. Total usable width = 12.0, gutter = 0.20 → 3.87 each.
    col_w = 3.87
    gutter = 0.20
    x0 = 0.62
    y0 = 2.80
    panel_h = 3.75
    panels = [
        ("01", "WHO", who_h, who_b, C_LAVENDER),
        ("02", "WHERE", where_h, where_b, C_LIME),
        ("03", "SO WHAT", sw_h, sw_b, C_MAGENTA),
    ]
    for i, (num, kicker, hd, body, accent) in enumerate(panels):
        x = x0 + i * (col_w + gutter)
        _rect(slide, x, y0, col_w, panel_h, RGBColor(0x14, 0x1F, 0x22))
        _text(slide, x + 0.35, y0 + 0.28, 1.0, 0.60, num,
              size=32, bold=True, color=accent, line_spacing=1.0)
        _text(slide, x + 0.35, y0 + 0.90, col_w - 0.60, 0.30, kicker,
              size=10, bold=True, color=accent, letter_spacing=0.14)
        _text(slide, x + 0.35, y0 + 1.30, col_w - 0.60, 0.90, hd,
              size=17, bold=True, color=C_CREAM, line_spacing=1.10)
        _text(slide, x + 0.35, y0 + 2.35, col_w - 0.60, panel_h - 2.50, body,
              size=10.5, color=C_MUTED, line_spacing=1.45)

    _text(slide, 0.62, 6.75, 12.0, 0.30, _source_line(p),
          size=8.5, color=C_MUTED2, letter_spacing=0.02)
    _footer(slide, 0, p["name"])


def _slide_day_in_life(prs: Presentation, p: dict):
    """A day-in-the-life narrative portrait, stitched from real signals.

    Only ships when the analyst brief returns a day_in_life narrative
    (needs LLM to be worth the page).

    Layout rewritten 2026-07-09: prior version dumped one 110-150 word
    paragraph into a single 3.80" tall text box — visually a wall of text
    that fought to be skimmed. Now renders as three time-of-day cards
    (MORNING / MID-DAY / EVENING) so a CMO can read the rhythm at a
    glance. Backward-compatible with the old string shape — auto-splits
    a single paragraph into three roughly-equal chunks."""
    raw = _brief_get(p, "day_in_life")
    if not raw:
        return
    subj = p["name"]

    morning = midday = evening = ""
    if isinstance(raw, dict):
        morning = _purge_pronouns(str(raw.get("morning") or "").strip(), subj)
        midday  = _purge_pronouns(str(raw.get("midday")
                                       or raw.get("afternoon")
                                       or "").strip(), subj)
        evening = _purge_pronouns(str(raw.get("evening")
                                       or raw.get("night") or "").strip(), subj)
    elif isinstance(raw, str):
        text = _purge_pronouns(raw.strip(), subj)
        # Auto-split legacy single-paragraph shape into thirds by sentence.
        sents = [s.strip() for s in re.split(r"(?<=[.!?])\s+", text) if s.strip()]
        if len(sents) >= 3:
            third = max(1, len(sents) // 3)
            morning = " ".join(sents[:third])
            midday  = " ".join(sents[third:2*third])
            evening = " ".join(sents[2*third:])
        elif len(sents) >= 1:
            morning = text

    # Need at least one substantive panel to bother shipping the slide.
    panels_with_text = sum(1 for t in (morning, midday, evening) if t)
    if panels_with_text < 1 or sum(len(t.split()) for t in
                                    (morning, midday, evening)) < 40:
        return

    slide = _blank_slide(prs)
    # Side image (right 5" strip) with darken so left-column cards stay clean
    _place_side_image(slide, p.get("_image_bytes"),
                      x=8.20, y=0, w=SLIDE_W_IN - 8.20, h=SLIDE_H_IN,
                      darken=0.35, subject_name=subj)
    _rect(slide, 0, 0, 8.40, SLIDE_H_IN, C_DARK)

    _section_eyebrow(slide, 0.52, subj, "DAY IN THE LIFE")
    _text(slide, 0.62, 0.82, 7.70, 1.10, "How they live it.",
          size=32, bold=True, color=C_CREAM, line_spacing=1.05)
    _text(slide, 0.62, 2.05, 7.70, 0.36,
          "The audience's rhythm, stitched from panel-observed behavior.",
          size=12, color=C_MUTED, line_spacing=1.4)

    # Three vertical time-of-day cards on the left ~7.70" column
    panels = [
        ("MORNING", morning, C_LIME),
        ("MID-DAY", midday, C_LAVENDER),
        ("EVENING", evening, C_MAGENTA),
    ]
    panel_y  = 2.85
    panel_h  = 1.20
    panel_gap = 0.10
    for i, (lbl, txt, accent) in enumerate(panels):
        py = panel_y + i * (panel_h + panel_gap)
        _rect(slide, 0.62, py, 7.70, panel_h, RGBColor(0x14, 0x1F, 0x22))
        _text(slide, 0.85, py + 0.13, 1.60, 0.28, lbl,
              size=10, bold=True, color=accent, letter_spacing=0.14)
        _text(slide, 0.85, py + 0.42, 7.25, panel_h - 0.48,
              txt or "(not observed)",
              size=11, color=C_CREAM, line_spacing=1.42)

    _text(slide, 0.62, 6.75, 8.0, 0.30, _source_line(p),
          size=8.5, color=C_MUTED2, letter_spacing=0.02)
    _footer(slide, 0, subj)


def _slide_whitespace(prs: Presentation, p: dict):
    """The whitespace slide: a single analytical call-out for the surface
    the audience over-indexes on but marketing isn't leveraging."""
    hd  = _purge_pronouns(_brief_get(p, "whitespace_headline") or "",
                          p["name"])
    body = _purge_pronouns(_brief_get(p, "whitespace_body") or "",
                           p["name"])
    if not hd or not body:
        return  # LLM didn't call it out — skip rather than fake it

    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], "WHITESPACE")
    _text(slide, 0.62, 0.82, 12.0, 0.32, "WHERE THE MONEY ISN'T YET",
          size=10.5, bold=True, color=C_LIME, letter_spacing=0.12)

    hd_pt = _fit_display_size(hd, 12.0, base_pt=44, min_pt=26, bold=True)
    hd_line_h = hd_pt / 72.0 * 1.05
    _text(slide, 0.62, 1.40, 12.0, hd_line_h + 0.30, hd,
          size=hd_pt, bold=True, color=C_CREAM, line_spacing=1.05)

    # Body panel positioned dynamically so it never rides up under the
    # headline when the LLM returns a long two-line hd.
    panel_top = max(3.20, 1.40 + hd_line_h + 0.40)
    panel_h   = min(3.20, 6.55 - panel_top)
    _rect(slide, 0.62, panel_top, 12.0, panel_h, RGBColor(0x14, 0x1F, 0x22))
    _text(slide, 0.90, panel_top + 0.20, 11.5, 0.32, "THE OPPORTUNITY",
          size=10, bold=True, color=C_MAGENTA, letter_spacing=0.10)
    _text(slide, 0.90, panel_top + 0.60, 11.5, panel_h - 0.75, body,
          size=15, color=C_CREAM, line_spacing=1.55)

    _text(slide, 0.62, 6.75, 12.0, 0.30, _source_line(p),
          size=8.5, color=C_MUTED2, letter_spacing=0.02)
    _footer(slide, 0, p["name"])


def _slide_media_plan(prs: Presentation, p: dict):
    """Recommended media plan: channel + % + rationale table.

    Consumes the analyst brief's media_plan array. Skips if the LLM
    didn't return a valid plan (rather than making one up)."""
    plan = _brief_get(p, "media_plan", default=[]) or []
    if not isinstance(plan, list):
        return
    # Normalize + sanity-check
    rows: list[tuple[str, int, str]] = []
    for row in plan[:6]:
        if not isinstance(row, dict):
            continue
        ch = str(row.get("channel") or "").strip()
        pct = int(_num(row.get("pct", 0)))
        rationale = _purge_pronouns(str(row.get("rationale") or "").strip(),
                                     p["name"])
        if ch and pct > 0 and rationale:
            rows.append((ch, pct, rationale))
    if len(rows) < 3:
        return
    # Normalize % to sum to 100 if the LLM was close but not exact
    total = sum(r[1] for r in rows)
    if 90 <= total <= 110 and total != 100:
        rows = [(c, round(p_ * 100 / total), r) for (c, p_, r) in rows]

    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], "RECOMMENDED MEDIA PLAN")
    _text(slide, 0.62, 0.82, 12.0, 1.10, "Where to spend, and why.",
          size=32, bold=True, color=C_CREAM, line_spacing=1.05)
    _text(slide, 0.62, 2.10, 12.0, 0.36,
          "Suggested paid-media weighting to match the audience's actual "
          "attention posture. Rebase per market as budget scales.",
          size=11, color=C_MUTED, line_spacing=1.4)

    # Column header
    hdr_y = 3.05
    _text(slide, 0.62, hdr_y, 3.20, 0.28, "CHANNEL",
          size=9, bold=True, color=C_LAVENDER, letter_spacing=0.08)
    _text(slide, 3.90, hdr_y, 1.20, 0.28, "% OF BUDGET",
          size=9, bold=True, color=C_LAVENDER, letter_spacing=0.08,
          align=PP_ALIGN.RIGHT)
    _text(slide, 5.30, hdr_y, 7.30, 0.28, "RATIONALE",
          size=9, bold=True, color=C_LAVENDER, letter_spacing=0.08)
    _hairline(slide, 0.62, hdr_y + 0.32, 12.0, C_STROKE)

    # Fit all rows between the header (y=3.55) and the source line
    # (y=6.75) with a 0.10" gap. Solve for row_h so the last row's
    # bottom edge lands at 6.65 exactly. Prevents the 5th row of a
    # 5-row media plan from overflowing the source line (slide 19
    # defect the user flagged 2026-07-09).
    y_start = hdr_y + 0.50
    y_end   = 6.65
    n_rows  = max(1, len(rows))
    row_gap = 0.05
    row_h   = min(0.72, (y_end - y_start - row_gap * (n_rows - 1)) / n_rows)
    row_h   = max(row_h, 0.42)  # readability floor
    y = y_start
    for i, (channel, pct, rationale) in enumerate(rows):
        _text(slide, 0.62, y, 3.20, row_h, channel,
              size=15, bold=True, color=C_CREAM,
              anchor=MSO_ANCHOR.MIDDLE)
        pct_color = (C_MAGENTA if pct >= 25 else
                     (C_LIME if pct >= 15 else C_CREAM))
        _text(slide, 3.90, y, 1.20, row_h, f"{pct}%",
              size=22, bold=True, color=pct_color,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, 5.30, y, 7.30, row_h, rationale,
              size=10.5, color=C_MUTED, line_spacing=1.35,
              anchor=MSO_ANCHOR.MIDDLE)
        if i < len(rows) - 1:
            _hairline(slide, 0.62, y + row_h, 12.0, C_STROKE)
        y += row_h + row_gap

    _text(slide, 0.62, 6.75, 12.0, 0.30, _source_line(p),
          size=8.5, color=C_MUTED2, letter_spacing=0.02)
    _footer(slide, 0, p["name"])


def _slide_signature_affinities(prs: Presentation, p: dict):
    """Compact roll-up of the audience's signature brand affinities across
    ALL categories.

    STRICT class-gated picks:
      * Reject SUBJECT itself (Green Day, Nimrods leak fix — 2026-07-09)
      * Reject MASS+LIFT / SCALE (those live on category deep-dives)
      * Reject NOISE (low reach + weak idx)
      * Reject freak-index (idx > 400) — panel low-N artifact
      * Only SIGNATURE (10-20% reach + >=150 idx) and EMERGING (10-20% +
        110-150 idx) belong here — the classes that actually define
        the audience's persona color without pretending to scale.
    """
    subject = str(p.get("name") or "").strip()

    def _accept(brand: str, pct: float, idx: float,
                 cat: str = "") -> bool:
        # 1) never surface the subject itself as its own signature affinity
        if _is_subject_brand(brand, subject):
            return False
        # 2) signal-class discipline — SIGNATURE + EMERGING only
        cls = _signal_class(pct, idx)
        if cls not in (_SIGCLASS_SIGNATURE, _SIGCLASS_EMERGING):
            return False
        # 3) freak-index guard for anything the classifier let through
        if idx > 400:
            return False
        return True

    brief_rows = _brief_get(p, "signature_affinities", default=[]) or []
    picks: list[dict] = []
    rejected_llm: list[str] = []
    if isinstance(brief_rows, list):
        for row in brief_rows[:24]:
            if not isinstance(row, dict):
                continue
            brand = str(row.get("brand") or "").strip()
            cat   = str(row.get("category") or "").strip()
            pct   = _num(row.get("pct", 0))
            idx   = int(_num(row.get("index", 0)))
            note  = _purge_pronouns(str(row.get("note") or "").strip(),
                                     p["name"])
            if not (brand and pct > 0 and idx > 0):
                continue
            if not _accept(brand, pct, idx, cat):
                rejected_llm.append(
                    f"{brand} pct={pct:.1f} idx={idx} "
                    f"class={_signal_class(pct, idx)}")
                continue
            picks.append({"brand": brand, "cat": cat, "pct": pct,
                           "index": idx, "note": note})

    # Backfill from data — same class discipline. Score = reach x sqrt(idx)
    # so higher-reach signature brands rank first.
    if len(picks) < 8:
        all_items = _all_behavioral_items(p)
        scored = []
        for it in all_items:
            pct = _num(it.get("pct", 0))
            idx = int(_num(it.get("index", 0)))
            brand = (it.get("name") or "").strip()
            if not brand:
                continue
            if not _accept(brand, pct, idx, it.get("_cat") or ""):
                continue
            scored.append((_signal_score(pct, idx), it))
        scored.sort(key=lambda t: -t[0])
        seen_brands = {r["brand"].upper() for r in picks}
        for _, it in scored:
            if len(picks) >= 10:
                break
            brand = (it.get("name") or "").strip()
            if brand.upper() in seen_brands:
                continue
            picks.append({
                "brand": brand,
                "cat":   _pretty_cat(it.get("_cat") or ""),
                "pct":   _num(it.get("pct", 0)),
                "index": int(_num(it.get("index", 0))),
                "note":  "",
            })
            seen_brands.add(brand.upper())

    if rejected_llm:
        _log(f"[signature_affinities] rejected {len(rejected_llm)} LLM picks "
              f"(subject/class/freak-index): "
              f"{'; '.join(rejected_llm[:6])}")
    if len(picks) < 4:
        return

    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], "SIGNATURE AFFINITIES")
    _text(slide, 0.62, 0.82, 12.0, 1.10, "The brands that color them.",
          size=32, bold=True, color=C_CREAM, line_spacing=1.05)
    _text(slide, 0.62, 2.05, 12.0, 0.60,
          "Signature affinities are the sharpest over-indexes across the "
          "file — the brands that give this audience its personality. "
          "They are NOT the mass-reach backbone (that lives on the "
          "category deep-dives). Use these for partner selection and "
          "cultural signaling, not for sizing the media buy.",
          size=10.5, color=C_MUTED, line_spacing=1.4)

    # Two-column layout. Cap at 4 rows per column (was 5) so the last
    # row's subtitle at y+0.32 + 0.28 = y+0.60 lands well above the
    # source line at 6.75. Prior 5-row layout ended at y=6.79, overlap.
    left_col  = picks[:4]
    right_col = picks[4:8]
    col_w = 5.85
    col_x = [0.62, 6.90]
    row_h = 0.78
    top_y = 3.15

    for col_idx, col_rows in enumerate([left_col, right_col]):
        cx = col_x[col_idx]
        for i, r in enumerate(col_rows):
            y = top_y + i * row_h
            _text(slide, cx, y, col_w - 2.00, 0.34,
                  r["brand"], size=13, bold=True, color=C_CREAM)
            sub_line = r["cat"] if r["cat"] else ""
            if r["note"]:
                sub_line = (sub_line + " · " + r["note"]) if sub_line else r["note"]
            _text(slide, cx, y + 0.32, col_w - 2.00, 0.28,
                  sub_line, size=9.5, color=C_MUTED, line_spacing=1.3)

            pct_color = C_LAVENDER
            idx_color = (C_MAGENTA if r["index"] >= 200 else
                          (C_LIME if r["index"] >= 130 else C_CREAM))
            _text(slide, cx + col_w - 1.85, y + 0.06, 1.00, 0.44,
                  f"{r['pct']:.1f}%",
                  size=16, bold=True, color=pct_color,
                  align=PP_ALIGN.RIGHT)
            _text(slide, cx + col_w - 0.85, y + 0.06, 0.85, 0.44,
                  f"{r['index']}",
                  size=16, bold=True, color=idx_color,
                  align=PP_ALIGN.RIGHT)
            if i < len(col_rows) - 1:
                _hairline(slide, cx, y + row_h - 0.03, col_w - 0.10, C_STROKE)

    _text(slide, 0.62, 6.75, 12.0, 0.30, _source_line(p),
          size=8.5, color=C_MUTED2, letter_spacing=0.02)
    _footer(slide, 0, p["name"])


# =============================================================================
#  Ethnographic portrait slides
# =============================================================================
#  These are the "who they actually are" slides — the depth the marketer
#  needs beyond the exec summary. Every builder self-skips gracefully
#  when the data + brief can't support it, so the deck stays clean on
#  thin profiles and gets rich on data-dense ones.
# =============================================================================


def _slide_portrait(prs: Presentation, p: dict):
    """4-quadrant ethnographic portrait: LIFE STAGE / MINDSET / CULTURAL
    POSTURE / DIFFERENTIATOR. Each quadrant is a mini-narrative panel."""
    pt = _brief_get(p, "portrait", default={}) or {}
    quadrants: list[tuple[str, str, str, RGBColor]] = []

    def _q(kicker: str, hd_key: str, body_key: str,
           fb_hd: str, fb_body: str, accent: RGBColor) -> None:
        hd = _purge_pronouns(str(pt.get(hd_key) or "").strip(), p["name"]) or fb_hd
        bd = _purge_pronouns(str(pt.get(body_key) or "").strip(), p["name"]) or fb_body
        quadrants.append((kicker, hd, bd, accent))

    # Data-driven fallbacks (each quadrant reads real numbers if the LLM
    # didn't return a brief — the deck must still be substantive).
    age = (p.get("demographics") or {}).get("age") or {}
    inc = (p.get("demographics") or {}).get("income") or {}
    eth = (p.get("demographics") or {}).get("ethnicity") or {}
    top_age = max(age.items(), key=lambda kv: _num(kv[1]), default=("", 0))
    hh_income_75plus = sum(_num(v) for k, v in inc.items()
                            if any(t in k for t in ("$75", "$100", "$150", "$200")))

    _q("LIFE STAGE", "life_stage_headline", "life_stage_body",
       f"Prime-earning {top_age[0]}." if top_age[0] else "Prime-earning audience.",
       (f"The core sits in {top_age[0]} ({_num(top_age[1]):.0f}% of the file). "
        f"About {hh_income_75plus:.0f}% earn $75K+, which sets the ceiling on "
        "premium-partner viability and the willingness to pay for direct-to-"
        "consumer subscription products."),
       C_LAVENDER)

    top_platform = ""
    top_platform_pct = 0.0
    for cat_names in [("SOCIAL MEDIA",), ("STREAMING/PLATFORM", "STREAMING VIDEO")]:
        items = _first_available(p, list(cat_names))
        if items:
            top_platform = (items[0].get("name") or "").strip()
            top_platform_pct = _num(items[0].get("pct", 0))
            break
    _q("MINDSET", "mindset_headline", "mindset_body",
       "Platform-native. Digitally-fluent.",
       (f"Their consumption posture is digital-first — {top_platform or 'their top platform'}"
        f" reaches {top_platform_pct:.0f}% of the audience vs. Gen Pop far below. "
        "That is a cohort that finds culture through algorithmic feed, "
        "not appointment TV or linear broadcast."),
       C_LIME)

    top_signals = sorted(
        _all_behavioral_items(p),
        key=lambda it: -(_num(it.get("pct", 0)) * max(1.0, _num(it.get("index", 0)) / 100.0)),
    )[:4]
    top_sig_names = ", ".join(
        f"{(it.get('name') or '').strip()} ({_num(it.get('pct', 0)):.0f}%)"
        for it in top_signals[:3] if (it.get("name") or "").strip()
    )
    _q("CULTURAL POSTURE", "cultural_posture_headline", "cultural_posture_body",
       "Culture-first early adopters.",
       (f"The brands and platforms that concentrate here — {top_sig_names or 'their top affinities'} — "
        "describe an audience that spends culture time on newer, "
        "creator-native surfaces. Traditional-media placements will "
        "under-perform this cohort at any given CPM."),
       C_MAGENTA)

    peak_idx = 0
    peak_bucket = ""
    peak_dim = ""
    for dim in ("age", "ethnicity", "income", "education"):
        idx_map = ((p.get("demographics_index") or {}).get(dim)) or {}
        for bucket, val in idx_map.items():
            if _num(val) > peak_idx:
                peak_idx = _num(val)
                peak_bucket = str(bucket)
                peak_dim = dim
    _q("DIFFERENTIATOR", "differentiator_headline", "differentiator_body",
       f"Sharpest edge: {peak_bucket}." if peak_bucket else "Sharpest edge vs Gen Pop.",
       (f"On {peak_dim}, the {peak_bucket} bucket over-indexes at {int(peak_idx)} — "
        f"{peak_idx/100:.1f}x Gen Pop. That single delta is the reason "
        "endemic-partner categories that skew to this bucket will "
        "outperform any broad-reach media buy on this file."),
       C_BLUE)

    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], "PORTRAIT")
    _text(slide, 0.62, 0.82, 12.0, 1.10, "Who they actually are.",
          size=32, bold=True, color=C_CREAM, line_spacing=1.05)
    _text(slide, 0.62, 2.05, 12.0, 0.36,
          "Four dimensions of the audience portrait, grounded in the panel.",
          size=11, color=C_MUTED, line_spacing=1.4)

    # 2x2 grid of narrative panels
    cell_w = 5.95
    cell_h = 2.10
    x0, y0 = 0.62, 2.65
    gutter = 0.20
    for i, (kicker, hd, bd, accent) in enumerate(quadrants[:4]):
        col = i % 2
        row = i // 2
        x = x0 + col * (cell_w + gutter)
        y = y0 + row * (cell_h + 0.18)
        _rect(slide, x, y, cell_w, cell_h, RGBColor(0x14, 0x1F, 0x22))
        _text(slide, x + 0.28, y + 0.18, cell_w - 0.50, 0.28, kicker,
              size=9.5, bold=True, color=accent, letter_spacing=0.14)
        _text(slide, x + 0.28, y + 0.50, cell_w - 0.50, 0.44, hd,
              size=16, bold=True, color=C_CREAM, line_spacing=1.08)
        _text(slide, x + 0.28, y + 1.02, cell_w - 0.50, cell_h - 1.12, bd,
              size=10.5, color=C_MUTED, line_spacing=1.42)

    _text(slide, 0.62, 6.75, 12.0, 0.30, _source_line(p),
          size=8.5, color=C_MUTED2, letter_spacing=0.02)
    _footer(slide, 0, p["name"])


def _slide_demographic_deep_dive(prs: Presentation, p: dict):
    """All demographic dimensions on one page with proportional mini-bars.
    No LLM needed — pure data visualization of the full demo shape."""
    demos = p.get("demographics") or {}
    gp    = p.get("demographics_gen_pop") or {}
    idxs  = p.get("demographics_index") or {}

    dims: list[tuple[str, dict, dict, dict]] = []
    for key, label in [
        ("age",       "AGE"),
        ("gender",    "GENDER"),
        ("ethnicity", "ETHNICITY"),
        ("income",    "INCOME"),
        ("education", "EDUCATION"),
    ]:
        d = demos.get(key) or {}
        if len(d) >= 2:
            dims.append((label, d, gp.get(key) or {}, idxs.get(key) or {}))
    if len(dims) < 3:
        return  # thin demo — skip; other slides cover it

    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], "DEMOGRAPHIC DEEP-DIVE")
    _text(slide, 0.62, 0.82, 12.0, 1.10, "The full demographic shape.",
          size=32, bold=True, color=C_CREAM, line_spacing=1.05)
    _text(slide, 0.62, 2.05, 12.0, 0.36,
          "Every dimension side-by-side with Gen Pop. Bar length = share "
          "of audience; the value on the right is the index vs U.S.",
          size=11, color=C_MUTED, line_spacing=1.4)

    # Layout: up to 5 dimensions, 2 columns
    show = dims[:6]
    col_w = 5.95
    x_cols = [0.62, 6.90]
    y_slots = [2.85, 4.55, 6.25]
    # We only have ~3.85 vertical inches for content between title area
    # and source line — so cap at 4 dimensions in a 2x2 grid.
    positions = [(0, 0), (0, 1), (1, 0), (1, 1)]
    for i, (label, d, d_gp, d_idx) in enumerate(show[:4]):
        col, row = positions[i][1], positions[i][0]
        x = x_cols[col]
        y = 2.85 + row * 1.95

        _text(slide, x, y, col_w, 0.28, label,
              size=10.5, bold=True, color=C_LAVENDER, letter_spacing=0.10)
        _hairline(slide, x, y + 0.32, col_w - 0.10, C_STROKE)

        # Sort buckets by canonical order if age/income, else by pct desc
        entries = sorted(d.items(), key=lambda kv: -_num(kv[1]))[:6]
        if label in ("AGE", "INCOME", "EDUCATION"):
            # Keep original data order (assumed ordinal) if provided
            entries = list(d.items())[:6]

        max_pct = max((_num(v) for _, v in entries), default=1.0)
        row_h = min(0.24, 1.45 / max(1, len(entries)))
        by = y + 0.42
        for bucket, val in entries:
            v = _num(val)
            width = 2.60 * (v / max_pct if max_pct > 0 else 0)
            idx = int(_num(d_idx.get(bucket, 0)))
            # Label
            _text(slide, x, by, 2.10, row_h,
                  _short_bucket_label(str(bucket)),
                  size=8.5, color=C_CREAM, anchor=MSO_ANCHOR.MIDDLE)
            # Bar
            _rect(slide, x + 2.15, by + 0.04, max(0.03, width),
                  row_h - 0.10,
                  C_MAGENTA if idx >= 130 else (C_LAVENDER if idx >= 100 else C_STROKE))
            # Pct value
            _text(slide, x + 4.85, by, 0.60, row_h, f"{v:.0f}%",
                  size=8.5, color=C_CREAM, align=PP_ALIGN.RIGHT,
                  anchor=MSO_ANCHOR.MIDDLE)
            # Index
            idx_color = (C_MAGENTA if idx >= 130 else
                          (C_LIME if idx >= 110 else
                           (C_MUTED if idx >= 90 else C_STROKE)))
            _text(slide, x + 5.45, by, 0.45, row_h, f"{idx}" if idx else "—",
                  size=8.5, bold=True, color=idx_color, align=PP_ALIGN.RIGHT,
                  anchor=MSO_ANCHOR.MIDDLE)
            by += row_h + 0.02

    _text(slide, 0.62, 6.75, 12.0, 0.30, _source_line(p),
          size=8.5, color=C_MUTED2, letter_spacing=0.02)
    _footer(slide, 0, p["name"])


def _slide_life_stage(prs: Presentation, p: dict):
    """Household + life-stage narrative. Ships when parental / relationship
    / education / income has real richness OR the LLM returned a life_stage
    block."""
    ls = _brief_get(p, "life_stage", default={}) or {}
    hd = _purge_pronouns(str(ls.get("headline") or "").strip(), p["name"])
    body = _purge_pronouns(str(ls.get("body") or "").strip(), p["name"])

    demos = p.get("demographics") or {}
    has_hh_data = any(len(demos.get(k) or {}) >= 2
                       for k in ("parental_status", "relationship", "education"))

    if not (hd and body) and not has_hh_data:
        return  # nothing meaningful to say

    if not hd or not body:
        # Data-driven fallback narrative
        top_parental = max(((demos.get("parental_status") or {}).items()),
                            key=lambda kv: _num(kv[1]),
                            default=("", 0))
        top_rel = max(((demos.get("relationship") or {}).items()),
                       key=lambda kv: _num(kv[1]),
                       default=("", 0))
        top_edu = max(((demos.get("education") or {}).items()),
                       key=lambda kv: _num(kv[1]),
                       default=("", 0))
        parts = []
        if top_parental[0]:
            parts.append(
                f"{_num(top_parental[1]):.0f}% describe their household as "
                f"{str(top_parental[0]).lower()}"
            )
        if top_rel[0]:
            parts.append(
                f"{_num(top_rel[1]):.0f}% report they are {str(top_rel[0]).lower()}"
            )
        if top_edu[0]:
            parts.append(
                f"{_num(top_edu[1]):.0f}% have {str(top_edu[0]).lower()}"
            )
        if not hd:
            hd = "The household context that drives the buy."
        if not body:
            body = (
                ". ".join(parts).capitalize() +
                ". This is the household shape a marketer plans against — "
                "who is in the buying decision, what daypart they open the "
                "phone, and what the messaging tone should be."
            ) if parts else ""

    if not hd or not body:
        return

    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], "LIFE STAGE & HOUSEHOLD")
    _text(slide, 0.62, 0.82, 12.0, 1.10, hd,
          size=32, bold=True, color=C_CREAM, line_spacing=1.05)

    # Big narrative body panel
    _rect(slide, 0.62, 2.30, 12.0, 4.20, RGBColor(0x14, 0x1F, 0x22))
    _text(slide, 0.90, 2.50, 11.5, 0.30, "HOUSEHOLD PORTRAIT",
          size=10, bold=True, color=C_LIME, letter_spacing=0.12)
    _text(slide, 0.90, 2.90, 11.5, 3.40, body,
          size=14, color=C_CREAM, line_spacing=1.55)

    _text(slide, 0.62, 6.75, 12.0, 0.30, _source_line(p),
          size=8.5, color=C_MUTED2, letter_spacing=0.02)
    _footer(slide, 0, p["name"])


def _slide_differ_from_genpop(prs: Presentation, p: dict):
    """The 6-8 sharpest deltas across ALL demo dimensions. Named by
    dimension so the CMO can see: 'ah, they differ most on ethnicity
    and age, not on income.'"""
    brief_rows = _brief_get(p, "differ_from_genpop", default=[]) or []
    picks: list[dict] = []
    if isinstance(brief_rows, list):
        for row in brief_rows[:8]:
            if not isinstance(row, dict):
                continue
            dim = str(row.get("dimension") or "").strip()
            bkt = str(row.get("bucket")    or "").strip()
            aud = _num(row.get("aud_pct", 0))
            gpv = _num(row.get("gp_pct",  0))
            idx = int(_num(row.get("index", 0)))
            note = _purge_pronouns(str(row.get("note") or "").strip(), p["name"])
            if dim and bkt and idx > 0:
                picks.append({"dim": dim, "bucket": bkt, "aud_pct": aud,
                               "gp_pct": gpv, "index": idx, "note": note})

    # Fallback: compute the sharpest deltas from demographics_index across
    # every dimension.
    if len(picks) < 5:
        idxs = p.get("demographics_index") or {}
        demos = p.get("demographics") or {}
        gp   = p.get("demographics_gen_pop") or {}
        scored = []
        for dim, dim_idx_map in idxs.items():
            if not isinstance(dim_idx_map, dict):
                continue
            for bucket, idx_val in dim_idx_map.items():
                idx = int(_num(idx_val))
                if idx <= 0:
                    continue
                if 90 <= idx <= 110:
                    continue  # not sharp enough
                aud = _num((demos.get(dim) or {}).get(bucket, 0))
                gpv = _num((gp.get(dim) or {}).get(bucket, 0))
                scored.append((abs(idx - 100), dim, bucket, aud, gpv, idx))
        scored.sort(key=lambda t: -t[0])
        seen_keys = {(r["dim"].lower(), r["bucket"].lower()) for r in picks}
        for _, dim, bucket, aud, gpv, idx in scored:
            if len(picks) >= 7:
                break
            key = (dim.lower(), bucket.lower())
            if key in seen_keys:
                continue
            picks.append({"dim": dim.title(), "bucket": bucket, "aud_pct": aud,
                           "gp_pct": gpv, "index": idx, "note": ""})
            seen_keys.add(key)

    if len(picks) < 3:
        return

    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], "HOW THEY DIFFER FROM GEN POP")
    _text(slide, 0.62, 0.82, 12.0, 1.10, "The sharpest deltas.",
          size=32, bold=True, color=C_CREAM, line_spacing=1.05)
    _text(slide, 0.62, 2.05, 12.0, 0.36,
          "Bucket-level differences vs. the U.S. adult baseline, ranked "
          "by absolute deviation. The 'index' is 100 = parity.",
          size=11, color=C_MUTED, line_spacing=1.4)

    # Column headers
    hdr_y = 2.90
    _text(slide, 0.62, hdr_y, 1.80, 0.28, "DIMENSION",
          size=9, bold=True, color=C_LAVENDER, letter_spacing=0.08)
    _text(slide, 2.55, hdr_y, 3.20, 0.28, "BUCKET",
          size=9, bold=True, color=C_LAVENDER, letter_spacing=0.08)
    _text(slide, 5.90, hdr_y, 1.10, 0.28, "AUD %",
          size=9, bold=True, color=C_LAVENDER, letter_spacing=0.08,
          align=PP_ALIGN.RIGHT)
    _text(slide, 7.10, hdr_y, 1.10, 0.28, "GEN POP %",
          size=9, bold=True, color=C_LAVENDER, letter_spacing=0.08,
          align=PP_ALIGN.RIGHT)
    _text(slide, 8.30, hdr_y, 0.90, 0.28, "INDEX",
          size=9, bold=True, color=C_LAVENDER, letter_spacing=0.08,
          align=PP_ALIGN.RIGHT)
    _text(slide, 9.35, hdr_y, 3.30, 0.28, "SO WHAT",
          size=9, bold=True, color=C_LAVENDER, letter_spacing=0.08)
    _hairline(slide, 0.62, hdr_y + 0.32, 12.0, C_STROKE)

    row_h = min(0.44, 3.15 / max(1, len(picks)))
    y = hdr_y + 0.44
    for i, r in enumerate(picks[:7]):
        _text(slide, 0.62, y, 1.80, row_h, r["dim"],
              size=11, bold=True, color=C_CREAM, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, 2.55, y, 3.20, row_h, r["bucket"],
              size=11, color=C_CREAM, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, 5.90, y, 1.10, row_h,
              f"{r['aud_pct']:.1f}%" if r["aud_pct"] else "—",
              size=11, color=C_CREAM, align=PP_ALIGN.RIGHT,
              anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, 7.10, y, 1.10, row_h,
              f"{r['gp_pct']:.1f}%" if r["gp_pct"] else "—",
              size=11, color=C_MUTED, align=PP_ALIGN.RIGHT,
              anchor=MSO_ANCHOR.MIDDLE)
        idx_color = (C_MAGENTA if r["index"] >= 150 else
                      (C_LIME if r["index"] >= 110 else
                       (C_LAVENDER if r["index"] < 80 else C_MUTED)))
        _text(slide, 8.30, y, 0.90, row_h, f"{r['index']}",
              size=13, bold=True, color=idx_color,
              align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, 9.35, y, 3.30, row_h, r["note"] or "",
              size=9.5, color=C_MUTED, anchor=MSO_ANCHOR.MIDDLE,
              line_spacing=1.30)
        if i < len(picks) - 1:
            _hairline(slide, 0.62, y + row_h, 12.0, C_STROKE)
        y += row_h + 0.02

    _text(slide, 0.62, 6.75, 12.0, 0.30, _source_line(p),
          size=8.5, color=C_MUTED2, letter_spacing=0.02)
    _footer(slide, 0, p["name"])


def _slide_commerce_persona(prs: Presentation, p: dict):
    """Cross-category shopping-personality narrative. Synthesizes MPB +
    RETAIL + COSMETICS + APPAREL + PERSONAL CARE into one prose portrait.

    STRICT reach-vs-index ordering: the narrative must LEAD with the
    highest-REACH brands (the mass-shopping backbone) and only mention
    high-index-low-reach picks in a second clause as SIGNATURE affinities.
    Never let a <20% reach brand anchor the story — that is analytically
    false because most of the audience does not shop that brand.
    """
    cp = _brief_get(p, "commerce_persona", default={}) or {}
    hd = _purge_pronouns(str(cp.get("headline") or "").strip(), p["name"])
    body = _purge_pronouns(str(cp.get("body") or "").strip(), p["name"])

    # Quality gate 1: if the LLM body has no data grounding (no brand+%
    # or brand+idx mentions), drop it so the fallback fires. Prevents
    # the "value and convenience" generic-slop leak.
    if body and not _body_has_data_grounding(body, min_mentions=2):
        _log(f"[commerce_persona] LLM body has no data grounding, using "
              f"fallback: {body[:120]!r}")
        body = ""
    # Quality gate 2: if the LLM body cites any sub-12% reach brand or
    # any >400 idx brand, drop it — the user flagged 'Killstar (9.9%,
    # idx 576)' and 'Urban Stems (7.9%, idx 399)' leaking into the
    # narrative even though the overall body had grounding.
    if body:
        has_banned, offending = _body_cites_banned_brand(body)
        if has_banned:
            _log(f"[commerce_persona] LLM body cites banned NOISE brands, "
                  f"using fallback: {offending}")
            body = ""

    # Collect commerce-side brands across MPB + RETAIL + APPAREL/FOOTWEAR
    # + COSMETICS + PERSONAL CARE so mass/signature split has enough surface.
    commerce_cats = ["MPB", "MOST PURCHASED BRANDS", "MOST_PURCHASED_BRANDS",
                     "RETAIL", "APPAREL/FOOTWEAR", "APPAREL_FOOTWEAR",
                     "APPAREL", "COSMETICS", "PERSONAL CARE",
                     "PERSONAL_CARE"]
    commerce_items: list[dict] = []
    seen_names: set[str] = set()
    for cat in commerce_cats:
        for it in (_first_available(p, [cat]) or []):
            nm = (it.get("name") or "").strip()
            key = nm.upper()
            if not nm or key in seen_names:
                continue
            seen_names.add(key)
            commerce_items.append(it)
    if len(commerce_items) < 4 and not (hd and body):
        return
    mass, sig = _split_by_reach(commerce_items, reach_threshold=20.0)
    mass_labels = [f"{(it.get('name') or '').strip()} "
                    f"({_num(it.get('pct',0)):.0f}%)"
                    for it in mass[:4]]
    sig_labels = [f"{(it.get('name') or '').strip()} "
                   f"(idx {int(round(_num(it.get('index',0))))})"
                   for it in sig[:3]]

    # If the LLM body leads with signature brands, repair it by prepending
    # a mass-brand anchor sentence. Catches the "Urban Stems (7.9%)" defect.
    if body and mass_labels:
        body = _repair_leading_signature(body, mass, sig)

    if not hd:
        if mass:
            top_reach = _num(mass[0].get("pct", 0))
            top_brand = (mass[0].get("name") or "").strip()
            hd = f"Mass-brand shoppers with signature over-indexes."
        else:
            hd = "Signature-brand shoppers, not mass buyers."
    if not body:
        if mass_labels and sig_labels:
            body = (
                f"Mass-brand commerce is anchored by "
                f"{', '.join(mass_labels[:3])}"
                + (f" and {mass_labels[3]}" if len(mass_labels) >= 4 else "")
                + ". That is the shopping backbone — how the buy will "
                "actually land. Layered on top are signature affinities "
                f"toward {', '.join(sig_labels)} — sharp over-indexes vs. "
                "Gen Pop that give the audience its personality without "
                "carrying the media weight. For the plan, size the buy "
                "against the mass brands; use the signatures for partner "
                "selection, co-branded activations, and cultural signaling."
            )
        elif mass_labels:
            body = (
                f"Mass-brand commerce is anchored by "
                f"{', '.join(mass_labels[:4])}. That is the shopping "
                "backbone — table-stakes brands the audience uses like "
                "the rest of America. Without a distinctive over-index "
                "layer, this reads as a mass-market shopper: value-tier "
                "promiscuity, no dominant category loyalty, standard "
                "big-box + Amazon digital rails."
            )
        elif sig_labels:
            body = (
                f"No single mass-brand anchor emerges above 20% reach — "
                f"the commerce fingerprint is defined by signature "
                f"affinities toward {', '.join(sig_labels)}. Read this as "
                "a distributed shopper with sharp niche taste, not a "
                "big-box loyalist. Partner selection matters more than "
                "media weight here."
            )
        else:
            return

    if not hd or not body:
        return

    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], "COMMERCE PERSONA")
    _text(slide, 0.62, 0.82, 12.0, 1.10, "How they shop.",
          size=32, bold=True, color=C_CREAM, line_spacing=1.05)
    _text(slide, 0.62, 2.05, 12.0, 0.36,
          "Cross-category shopping personality, synthesized across MPB, "
          "apparel, personal care, and retail signals.",
          size=11, color=C_MUTED, line_spacing=1.4)

    # Panel with headline + body
    _rect(slide, 0.62, 2.80, 12.0, 3.70, RGBColor(0x14, 0x1F, 0x22))
    _text(slide, 0.90, 3.00, 11.5, 0.30, "THE SHOPPING PERSONALITY",
          size=10, bold=True, color=C_LIME, letter_spacing=0.12)
    _text(slide, 0.90, 3.40, 11.5, 0.80, hd,
          size=22, bold=True, color=C_CREAM, line_spacing=1.10)
    _text(slide, 0.90, 4.35, 11.5, 2.05, body,
          size=13, color=C_CREAM, line_spacing=1.55)

    _text(slide, 0.62, 6.75, 12.0, 0.30, _source_line(p),
          size=8.5, color=C_MUTED2, letter_spacing=0.02)
    _footer(slide, 0, p["name"])


def _slide_cultural_interests(prs: Presentation, p: dict):
    """Integrated cultural-time narrative pulling from MUSIC / GAMES /
    FESTIVALS / TICKETING / SPORTS TEAM / STREAMING signals."""
    ci = _brief_get(p, "cultural_interests", default={}) or {}
    hd = _purge_pronouns(str(ci.get("headline") or "").strip(), p["name"])
    body = _purge_pronouns(str(ci.get("body") or "").strip(), p["name"])

    # Quality gate 1: kill generic LLM bodies with no brand names / numbers.
    if body and not _body_has_data_grounding(body, min_mentions=2):
        _log(f"[cultural_interests] LLM body has no data grounding, using "
              f"fallback: {body[:120]!r}")
        body = ""
    # Quality gate 2: reject bodies citing NOISE-class brands.
    if body:
        has_banned, offending = _body_cites_banned_brand(body)
        if has_banned:
            _log(f"[cultural_interests] LLM body cites banned NOISE brands, "
                  f"using fallback: {offending}")
            body = ""

    # Fallback: check whether the profile has enough cultural-time
    # categories to justify the slide.
    beh = p.get("behavioral") or {}
    cultural_cats = ["MUSIC", "GAMES", "GAMING", "FESTIVALS",
                      "TICKETING PLATFORMS", "SPORTS TEAM",
                      "STREAMING VIDEO", "STREAMING/PLATFORM"]
    cat_hits = [c for c in cultural_cats
                if any(c.upper() == k.upper() and v for k, v in beh.items())]
    if not (hd and body) and len(cat_hits) < 2:
        return

    # Collect cultural-side brands with dedupe so mass/signature split
    # sees the full cross-category surface.
    cultural_items: list[dict] = []
    seen_names: set[str] = set()
    for cat in cultural_cats:
        for it in (_first_available(p, [cat]) or []):
            nm = (it.get("name") or "").strip()
            key = nm.upper()
            if not nm or key in seen_names:
                continue
            seen_names.add(key)
            cultural_items.append(it)
    # Cultural surfaces skew smaller (Spotify 40% is 'big' here) so relax
    # the mass threshold from 20% to 15%.
    mass, sig = _split_by_reach(cultural_items, reach_threshold=15.0)
    mass_labels = [f"{(it.get('name') or '').strip()} "
                    f"({_num(it.get('pct',0)):.0f}%)"
                    for it in mass[:4]]
    sig_labels = [f"{(it.get('name') or '').strip()} "
                   f"(idx {int(round(_num(it.get('index',0))))})"
                   for it in sig[:3]]

    if body and (mass or sig):
        body = _repair_leading_signature(body, mass, sig,
                                         reach_threshold=15.0)

    if not hd:
        hd = "Where they spend cultural time."
    if not body:
        if mass_labels and sig_labels:
            body = (
                f"Cultural time concentrates in {', '.join(mass_labels[:3])}"
                + (f" and {mass_labels[3]}" if len(mass_labels) >= 4 else "")
                + ". Those are the surfaces where the audience actually "
                "spends time and where sponsor properties will earn "
                f"attention. Beyond that scale layer, signature interests "
                f"in {', '.join(sig_labels)} give the audience its "
                "distinctive cultural fingerprint — sharp over-indexes "
                "vs. Gen Pop that make for authentic cultural partners "
                "without carrying the reach numbers themselves."
            )
        elif mass_labels:
            body = (
                f"Cultural time concentrates in {', '.join(mass_labels[:4])}. "
                "Those are the surfaces where sponsor properties will earn "
                "attention. The distribution is mass-shaped rather than "
                "signature-shaped — activation should ride the reach here, "
                "not chase niche indexes."
            )
        elif sig_labels:
            body = (
                f"No dominant mass cultural surface anchors the audience — "
                f"cultural interest is defined by signature affinities: "
                f"{', '.join(sig_labels)}. This is a distributed cultural "
                "consumer with sharp niche taste; sponsor selection matters "
                "more than sheer reach when picking properties."
            )
        else:
            return

    if not hd or not body:
        return

    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], "CULTURAL INTERESTS")
    _text(slide, 0.62, 0.82, 12.0, 1.10, "How they spend cultural time.",
          size=32, bold=True, color=C_CREAM, line_spacing=1.05)
    _text(slide, 0.62, 2.05, 12.0, 0.36,
          "Music, games, sports, festivals, streaming — the integrated "
          "cultural-time portrait for this audience.",
          size=11, color=C_MUTED, line_spacing=1.4)

    _rect(slide, 0.62, 2.80, 12.0, 3.70, RGBColor(0x14, 0x1F, 0x22))
    _text(slide, 0.90, 3.00, 11.5, 0.30, "THE CULTURAL PROFILE",
          size=10, bold=True, color=C_MAGENTA, letter_spacing=0.12)
    _text(slide, 0.90, 3.40, 11.5, 0.80, hd,
          size=22, bold=True, color=C_CREAM, line_spacing=1.10)
    _text(slide, 0.90, 4.35, 11.5, 2.05, body,
          size=13, color=C_CREAM, line_spacing=1.55)

    _text(slide, 0.62, 6.75, 12.0, 0.30, _source_line(p),
          size=8.5, color=C_MUTED2, letter_spacing=0.02)
    _footer(slide, 0, p["name"])


def _slide_methodology(prs: Presentation, p: dict):
    """KD slide 23 pattern: 'How this brief was built.' — 4 subsections
    (Sample, Measurement, Benchmark, Reading the numbers)."""
    slide = _blank_slide(prs)
    _section_eyebrow(slide, 0.52, p["name"], "METHODOLOGY")
    _text(slide, 0.62, 0.82, 12.0, 1.10, "How this brief was built.",
          size=28, bold=True, color=C_CREAM, line_spacing=1.05)

    sample = p["sample_size"]
    proj   = p["projected_us"]
    dr     = p.get("date_range") or ""
    subj   = p["name"]

    sections = [
        ("SAMPLE",
         (f"{_fmt_int(sample) if sample else 'The panel'} zero-party opted-in "
          f"U.S. panelists who exhibited {subj} engagement behavior "
          f"{('between ' + dr) if dr else 'in the observation window'}. "
          "Census-based panel of 10M respondents projected to a total U.S. "
          "digital audience of "
          f"{_fmt_int(proj) if proj else '329.9M adults'}.\n\n"
          "Engagement vectors include: Search, Social, Owned Channels, "
          "Commerce, and Media.")),
        ("MEASUREMENT",
         ("Crosswalk Profile IQ. Behavioral observations are deterministic and "
          "observed across the open web at the individual panelist level. "
          "Demographic data is self-reported at registration.")),
        ("BENCHMARK",
         ("Comparison to a General Population (Gen. Pop.) sample projected to "
          "329.9M U.S. adults, drawn from the same observation window.\n\n"
          "Index = (audience share) / (Gen Pop share) x 100. "
          "Index of 100 = parity (national norm).")),
        ("READING THE NUMBERS",
         ("Charts display talent/brand/category percent penetration - the % "
          "of the audience exhibiting that behavior. Every panelist voluntarily "
          "opted in. Clean consent lineage. Audit-ready provenance.")),
    ]

    col_w = 5.90
    col_gap = 0.30
    row_h = 2.05
    positions = [
        (0.62, 2.30),
        (0.62 + col_w + col_gap, 2.30),
        (0.62, 2.30 + row_h + 0.20),
        (0.62 + col_w + col_gap, 2.30 + row_h + 0.20),
    ]
    for (title, body), (x, y) in zip(sections, positions):
        _text(slide, x, y, col_w, 0.32, title,
              size=11, bold=True, color=C_LAVENDER, letter_spacing=0.10)
        _text(slide, x, y + 0.38, col_w, row_h - 0.42, body,
              size=10.5, color=C_MUTED, line_spacing=1.4)

    _footer(slide, 99, p["name"])


def _slide_about(prs: Presentation, p: dict):
    """KD slide 24 pattern: 'Behavioral intelligence at scale.' — 4
    Crosswalk capability stats (events, panelists, freshness, archive)."""
    slide = _blank_slide(prs, bg=C_DARK)
    _text(slide, 0.62, 0.62, 12.0, 0.30, "ABOUT",
          size=10.5, color=C_LIME, letter_spacing=0.10)
    _text(slide, 0.62, 0.98, 12.0, 1.20, "Behavioral intelligence at scale.",
          size=36, bold=True, color=C_CREAM, line_spacing=1.05)
    _text(slide, 0.62, 2.55, 12.0, 0.40,
          "We don't infer behavior. We observe it.",
          size=15, color=C_MUTED, line_spacing=1.4)

    stats = [
        ("8B+", "DAILY BEHAVIORAL EVENTS",
         "Real-time observation across the open web — purchase, attention, "
         "subscription, media."),
        ("30M", "PERMISSIONED PANELISTS",
         "Zero-party consented individuals, longitudinally observed for up to "
         "five years. Demographic-grounded. 10M census-based Gen. Pop. subset."),
        ("T+1", "DAILY DATA",
         "Yesterday's behavior, not last year's. Live attention migration and "
         "consumer shifts available at enterprise scale."),
        ("ARCHIVE", "ACTIVE HISTORICAL",
         "Unprecedented access to historical data, anytime."),
    ]
    n = len(stats)
    slot_w = (SLIDE_W_IN - 1.24 - 0.15 * (n - 1)) / n
    x = 0.62
    for value, label, note in stats:
        _text(slide, x, 3.60, slot_w, 1.20, value,
              size=56, bold=True, color=C_LAVENDER, line_spacing=1.0)
        _text(slide, x, 4.85, slot_w, 0.30, label,
              size=10, bold=True, color=C_CREAM, letter_spacing=0.10)
        _text(slide, x, 5.22, slot_w, 1.30, note,
              size=10.5, color=C_MUTED, line_spacing=1.45)
        x += slot_w + 0.15

    _footer(slide, 99, p["name"])


# =============================================================================
#  Slide builders  ·  COMBINED / COMPARATIVE
# =============================================================================


def _slide_combined_cover(prs: Presentation, profiles: list[dict]):
    """Comparison cover: right-side vertical image collage of the profiles,
    left side big two-line comparison title + editorial tagline. Falls back
    to placeholder gradient if none of the profiles resolved images."""
    slide = _blank_slide(prs, bg=C_DARK)

    # Right-side image column split evenly by profile count. Each image is
    # darkened + center-cropped to fill its slot.
    imgs = [pr for pr in profiles if pr.get("_image_bytes")]
    if imgs:
        slot_h = SLIDE_H_IN / len(imgs)
        for i, pr in enumerate(imgs):
            try:
                _place_prepared_image(
                    slide, pr["_image_bytes"],
                    x=8.05, y=i * slot_h,
                    w=SLIDE_W_IN - 8.05, h=slot_h,
                    darken=0.68,
                )
            except Exception:
                _rect(slide, 8.05, i * slot_h,
                      SLIDE_W_IN - 8.05, slot_h,
                      RGBColor(0x12, 0x1D, 0x1F))
    else:
        _rect(slide, 8.05, 0, SLIDE_W_IN - 8.05, SLIDE_H_IN,
              RGBColor(0x12, 0x1D, 0x1F))
    # Left-side dark scrim
    _rect(slide, 0, 0, 8.05, SLIDE_H_IN, C_DARK)

    year_label = datetime.now().strftime("%Y")
    _text(slide, 0.62, 0.62, 10.0, 0.30,
          f"A BEHAVIORAL AUDIENCE COMPARISON   \u00b7   {year_label}",
          size=10, color=C_MUTED, letter_spacing=0.06)

    if len(profiles) == 2:
        title_a = profiles[0]["name"] + ","
        title_b = f"vs. {profiles[1]['name']}."
    elif len(profiles) == 3:
        title_a = f"{profiles[0]['name']}, {profiles[1]['name']},"
        title_b = f"and {profiles[2]['name']}."
    else:
        title_a = f"{len(profiles)} audiences."
        title_b = "One dataset."
    _text(slide, 0.62, 2.60, 7.30, 1.30, title_a,
          size=44, bold=True, color=C_CREAM, line_spacing=1.05)
    _text(slide, 0.62, 3.80, 7.30, 1.30, title_b,
          size=44, bold=True, color=C_CREAM, line_spacing=1.05)

    subj_list = ", ".join(pr["name"] for pr in profiles[:4])
    tag = (
        f"Where do their fans overlap, and where do they diverge? "
        f"A side-by-side read of {subj_list} on scale, demographics, "
        "category signal, and the brands they share."
    )
    _text(slide, 0.62, 5.30, 7.0, 1.60, tag,
          size=12, color=C_MUTED, line_spacing=1.45)

    _text(slide, 0.62, 7.00, 4.0, 0.28, "CONFIDENTIAL",
          size=9, color=C_MUTED2, letter_spacing=0.06)
    _text(slide, 8.11, 7.00, 4.60, 0.28, "CROSSWALK  PROFILE IQ",
          size=9, color=C_MUTED2, align=PP_ALIGN.RIGHT, letter_spacing=0.06)


def _slide_combined_profile_spotlight(prs: Presentation, pr: dict):
    """One slide per profile in the combined deck. Shows the subject image,
    scale, top demo call-out, and the profile's #1 over-index brand. Gives
    the deck a rhythm of comparison -> spotlight -> comparison -> spotlight."""
    slide = _blank_slide(prs, bg=C_DARK)
    # Right-side subject image (like LISA identity)
    _place_side_image(slide, pr.get("_image_bytes"), x=8.05, y=0.0,
                      w=SLIDE_W_IN - 8.05, h=SLIDE_H_IN,
                      darken=0.72, subject_name=pr["name"])

    _text(slide, 0.62, 0.62, 12.0, 0.30,
          f"AUDIENCE SPOTLIGHT  \u00b7  {pr['name'].upper()}",
          size=10, color=C_LAVENDER, letter_spacing=0.10)
    _text(slide, 0.60, 0.98, 8.0, 1.20, pr["name"] + ".",
          size=40, bold=True, color=C_CREAM, line_spacing=1.05)

    # Three-stat mini "by the numbers"
    proj = pr["projected_us"] or 0
    top = _top_over_index_brand(pr)
    peak = _peak_index(pr)
    demos = pr.get("demographics") or {}
    prime = sum(_num(v) for k, v in demos.get("age", {}).items()
                if str(k).replace("–", "-") in ("18-24", "25-34", "35-44"))

    stats = [
        (_fmt_int_compact(proj) or "—", "U.S. DIGITAL AUDIENCE"),
        (f"{prime:.1f}%" if prime else "—", "AGE 18-44 SHARE"),
        (f"{int(peak)}" if peak else "—", "TOP INDEX vs GEN POP"),
    ]
    y = 2.65
    for value, label in stats:
        _text(slide, 0.62, y, 3.4, 1.10, value,
              size=44, bold=True, color=C_CREAM, line_spacing=1.0)
        _text(slide, 0.62, y + 1.10, 3.4, 0.30, label,
              size=9.5, color=C_MUTED, letter_spacing=0.10)
        y += 1.50

    if top:
        _text(slide, 4.60, 2.65, 3.3, 0.30, "STRONGEST BRAND SIGNAL",
              size=9.5, bold=True, color=C_LIME, letter_spacing=0.10)
        _text(slide, 4.60, 3.00, 3.3, 0.50, top["name"],
              size=22, bold=True, color=C_CREAM)
        _text(slide, 4.60, 3.55, 3.3, 0.30,
              f"{_num(top.get('pct', 0)):.1f}% penetration",
              size=11, color=C_CREAM)
        _text(slide, 4.60, 3.85, 3.3, 0.30,
              f"{int(_num(top.get('index', 0)))} index vs Gen Pop",
              size=11, bold=True, color=C_MAGENTA)
        _text(slide, 4.60, 4.25, 3.3, 0.30,
              f"Category: {_pretty_cat(top.get('_cat', ''))}",
              size=10, color=C_MUTED)

    _footer(slide, 99, "Comparison")


def _slide_combined_overview(prs: Presentation, profiles: list[dict]):
    slide = _blank_slide(prs)
    _text(slide, 0.62, 0.62, 12.0, 0.30,
          "AUDIENCE COMPARISON  ·  OVERVIEW",
          size=10, color=C_MUTED, letter_spacing=0.06)
    _title_block(slide, 1.0, "What's inside.")

    body = (
        f"We stacked {len(profiles)} audiences against each other and the U.S. baseline. "
        f"Every slide answers one question: who's biggest, who over-indexes with whom, and "
        f"which brands live in every fandom."
    )
    _text(slide, 0.62, 2.35, 8.0, 1.4, body, size=13, color=C_MUTED,
          line_spacing=1.45)

    chapters = [
        ("01", "Scale", "How large is each digital audience, side by side?"),
        ("02", "Demographics", "Age, gender, ethnicity, income (who's in each room)."),
        ("03", "Category winners", "For each behavioral category, which audience over-indexes most."),
        ("04", "Shared brands", "Which brands live in every audience: the connective tissue."),
    ]
    left_x, top_y = 0.62, 4.05
    card_w, card_h = 2.94, 2.65
    gap = 0.18
    for i, (num, title, subtitle) in enumerate(chapters):
        x = left_x + i * (card_w + gap)
        _rect(slide, x, top_y, card_w, card_h, RGBColor(0x14, 0x1F, 0x22))
        _text(slide, x + 0.32, top_y + 0.30, 2.4, 0.32, num,
              size=11, bold=True, color=C_LAVENDER, letter_spacing=0.05)
        _text(slide, x + 0.32, top_y + 0.72, 2.4, 1.0, title,
              size=17, bold=True, color=C_CREAM, line_spacing=1.15)
        _text(slide, x + 0.32, top_y + 1.65, 2.4, 0.95, subtitle,
              size=10.5, color=C_MUTED, line_spacing=1.35)

    _footer(slide, 2, "Comparison")


def _slide_combined_scale(prs: Presentation, profiles: list[dict]):
    slide = _blank_slide(prs)
    _text(slide, 0.62, 0.62, 12.0, 0.30,
          "AUDIENCE COMPARISON  ·  SCALE",
          size=10, color=C_MUTED, letter_spacing=0.06)
    _title_block(slide, 1.0, "How large is each audience.",
                 "Projected to U.S. digital adults.")

    _text(slide, 0.62, 2.35, 11.5, 0.6,
          "Every projection is the audience size that the Crosswalk panel scales to. "
          "Bars are proportional; index vs Gen Pop shown as a chip.",
          size=11, color=C_MUTED, line_spacing=1.4)

    # Horizontal bar chart
    max_proj = max((p["projected_us"] for p in profiles), default=0) or 1
    row_h = min(0.85, 3.4 / max(len(profiles), 1))
    y0 = 3.15
    label_w = 2.8
    bar_start = 0.62 + label_w
    bar_max_w = SLIDE_W_IN - bar_start - 2.6

    for i, p in enumerate(profiles):
        y = y0 + i * (row_h + 0.18)
        _text(slide, 0.62, y, label_w, row_h - 0.10,
              p["name"], size=15, bold=True, color=C_CREAM,
              anchor=MSO_ANCHOR.MIDDLE)
        proj = p["projected_us"] or 0
        pct_of_max = proj / max_proj
        color = CYCLE_ACCENT[i % len(CYCLE_ACCENT)]
        _rect(slide, bar_start, y + 0.05, max(bar_max_w * pct_of_max, 0.05),
              row_h - 0.20, color)
        _text(slide, bar_start + bar_max_w + 0.15, y + 0.05, 2.4, row_h - 0.20,
              _fmt_int_compact(proj) or "—",
              size=18, bold=True, color=C_CREAM,
              anchor=MSO_ANCHOR.MIDDLE)

    _footer(slide, 3, "Comparison")


def _slide_combined_demographics(prs: Presentation, profiles: list[dict]):
    slide = _blank_slide(prs)
    _text(slide, 0.62, 0.62, 12.0, 0.30,
          "AUDIENCE COMPARISON  ·  DEMOGRAPHICS",
          size=10, color=C_MUTED, letter_spacing=0.06)
    _title_block(slide, 1.0, "Who's in each room.")

    _text(slide, 0.62, 2.35, 11.5, 0.6,
          "Age and gender skew for every audience, side by side. Percentages are share of that audience.",
          size=11, color=C_MUTED, line_spacing=1.4)

    # Two panels: AGE (grouped bars per audience) and GENDER
    _combined_demo_panel(slide, 0.62, 3.05, 6.05, 3.9,
                         "AGE", profiles,
                         lambda p: p["demographics"].get("age", {}))
    _combined_demo_panel(slide, 6.72, 3.05, 6.05, 3.9,
                         "GENDER", profiles,
                         lambda p: p["demographics"].get("gender", {}))

    _footer(slide, 4, "Comparison")


def _combined_demo_panel(slide, x, y, w, h, label, profiles, getter):
    _text(slide, x, y, w, 0.30, label,
          size=10, bold=True, color=C_LAVENDER, letter_spacing=0.06)
    # Get union of buckets across profiles (preserve order from first)
    buckets = []
    seen = set()
    for p in profiles:
        for k in getter(p).keys():
            if k not in seen:
                seen.add(k)
                buckets.append(k)
    if not buckets:
        _text(slide, x, y + 0.4, w, 0.4, "No data.",
              size=10, color=C_MUTED)
        return
    if len(buckets) > 6:
        buckets = buckets[:6]

    # Grouped bars: one column per bucket, N side-by-side bars per column
    panel_top = y + 0.6
    panel_h = h - 0.9
    slot_w = w / len(buckets)
    n_prof = len(profiles)
    bar_group_w = slot_w * 0.68
    bar_w = bar_group_w / max(n_prof, 1)

    for bi, b in enumerate(buckets):
        cx = x + bi * slot_w
        # Bucket label
        _text(slide, cx, panel_top + panel_h + 0.02, slot_w - 0.05, 0.24,
              _short_bucket_label(b), size=8.5, color=C_MUTED,
              align=PP_ALIGN.CENTER)
        # Bars
        max_val = max(_num(getter(p).get(b, 0)) for p in profiles)
        max_val = max(max_val * 1.1, 1)
        for pi, p in enumerate(profiles):
            val = _num(getter(p).get(b, 0))
            bar_h = panel_h * (val / max_val)
            bx = cx + (slot_w - bar_group_w) / 2 + pi * bar_w
            color = CYCLE_ACCENT[pi % len(CYCLE_ACCENT)]
            _rect(slide, bx, panel_top + (panel_h - bar_h),
                  bar_w * 0.85, bar_h, color)

    # Mini legend under panel
    _draw_legend(slide, x, y + h - 0.30, profiles)


def _draw_legend(slide, x, y, profiles):
    chip_x = x
    for i, p in enumerate(profiles):
        color = CYCLE_ACCENT[i % len(CYCLE_ACCENT)]
        _rect(slide, chip_x, y + 0.06, 0.16, 0.14, color)
        # Estimate width from name length to avoid overlap
        est_w = 0.10 + 0.06 * len(p["name"])
        _text(slide, chip_x + 0.22, y, est_w + 0.4, 0.26,
              p["name"], size=9, color=C_CREAM)
        chip_x += 0.32 + est_w


def _slide_combined_category_winners(prs: Presentation, profiles: list[dict]):
    slide = _blank_slide(prs)
    _text(slide, 0.62, 0.62, 12.0, 0.30,
          "AUDIENCE COMPARISON  ·  CATEGORY WINNERS",
          size=10, color=C_MUTED, letter_spacing=0.06)
    _title_block(slide, 1.0, "Who owns each category.")

    _text(slide, 0.62, 2.35, 11.5, 0.6,
          "For each behavioral category, we identify which audience over-indexes hardest and what "
          "their #1 brand is. A dash means the audience has no signal in that category.",
          size=11, color=C_MUTED, line_spacing=1.45)

    # Get union of categories across profiles
    cats = _combined_category_ranking(profiles)[:8]
    if not cats:
        _text(slide, 0.62, 4.0, 11.5, 0.5,
              "No overlapping behavioral categories to compare.",
              size=12, color=C_MUTED)
        _footer(slide, 5, "Comparison")
        return

    # Table header row
    y = 3.15
    row_h = 0.42
    col_x = [0.62, 3.6, 6.8, 9.4]
    col_w = [2.9, 3.1, 2.5, 3.5]

    _text(slide, col_x[0], y, col_w[0], 0.28, "CATEGORY",
          size=9, bold=True, color=C_LAVENDER, letter_spacing=0.06)
    _text(slide, col_x[1], y, col_w[1], 0.28, "TOP-INDEX AUDIENCE",
          size=9, bold=True, color=C_LAVENDER, letter_spacing=0.06)
    _text(slide, col_x[2], y, col_w[2], 0.28, "INDEX",
          size=9, bold=True, color=C_LAVENDER, letter_spacing=0.06)
    _text(slide, col_x[3], y, col_w[3], 0.28, "THEIR TOP BRAND",
          size=9, bold=True, color=C_LAVENDER, letter_spacing=0.06)
    _hairline(slide, 0.62, y + 0.30, 12.10, C_STROKE)

    y += 0.42
    for cat_row in cats:
        cat, winner_name, winner_idx, brand = cat_row
        _text(slide, col_x[0], y, col_w[0], row_h - 0.10,
              cat.title(), size=12, color=C_CREAM,
              anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, col_x[1], y, col_w[1], row_h - 0.10,
              winner_name or "—", size=12, bold=True, color=C_CREAM,
              anchor=MSO_ANCHOR.MIDDLE)
        idx_color = C_MAGENTA if winner_idx >= 150 else C_LIME
        _text(slide, col_x[2], y, col_w[2], row_h - 0.10,
              f"{int(winner_idx)}" if winner_idx else "—",
              size=14, bold=True, color=idx_color,
              anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, col_x[3], y, col_w[3], row_h - 0.10,
              brand or "—", size=12, color=C_CREAM,
              anchor=MSO_ANCHOR.MIDDLE)
        _hairline(slide, 0.62, y + row_h - 0.05, 12.10, C_STROKE)
        y += row_h

    _footer(slide, 5, "Comparison")


def _slide_combined_brand_overlap(prs: Presentation, profiles: list[dict]):
    slide = _blank_slide(prs)
    _text(slide, 0.62, 0.62, 12.0, 0.30,
          "AUDIENCE COMPARISON  ·  SHARED BRANDS",
          size=10, color=C_MUTED, letter_spacing=0.06)
    _title_block(slide, 1.0, "The connective tissue.",
                 "Brands that live in every fandom.")

    _text(slide, 0.62, 2.35, 11.5, 0.6,
          "Brands where every profile in this deck posts a meaningful over-index (≥ 110) vs. Gen Pop. "
          "The higher the row, the tighter the shared behavior.",
          size=11, color=C_MUTED, line_spacing=1.4)

    overlap = _brand_overlap_rows(profiles, min_index=110, top_n=8)
    if not overlap:
        _text(slide, 0.62, 4.0, 11.5, 0.6,
              "No brands over-index across all profiles at the 110+ threshold. "
              "These audiences don't share a single common brand. That is a signal in itself.",
              size=13, color=C_MUTED, line_spacing=1.45)
        _footer(slide, 6, "Comparison")
        return

    # Header
    y = 3.15
    row_h = 0.42
    per_profile_w = 1.3
    n = len(profiles)
    start_x = SLIDE_W_IN - 0.62 - per_profile_w * n
    _text(slide, 0.62, y, 4.0, 0.28, "BRAND",
          size=9, bold=True, color=C_LAVENDER, letter_spacing=0.06)
    _text(slide, 4.7, y, 2.5, 0.28, "CATEGORY",
          size=9, bold=True, color=C_LAVENDER, letter_spacing=0.06)
    for i, p in enumerate(profiles):
        _text(slide, start_x + i * per_profile_w, y, per_profile_w - 0.1, 0.28,
              p["name"], size=9, bold=True, color=C_LAVENDER,
              align=PP_ALIGN.RIGHT, letter_spacing=0.05)
    _hairline(slide, 0.62, y + 0.30, 12.10, C_STROKE)

    y += 0.42
    for row in overlap:
        _text(slide, 0.62, y, 4.0, row_h - 0.10, row["brand"],
              size=12, bold=True, color=C_CREAM, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, 4.7, y, 2.5, row_h - 0.10, row["category"].title(),
              size=11, color=C_MUTED, anchor=MSO_ANCHOR.MIDDLE)
        for i, p in enumerate(profiles):
            idx = row["indices"].get(p["name"], 0)
            color = C_MAGENTA if idx >= 150 else (C_LIME if idx >= 120 else C_CREAM)
            _text(slide, start_x + i * per_profile_w, y,
                  per_profile_w - 0.1, row_h - 0.10,
                  f"{int(idx)}" if idx else "—",
                  size=14, bold=True, color=color, align=PP_ALIGN.RIGHT,
                  anchor=MSO_ANCHOR.MIDDLE)
        _hairline(slide, 0.62, y + row_h - 0.05, 12.10, C_STROKE)
        y += row_h

    _footer(slide, 6, "Comparison")


def _slide_combined_final_insight(prs: Presentation, profiles: list[dict]):
    slide = _blank_slide(prs)
    _text(slide, 0.62, 0.62, 12.0, 0.30,
          "AUDIENCE COMPARISON  ·  THE FINAL INSIGHT",
          size=10, color=C_MUTED, letter_spacing=0.06)

    largest = max(profiles, key=lambda p: p["projected_us"] or 0)
    tightest = max(profiles, key=lambda p: _peak_index(p))

    _text(slide, 0.62, 1.6, 12.0, 1.4, "Scale sells the room.",
          size=44, bold=True, color=C_CREAM, line_spacing=1.05)
    _text(slide, 0.62, 2.5, 12.0, 1.4, "Index sells the message.",
          size=44, bold=True, color=C_LAVENDER, line_spacing=1.05)

    _hairline(slide, 0.62, 4.05, 12.10, C_STROKE)

    body_lines = [
        f"{largest['name']} reaches the largest audience: "
        f"{_fmt_int_compact(largest['projected_us'])} U.S. digital adults.",
        f"{tightest['name']} owns the tightest signal: peak index of "
        f"{int(_peak_index(tightest))} vs. Gen Pop.",
        "Match the buy to the goal: reach-first for one, resonance-first for the other.",
    ]
    _text(slide, 0.62, 4.35, 11.5, 2.6, "\n\n".join(body_lines),
          size=14, color=C_MUTED, line_spacing=1.55)

    _footer(slide, 7, "Comparison")


# =============================================================================
#  Data helpers
# =============================================================================


def _fmt_int(n) -> str:
    try:
        n = int(round(float(n)))
    except (TypeError, ValueError):
        return ""
    if n <= 0:
        return ""
    return f"{n:,}"


def _fmt_int_compact(n) -> str:
    """Human-readable: 23,100,000 -> '23.1M'."""
    try:
        n = float(n)
    except (TypeError, ValueError):
        return ""
    if not n or n <= 0:
        return ""
    if n >= 1_000_000_000:
        return f"{n/1_000_000_000:.1f}B"
    if n >= 1_000_000:
        return f"{n/1_000_000:.1f}M"
    if n >= 1_000:
        return f"{n/1_000:.0f}K"
    return f"{int(n)}"


def _short_bucket_label(b: str) -> str:
    """Shorten long demographic bucket labels for tight chart columns."""
    s = str(b or "").strip()
    if not s:
        return ""
    if len(s) <= 12:
        return s
    # Some canonical shortenings
    repl = {
        "Bachelor's Degree": "Bachelor's",
        "Master's Degree":   "Master's",
        "Doctoral Degree":   "PhD",
        "Some College":      "Some Coll.",
        "High School":       "HS",
        "Prefer Not to Say": "N/A",
        "Not Hispanic or Latino": "Non-Hisp.",
        "Hispanic or Latino":     "Hispanic",
        "Black or African American": "Black",
        "American Indian or Alaska Native": "Am. Indian",
        "Native Hawaiian or Other Pacific Islander": "NHPI",
        "Two or More Races": "Mixed",
        "White":             "White",
        "Asian":             "Asian",
    }
    if s in repl:
        return repl[s]
    return s[:14] + "…"


def _year_label(date_range: str) -> str:
    if not date_range:
        return ""
    m = re.search(r"(20\d{2})", date_range)
    return m.group(1) if m else ""


def _all_behavioral_items(p: dict) -> list[dict]:
    out = []
    for cat, items in (p.get("behavioral") or {}).items():
        for it in items:
            it2 = dict(it)
            it2["_cat"] = cat
            out.append(it2)
    return out


def _top_over_index_brand(p: dict, *,
                          min_pct: float = 8.0,
                          max_index: float = 350.0,
                          min_index: float = 130.0) -> Optional[dict]:
    """Strongest over-indexing brand with real reach.

    ``max_index`` caps the search to weed out low-N artifacts (a 5%-reach
    brand at 800 idx is almost always panel noise, not audience signal —
    see 'THE FANTOM WALLET' defect the user flagged). Falls back to a
    looser filter, then to the raw top, so this never returns None as long
    as any behavioral item exists.

    Callers can pass ``min_pct=1.0`` etc. if they specifically want the
    old lax behavior (still used by combined-deck spotlight)."""
    all_items = _all_behavioral_items(p)
    # Tier 1: strict — real reach + meaningful (but not freak) over-index
    tier1 = [it for it in all_items
             if _num(it.get("pct", 0)) >= min_pct
             and min_index <= _num(it.get("index", 0)) <= max_index]
    if tier1:
        tier1.sort(key=lambda it: (-_num(it.get("index", 0)),
                                    -_num(it.get("pct", 0))))
        return tier1[0]
    # Tier 2: allow higher-index outliers ONLY when they still have some reach
    tier2 = [it for it in all_items
             if _num(it.get("pct", 0)) >= max(min_pct * 0.6, 5.0)
             and _num(it.get("index", 0)) >= min_index]
    if tier2:
        tier2.sort(key=lambda it: (-_num(it.get("index", 0)),
                                    -_num(it.get("pct", 0))))
        return tier2[0]
    # Tier 3: whatever we have
    tier3 = [it for it in all_items
             if _num(it.get("index", 0)) > 0
             and _num(it.get("pct", 0)) > 0]
    if not tier3:
        return None
    tier3.sort(key=lambda it: (-_num(it.get("index", 0)),
                                -_num(it.get("pct", 0))))
    return tier3[0]


# =============================================================================
#  Signal class framework — reach vs index nuance
# =============================================================================
#
# Every brand behavioral row sits somewhere on a 2x2 of REACH (% of the
# audience that engages) x INDEX (engagement rate vs Gen Pop, 100 =
# parity). The two dimensions carry different analytical weight:
#
#   * REACH governs SCALE — how many people you can actually reach.
#     A brand at 3% reach + 300 index is a persona characteristic; it
#     cannot deliver a media buy at scale, no matter how sharp the index.
#   * INDEX governs DIFFERENTIATION — how uniquely THIS audience over-
#     consumes the brand vs. the average American. A brand at 60% reach
#     + 105 index is table-stakes, not a signature.
#
# The five classes below encode the useful cells of that 2x2. Every
# narrative fallback in this module respects them: SIGNATURE brands
# never lead a sentence; MASS+LIFT brands always do.
# =============================================================================

_SIGCLASS_MASS_LIFT = "MASS_LIFT"      # reach >=20 and index >=130 — actionable
_SIGCLASS_SCALE     = "SCALE"          # reach >=20 and 90 <= index < 130 — table stakes
_SIGCLASS_UNDER     = "UNDER"          # reach >=20 and index < 90 — under-indexes
_SIGCLASS_SIGNATURE = "SIGNATURE"      # reach 8-20 and index >=150 — persona-only
_SIGCLASS_EMERGING  = "EMERGING"       # reach 10-20 and index 110-150 — moderate lift
_SIGCLASS_NOISE     = "NOISE"          # low reach + not high index, OR freak index

_SIGCLASS_LABEL = {
    _SIGCLASS_MASS_LIFT: "MASS + LIFT",
    _SIGCLASS_SCALE:     "SCALE",
    _SIGCLASS_UNDER:     "UNDER-INDEX",
    _SIGCLASS_SIGNATURE: "SIGNATURE",
    _SIGCLASS_EMERGING:  "EMERGING",
    _SIGCLASS_NOISE:     "NOISE",
}

# What each class means in analyst terms — used in captions + prompt.
_SIGCLASS_MEANING = {
    _SIGCLASS_MASS_LIFT: "Big reach and clear over-index — the actionable sweet spot for a partner buy.",
    _SIGCLASS_SCALE:     "Big reach but at Gen-Pop parity — table-stakes for the media buy, not a differentiator.",
    _SIGCLASS_SIGNATURE: "Small reach but sharp over-index — persona color, not a media buy at scale.",
    _SIGCLASS_EMERGING:  "Moderate reach with meaningful lift — worth a test but not a headline signal.",
    _SIGCLASS_UNDER:     "Reached but under-indexed — de-prioritize vs. the audience's over-index picks.",
    _SIGCLASS_NOISE:     "Too little reach or a freak index — probably panel artifact, omit.",
}


def _signal_class(pct, index) -> str:
    """Classify a (reach, index) pair. Reach in percent (0-100), index
    with 100 = parity.

    Thresholds tightened 2026-07-09 (round 2) after user flagged
    'Killstar (9.9%, idx 576)' and 'Urban Stems (7.9%, idx 399)' still
    leaking into the commerce narrative — 'less than 10% of the audience'
    is too thin to earn a mention in the persona story regardless of how
    sharp the index reads.

      * SIGNATURE minimum reach raised from 10% to 12% (was 5% with a
        'borderline signature' bucket that let 7-9% brands ship)
      * EMERGING minimum reach raised from 10% to 12%
      * Freak-index ceiling holds: idx > 500 downgrades regardless of
        reach; idx > 400 with reach < 15% becomes NOISE (Killstar 9.9%
        at 576 → NOISE, not SIGNATURE)
    """
    p = _num(pct)
    i = _num(index)
    if p <= 0 or i <= 0:
        return _SIGCLASS_NOISE
    if i > 500:
        return _SIGCLASS_SCALE if p >= 20 else _SIGCLASS_NOISE
    if p >= 20:
        if i >= 130:
            return _SIGCLASS_MASS_LIFT
        if i >= 90:
            return _SIGCLASS_SCALE
        return _SIGCLASS_UNDER
    if p >= 12:
        if i >= 150:
            return _SIGCLASS_SIGNATURE
        if i >= 110:
            return _SIGCLASS_EMERGING
        return _SIGCLASS_NOISE
    return _SIGCLASS_NOISE


def _signal_score(pct, index) -> float:
    """Weighted score for ranking brands by MARKETING IMPACT — how much
    audience you can actually reach times how differentiated that reach
    is. Reach dominates; the sqrt(index) term tempers freak-index bias
    so a 3%-reach-at-800-index brand can't jump above a 25%-at-140.
    """
    p = _num(pct)
    i = min(_num(index), 300.0)
    if p <= 0 or i <= 0:
        return 0.0
    return p * ((i / 100.0) ** 0.5)


def _log(msg: str) -> None:
    """Lightweight logger — writes to stderr with a module tag so debug
    output shows up in Render / local logs but never surfaces on-slide."""
    try:
        print(f"[profile_iq_deck_builder] {msg}", file=sys.stderr, flush=True)
    except Exception:
        pass


# Aggressive tokeniser used by _is_subject_brand — strips punctuation,
# common determiners, and film/tv-boilerplate words so 'The Nimrods: A
# Green Day Comedy' tokenises to {NIMRODS, GREEN, DAY, COMEDY} which
# matches the brand 'Green Day' at 100% overlap.
_SUBJECT_STOPWORDS = {
    "THE", "A", "AN", "AND", "&", "OF", "IN", "ON", "AT", "FOR", "TO",
    "OR", "MOVIE", "FILM", "SHOW", "SEASON", "EPISODE", "SERIES", "PART",
    "VOLUME", "VOL", "MR", "MS", "MRS", "DR",
}


def _subject_tokens(subject: str) -> set[str]:
    if not subject:
        return set()
    cleaned = re.sub(r"[^A-Za-z0-9 ]+", " ", subject).upper()
    return {t for t in cleaned.split()
            if t and t not in _SUBJECT_STOPWORDS and len(t) >= 2}


def _is_subject_brand(brand: str, subject: str) -> bool:
    """True if ``brand`` is materially the subject (or a fragment of it).

    Catches the Green Day / Nimrods leak: the subject 'The Nimrods A Green
    DAY Comedy' has tokens {NIMRODS, GREEN, DAY, COMEDY}. The brand
    'Green Day' tokenises to {GREEN, DAY} which is a full subset — subject
    match. 'Weezer' has no token overlap — not the subject.

    Guard rules:
      * Exact string match (case + punctuation insensitive) always fires
      * >=2 subject-token overlap fires
      * Single-token brand that IS a subject token fires (Nimrods == Nimrods)
    """
    if not brand or not subject:
        return False
    b = re.sub(r"[^A-Za-z0-9 ]+", " ", brand).upper().strip()
    s = re.sub(r"[^A-Za-z0-9 ]+", " ", subject).upper().strip()
    if not b or not s:
        return False
    if b == s or b in s or s in b:
        return True
    subj_toks = _subject_tokens(subject)
    brand_toks = {t for t in b.split()
                  if t and t not in _SUBJECT_STOPWORDS and len(t) >= 2}
    if not brand_toks or not subj_toks:
        return False
    # Full subset match — every brand token is a subject token (Green Day
    # inside Nimrods A Green Day Comedy)
    if brand_toks.issubset(subj_toks):
        return True
    # Single-token brand that IS a subject token
    if len(brand_toks) == 1 and next(iter(brand_toks)) in subj_toks:
        return True
    return False


_REACH_PATTERN = re.compile(r"([A-Z][A-Za-z0-9&'\-\.\+\/ ]{1,40})\s*\(\s*([0-9]{1,2}(?:\.[0-9])?)\s*%\s*(?:reach)?\s*\)")


def _body_has_data_grounding(body: str, min_mentions: int = 2) -> bool:
    """True if the narrative body mentions at least ``min_mentions``
    specific brand or bucket references with numeric context (pct or idx).

    This is the quality gate that catches generic LLM outputs like 'The
    audience exhibits strong preferences for brands that offer value and
    convenience' — zero brand names, zero numbers, useless for a CMO.
    Anything with fewer than 2 grounded references falls back to the
    data-driven template."""
    if not body:
        return False
    reach_hits = _REACH_PATTERN.findall(body)
    idx_hits   = re.findall(r"\b(?:idx|index)\s*[:=]?\s*[0-9]{2,4}\b",
                             body, flags=re.IGNORECASE)
    paren_idx  = re.findall(r"\(\s*idx\s*[0-9]{2,4}\s*\)",
                             body, flags=re.IGNORECASE)
    return (len(reach_hits) + len(idx_hits) + len(paren_idx)) >= min_mentions


def _body_cites_banned_brand(body: str,
                              signature_min_pct: float = 12.0,
                              max_allowed_idx: float = 400.0) -> tuple[bool, list]:
    """Detect LLM bodies that cite brands the framework has banned as
    NOISE — anything under ``signature_min_pct`` reach (2026-07-09
    round-2 threshold) or with an index above ``max_allowed_idx`` (freak
    signal, subject-adjacent panel skew).

    Returns (has_banned, offending_mentions) where offending_mentions is
    a list of (brand, pct, idx_or_None) tuples for logging. Callers
    should DISCARD the body and fire the data-driven fallback rather
    than trying to surgically strip the offending phrase — regex
    surgery on natural language leaves grammatically-broken sentences.
    """
    if not body:
        return False, []
    offending: list[tuple[str, float, Optional[float]]] = []
    # Match "BRAND (9.9%, index 576)" or "BRAND (9.9%)" or "BRAND (9.9% reach)"
    combo_pat = re.compile(
        r"([A-Z][A-Za-z0-9&'\-\.\+\/ ]{1,40})\s*\(\s*"
        r"([0-9]{1,2}(?:\.[0-9])?)\s*%[^)]*?(?:index\s*|idx\s*)?"
        r"([0-9]{2,4})?[^)]*\)"
    )
    for m in combo_pat.finditer(body):
        brand = (m.group(1) or "").strip().strip(",.:;")
        try:
            pct = float(m.group(2))
        except (TypeError, ValueError):
            continue
        try:
            idx = float(m.group(3)) if m.group(3) else None
        except (TypeError, ValueError):
            idx = None
        if len(brand) < 2:
            continue
        if pct < signature_min_pct:
            offending.append((brand, pct, idx))
            continue
        if idx is not None and idx > max_allowed_idx:
            offending.append((brand, pct, idx))
    return (len(offending) > 0, offending)


def _extract_narrated_reaches(body: str) -> list[tuple[str, float]]:
    """Return the (brand, pct) mentions in narrative order — used to
    detect when the LLM led a sentence with sub-mass brands."""
    if not body:
        return []
    out: list[tuple[str, float]] = []
    for m in _REACH_PATTERN.finditer(body):
        nm = (m.group(1) or "").strip().strip(",.:;")
        try:
            pct = float(m.group(2))
        except (TypeError, ValueError):
            continue
        # Skip generic "audience" / "them" tokens that pattern-matched.
        if len(nm) < 2 or nm.lower() in {"the audience", "and", "with"}:
            continue
        out.append((nm, pct))
    return out


def _repair_leading_signature(body: str, mass: list, sig: list,
                                reach_threshold: float = 20.0) -> str:
    """If the LLM body opens by naming brands under the mass threshold,
    prepend a mass-anchor sentence so the narrative reads correctly.

    This is defence-in-depth: the system prompt teaches the framework
    but a smaller model can still slip. Rather than override the LLM
    body entirely, we prepend a mass anchor and let the LLM's persona
    read as the SECOND half of the story — exactly its rightful role.
    """
    if not body or not mass:
        return body
    hits = _extract_narrated_reaches(body)
    if not hits:
        return body
    first_three = hits[:3]
    # If any of the first three named brands is above threshold, the LLM
    # got the ordering right — leave it alone.
    if any(pct >= reach_threshold for _, pct in first_three):
        return body
    # Otherwise, all three lead brands are sub-mass. Prepend an anchor.
    top = mass[:3]
    anchor_bits = [
        f"{(it.get('name') or '').strip()} "
        f"({_num(it.get('pct',0)):.0f}%)"
        for it in top
    ]
    if not anchor_bits:
        return body
    anchor = ("Mass-brand shopping is anchored by "
               + (", ".join(anchor_bits[:-1]) + ", and " + anchor_bits[-1]
                  if len(anchor_bits) > 1 else anchor_bits[0])
               + ". ")
    return anchor + body


def _split_by_reach(items: list, reach_threshold: float = 20.0,
                     signature_min_pct: float = 12.0,
                     signature_max_idx: float = 400.0,
                     ) -> tuple[list, list]:
    """Return (mass_brands, signature_brands).

    Thresholds tightened 2026-07-09 (round 2):
      * mass: reach >= reach_threshold (default 20%)
      * signature: reach in [signature_min_pct, reach_threshold) AND
        index between 130 and signature_max_idx.
          - min pct raised from 5% to 12% — 'less than 10% of the
            audience does this' is too thin to earn a mention in the
            persona narrative regardless of how sharp the index reads
          - max idx capped at 400 — freak-index brands (Killstar at
            9.9%/576) are almost always subject-adjacent panel skew,
            not real actionable persona signal
    """
    mass, sig = [], []
    for it in items or []:
        p = _num(it.get("pct", 0))
        i = _num(it.get("index", 0))
        if p <= 0:
            continue
        if p >= reach_threshold:
            mass.append(it)
        elif (p >= signature_min_pct
              and 130 <= i <= signature_max_idx):
            sig.append(it)
    mass.sort(key=lambda it: -_num(it.get("pct", 0)))
    sig.sort(key=lambda it: -_signal_score(it.get("pct", 0),
                                             it.get("index", 0)))
    return mass, sig


def _top_owned_reach(p: dict) -> float:
    items = _all_behavioral_items(p)
    if not items:
        return 0
    return max((_num(it.get("projection", 0)) for it in items), default=0)


def _first_available(p: dict, cat_names: list[str]) -> list[dict]:
    """Return the first matching category's item list."""
    beh = p.get("behavioral") or {}
    # Case-insensitive match
    lut = {k.upper(): v for k, v in beh.items()}
    for name in cat_names:
        v = lut.get(name.upper())
        if v:
            return list(v)
    return []


def _pick_overview_chapters(p: dict) -> list[tuple[str, str, str]]:
    """Choose the four chapter cards to render on the OVERVIEW slide, based
    on which behavioral categories actually have data."""
    subj = p["name"]
    chapters = [
        ("01", "Who they are",
         f"Demographic composition of the {subj} audience vs. U.S. Gen Pop."),
        ("02", "How they signal",
         "The category over-indexers — where this audience concentrates."),
    ]
    has_beauty = _first_available(p, ["COSMETICS", "BEAUTY", "PERSONAL CARE",
                                       "PERSONAL_CARE"])
    has_lux    = _first_available(p, ["LUXURY FASHION", "LUXURY_FASHION",
                                       "APPAREL/FOOTWEAR", "APPAREL_FOOTWEAR",
                                       "MPB"])
    has_media  = _first_available(p, ["STREAMING/PLATFORM", "STREAMING_PLATFORM",
                                       "STREAMING VIDEO", "MEDIA", "SOCIAL MEDIA"])
    if has_beauty:
        chapters.append((f"{len(chapters)+1:02d}",
                         "Beauty & prestige",
                         "Cosmetics, skincare, and fragrance — where the audience actually shops."))
    if has_lux and len(chapters) < 4:
        chapters.append((f"{len(chapters)+1:02d}",
                         "Fashion & apparel",
                         "Luxury and everyday apparel where this audience actually spends."))
    if has_media and len(chapters) < 4:
        chapters.append((f"{len(chapters)+1:02d}",
                         "Where attention lives",
                         "Streaming, social, and platform surfaces to reach the audience."))
    filler = [
        ("Reach quality",     "Beyond raw scale — engagement intensity by category."),
        ("Whitespace & activation", "The strongest, most-actionable partnership signals."),
    ]
    while len(chapters) < 4 and filler:
        t, s = filler.pop(0)
        chapters.append((f"{len(chapters)+1:02d}", t, s))
    return chapters[:4]


PERSONA_CATEGORY_HEADLINES: dict[str, tuple[str, str]] = {
    "COSMETICS":              ("Beauty retail",              "dominates."),
    "PERSONAL CARE":          ("A personal-care audience,",  "already spending."),
    "APPAREL/FOOTWEAR":       ("A clear lens on",            "how the audience dresses."),
    "MPB":                    ("Most-purchased brands.",     "Where the audience actually shops."),
    "MOST_PURCHASED_BRANDS":  ("Most-purchased brands.",     "Where the audience actually shops."),
    "LUXURY FASHION":         ("A clear lens on luxury.",    ""),
    "QSR":                    ("How the audience eats,",     "when they're out."),
    "RETAIL":                 ("Where the audience buys,",   "beyond luxury."),
    "TRAVEL":                 ("How the audience travels.",  "And where."),
    "BANKING":                ("Where the money lives.",     ""),
    "AUTOMOBILE":             ("What the audience drives.",  ""),
    "TALENT/ACTOR":           ("The talent overlap.",        ""),
    "SOCIAL MEDIA":           ("Social platform-native.",    "Content curious."),
    "STREAMING/PLATFORM":     ("How the audience watches,",  "when they choose."),
    "STREAMING VIDEO":        ("How the audience watches,",  "when they choose."),
    "GAMES":                  ("The gaming signal.",         "Where fandom converts to session time."),
    "GAMING":                 ("The gaming signal.",         "Where fandom converts to session time."),
    "TICKETING PLATFORMS":    ("Live-event buyers.",         ""),
    "FESTIVALS":              ("Festival culture.",          ""),
    "MUSIC":                  ("Music discovery habits.",    ""),
    "ARTIST":                 ("The artist overlap.",        ""),
    "ARTISTS":                ("The artist overlap.",        ""),
    "CULTURE":                ("Cultural affinity map.",     ""),
    "INSURANCE":              ("Insurance signal.",          ""),
    "PHARMACY":               ("Where the audience refills.",""),
    "TECHNOLOGY DEVICE":      ("The device stack.",          ""),
    "APP/PLATFORM":           ("The app stack.",             "Where the audience opens the phone."),
    "SEARCH ENGINE/AI":       ("Search and AI, blended.",    ""),
}


def _persona_category_headline(cat: str) -> tuple[str, str]:
    return PERSONA_CATEGORY_HEADLINES.get(cat.upper(), (f"{_pretty_cat(cat)}.", "The audience signal."))


# Category label prettifier so titles never render as "Qsr" or "Mpb".
_CAT_LABEL_OVERRIDES = {
    "QSR": "QSR", "MPB": "Most-Purchased Brands", "AI": "AI",
    "APP/PLATFORM": "App / Platform", "APPAREL/FOOTWEAR": "Apparel & Footwear",
    "STREAMING/PLATFORM": "Streaming & Platform",
    "STREAMING VIDEO": "Streaming Video", "STREAMING MUSIC": "Streaming Music",
    "SEARCH ENGINE/AI": "Search Engine & AI",
    "SOCIAL MEDIA": "Social Media", "MOST_PURCHASED_BRANDS": "Most-Purchased Brands",
    "PERSONAL CARE": "Personal Care", "LUXURY FASHION": "Luxury Fashion",
    "VMVPD/FAST": "vMVPD / FAST", "VIRTUAL MVPD/FAST": "Virtual MVPD / FAST",
}


def _pretty_cat(cat: str) -> str:
    """Return a title-cased version of a category name, respecting overrides."""
    if not cat:
        return ""
    key = cat.strip().upper()
    if key in _CAT_LABEL_OVERRIDES:
        return _CAT_LABEL_OVERRIDES[key]
    # Otherwise title-case, but keep known acronyms upper
    words = re.split(r"[\s_/]+", cat)
    out = []
    for w in words:
        if not w:
            continue
        if w.upper() in {"QSR", "MPB", "AI", "TV", "MLB", "NBA", "NFL",
                         "MLS", "WNBA", "MILB", "USPS", "CFB", "EPL"}:
            out.append(w.upper())
        else:
            out.append(w.title())
    return " ".join(out)


# Categories that make good persona deep-dive slides (skip social/streaming
# which get their own slide, skip huge structural buckets like AGE etc.)
_PERSONA_SLIDE_ORDER = [
    "COSMETICS", "PERSONAL CARE", "PERSONAL_CARE",
    "APPAREL/FOOTWEAR", "APPAREL_FOOTWEAR",
    "LUXURY FASHION", "LUXURY_FASHION",
    "RETAIL", "QSR", "TRAVEL", "BANKING", "AUTOMOBILE",
    "GAMES", "GAMING", "TICKETING PLATFORMS", "FESTIVALS",
    "CULTURE", "MUSIC", "ARTIST", "ARTISTS",
]

_SKIP_PERSONA_CATS = {"SOCIAL MEDIA", "SOCIAL_MEDIA", "STREAMING/PLATFORM",
                      "STREAMING_PLATFORM", "STREAMING VIDEO",
                      "STREAMING_VIDEO", "MEDIA",
                      # MPB gets its own dedicated ranked-table slide via
                      # _slide_mpb_table, so don't also spend a persona slot on it.
                      "MPB", "MOST PURCHASED BRANDS", "MOST_PURCHASED_BRANDS"}


def _category_earns_deep_dive(items: list) -> bool:
    """A category earns its own deep-dive slide only if it clears the
    McKinsey reach bar: top brand >=15% reach AND at least 2 brands >=10%
    reach. Everything else is a wasted page (Marni-at-3.6% problem) and
    rolls up into the signature_affinities compact slide instead."""
    if not items:
        return False
    reaches = sorted(
        (_num(it.get("pct", 0)) for it in items),
        reverse=True,
    )
    if not reaches or reaches[0] < 15.0:
        return False
    strong_count = sum(1 for r in reaches if r >= 10.0)
    return strong_count >= 2


def _pick_persona_category_slides(p: dict, max_slides: int = 2) -> list[tuple[str, list]]:
    """Pick up to N behavioral categories to spotlight as their own slides.

    Applies the reach gate (_category_earns_deep_dive) so a category with a
    3.6%-reach top brand can never eat a page. If the LLM analyst brief
    nominated categories via ``category_why``, prefer those (in the order
    given by the brief), otherwise fall back to canonical persona order.
    """
    beh = p.get("behavioral") or {}
    picked: list[tuple[str, list]] = []
    seen = set()

    # 1. If the analyst brief called out specific categories in category_why,
    #    honor its picks first — the LLM already applied the reach + story
    #    filter, and its ordering encodes narrative priority.
    cat_why = _brief_get(p, "category_why", default={}) or {}
    brief_cats = [k for k in cat_why.keys() if k]
    for pref in brief_cats:
        if len(picked) >= max_slides:
            break
        for cat, items in beh.items():
            if cat in seen or cat.upper() in _SKIP_PERSONA_CATS:
                continue
            if (cat.upper() == pref.upper() and items
                    and _category_earns_deep_dive(items)):
                picked.append((cat, items))
                seen.add(cat)
                break

    # 2. Canonical persona order, still reach-gated.
    for pref in _PERSONA_SLIDE_ORDER:
        if len(picked) >= max_slides:
            break
        for cat, items in beh.items():
            if cat in seen:
                continue
            if (cat.upper() == pref.upper() and items
                    and _category_earns_deep_dive(items)):
                picked.append((cat, items))
                seen.add(cat)
                break

    # 3. Any other high-signal reach-gated category to backfill.
    if len(picked) < max_slides:
        candidates = []
        for cat, items in beh.items():
            if cat in seen or cat.upper() in _SKIP_PERSONA_CATS:
                continue
            if not items or not _category_earns_deep_dive(items):
                continue
            peak = max((_num(it.get("index", 0)) for it in items), default=0)
            if peak >= 130:
                candidates.append((peak, cat, items))
        candidates.sort(key=lambda t: -t[0])
        for _, cat, items in candidates[: max_slides - len(picked)]:
            picked.append((cat, items))
    return picked


def _generate_tagline(p: dict) -> str:
    """One-line tagline for the cover, driven by the analyst brief when
    available, otherwise by the strongest durable audience signal (with
    the low-reach / freak-index filter applied so we don't lead with a
    5%-reach brand at 800 idx)."""
    brief_t = _brief_get(p, "cover_tagline")
    if brief_t:
        return _purge_pronouns(brief_t, p["name"])
    top = _top_over_index_brand(p, min_pct=10.0, max_index=300.0)
    proj = _fmt_int_compact(p["projected_us"] or 0)
    proj_phrase = (f"{proj} U.S. digital adults" if proj
                    else "the U.S. digital audience")
    arch = p.get("_archetype") or _ARCHETYPE_GENERIC
    if arch == _ARCHETYPE_MOVIE:
        arch_phrase = f"US adults with a digital touchpoint to {p['name']}"
    elif arch == _ARCHETYPE_TV_SHOW:
        arch_phrase = f"US adults engaged with {p['name']}"
    else:
        arch_phrase = f"the {p['name']} audience"
    if top:
        cat = _pretty_cat(top.get("_cat") or "").lower()
        return (
            f"A behavioral read on {arch_phrase} — {proj_phrase} — with "
            f"durable brand affinities like {top['name']} "
            f"({int(_num(top['index']))} index in {cat}) telling a real story "
            "about how this audience buys, watches, and shows up."
        )
    return (
        f"A behavioral profile of {arch_phrase}. {proj_phrase}, drawn from "
        "Crosswalk's zero-party panel."
    )


def _combined_category_ranking(profiles: list[dict]) -> list[tuple[str, str, float, str]]:
    """Return rows of (category, winner_audience, winner_peak_index, winner_top_brand)."""
    cats = set()
    for p in profiles:
        for c in (p.get("behavioral") or {}).keys():
            if c.upper() not in _SKIP_PERSONA_CATS:
                cats.add(c)
    rows = []
    for c in cats:
        best = None
        for p in profiles:
            items = (p.get("behavioral") or {}).get(c) or []
            if not items:
                # Try case-variant match
                lut = {k.upper(): v for k, v in p["behavioral"].items()}
                items = lut.get(c.upper()) or []
            if not items:
                continue
            peak = max(items, key=lambda it: _num(it.get("index", 0)))
            peak_idx = _num(peak.get("index", 0))
            if not best or peak_idx > best[2]:
                best = (c, p["name"], peak_idx, peak.get("name", ""))
        if best and best[2] >= 100:
            rows.append(best)
    rows.sort(key=lambda r: -r[2])
    return rows


def _brand_overlap_rows(profiles: list[dict], *,
                        min_index: float = 110, top_n: int = 8) -> list[dict]:
    """Return brands that every profile posts an index >= min_index for."""
    if not profiles:
        return []
    # Build per-profile { normalized_brand -> (raw_name, category, index) }
    def norm(s: str) -> str:
        return re.sub(r"[^a-z0-9]+", "", (s or "").lower())
    per: list[dict] = []
    for p in profiles:
        d = {}
        for cat, items in (p.get("behavioral") or {}).items():
            if cat.upper() in _SKIP_PERSONA_CATS:
                continue
            for it in items:
                nm = norm(it.get("name", ""))
                if not nm:
                    continue
                idx = _num(it.get("index", 0))
                if idx < min_index:
                    continue
                cur = d.get(nm)
                if not cur or idx > cur[2]:
                    d[nm] = (it.get("name", ""), cat, idx)
        per.append(d)
    if not per:
        return []
    shared = set(per[0].keys())
    for m in per[1:]:
        shared &= set(m.keys())
    rows = []
    for k in shared:
        rec = {
            "brand":    per[0][k][0],
            "category": per[0][k][1],
            "indices":  {p["name"]: per[i][k][2] for i, p in enumerate(profiles) if k in per[i]},
        }
        rec["_avg_idx"] = sum(rec["indices"].values()) / max(len(rec["indices"]), 1)
        rows.append(rec)
    rows.sort(key=lambda r: -r["_avg_idx"])
    return rows[:top_n]


def _peak_index(p: dict) -> float:
    peak = 0.0
    for it in _all_behavioral_items(p):
        v = _num(it.get("index", 0))
        if v > peak:
            peak = v
    return peak


# =============================================================================
#  Data helpers for KD-style slides (by-the-numbers, index story, geography,
#  case studies, takeaways, MPB table, methodology, sub-section headers)
# =============================================================================

# Rough Gen Pop shares used for index computations when the payload does not
# ship pre-computed index maps. Approximates U.S. adult digital census (2025).
_APPROX_GENPOP_AGE = {
    "17 & under": 8.0, "17 and under": 8.0, "under 18": 8.0,
    "18-24": 12.0, "18–24": 12.0,
    "25-34": 17.0, "25–34": 17.0,
    "35-44": 16.5, "35–44": 16.5,
    "45-54": 15.5, "45–54": 15.5,
    "55-64": 14.0, "55–64": 14.0,
    "65+": 17.0,
}
_APPROX_GENPOP_ETH = {
    "White": 59.0, "Black": 12.5, "Black or African American": 12.5,
    "Hispanic": 18.5, "Hispanic or Latino": 18.5,
    "Asian": 6.5, "Two or More Races": 2.5, "Another": 1.5,
    "American Indian or Alaska Native": 0.7,
    "Native Hawaiian or Other Pacific Islander": 0.3,
}
_APPROX_GENPOP_GENDER = {"Female": 51.0, "Male": 49.0}


def _lookup_idx(p: dict, dim: str, bucket: str) -> float:
    """Return audience index for a demo bucket; fall back to approximation."""
    idx_map = (p.get("demographics_index") or {}).get(dim, {}) or {}
    v = _num(idx_map.get(bucket, 0))
    if v:
        return v
    share = _num((p.get("demographics") or {}).get(dim, {}).get(bucket, 0))
    if not share:
        return 0
    gp_source = {"age": _APPROX_GENPOP_AGE,
                 "ethnicity": _APPROX_GENPOP_ETH,
                 "gender": _APPROX_GENPOP_GENDER}.get(dim, {})
    gp = _num(gp_source.get(bucket, 0))
    if not gp:
        return 0
    return round(share / gp * 100)


def _source_line(p: dict) -> str:
    """Canonical KD-style source footer line."""
    n = _fmt_int(p.get("sample_size") or 0)
    dr = (p.get("date_range") or "").strip()
    parts = ["Source: Crosswalk Profile IQ"]
    if dr:
        parts.append(dr)
    if n:
        parts.append(f"n = {n} panelists")
    parts.append("Index = audience share vs Gen. Pop.")
    return "  /  ".join(parts)


def _by_the_numbers_stats(p: dict) -> list[tuple[str, str, str, str]]:
    """Pick the 4 most-quotable stats for the KD 'By the numbers' slide.
    Returns [(value, label, note, tone), ...] where tone is
    'cream' / 'lavender' / 'magenta' for color choice."""
    out: list[tuple[str, str, str, str]] = []
    proj = p.get("projected_us") or 0
    sample = p.get("sample_size") or 0

    # Stat 1: audience scale
    out.append((
        _fmt_int_compact(proj) or "—",
        "U.S. DIGITAL AUDIENCE",
        (f"Insights based on {_fmt_int(sample) or 'panel'} panelists" +
         (f", tracked {p.get('date_range')}" if p.get("date_range") else "")
         + "."),
        "cream",
    ))

    # Stat 2: primary age band (combine 18-24, 25-34, 35-44 if present)
    age = (p.get("demographics") or {}).get("age", {}) or {}
    prime = 0.0
    for k, v in age.items():
        key = str(k).strip().lower().replace("–", "-")
        if key in ("18-24", "25-34", "35-44"):
            prime += _num(v)
    if prime:
        prime_idx = round((prime / 41.5) * 100)  # ~41.5% Gen Pop 18-44
        out.append((
            f"{prime:.1f}%",
            "AGE 18-44",
            (f"Driving culture with a core audience. "
             f"Index {prime_idx} vs ~41.5% of Gen. Pop."),
            "lavender" if prime_idx >= 130 else "cream",
        ))

    # Stat 3: top ethnicity over-index
    eth = (p.get("demographics") or {}).get("ethnicity", {}) or {}
    top_eth = None
    best_idx = 0
    for k, v in eth.items():
        idx = _lookup_idx(p, "ethnicity", k)
        if idx > best_idx and idx > 110:
            top_eth = (k, _num(v), idx)
            best_idx = idx
    if top_eth:
        k, share, idx = top_eth
        gp = _APPROX_GENPOP_ETH.get(k, 0)
        note = (f"{share:.1f}% of the audience is {k}"
                + (f" vs {gp:.1f}% of Gen Pop." if gp else ".")
                + f" A significant anchor of this profile.")
        out.append((f"{int(idx)}", f"INDEX — {k.upper()} AUDIENCE",
                    note, "magenta"))

    # Stat 4: income $75k+ combined
    inc = (p.get("demographics") or {}).get("income", {}) or {}
    high = 0.0
    for k, v in inc.items():
        key = str(k).lower()
        if any(t in key for t in ("$75", "$100", "$150", "$200", "$250",
                                  "75-", "100-", "150-", "200+", "250+")):
            # Avoid double-counting narrow bands: keep only $75k+ patterns
            if ("75" in key or "100" in key or "150" in key or
                    "200" in key or "250" in key):
                high += _num(v)
    if high:
        out.append((
            f"{high:.1f}%",
            "INCOME $75K PLUS",
            ("Bifurcated commercial opportunities that run the gamut "
             "of household income levels."),
            "cream",
        ))

    return out[:4]


def _identity_headline(p: dict) -> str:
    """Sentence-case shorthand of who this audience is.

    Prefers the analyst brief. Falls back to a heuristic that respects the
    Gen Pop split — 54% male at index 113 is now correctly named 'Male-
    leaning.' rather than 'Balanced-gender.'."""
    brief_h = _brief_get(p, "identity_headline")
    if brief_h:
        return _purge_pronouns(brief_h, p["name"])
    gender = (p.get("demographics") or {}).get("gender", {}) or {}
    fem = _num(gender.get("Female", 0))
    mal = _num(gender.get("Male", 0))
    parts: list[str] = []
    if fem >= 60:
        parts.append("Female.")
    elif mal >= 60:
        parts.append("Male.")
    elif fem >= 54:
        parts.append("Female-leaning.")
    elif mal >= 54:
        parts.append("Male-leaning.")
    elif fem or mal:
        parts.append("Balanced-gender.")
    age = (p.get("demographics") or {}).get("age", {}) or {}
    prime = sum(_num(v) for k, v in age.items()
                if str(k).strip().replace("–", "-") in ("18-24", "25-34", "35-44"))
    if prime >= 50:
        parts.append("Prime-age.")
    for k in ("Black", "Black or African American", "Hispanic",
              "Hispanic or Latino", "Asian"):
        if _lookup_idx(p, "ethnicity", k) >= 140:
            parts.append("Multicultural.")
            break
    inc = (p.get("demographics") or {}).get("income", {}) or {}
    high = sum(_num(v) for k, v in inc.items()
               if any(t in str(k) for t in ("$75", "$100", "$150", "$200", "$250")))
    if high >= 45:
        parts.append("Mid-income+.")
    return " ".join(parts) or "The audience, unfiltered."


def _identity_stat_rows(p: dict) -> list[tuple[str, str, float]]:
    """Return up to 4 (pct_str, label, idx_val) rows for the IDENTITY slide,
    prioritized by highest audience share within a category (top of each dim)."""
    out: list[tuple[str, str, float]] = []

    def _top(dim: str, dim_name: str):
        d = (p.get("demographics") or {}).get(dim, {}) or {}
        if not d:
            return None
        top_k = max(d.keys(), key=lambda k: _num(d[k]))
        top_v = _num(d[top_k])
        if not top_v:
            return None
        idx = _lookup_idx(p, dim, top_k)
        return (f"{top_v:.1f}%", f"{dim_name}: {top_k.upper()}", idx)

    for dim, name in [
        ("gender", "GENDER"),
        ("age", "AGE"),
        ("ethnicity", "ETHNICITY"),
        ("income", "INCOME"),
    ]:
        row = _top(dim, name)
        if row:
            out.append(row)
    return out[:4]


def _identity_why_it_matters(p: dict) -> str:
    brief_w = _brief_get(p, "identity_why")
    if brief_w:
        return _purge_pronouns(brief_w, p["name"])
    subj = p["name"]
    idx_hits = [(k, _lookup_idx(p, "ethnicity", k))
                for k in ("Black or African American",
                          "Hispanic or Latino", "Asian")
                if _lookup_idx(p, "ethnicity", k) >= 130]
    if idx_hits:
        top = max(idx_hits, key=lambda t: t[1])
        return (f"The {subj} audience is anchored by strong over-representation "
                f"in {top[0]} ({int(top[1])} index) — a cultural anchor that "
                "shapes brand fit, media plan, and creative direction.")
    return (f"The {subj} audience shape is an ideal fit for brand "
            "partnerships, media alignment, and category activation.")


def _final_headline(p: dict) -> tuple[str, str]:
    """Two-line editorial closer for the FINAL INSIGHT slide.

    Prefers the analyst brief. Fallback avoids the 'reads it as taste'
    cliché the user flagged."""
    brief_a = _brief_get(p, "final_line_a")
    brief_b = _brief_get(p, "final_line_b")
    if brief_a and brief_b:
        return (_purge_pronouns(brief_a, p["name"]),
                _purge_pronouns(brief_b, p["name"]))
    subj = p["name"]
    arch = p.get("_archetype") or _ARCHETYPE_GENERIC
    peak = _peak_index(p)
    if arch == _ARCHETYPE_MOVIE:
        return (f"Where {subj} lives online,",
                "the audience is already waiting.")
    if arch == _ARCHETYPE_TV_SHOW:
        return (f"The {subj} audience already lives",
                "inside the right media surfaces.")
    if peak >= 300:
        return (f"This is not a generic fanbase.",
                "It concentrates habit and spend.")
    return (f"The {subj} audience is a signal,",
            "not just a headcount.")


def _final_body(p: dict, top) -> str:
    brief_body = _brief_get(p, "final_body")
    if brief_body:
        return _purge_pronouns(brief_body, p["name"])
    subj = p["name"]
    arch = p.get("_archetype") or _ARCHETYPE_GENERIC
    n_over = sum(1 for it in _all_behavioral_items(p)
                 if _num(it.get("index", 0)) >= 150)
    parts: list[str] = []
    parts.append(
        f"Across {n_over}+ behavioral signals, the {subj} audience over-indexes "
        "150+ vs. the U.S. digital baseline — a concentration of habit that "
        "reads as durable purchase intent, not passing attention."
    )
    if top and _num(top.get("pct", 0)) >= 8.0:
        cat = _pretty_cat(top.get("_cat", ""))
        parts.append(
            f"The strongest single durable read is {top['name']} at "
            f"{int(_num(top['index']))} in {cat} — a ready-made "
            "co-signer for any marketer trying to reach this school of fish."
        )
    if arch == _ARCHETYPE_MOVIE:
        parts.append(
            "For a theatrical or streaming release, prioritize media surfaces "
            "where the audience already indexes 130+ and lean into "
            "co-marketing with the highest-fit brand partners identified above."
        )
    elif arch == _ARCHETYPE_TV_SHOW:
        parts.append(
            "For launch or renewal marketing, weight social + streaming "
            "surfaces the audience already prefers, and pursue partner-brand "
            "co-marketing where the index is durable."
        )
    return "  ".join(parts)


def _top_demo_index_callouts(p: dict, n: int = 4) -> list[tuple[float, str, str]]:
    """Rank demographic buckets by index vs Gen Pop, return top ``n`` for
    the KD 'Where the audience over-indexes most' slide.

    Rules (added 2026-07-09 after the 'LESS THAN $25,000' + '$25-49K'
    both showing as top over-indexes for the Nimrods comedy audience —
    two income buckets crowding out ethnicity + education):
      * ONE bucket per dimension max — dedupe so we don't blast two
        income buckets as headline stats
      * Skip fringe buckets (share < 5%)
      * Skip weak over-indexes (idx < 120)
      * Skip label-noise buckets (empty string, 'Prefer not to say', etc.)

    Each row: (index_value, big_label, sub_context)."""
    NOISE_LABELS = {"", "PREFER NOT TO SAY", "UNKNOWN", "OTHER",
                    "NON-BINARY", "TRANS MALE", "TRANS FEMALE"}
    per_dim_best: dict[str, tuple[float, str, str]] = {}
    for dim in ["ethnicity", "age", "income", "education", "gender"]:
        d = (p.get("demographics") or {}).get(dim, {}) or {}
        for k, share in d.items():
            share_n = _num(share)
            if share_n < 5:
                continue
            k_upper = str(k).upper().strip()
            if k_upper in NOISE_LABELS:
                continue
            idx = _lookup_idx(p, dim, k)
            if idx < 120:
                continue
            gp = 0.0
            if dim == "ethnicity":
                gp = _APPROX_GENPOP_ETH.get(k, 0.0)
            elif dim == "age":
                gp = _APPROX_GENPOP_AGE.get(str(k).replace("–", "-"), 0.0)
            elif dim == "gender":
                gp = _APPROX_GENPOP_GENDER.get(k, 0.0)
            sub = (f"{share_n:.1f}% of the audience"
                   + (f" vs {gp:.1f}% of Gen Pop" if gp else ""))
            candidate = (idx, k_upper, sub)
            # Keep the highest-index bucket in this dimension
            if dim not in per_dim_best or candidate[0] > per_dim_best[dim][0]:
                per_dim_best[dim] = candidate
    picks = list(per_dim_best.values())
    picks.sort(key=lambda t: -t[0])
    return picks[:n]


def _top_dmas_ranked(p: dict, n: int = 10) -> list[tuple[str, float, float]]:
    """Return top DMAs from ``p['locations']`` sorted by audience share,
    with (name, pct, idx). Handles both list-of-dicts and dict shapes."""
    locs = p.get("locations") or []
    rows: list[tuple[str, float, float]] = []
    if isinstance(locs, list):
        for item in locs:
            if isinstance(item, dict):
                name = (item.get("name") or item.get("dma")
                        or item.get("market") or "").strip()
                pct = _num(item.get("pct") or item.get("percent")
                           or item.get("share"))
                idx = _num(item.get("index"))
                if name and pct:
                    rows.append((name, pct, idx))
            elif isinstance(item, (list, tuple)) and len(item) >= 2:
                name = str(item[0]).strip()
                pct = _num(item[1])
                idx = _num(item[2]) if len(item) > 2 else 0
                if name and pct:
                    rows.append((name, pct, idx))
    elif isinstance(locs, dict):
        for k, v in locs.items():
            pct = _num(v.get("pct") if isinstance(v, dict) else v)
            idx = _num(v.get("index") if isinstance(v, dict) else 0)
            if k and pct:
                rows.append((str(k).strip(), pct, idx))
    rows.sort(key=lambda r: -r[1])
    return rows[:n]


def _why_it_matters_for_category(cat: str, p: dict, top: list) -> str:
    """One-liner narrative for a category slide's WHY IT MATTERS strip.

    Prefers the LLM analyst brief (which delivers per-category tailored
    'why it matters' one-liners). Falls back to an archetype-aware
    templated line."""
    # Try the brief first — analyst may have written a specific line
    cat_map = _brief_get(p, "category_why", default={}) or {}
    if isinstance(cat_map, dict):
        cat_key = cat.upper().strip()
        # Try exact key, then pretty variants
        for k in (cat_key, cat.upper(),
                  _pretty_cat(cat).upper(),
                  cat_key.replace("_", " "),
                  cat_key.replace("/", " / ")):
            v = cat_map.get(k)
            if v and isinstance(v, str):
                return _purge_pronouns(v, p["name"])

    if not top:
        return "No standout brand signal yet in this category."
    lead = top[0]
    name = lead.get("name", "").strip()
    idx  = int(_num(lead.get("index", 0)))
    pct  = _num(lead.get("pct", 0))
    cat_pretty = _pretty_cat(cat)
    arch = p.get("_archetype") or _ARCHETYPE_GENERIC
    lead_class = _signal_class(pct, idx)

    # Extreme-index low-reach guard — never lead with a freak-index brand
    if idx >= 400 and pct < 8:
        return (f"The audience spreads across mainstream "
                f"{cat_pretty.lower()} choices with meaningful over-index; "
                "scale is present, no single brand dominates the signal.")

    if lead_class == _SIGCLASS_MASS_LIFT:
        # Real MASS + LIFT — the sweet spot. Frame it as actionable scale.
        base = (f"{name} anchors this category with {pct:.1f}% reach at a "
                f"{idx} index — real scale AND real over-index. This is the "
                f"actionable pick for a {cat_pretty.lower()} partnership.")
    elif lead_class == _SIGCLASS_SCALE:
        base = (f"{name} reaches {pct:.1f}% of the audience — real scale, but "
                f"at a {idx} index it reads as table stakes rather than "
                f"differentiation. Sizing follows scale, not signature.")
    elif lead_class in (_SIGCLASS_SIGNATURE, _SIGCLASS_EMERGING):
        # Signature-tier lead: caveat that this is persona color, not scale.
        base = (f"The category is signature-tier: {name} over-indexes at {idx} "
                f"but only reaches {pct:.1f}% of the audience. Read this as "
                f"persona color for {cat_pretty.lower()} — sharp taste, "
                "not a channel that will scale a buy on its own.")
    elif idx >= 130:
        base = (f"{name} leads with {pct:.1f}% penetration ({idx} index) — a "
                f"reliable {cat_pretty.lower()} beachhead.")
    else:
        return (f"The audience distributes across mainstream "
                f"{cat_pretty.lower()} choices; scale is present, "
                "differentiation is soft.")

    # Archetype-specific tag — but only if the lead is scaled enough to
    # actually stand up as a partner brand. Signature/emerging leads get a
    # SIGNATURE-appropriate tag: partner selection, not media weight.
    if lead_class in (_SIGCLASS_MASS_LIFT, _SIGCLASS_SCALE):
        if arch == _ARCHETYPE_MOVIE:
            base += (f" For the release-marketing plan, {name} is a natural "
                     "co-marketing / partner brand.")
        elif arch == _ARCHETYPE_TV_SHOW:
            base += (f" For launch or renewal marketing, {name} is a high-fit "
                     "sponsor and audience-lookalike seed.")
    else:
        if arch == _ARCHETYPE_MOVIE:
            base += (f" Best used as a signature partner or cultural signal "
                     "in the release marketing, not as a scale channel.")
        elif arch == _ARCHETYPE_TV_SHOW:
            base += (f" Best used as a taste-maker partner, not as a scale "
                     "sponsor.")
    return base


def _sub_section_headers(cat: str, n_cols: int) -> list[str]:
    """Sub-section header labels for the multi-column persona slide.

    Previous version had hardcoded conceptual labels (e.g. Travel →
    AIRLINES/HOTELS/OTA/EXPERIENCES) that the row assignments did NOT
    respect — a hotel brand ended up under 'AIRLINES', a lingerie brand
    under 'STREETWEAR', etc. Since we do not have a real sub-category
    signal in the raw payload, we now always use tier labels driven by
    index-rank order. This is honest and never misleading.

    Categories like MPB where the split IS just rank-based keep their
    'TOP 5 / 6-10 / ...' style labels since those describe reality."""
    cat_key = cat.upper().replace("_", " ").replace("/", " / ")
    if cat_key in ("MPB", "MOST PURCHASED BRANDS"):
        return ["TOP 5", "6-10", "11-15", "16-20"][:n_cols]
    # Tier labels — the columns are populated top-of-index first, so tier
    # names reflect actual ranking order and don't fabricate a taxonomy.
    tiers = ["STRONGEST SIGNAL", "STRONG SIGNAL",
             "NOTABLE SIGNAL", "EMERGING SIGNAL"]
    return tiers[:n_cols]


def _pick_case_studies(p: dict, n: int = 3) -> list[dict]:
    """Pick top ``n`` distinct brand case studies — each becomes one slide.

    Priority order:
      1. If the LLM analyst brief provided case studies, honor its picks
         (matched back to the real data so stats stay accurate).
      2. Otherwise, filter noise: require pct >= 8% AND index >= 130 AND
         (if idx > 350 then also pct >= 15%). Rank by pct * log(index/100+1)
         which favors real reach over freak-index outliers. This is the
         fix for the 'FANTOM WALLET at 849 with 5.7% reach' defect the
         user flagged.
      3. Spread across categories so we don't do 3 case studies from the
         same category."""
    brief_cs = _brief_get(p, "case_studies")
    if brief_cs and isinstance(brief_cs, list):
        all_items = _all_behavioral_items(p)
        # Lookup by upper-cased brand name
        item_lut = {(it.get("name") or "").strip().upper(): it
                     for it in all_items}
        picked: list[dict] = []
        for bcs in brief_cs:
            if not isinstance(bcs, dict):
                continue
            brand_key = (bcs.get("brand") or "").strip().upper()
            if not brand_key:
                continue
            it = item_lut.get(brand_key)
            if not it:
                # Partial-match fallback (brand names sometimes drift)
                for k, v in item_lut.items():
                    if brand_key in k or (len(k) >= 4 and k in brand_key):
                        it = v
                        break
            if not it:
                continue
            picked.append(_case_study_from_item(
                p, it,
                headline_a=_purge_pronouns(str(bcs.get("headline_a") or ""),
                                           p["name"]),
                headline_b=_purge_pronouns(str(bcs.get("headline_b") or ""),
                                           p["name"]),
                narrative=_purge_pronouns(str(bcs.get("narrative") or ""),
                                          p["name"]),
            ))
            if len(picked) >= n:
                break
        if picked:
            return picked

    # ── Deterministic fallback path ─────────────────────────────────────
    all_items = _all_behavioral_items(p)
    scored: list[tuple[float, dict]] = []
    for it in all_items:
        pct = _num(it.get("pct", 0))
        idx = _num(it.get("index", 0))
        if pct < 8.0 or idx < 130:
            continue
        # Freak-index low-reach filter (Rule #16 in the pipeline rules —
        # panel-reality calibration): a 5%-reach brand at 800 idx is much
        # more likely a low-N panel artifact than a real audience signal.
        if idx >= 350 and pct < 15:
            continue
        if len((it.get("name") or "").strip()) < 2:
            continue
        # Score weights reach linearly + log(index) so we don't pick 3
        # tiny freak-index brands over 3 meaningful mid-index brands.
        score = pct * math.log((idx / 100.0) + 1.0)
        scored.append((score, it))
    scored.sort(key=lambda t: -t[0])
    picked_items: list[dict] = []
    seen_cats: set[str] = set()
    for _, it in scored:
        cat = (it.get("_cat") or "").upper()
        if cat in seen_cats and len(picked_items) >= 1:
            continue
        picked_items.append(it)
        seen_cats.add(cat)
        if len(picked_items) >= n:
            break
    return [_case_study_from_item(p, it) for it in picked_items]


def _case_study_from_item(p: dict, it: dict, *,
                          headline_a: str = "",
                          headline_b: str = "",
                          narrative: str = "") -> dict:
    """Build the case-study payload the slide renderer expects from a raw
    behavioral item. If any editorial string is missing, a smart archetype-
    aware fallback fills it in — but the LLM brief's copy always wins when
    provided."""
    brand = (it.get("name") or "").strip()
    pct   = _num(it.get("pct", 0))
    idx   = int(_num(it.get("index", 0)))
    gp    = _num(it.get("genPopPct", 0))
    proj  = _num(it.get("projection", 0))
    cat   = (it.get("_cat") or "")
    if not gp and idx:
        gp = round(pct / (idx / 100), 1) if idx else 0

    if not headline_a or not headline_b:
        if idx >= 250:
            headline_a, headline_b = "All-in on", f"{brand}."
        elif idx >= 180:
            headline_a, headline_b = "Wired to buy", f"{brand}."
        elif idx >= 140:
            headline_a, headline_b = "Pre-sold on", f"{brand}."
        else:
            headline_a, headline_b = "Meaningful lift for", f"{brand}."

    if not narrative:
        subj = p["name"]
        arch = p.get("_archetype") or _ARCHETYPE_GENERIC
        base = (
            f"The {subj} audience over-indexes {idx} on {brand} — "
            f"{pct:.1f}% engage in-window vs. {gp:.1f}% of the U.S. general "
            "population. That gap is a durable purchase signal, not a "
            "one-off spike."
        )
        if arch == _ARCHETYPE_MOVIE:
            tail = (f"For the release marketing plan, {brand} is a natural "
                    "co-marketing partner and a lookalike-audience seed "
                    "worth activating against.")
        elif arch == _ARCHETYPE_TV_SHOW:
            tail = (f"For launch and renewal marketing, {brand} is a high-"
                    "fit sponsor and a lookalike-seed audience worth "
                    "activating against.")
        elif arch in (_ARCHETYPE_MUSICIAN, _ARCHETYPE_ACTOR,
                       _ARCHETYPE_ATHLETE, _ARCHETYPE_TEAM):
            tail = (f"For endorsement or partnership planning, {brand} is a "
                    "high-fit brand — the affinity is already there without "
                    "a paid buy.")
        else:
            tail = (f"For any partner brand aligned with the audience, "
                    f"{brand} is a ready-made co-signer.")
        narrative = base + "\n\n" + tail

    stats = [
        (f"{brand.upper()} INDEX", f"{idx}", "vs Gen Pop"),
        ("% PENETRATION",          f"{pct:.1f}%",
         (f"vs {gp:.1f}% Gen. Pop." if gp else "of the audience")),
        ("AUDIENCE REACH",         (_fmt_int_compact(proj) or "—"),
         "U.S. digital adults engaged"),
    ]
    return {
        "name": brand, "category": cat,
        "pct": pct, "index": idx, "gen_pop_pct": gp,
        "projection": proj,
        "narrative": narrative,
        "headline_a": headline_a, "headline_b": headline_b,
        "stats": stats,
    }


def _mpb_top_brands(p: dict, n: int = 20) -> list[dict]:
    """Return top MPB brands. Prefers the 'MPB' or 'MOST PURCHASED BRANDS'
    category; falls back to a flat top-N across all behavioral items."""
    beh = p.get("behavioral") or {}
    for key in ("MPB", "MOST PURCHASED BRANDS", "MOST_PURCHASED_BRANDS",
                "Most Purchased Brands"):
        if key in beh and beh[key]:
            items = list(beh[key])
            items.sort(key=lambda it: -_num(it.get("pct", 0)))
            return items[:n]
    # Fallback: all-cat flat rank by pct
    flat = [it for it in _all_behavioral_items(p) if _num(it.get("pct", 0)) > 0]
    flat.sort(key=lambda it: -_num(it.get("pct", 0)))
    return flat[:n]


def _pick_takeaways(p: dict) -> list[tuple[str, str, str]]:
    """Return up to 6 (number, title, body) tuples for the KEY TAKEAWAYS
    slide, built dynamically from the strongest signals.

    Bodies are always readable sentences (not bullet lists) so they render
    cleanly in a 3.55"-wide tile without pipe-character artifacts."""
    subj      = p["name"]
    proj      = p["projected_us"] or 0
    demos     = p.get("demographics") or {}
    top_brand = _top_over_index_brand(p, min_pct=10.0, max_index=300.0)
    picks: list[tuple[str, str, str]] = []

    # 01 — Who they are
    facts = []
    if proj:
        facts.append(f"{_fmt_int_compact(proj)} U.S. digital fans")
    gender = demos.get("gender", {})
    for k in ("Male", "Female"):
        v = _num(gender.get(k, 0))
        if v >= 54:
            facts.append(f"{v:.0f}% {k.lower()}-leaning")
            break
    prime = sum(_num(v) for k, v in demos.get("age", {}).items()
                if str(k).replace("–", "-") in ("18-24", "25-34", "35-44"))
    if prime:
        facts.append(f"{prime:.0f}% aged 18-44")
    for k in ("Black or African American", "Hispanic or Latino", "Asian"):
        idx = _lookup_idx(p, "ethnicity", k)
        if idx >= 130:
            facts.append(f"{k} index {int(idx)}")
            break
    body_01 = "Scale meets resonance. " + "; ".join(facts[:4]) + "."
    picks.append(("01", "Who they are.", body_01))

    # 02 — Where they live
    dmas = _top_dmas_ranked(p, n=3)
    if dmas:
        dma_line = ", ".join(f"{name} ({pct:.1f}%)" for name, pct, _ in dmas)
        picks.append(("02", "Where they live.",
                      f"Concentrated in major digital markets: {dma_line}. "
                      "Prioritize these DMAs for high-density media buys."))

    # 03 — Live events (only if that data exists)
    live = _first_available(p, ["TICKETING PLATFORM", "TICKETING PLATFORMS",
                                "FESTIVAL", "FESTIVALS", "LIVE EVENTS",
                                "EVENTS"])
    if live:
        top3 = live[:3]
        names = ", ".join(
            f"{it['name']} ({int(_num(it.get('index',0)))})" for it in top3
        )
        picks.append(("03", "They show up.",
                      f"Live-event and ticketing behavior far exceeds Gen Pop: "
                      f"{names}. Sponsorship and on-site activation are "
                      "high-fit for this audience."))

    # 04 — Where attention lives. Rank media surfaces by SIGNAL SCORE
    # (reach with a lift kicker), not by index alone — a 2.8%-reach
    # freak-index streamer is not a "highest-fit media surface" no matter
    # how sharp the index reads.
    media = _first_available(p, ["STREAMING/PLATFORM", "STREAMING VIDEO",
                                 "MEDIA", "SOCIAL MEDIA"])
    if media:
        ranked = sorted(media,
                        key=lambda it: -_signal_score(it.get("pct", 0),
                                                        it.get("index", 0)))
        top3 = ranked[:3]
        names = ", ".join(
            f"{it['name']} ({_num(it.get('pct',0)):.0f}%, "
            f"{int(_num(it.get('index',0)))} idx)" for it in top3
        )
        picks.append(("04", "Where attention lives.",
                      f"Highest-fit media surfaces balancing reach and "
                      f"over-index: {names}. Prioritize these over generic "
                      "linear or broad social buys."))

    # 05 — Commerce / financial (if present)
    fin = _first_available(p, ["BANKING", "DIGITAL BANKING", "INVESTMENTS",
                               "CREDIT PROVIDER"])
    if fin:
        top3 = fin[:3]
        names = ", ".join(
            f"{it['name']} ({int(_num(it.get('index',0)))})" for it in top3
        )
        picks.append(("05", "How they transact.",
                      f"Financial-services signature: {names}. "
                      "Category-fit for co-marketing and sponsor plays."))

    # 06 — Whitespace / proof point
    ws = _brief_get(p, "whitespace")
    if ws:
        picks.append(("06", "The whitespace.",
                      _purge_pronouns(str(ws), subj)))
    elif top_brand:
        picks.append(("06", "The taste anchor.",
                      f"{top_brand['name']} indexes at "
                      f"{int(_num(top_brand['index']))} — the strongest "
                      "durable brand signal in the file. Any partnership "
                      "activation should key off this affinity."))

    # Renumber 01..N in produced order
    picks = [(f"{i+1:02d}", t, b) for i, (_, t, b) in enumerate(picks)][:6]
    return picks


# =============================================================================
#  Optional: entry point for CLI testing
# =============================================================================


if __name__ == "__main__":
    import json
    import sys
    if len(sys.argv) < 2:
        print("usage: python3 profile_iq_deck_builder.py <profile_json> [more_profiles...]")
        print("       last arg can be --combined to build a combined deck")
        sys.exit(2)
    combined = False
    args = list(sys.argv[1:])
    if args and args[-1] == "--combined":
        combined = True
        args.pop()
    profiles = []
    for a in args:
        with open(a) as fh:
            profiles.append(json.load(fh))
    if combined and len(profiles) > 1:
        data = build_combined_deck(profiles)
        fname = suggested_combined_filename(profiles)
    else:
        data = build_deck(profiles[0])
        fname = suggested_filename(profiles[0])
    with open(fname, "wb") as fh:
        fh.write(data)
    print(f"Wrote {fname} ({len(data):,} bytes)")

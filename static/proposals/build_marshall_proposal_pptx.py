#!/usr/bin/env python3
"""Generate A. Marshall Hospitality × Crosswalk proposal PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageOps
import os

BASE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(BASE, "images")
OUT = os.path.join(BASE, "A-Marshall-Hospitality-Crosswalk-Proposal.pptx")
LOGO_CW = os.path.join(BASE, "..", "crosswalk-logo-deck.png")
LOGO_CW_WHITE = os.path.join(IMG, "crosswalk-logo-white.png")
LOGO_PU = os.path.join(IMG, "pucketts-logo-ppt.png")
LOGO_DE = os.path.join(IMG, "deacons-logo-ppt.png")
LOGO_SC = os.path.join(IMG, "scouts-pub-logo-ppt.png")
PHOTO_HERO = os.path.join(IMG, "pgr-mule-room_jpg.jpg")
PHOTO_CULLMAN = os.path.join(IMG, "pucketts-cullman-exterior_jpg.jpg")
PHOTO_RIBEYE = os.path.join(IMG, "deacons-ribeye_jpg.jpg")
PHOTO_FARM = os.path.join(IMG, "pucketts-whole-farm_jpg.jpg")
PHOTO_PIGEON = os.path.join(IMG, "pucketts-pigeon-forge-dusk_jpg.jpg")
PHOTO_COCKTAIL = os.path.join(IMG, "deacons-cocktail-hero.png")
BLEND_DIR = os.path.join(IMG, "blends")

CREAM_RGB = (1, 83, 93)  # Valspar Deep Water #01535D

CREAM = RGBColor(0x01, 0x53, 0x5D)
TEXT = RGBColor(0xF5, 0xF8, 0xF9)
MUTED = RGBColor(0xA8, 0xC4, 0xC8)
LIME = RGBColor(0xB8, 0xD9, 0x45)
MAGENTA = RGBColor(0xC0, 0x26, 0xA3)
PURPLE = RGBColor(0xA5, 0xB4, 0xFC)
LIGHT = RGBColor(0xA8, 0xC4, 0xC8)


def build_organic_blend_image(photo_path, out_path, width=1920, height=1080, photo_width_pct=0.58, centering=(0.5, 0.42)):
    """Photo fades organically into Deep Water deck background — no hard crop box."""
    if not photo_path or not os.path.exists(photo_path):
        return None
    canvas = Image.new("RGB", (width, height), CREAM_RGB)
    photo = Image.open(photo_path).convert("RGB")
    photo_w = int(width * photo_width_pct)
    photo = ImageOps.fit(photo, (photo_w, height), Image.LANCZOS, centering=centering)

    mask = Image.new("L", (photo_w, height), 0)
    draw = ImageDraw.Draw(mask)
    fade_end = int(photo_w * 0.42)
    for x in range(photo_w):
        if x < fade_end:
            t = x / max(fade_end, 1)
            alpha = int(255 * (t ** 1.45))
        else:
            alpha = 255
        draw.line([(x, 0), (x, height)], fill=alpha)

    for y in range(height):
        edge = min(y, height - 1 - y) / (height * 0.12)
        edge = min(1.0, max(0.0, edge))
        for x in range(photo_w):
            current = mask.getpixel((x, y))
            mask.putpixel((x, y), int(current * edge))

    canvas.paste(photo, (width - photo_w, 0), mask)
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    canvas.save(out_path, "JPEG", quality=93)
    return out_path


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def set_cream_bg(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = CREAM

    def add_footer(slide, num):
        box = slide.shapes.add_textbox(Inches(0.45), Inches(7.05), Inches(8), Inches(0.3))
        p = box.text_frame.paragraphs[0]
        p.text = "A. MARSHALL HOSPITALITY / PARTNERSHIP PROPOSAL"
        p.font.size = Pt(7)
        p.font.color.rgb = MUTED
        p.font.name = "Arial"
        box2 = slide.shapes.add_textbox(Inches(12.2), Inches(7.05), Inches(0.8), Inches(0.3))
        p2 = box2.text_frame.paragraphs[0]
        p2.text = f"{num:02d}"
        p2.font.size = Pt(7)
        p2.font.color.rgb = MUTED
        p2.alignment = PP_ALIGN.RIGHT

    def add_crosswalk_logo(slide):
        path = LOGO_CW_WHITE if os.path.exists(LOGO_CW_WHITE) else LOGO_CW
        if os.path.exists(path):
            slide.shapes.add_picture(path, Inches(0.45), Inches(0.35), width=Inches(1.35))

    def add_picture_safe(slide, path, left, top, width=None, height=None):
        if not os.path.exists(path):
            return
        if width and height:
            slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width), height=Inches(height))
        elif width:
            slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))
        elif height:
            slide.shapes.add_picture(path, Inches(left), Inches(top), height=Inches(height))
        else:
            slide.shapes.add_picture(path, Inches(left), Inches(top))

    def add_label(slide, text, left, top, width=4, color=MUTED, size=7, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.35))
        p = box.text_frame.paragraphs[0]
        p.text = text.upper()
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.font.bold = True
        if align == PP_ALIGN.RIGHT:
            p.alignment = PP_ALIGN.RIGHT

    def add_headline(slide, text, left, top, width=8, size=32, color=TEXT):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(1.2))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Georgia"
        p.font.bold = True

    def apply_hero_background(slide, photo_path, cache_name, centering=(0.5, 0.42)):
        if not photo_path or not os.path.exists(photo_path):
            set_cream_bg(slide)
            return
        os.makedirs(BLEND_DIR, exist_ok=True)
        blend_path = os.path.join(BLEND_DIR, f"{cache_name}.jpg")
        build_organic_blend_image(photo_path, blend_path, centering=centering)
        set_cream_bg(slide)
        if os.path.exists(blend_path):
            add_picture_safe(slide, blend_path, 0, 0, width=13.333, height=7.5)

    def add_brand_logos_top_right(slide):
        logo_specs = [
            (LOGO_PU, 7.75, 1.28, 0.24),
            (LOGO_SC, 9.15, 1.18, 0.28),
            (LOGO_DE, 10.85, 1.22, 0.22),
        ]
        for path, left, w, h in logo_specs:
            if os.path.exists(path):
                slide.shapes.add_picture(path, Inches(left), Inches(0.52), width=Inches(w), height=Inches(h))

    def add_body(slide, text, left, top, width, height, size=11, italic=False, color=TEXT):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.font.italic = italic

    def add_bullet_list(slide, items, left, top, width, size=10, spacing=0.34):
        for i, item in enumerate(items):
            add_body(slide, f"•  {item}", left, top + i * spacing, width, 0.32, size=size)

    def add_why_it_matters(slide, text, top=1.85, width=5.5):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(top), Inches(0.04), Inches(0.85))
        bar.fill.solid()
        bar.fill.fore_color.rgb = MAGENTA
        bar.line.fill.background()
        add_label(slide, "Why It Matters", 0.6, top - 0.05, 2, MAGENTA)
        add_body(slide, text, 0.6, top + 0.18, width, 1.0, size=11, italic=True)

    def start_hero_slide(photo_path, cache_name, centering=(0.5, 0.42)):
        slide = prs.slides.add_slide(blank)
        apply_hero_background(slide, photo_path, cache_name, centering=centering)
        add_crosswalk_logo(slide)
        return slide

    # Slide 1 — Title
    s = prs.slides.add_slide(blank)
    apply_hero_background(s, PHOTO_HERO, "slide-01-title", centering=(0.5, 0.5))
    add_crosswalk_logo(s)
    add_label(s, "A Partnership Proposal / June 2026", 8.5, 0.4, 4.3, MAGENTA, size=7, align=PP_ALIGN.RIGHT)
    add_brand_logos_top_right(s)
    add_headline(s, "A. Marshall", 0.45, 1.45, 6.2, size=40)
    add_headline(s, "Hospitality", 0.45, 2.05, 6.2, size=40)
    add_body(
        s,
        "behavioral intelligence: making visitors feel like locals",
        0.45,
        3.05,
        5.8,
        0.5,
        size=13,
        color=MUTED,
        italic=True,
    )
    add_label(s, "Presented by Crosswalk Technologies Inc", 0.45, 5.85, 6, MUTED, size=7)
    add_footer(s, 1)

    # Slide 2 — What's inside
    s = prs.slides.add_slide(blank)
    set_cream_bg(s)
    add_crosswalk_logo(s)
    add_label(s, "Overview", 10.5, 0.4, 2.3, align=PP_ALIGN.RIGHT)
    add_headline(s, "What's inside.", 0.45, 1.0, 5.5, size=36)
    intro = (
        "A. Marshall Hospitality operates distinctive brands across multiple markets — Puckett's, Deacon's, and Scout's Pub. "
        "Crosswalk proposes a twelve-month partnership to deliver location-level Profile IQ, competitive benchmarking, and "
        "digital journey intelligence so each restaurant can act on guest behaviors and conversion motivators."
    )
    add_body(s, intro, 0.45, 2.0, 5.2, 2.8, size=11)
    add_label(s, "Contents", 7.0, 1.0, 2, MAGENTA)
    toc = [
        ("01", "Profile IQ by Location"),
        ("02", "Competitive Benchmarking"),
        ("03", "Digital Journeys & Media Conversion"),
        ("04", "Your Restaurant Portfolio"),
        ("05", "Investment Options"),
        ("06", "Next Steps"),
    ]
    y = 1.45
    for num, title in toc:
        add_body(s, num, 7.0, y, 0.5, 0.3, size=12, color=MAGENTA)
        add_body(s, title, 7.55, y, 5, 0.35, size=11)
        y += 0.42
    add_footer(s, 2)

    # Slide 3 — Profile IQ
    s = start_hero_slide(PHOTO_CULLMAN, "slide-03-profile-iq", centering=(0.55, 0.35))
    add_label(s, "A. Marshall Hospitality / Profile IQ", 7.5, 0.4, 5.3, align=PP_ALIGN.RIGHT)
    add_headline(s, "Profile IQ for Every Location", 0.45, 1.0, 10, size=28)
    add_body(s, "A dedicated audience portrait for each restaurant location — refreshed on your cadence.", 0.45, 1.55, 10, 0.4, size=11, color=MUTED, italic=True)
    add_why_it_matters(
        s,
        "Each location draws a different guest — different ages, incomes, dining occasions, and media habits. "
        "Location-level Profile IQ turns that variance into actionable strategy.",
    )
    add_label(s, "Deliverables per location", 0.45, 3.0, 4)
    add_bullet_list(
        s,
        [
            "Demographic profile (age, gender, income, ethnicity, and more)",
            "Behavioral affinities and interests",
            "Geographic concentration and DMA views",
            "Crosswalk dashboard access for your team",
            "Deck-ready exports for internal and external sharing",
        ],
        0.45,
        3.35,
        6.5,
    )
    add_footer(s, 3)

    # Slide 4 — Competitors
    s = start_hero_slide(PHOTO_RIBEYE, "slide-04-competitors", centering=(0.45, 0.5))
    add_label(s, "A. Marshall Hospitality / Competitive Intel", 6.8, 0.4, 6, align=PP_ALIGN.RIGHT)
    add_headline(s, "Competitors Loaded by Location", 0.45, 1.0, 10, size=28)
    add_body(s, "You identify the competitive set; we benchmark each location against them in the Crosswalk dashboard.", 0.45, 1.55, 10, 0.4, size=11, color=MUTED, italic=True)
    add_why_it_matters(
        s,
        "Each market competes differently. We upload your named competitors per location, reflecting indexing and gap analysis "
        "to show real market dynamics, rather than generic category averages.",
    )
    add_label(s, "We handle", 0.45, 3.05, 2)
    add_bullet_list(
        s,
        [
            "Competitor profile upload per location",
            "Side-by-side index and share comparison",
        ],
        0.45,
        3.4,
        5.5,
    )
    add_label(s, "You provide", 6.8, 3.05, 2)
    add_bullet_list(
        s,
        [
            "Location list across your portfolio",
            "2–3 named competitors per site",
            "Quarterly refresh of comp lists",
        ],
        6.8,
        3.4,
        5.5,
    )
    add_footer(s, 4)

    # Slide 5 — Digital Journeys
    s = start_hero_slide(PHOTO_FARM, "slide-05-digital", centering=(0.5, 0.45))
    add_label(s, "A. Marshall Hospitality / Digital Journeys", 6.5, 0.4, 6.3, align=PP_ALIGN.RIGHT)
    add_headline(s, "Digital Journeys & Media Conversion", 0.45, 1.0, 10, size=28)
    add_body(s, "Per-location path-to-plate intelligence — from discovery to reservation to delivery.", 0.45, 1.55, 10, 0.4, size=11, color=MUTED, italic=True)
    add_why_it_matters(
        s,
        "Guests discover you on Yelp, book on OpenTable, and order on DoorDash or Uber Eats, often before they ever walk through the door. "
        "Digital journey maps show where each location wins, leaks, and converts across the platforms that drive covers and checks.",
    )
    add_label(s, "Platforms in scope (per location)", 0.45, 3.0, 5)
    add_bullet_list(
        s,
        ["OpenTable", "Yelp", "DoorDash", "Uber Eats", "Google / Maps", "Meta and programmatic media"],
        0.45,
        3.35,
        5.5,
        spacing=0.3,
    )
    add_label(s, "Included with each digital journey", 0.45, 4.85, 5)
    add_bullet_list(
        s,
        [
            "Media conversion analysis per location",
            "Platform affinity index scores",
            "Full journey stage mapping (discover → book → order → visit)",
        ],
        0.45,
        5.2,
        6.5,
        spacing=0.32,
    )
    add_footer(s, 5)

    # Slide 6 — Portfolio
    s = start_hero_slide(PHOTO_PIGEON, "slide-06-portfolio", centering=(0.6, 0.4))
    add_label(s, "A. Marshall Hospitality / Portfolio", 7.8, 0.4, 5, align=PP_ALIGN.RIGHT)
    add_headline(s, "Your Restaurant Portfolio", 0.45, 1.0, 5.8, size=28)
    add_body(s, "One partnership covering every location across Puckett's, Deacon's, and Scout's Pub.", 0.45, 1.55, 5.8, 0.5, size=11, color=MUTED, italic=True)
    add_label(s, "What's included at every location", 0.45, 2.2, 5, MAGENTA)
    add_bullet_list(
        s,
        [
            "Profile IQ — full demographic and behavioral portrait",
            "Competitive benchmarking with your named comp set",
            "Digital journey mapping across reservation, delivery, and discovery platforms",
            "Media conversion analysis tied to each journey",
            "Crosswalk dashboard access and deck-ready exports",
        ],
        0.45,
        2.55,
        5.8,
        size=11,
        spacing=0.38,
    )
    add_footer(s, 6)

    # Slide 7 — Investment (cocktail hero)
    cocktail_photo = PHOTO_COCKTAIL if os.path.exists(PHOTO_COCKTAIL) else PHOTO_RIBEYE
    s = start_hero_slide(cocktail_photo, "slide-07-investment", centering=(0.35, 0.55))
    add_label(s, "A. Marshall Hospitality / Investment", 7.8, 0.4, 5, align=PP_ALIGN.RIGHT)
    add_headline(s, "Investment Options", 0.45, 1.0, 10, size=28)
    add_body(
        s,
        "Twelve-month contract · All tiers include full Profile IQ, competitive uploads, digital journeys and media conversion per location.",
        0.45,
        1.55,
        10,
        0.5,
        size=11,
        color=MUTED,
        italic=True,
    )
    add_body(
        s,
        "Panel participation lowers your cash investment while enriching the behavioral dataset that powers sharper profiles. "
        "The more locations participate, the better the intelligence and the lower your cost.",
        0.45,
        2.15,
        6.5,
        1.0,
        size=11,
        italic=True,
    )
    tiers = [
        ("OPTION A — ALL CASH", "Full service, no panel obligations.", "$10,000/mo", False),
        ("OPTION B — PANEL INCENTIVE", "One panel ingestion incentive per location per month.", "$5,000/mo", False),
        ("OPTION C — FULL PANEL PARTNERSHIP", "Weekly incentives per location per month.", "No Charge", True),
    ]
    y = 3.35
    for opt, desc, price, highlight in tiers:
        add_label(s, opt, 0.45, y, 6, MAGENTA if highlight else MUTED, size=8)
        add_body(s, desc, 0.45, y + 0.22, 6.5, 0.35, size=10)
        pb = s.shapes.add_textbox(Inches(9.5), Inches(y), Inches(3), Inches(0.5))
        pp = pb.text_frame.paragraphs[0]
        pp.text = price
        pp.font.size = Pt(22)
        pp.font.name = "Georgia"
        pp.font.bold = True
        pp.font.color.rgb = MAGENTA if highlight else TEXT
        pp.alignment = PP_ALIGN.RIGHT
        y += 0.85
    add_body(s, "All options: 12-month initial term", 0.45, 6.35, 12, 0.4, size=9, color=MUTED, italic=True)
    add_footer(s, 7)

    # Slide 8 — Next steps
    s = prs.slides.add_slide(blank)
    set_cream_bg(s)
    add_crosswalk_logo(s)
    add_label(s, "A. Marshall Hospitality / Next Steps", 7.5, 0.4, 5.3, align=PP_ALIGN.RIGHT)
    add_headline(s, "Let's get started.", 0.45, 1.2, 8, size=36)
    add_body(
        s,
        "We're ready to scope locations, load your competitive sets, and stand up your first Profile IQ dashboards within weeks of contract execution.",
        0.45,
        2.1,
        6.5,
        0.9,
        size=11,
    )
    add_bullet_list(
        s,
        [
            "Confirm location list across your portfolio",
            "Select investment tier and panel cadence",
            "Kickoff: competitor uploads and journey mapping",
            "First location profiles delivered to dashboard",
        ],
        0.45,
        3.1,
        6.5,
        size=11,
        spacing=0.38,
    )
    add_label(s, "Crosswalk Technologies Inc", 8.5, 5.5, 4.3, align=PP_ALIGN.RIGHT)
    box = s.shapes.add_textbox(Inches(8.5), Inches(5.85), Inches(4.3), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = "jessie@crosswalknyc.com"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT
    p.alignment = PP_ALIGN.RIGHT
    add_footer(s, 8)

    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build_deck()

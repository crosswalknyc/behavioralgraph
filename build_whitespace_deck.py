#!/usr/bin/env python3
"""Generate 'Adult Animation: The Female White Space' PowerPoint deck."""

import sys, os
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), ".pip_packages"))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Adult_Animation_Female_White_Space.pptx")

CYAN = RGBColor(0x4D, 0xC9, 0xD4)
ORANGE = RGBColor(0xF2, 0x8C, 0x28)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
LIGHT_CYAN_BG = RGBColor(0xE8, 0xF8, 0xF9)
LIGHT_ORANGE_BG = RGBColor(0xFE, 0xF0, 0xE0)
CYAN_DARK = RGBColor(0x37, 0x9E, 0xA7)
ORANGE_DARK = RGBColor(0xD4, 0x74, 0x10)
MID_BLUE = RGBColor(0x5B, 0x9B, 0xD5)

SHOWS = (
    "The Simpsons, Family Guy, South Park, Bob's Burgers, American Dad!, "
    "Krapopolis, Grimsburg, Rick and Morty, Smiling Friends, Common Side Effects, "
    "Royal Crackers, Harley Quinn, Lazarus, Invincible, Hazbin Hotel, "
    "The Legend of Vox Machina, Solar Opposites, Hit-Monkey, X-Men '97, Big Mouth"
)


def white_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE


def add_text(slide, text, left, top, width, height, size=11,
             color=DARK, bold=False, italic=False, align=PP_ALIGN.LEFT,
             font="Arial", anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = font
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = align
    return box


def add_stat_box(slide, value, label, left, top, width, val_color, bg_color):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.9),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg_color
    rect.line.fill.background()
    rect.shadow.inherit = False
    add_text(slide, value, left, top + 0.08, width, 0.45,
             size=28, color=val_color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, left, top + 0.55, width, 0.3,
             size=9, color=MUTED, align=PP_ALIGN.CENTER)


def style_chart_text(chart, cat_font_size=Pt(9), val_font_size=Pt(8)):
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = cat_font_size
    cat_axis.tick_labels.font.color.rgb = DARK
    cat_axis.tick_labels.font.name = "Arial"
    cat_axis.format.line.fill.background()
    cat_axis.major_tick_mark = 2  # XL_TICK_MARK.NONE

    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = val_font_size
    val_axis.tick_labels.font.color.rgb = MUTED
    val_axis.tick_labels.font.name = "Arial"
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.color.rgb = LIGHT_GRAY
    val_axis.major_gridlines.format.line.width = Pt(0.5)
    val_axis.major_tick_mark = 2


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── SLIDE 1: TITLE + STATS ──
    s1 = prs.slides.add_slide(blank)
    white_bg(s1)

    accent_bar = s1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.06),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = CYAN
    accent_bar.line.fill.background()

    add_text(s1, "Adult Animation", 0.7, 0.5, 10, 0.7,
             size=38, color=DARK, bold=True, font="Georgia")
    add_text(s1, "The Female White Space", 0.7, 1.1, 10, 0.6,
             size=32, color=CYAN, bold=True, font="Georgia")

    add_text(s1, (
        f"The Adult Animation audience reflects viewers of at least one of these "
        f"20 titles: {SHOWS}."
    ), 0.7, 2.0, 8.5, 1.0, size=10, color=MUTED, italic=True)

    add_text(s1,
        "Compared against the proven female comedy audiences of Tina Fey and Amy Poehler.",
        0.7, 2.85, 8.5, 0.4, size=10, color=MUTED, italic=True)

    stat_y = 3.8
    stat_w = 2.6
    gap = 0.35
    x0 = 0.7
    add_stat_box(s1, "65%", "ADULT ANIM. MALE", x0, stat_y, stat_w, MUTED, LIGHT_ORANGE_BG)
    add_stat_box(s1, "33%", "ADULT ANIM. FEMALE", x0 + stat_w + gap, stat_y, stat_w, ORANGE, LIGHT_ORANGE_BG)
    add_stat_box(s1, "65%", "TINA FEY FEMALE", x0 + 2 * (stat_w + gap), stat_y, stat_w, CYAN, LIGHT_CYAN_BG)
    add_stat_box(s1, "59%", "AMY POEHLER FEMALE", x0 + 3 * (stat_w + gap), stat_y, stat_w, CYAN, LIGHT_CYAN_BG)

    callout = s1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(5.2), Inches(11.9), Inches(1.5),
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = LIGHT_CYAN_BG
    callout.line.fill.background()
    callout.shadow.inherit = False

    bar_l = s1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7), Inches(5.2), Inches(0.06), Inches(1.5),
    )
    bar_l.fill.solid()
    bar_l.fill.fore_color.rgb = CYAN
    bar_l.line.fill.background()

    add_text(s1, "The Opportunity", 1.05, 5.3, 5, 0.3,
             size=13, color=DARK, bold=True)
    add_text(s1, (
        "Women 35-54 are a proven, passionate comedy audience — and adult animation "
        "has barely spoken to them. Across 20 titles, adult animation draws just 33% "
        "women (Index 66). Meanwhile Tina Fey and Amy Poehler command 60-65% female "
        "audiences. An animated comedy channeling their sensibility could unlock a "
        "massive, under-served demographic already primed for the format."
    ), 1.05, 5.65, 11.2, 1.0, size=10, color=MUTED)

    # ── SLIDE 2: GENDER COMPOSITION ──
    s2 = prs.slides.add_slide(blank)
    white_bg(s2)

    add_text(s2, "Who's Watching: Gender Split", 0.7, 0.4, 10, 0.6,
             size=28, color=DARK, bold=True, font="Georgia")
    add_text(s2, (
        "Profile % of each audience by gender. Adult Animation is nearly 2:1 male. "
        "Fey and Poehler audiences flip that ratio entirely."
    ), 0.7, 1.05, 9, 0.5, size=11, color=MUTED)

    chart_data = CategoryChartData()
    chart_data.categories = ["Female", "Male"]
    chart_data.add_series("Adult Animation", (32.87, 65.30))
    chart_data.add_series("Tina Fey",        (65.43, 33.72))
    chart_data.add_series("Amy Poehler",     (59.27, 39.10))

    chart_frame = s2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(1.7),
        Inches(11.9), Inches(5.2), chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.legend.font.color.rgb = DARK

    s2_colors = [ORANGE, CYAN, CYAN_DARK]
    for i, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = s2_colors[i]
        series.format.line.fill.background()
        dl = series.data_labels
        dl.show_value = True
        dl.font.size = Pt(10)
        dl.font.color.rgb = DARK
        dl.font.bold = True
        dl.number_format = '0.0"%"'
        dl.number_format_is_linked = False

    style_chart_text(chart)
    chart.value_axis.maximum_scale = 80

    # ── SLIDE 3: FEMALE INDEX ──
    s3 = prs.slides.add_slide(blank)
    white_bg(s3)

    add_text(s3, "Female Engagement Index", 0.7, 0.4, 10, 0.6,
             size=28, color=DARK, bold=True, font="Georgia")
    add_text(s3, (
        "How much each audience over- or under-represents women relative to "
        "the U.S. population. 100 = expected share; Adult Animation sits at half that."
    ), 0.7, 1.05, 9, 0.5, size=11, color=MUTED)

    chart_data2 = CategoryChartData()
    chart_data2.categories = ["Adult Animation", "Tina Fey", "Amy Poehler"]
    chart_data2.add_series("Female Index", (66, 132, 120))

    chart_frame2 = s3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.5), Inches(1.7),
        Inches(10.3), Inches(5.2), chart_data2,
    )
    chart2 = chart_frame2.chart
    chart2.has_legend = False

    bar_colors = [ORANGE, CYAN, CYAN_DARK]
    series2 = chart2.series[0]
    for idx, color in enumerate(bar_colors):
        pt = series2.points[idx]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = color

    series2.format.line.fill.background()
    dl2 = series2.data_labels
    dl2.show_value = True
    dl2.font.size = Pt(14)
    dl2.font.color.rgb = DARK
    dl2.font.bold = True

    style_chart_text(chart2, cat_font_size=Pt(11))
    chart2.value_axis.maximum_scale = 160

    # ── SLIDE 4: 35-54 SWEET SPOT ──
    s4 = prs.slides.add_slide(blank)
    white_bg(s4)

    add_text(s4, "The 35-54 Sweet Spot", 0.7, 0.4, 10, 0.6,
             size=28, color=DARK, bold=True, font="Georgia")
    add_text(s4, (
        "Index vs. U.S. population for the two age brackets where Fey and Poehler "
        "over-index most aggressively — and where Adult Animation already has strong "
        "viewership, but driven almost entirely by men."
    ), 0.7, 1.05, 9, 0.5, size=11, color=MUTED)

    chart_data3 = CategoryChartData()
    chart_data3.categories = ["35-44", "45-54"]
    chart_data3.add_series("Adult Animation", (170, 121))
    chart_data3.add_series("Tina Fey",        (160, 187))
    chart_data3.add_series("Amy Poehler",     (227, 200))

    chart_frame3 = s4.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.5), Inches(1.7),
        Inches(10.3), Inches(5.2), chart_data3,
    )
    chart3 = chart_frame3.chart
    chart3.has_legend = True
    chart3.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart3.legend.include_in_layout = False
    chart3.legend.font.size = Pt(10)
    chart3.legend.font.color.rgb = DARK

    for i, series in enumerate(chart3.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = s2_colors[i]
        series.format.line.fill.background()
        dl = series.data_labels
        dl.show_value = True
        dl.font.size = Pt(10)
        dl.font.color.rgb = DARK
        dl.font.bold = True

    style_chart_text(chart3, cat_font_size=Pt(11))
    chart3.value_axis.maximum_scale = 260

    # ── SLIDE 5: TAKEAWAY ──
    s5 = prs.slides.add_slide(blank)
    white_bg(s5)

    accent_bar2 = s5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.06),
    )
    accent_bar2.fill.solid()
    accent_bar2.fill.fore_color.rgb = ORANGE
    accent_bar2.line.fill.background()

    add_text(s5, "The White Space", 0.7, 0.7, 10, 0.6,
             size=34, color=DARK, bold=True, font="Georgia")
    add_text(s5, "is a Female 35-54 Audience", 0.7, 1.3, 10, 0.5,
             size=26, color=CYAN, bold=True, font="Georgia")

    findings = [
        ("33%", "of the adult animation audience is female — Index 66, "
         "roughly half the expected share of women vs. the U.S. population."),
        ("65%", "of Tina Fey's audience is female (Index 132) and 59% "
         "of Amy Poehler's (Index 120) — proving women 35-54 are a deeply "
         "engaged comedy demo."),
        ("227", "is Amy Poehler's index among 35-44 year-olds. Adult Animation "
         "indexes at 170 in the same bracket — but those viewers are overwhelmingly male."),
    ]

    for i, (num, desc) in enumerate(findings):
        y = 2.4 + i * 1.35
        circ = s5.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.7), Inches(y), Inches(0.9), Inches(0.9),
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = CYAN if i > 0 else ORANGE
        circ.line.fill.background()
        circ.shadow.inherit = False
        add_text(s5, num, 0.7, y + 0.15, 0.9, 0.6,
                 size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s5, desc, 1.85, y + 0.1, 10, 0.8,
                 size=13, color=DARK)

    add_text(s5, (
        "An animated comedy that channels the sensibility of Fey or Poehler "
        "could unlock a massive, under-served demographic — one that already "
        "watches comedy, already indexes above average in the category's core "
        "age range, and simply has not been given content made for them."
    ), 0.7, 6.2, 11.9, 0.8, size=12, color=MUTED, italic=True)

    prs.save(OUT)
    print(f"Saved → {OUT}")


if __name__ == "__main__":
    build_deck()

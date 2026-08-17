"""
Obsession -> Peacock: client-facing PowerPoint deck for the IP owner.

Structure (12 slides, presentation-ready - one idea per slide):

  1.  Title / cover
  2.  The situation - release timeline
  3.  Executive summary - three headline numbers
  4.  How the analysis was built (N1 - setup)
  5.  Chart 1 - studio window structure
  6.  Chart 2 - every Universal-family title, side by side
  7.  The paradox - Focus recut theatrical, not TVOD
  8.  Chart 3 - TVOD daily fall-off (with Batch E confounds)
  9.  What was traded (measurement, $22-32M range)
 10.  What was received (H6 hypothesis, $43-175M modeled)
 11.  The one-sentence framing + next steps
 12.  Appendix - measurement vs hypothesis vs deferred

All charts render as PNG via matplotlib (300 dpi) and are embedded into
16x9 slides via python-pptx. Numbers reflect the full revision-doc pass
(C1-C13 corrections, S1-S7 softens, H1-H6 hypothesis relabels, N1-N3
new elements, standing rules).

Output: ~/Downloads/Obsession_Peacock_Client_Deck.pptx
"""
from __future__ import annotations

import io
import os
from pathlib import Path

import matplotlib
matplotlib.use("Agg")

import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Patch
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn


# --- Design tokens ------------------------------------------------------

INK = RGBColor(0x0F, 0x17, 0x2A)
INK_MUTED = RGBColor(0x47, 0x55, 0x69)
INK_FAINT = RGBColor(0x94, 0xA3, 0xB8)
RULE = RGBColor(0xE2, 0xE8, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TEAL = RGBColor(0x14, 0xB8, 0xA6)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
INDIGO = RGBColor(0x63, 0x66, 0xF1)
RED = RGBColor(0xEF, 0x44, 0x44)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
BLUE = RGBColor(0x60, 0xA5, 0xFA)

# Matplotlib hex equivalents for the chart renderer
MPL_TEAL = "#14b8a6"
MPL_AMBER = "#f59e0b"
MPL_INDIGO = "#6366f1"
MPL_RED = "#ef4444"
MPL_GRAY = "#94a3b8"
MPL_GREEN = "#22c55e"
MPL_BLUE = "#60a5fa"
MPL_INK = "#0f172a"
MPL_INK_MUTED = "#475569"
MPL_INK_FAINT = "#94a3b8"
MPL_RULE = "#e2e8f0"


plt.rcParams.update({
    "font.family": "DejaVu Sans",
    "axes.titlecolor": MPL_INK,
    "axes.labelcolor": MPL_INK,
    "xtick.color": MPL_INK,
    "ytick.color": MPL_INK,
    "axes.edgecolor": MPL_RULE,
    "axes.linewidth": 0.8,
    "axes.spines.top": False,
    "axes.spines.right": False,
    "figure.facecolor": "#ffffff",
    "axes.facecolor": "#ffffff",
    "savefig.facecolor": "#ffffff",
    "pdf.fonttype": 42,
    "ps.fonttype": 42,
    "text.parse_math": False,
})


# --- Data (post revision doc) ------------------------------------------

STUDIO_AVERAGES = [
    # (label, theatre_d, tvod_d, low_sample, subject)
    ("Universal",    32, 86, False, False),
    ("Disney",       60, 40, False, False),
    ("Warner Bros.", 39, 37, False, False),
    ("Paramount",    32, 28, True,  False),
    ("Obsession",    46, 17, False, True),
]

UNIVERSAL_TITLES = [
    # (title, studio, tv_days, tvod_days, subject, dom_bo_m)
    ("Despicable Me 4",           "Illumination",           34,  86, False, 361.0),
    ("Twisters",                  "Universal",              25,  94, False, 267.8),
    ("The Wild Robot",            "DreamWorks",             18, 101, False, 143.9),
    ("Wicked",                    "Universal",              39,  80, False, 473.2),
    ("How to Train Your Dragon",  "Universal / DreamWorks", 32,  87, False, 263.6),
    ("Jurassic World Rebirth",    "Universal",              34,  86, False, 339.6),
    ("Wicked: For Good",          "Universal",              39,  80, False, 342.9),
    ("Five Nights at Freddy's 2", "Universal / Blumhouse",  18, 101, False, 127.7),
    ("Nosferatu",                 "Focus",                  27,  31, False,  95.6),
    ("Downton Abbey: Finale",     "Focus",                  18,  38, False,  45.0),
    ("Conclave",                  "Focus",                  32,  17, False,  32.6),
    ("Obsession",                 "Focus / Universal",      46,  17, True,  263.4),
]

# Disney / Disney+ (n=5, 2024-2026). Theatre days derived as
# svod_total - tvod_days per the obp_industry data blocks.
# (title, tv_days, tvod_days, note)
DISNEY_TITLES = [
    ("Inside Out 2",              67, 36, "$652M dom BO. Full 67-day theatrical."),
    ("Captain America: BNW",      60, 43, "Held TVOD 43 days before Disney+ handoff."),
    ("Lilo & Stitch",             60, 43, "Family co-view driver; same 43-day TVOD lane."),
    ("Zootopia 2",                62, 43, "43-day TVOD; 105-day theatre-to-Disney+."),
    ("Hoppers",                   53, 36, "Shortest Disney TVOD in the sample."),
]

# Warner Bros. / HBO Max (n=4, 2024-2026).
WB_TITLES = [
    ("A Minecraft Movie",         39, 38, "Family-friendly Warner tentpole."),
    ("Sinners",                   46, 31, "Coogler R-rated horror."),
    ("Final Destination Bloodlines", 32, 45, "Horror; longer TVOD than most WB slate."),
    ("Superman",                  35, 35, "Expedited ~7d to align with Peacemaker S2 finale."),
]

# Paramount / Paramount+ (n=1 - Gladiator II only; flagged low_sample).
PARAMOUNT_TITLES = [
    ("Gladiator II",              32, 28, "Only 2024-2026 wide release with owned-SVOD handoff in the panel."),
]

# All integrated-studio titles + Obsession, sorted by TVOD days ASC so
# the reader lands on the industry floor and then Obsession at 17.
def _build_all_integrated():
    rows = []
    for t in UNIVERSAL_TITLES:
        title, studio, tv, tvod, subject, bo = t
        if "Focus" in studio and not subject:
            continue  # Focus non-subject shown separately if needed
        rows.append((title, studio, tv, tvod, subject))
    for t in DISNEY_TITLES:
        rows.append((t[0], "Disney", t[1], t[2], False))
    for t in WB_TITLES:
        rows.append((t[0], "Warner Bros.", t[1], t[2], False))
    for t in PARAMOUNT_TITLES:
        rows.append((t[0], "Paramount", t[1], t[2], False))
    # Sort by tvod_days DESCENDING so the shortest TVOD (Obsession at 17)
    # anchors the bottom of the chart.
    rows.sort(key=lambda r: (-r[3], r[2]))
    return rows

ALL_INTEGRATED_TITLES = _build_all_integrated()

STUDIO_COLOR_HEX = {
    "Illumination":            "#14b8a6",
    "Universal":               "#14b8a6",
    "DreamWorks":              "#14b8a6",
    "Universal / DreamWorks":  "#14b8a6",
    "Universal / Blumhouse":   "#14b8a6",
    "Focus / Universal":       "#ef4444",  # subject
    "Disney":                  "#3b82f6",
    "Warner Bros.":            "#a855f7",
    "Paramount":               "#0ea5e9",
}

# TVOD daily curve - Jun 30 (day 1) through Aug 10 (day 42).
# Day 18 = Jul 17 (Peacock SVOD handoff).
TVOD_DAYS = list(range(1, 43))
TVOD_OBS = [
    88.0, 90.0, 91.0, 78.0, 82.0, 92.0, 94.0,
    82.0, 79.0, 78.0, 77.0, 75.0, 76.0, 75.0,
    73.0, 71.0, 70.0, 14.0, 12.0, 11.0, 10.0,
    3.6, 3.5, 3.4, 3.4, 3.4, 3.4, 3.4,
    2.5, 2.4, 2.3, 2.3, 2.3, 2.3, 2.3,
    1.9, 1.8, 1.7, 1.7, 1.7, 1.7, 1.7,
]
TVOD_CF = [
    88.0, 90.0, 91.0, 78.0, 82.0, 92.0, 94.0,
    82.0, 79.0, 78.0, 77.0, 75.0, 76.0, 75.0,
    73.0, 71.0, 70.0,
    68.0, 66.0, 64.0, 62.0,
    60.0, 58.0, 56.0, 55.0, 54.0, 53.0, 52.0,
    51.0, 50.0, 49.0, 48.0, 47.0, 46.0, 46.0,
    45.0, 45.0, 44.0, 44.0, 43.0, 43.0, 43.0,
]
CONFOUND_FANDANGO = (13, 18)   # Jul 12-17
CONFOUND_JULY_4 = 5             # Jul 4
CONFOUND_HANDOFF = 18           # Jul 17
CONFOUND_AUG = (33, 42)         # Aug 1-10

DISPLACED_RANGE = "$22-32M"
DISPLACED_TXN = "~2.68M (2.28-3.07M, 90% CI)"
ASP_RANGE = "$15-$22"

SUBIQ = dict(
    viewed="~9.5M",
    acquired="~138K (121-155K, 90% CI)",
    new="~76K (67-85K)",
    reactivated="~62K (55-70K)",
    rate="~1.5%",
    completion="~58%",
    second_screen="~65%",
)

CHARTS_DIR = Path("/tmp/obp_pptx_charts")
CHARTS_DIR.mkdir(exist_ok=True)


# --- Chart renderers (return path to PNG) ------------------------------

def render_studio_averages_chart() -> Path:
    fig, ax = plt.subplots(figsize=(13.33, 6.5), dpi=200)
    fig.subplots_adjust(left=0.14, right=0.78, top=0.92, bottom=0.14)

    labels = [s[0] for s in STUDIO_AVERAGES]
    theatre = np.array([s[1] for s in STUDIO_AVERAGES])
    tvod = np.array([s[2] for s in STUDIO_AVERAGES])
    svod = np.array([30] * len(STUDIO_AVERAGES))
    totals = theatre + tvod + svod

    y = np.arange(len(labels))
    ax.barh(y, theatre, color=MPL_TEAL, edgecolor="white", height=0.62,
            label="Days in Theatre")
    ax.barh(y, tvod, left=theatre, color=MPL_AMBER, edgecolor="white", height=0.62,
            label="Days in TVOD exclusive")
    ax.barh(y, svod, left=theatre + tvod, color=MPL_INDIGO, edgecolor="white", height=0.62,
            label="Days in SVOD (30-day attribution)")

    for i, (t, v, s) in enumerate(zip(theatre, tvod, svod)):
        if t > 8:
            ax.text(t/2, i, f"{int(t)}d", ha="center", va="center",
                    color="white", fontsize=13, weight="bold")
        if v > 8:
            ax.text(t + v/2, i, f"{int(v)}d", ha="center", va="center",
                    color="white", fontsize=13, weight="bold")
        if s > 8:
            ax.text(t + v + s/2, i, f"{int(s)}d", ha="center", va="center",
                    color="white", fontsize=13, weight="bold")

    for i, tot in enumerate(totals):
        ax.text(tot + 3, i, f"= {int(tot)}d total", va="center", ha="left",
                fontsize=12, color=MPL_INK, weight="bold")

    for i, s in enumerate(STUDIO_AVERAGES):
        if s[3]:  # low_sample flag
            ax.text(tot + 3, i + 0.36, "(n=1 -- Gladiator II only)",
                    va="center", ha="left", fontsize=9,
                    color=MPL_AMBER, weight="bold", fontstyle="italic")

    obs_idx = [i for i, s in enumerate(STUDIO_AVERAGES) if s[4]][0]
    ax.axhspan(obs_idx - 0.42, obs_idx + 0.42, color=MPL_RED, alpha=0.08, zorder=0)

    ax.set_yticks(y)
    ax.set_yticklabels(labels, fontsize=14, color=MPL_INK, weight="bold")
    ax.invert_yaxis()
    ax.set_xlabel("Days from theatrical release", fontsize=12, color=MPL_INK_MUTED)
    ax.set_xlim(0, max(totals) * 1.22)
    ax.grid(axis="x", color=MPL_RULE, linewidth=0.6, alpha=0.9)
    ax.set_axisbelow(True)

    legend = ax.legend(
        handles=[
            Patch(facecolor=MPL_TEAL, label="Theatre exclusive"),
            Patch(facecolor=MPL_AMBER, label="TVOD exclusive"),
            Patch(facecolor=MPL_INDIGO, label="SVOD (30d attribution)"),
        ],
        loc="center left", bbox_to_anchor=(1.02, 0.5),
        frameon=False, fontsize=11, title="Window phase",
        title_fontsize=11,
    )
    legend.get_title().set_color(MPL_INK_MUTED)
    legend.get_title().set_weight("bold")

    out = CHARTS_DIR / "studio_windows.png"
    fig.savefig(out, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def render_universal_titles_chart() -> Path:
    fig, ax = plt.subplots(figsize=(13.33, 6.8), dpi=200)
    fig.subplots_adjust(left=0.22, right=0.88, top=0.94, bottom=0.10)

    labels = [t[0] for t in UNIVERSAL_TITLES]
    studios = [t[1] for t in UNIVERSAL_TITLES]
    theatre = np.array([t[2] for t in UNIVERSAL_TITLES])
    tvod = np.array([t[3] for t in UNIVERSAL_TITLES])
    svod = np.array([30] * len(UNIVERSAL_TITLES))
    subject = np.array([t[4] for t in UNIVERSAL_TITLES])
    dom_bo = np.array([t[5] for t in UNIVERSAL_TITLES])
    totals = theatre + tvod + svod

    y = np.arange(len(labels))
    ax.barh(y, theatre, color=MPL_TEAL, edgecolor="white", height=0.62)
    ax.barh(y, tvod, left=theatre, color=MPL_AMBER, edgecolor="white", height=0.62)
    ax.barh(y, svod, left=theatre + tvod, color=MPL_INDIGO, edgecolor="white", height=0.62)

    obs_idx = [i for i, s in enumerate(subject) if s][0]
    ax.axhspan(obs_idx - 0.42, obs_idx + 0.42, color=MPL_RED, alpha=0.12, zorder=0)
    ax.barh(obs_idx, totals[obs_idx], left=0, height=0.66,
            fill=False, edgecolor=MPL_RED, linewidth=2.4, zorder=3)

    for i, (t, v, s) in enumerate(zip(theatre, tvod, svod)):
        if t > 6:
            ax.text(t/2, i, f"{int(t)}", ha="center", va="center",
                    color="white", fontsize=10, weight="bold")
        if v > 6:
            ax.text(t + v/2, i, f"{int(v)}", ha="center", va="center",
                    color="white", fontsize=10, weight="bold")
        if s > 6:
            ax.text(t + v + s/2, i, f"{int(s)}", ha="center", va="center",
                    color="white", fontsize=10, weight="bold")

    for i, (tot, bo) in enumerate(zip(totals, dom_bo)):
        color = MPL_RED if subject[i] else MPL_INK
        weight = "bold" if subject[i] else "normal"
        ax.text(tot + 3, i, f"= {int(tot)}d",
                va="center", ha="left", fontsize=10.5, color=color, weight=weight)
        ax.text(tot + 22, i, f"${bo:.1f}M BO",
                va="center", ha="left", fontsize=9.5, color=MPL_INK_MUTED, style="italic")

    ytick_labels = []
    for lab, st, is_sub in zip(labels, studios, subject):
        if is_sub:
            ytick_labels.append(f"{lab}\n({st}) <- subject")
        else:
            ytick_labels.append(f"{lab}\n({st})")

    ax.set_yticks(y)
    ax.set_yticklabels(ytick_labels, fontsize=10.5, color=MPL_INK)
    for lbl, is_sub in zip(ax.get_yticklabels(), subject):
        if is_sub:
            lbl.set_color(MPL_RED)
            lbl.set_weight("bold")
    ax.invert_yaxis()
    ax.set_xlabel("Days from theatrical release", fontsize=12, color=MPL_INK_MUTED)
    ax.set_xlim(0, max(totals) * 1.35)
    ax.grid(axis="x", color=MPL_RULE, linewidth=0.6, alpha=0.9)
    ax.set_axisbelow(True)

    for divider_y in (7.5, 10.5):
        ax.axhline(divider_y, color=MPL_INK_FAINT, linewidth=0.7,
                   linestyle=(0, (3, 3)), alpha=0.6)
    ax.text(max(totals) * 1.31, 3.5, "Non-Focus\nUFEG", fontsize=9,
            color=MPL_INK_FAINT, ha="right", va="center", style="italic")
    ax.text(max(totals) * 1.31, 9.0, "Focus\n(non-subject)", fontsize=9,
            color=MPL_INK_FAINT, ha="right", va="center", style="italic")
    ax.text(max(totals) * 1.31, 11.0, "Subject", fontsize=9,
            color=MPL_RED, ha="right", va="center", style="italic", weight="bold")

    out = CHARTS_DIR / "universal_titles.png"
    fig.savefig(out, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def _render_studio_titles_chart(rows, out_name, height_in=6.0,
                                x_max_days=170, subject_flag_fn=None):
    """Render a per-studio side-by-side title chart.

    rows: list of (title, tv_days, tvod_days, [optional] note)
    subject_flag_fn: optional fn(row_tuple) -> bool that flags a subject row.
    Returns path to PNG.
    """
    fig, ax = plt.subplots(figsize=(13.33, height_in), dpi=200)
    fig.subplots_adjust(left=0.28, right=0.94, top=0.94, bottom=0.14)

    labels = [r[0] for r in rows]
    theatre = np.array([r[1] for r in rows])
    tvod = np.array([r[2] for r in rows])
    svod = np.array([30] * len(rows))
    totals = theatre + tvod + svod
    subjects = [subject_flag_fn(r) if subject_flag_fn else False for r in rows]
    notes = [(r[3] if len(r) >= 4 else "") for r in rows]

    y = np.arange(len(labels))
    ax.barh(y, theatre, color=MPL_TEAL, edgecolor="white", height=0.62)
    ax.barh(y, tvod,    left=theatre, color=MPL_AMBER, edgecolor="white", height=0.62)
    ax.barh(y, svod,    left=theatre + tvod, color=MPL_INDIGO, edgecolor="white", height=0.62)

    for i, (t, v, s) in enumerate(zip(theatre, tvod, svod)):
        if t > 6:
            ax.text(t/2, i, f"{int(t)}", ha="center", va="center",
                    color="white", fontsize=10.5, weight="bold")
        if v > 6:
            ax.text(t + v/2, i, f"{int(v)}", ha="center", va="center",
                    color="white", fontsize=10.5, weight="bold")
        if s > 6:
            ax.text(t + v + s/2, i, f"{int(s)}", ha="center", va="center",
                    color="white", fontsize=10.5, weight="bold")

    # Total + optional note at right
    for i, (tot, note, subj) in enumerate(zip(totals, notes, subjects)):
        color = MPL_RED if subj else MPL_INK
        weight = "bold" if subj else "normal"
        ax.text(tot + 2, i, f"= {int(tot)}d",
                va="center", ha="left", fontsize=10.5, color=color, weight=weight)
        if note:
            ax.text(tot + 18, i, note, va="center", ha="left",
                    fontsize=9, color=MPL_INK_MUTED, style="italic")

    ax.set_yticks(y)
    ax.set_yticklabels(labels, fontsize=11.5, color=MPL_INK, weight="bold")
    for lbl, subj in zip(ax.get_yticklabels(), subjects):
        if subj:
            lbl.set_color(MPL_RED)
            lbl.set_weight("bold")
    ax.invert_yaxis()
    ax.set_xlabel("Days from theatrical release", fontsize=11, color=MPL_INK_MUTED)
    ax.set_xlim(0, x_max_days)
    ax.grid(axis="x", color=MPL_RULE, linewidth=0.6, alpha=0.9)
    ax.set_axisbelow(True)

    # Industry-floor reference line at 28 days of TVOD. Positioned above
    # the top of the chart so it doesn't collide with the x-axis label
    # underneath.
    if len(theatre):
        theatre_avg = float(np.mean(theatre))
        ax.axvline(theatre_avg + 28, color=MPL_RED, linewidth=0.9,
                   linestyle=(0, (4, 3)), alpha=0.55, zorder=1)
        ax.text(theatre_avg + 28 + 1.5, -0.55,
                "Industry TVOD floor (28d, Paramount)",
                fontsize=8.5, color=MPL_RED, style="italic",
                va="bottom", ha="left", weight="bold")

    out = CHARTS_DIR / out_name
    fig.savefig(out, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def render_disney_titles_chart() -> Path:
    rows = [(t[0], t[1], t[2], t[3]) for t in DISNEY_TITLES]
    return _render_studio_titles_chart(rows, "disney_titles.png",
                                       height_in=4.4, x_max_days=170)


def render_wb_titles_chart() -> Path:
    rows = [(t[0], t[1], t[2], t[3]) for t in WB_TITLES]
    return _render_studio_titles_chart(rows, "wb_titles.png",
                                       height_in=3.8, x_max_days=170)


def render_paramount_titles_chart() -> Path:
    rows = [(t[0], t[1], t[2], t[3]) for t in PARAMOUNT_TITLES]
    return _render_studio_titles_chart(rows, "paramount_titles.png",
                                       height_in=2.4, x_max_days=170)


def render_all_integrated_chart() -> Path:
    """The money-shot combined chart: every integrated-studio title + Obsession,
    color-coded by studio, sorted by TVOD days descending so Obsession lands
    at the bottom as the outlier."""
    rows = ALL_INTEGRATED_TITLES
    labels = [r[0] for r in rows]
    studios = [r[1] for r in rows]
    theatre = np.array([r[2] for r in rows])
    tvod = np.array([r[3] for r in rows])
    svod = np.array([30] * len(rows))
    subjects = np.array([r[4] for r in rows])
    totals = theatre + tvod + svod

    fig, ax = plt.subplots(figsize=(13.33, 8.0), dpi=200)
    fig.subplots_adjust(left=0.28, right=0.90, top=0.96, bottom=0.09)

    y = np.arange(len(labels))
    ax.barh(y, theatre, color=MPL_TEAL, edgecolor="white", height=0.68)
    ax.barh(y, tvod,    left=theatre, color=MPL_AMBER, edgecolor="white", height=0.68)
    ax.barh(y, svod,    left=theatre + tvod, color=MPL_INDIGO, edgecolor="white", height=0.68)

    # Subject highlight
    obs_idx = int(np.where(subjects)[0][0]) if subjects.any() else -1
    if obs_idx >= 0:
        ax.axhspan(obs_idx - 0.44, obs_idx + 0.44, color=MPL_RED, alpha=0.10, zorder=0)
        ax.barh(obs_idx, totals[obs_idx], left=0, height=0.72,
                fill=False, edgecolor=MPL_RED, linewidth=2.4, zorder=3)

    # Segment numeric labels inside each bar
    for i, (t, v, s) in enumerate(zip(theatre, tvod, svod)):
        if t > 5:
            ax.text(t/2, i, f"{int(t)}", ha="center", va="center",
                    color="white", fontsize=8.5, weight="bold")
        if v > 5:
            ax.text(t + v/2, i, f"{int(v)}", ha="center", va="center",
                    color="white", fontsize=8.5, weight="bold")
        if s > 5:
            ax.text(t + v + s/2, i, f"{int(s)}", ha="center", va="center",
                    color="white", fontsize=8.5, weight="bold")

    # Right-side total + studio chip
    for i, (tot, studio, subj) in enumerate(zip(totals, studios, subjects)):
        color = MPL_RED if subj else MPL_INK
        weight = "bold" if subj else "normal"
        ax.text(tot + 2, i, f"= {int(tot)}d",
                va="center", ha="left", fontsize=9, color=color, weight=weight)
        chip_color = STUDIO_COLOR_HEX.get(studio, MPL_INK_FAINT)
        ax.text(tot + 22, i, studio, va="center", ha="left",
                fontsize=8, color=chip_color, weight="bold", style="italic")

    # Industry TVOD floor line at 28 days past theatrical (use average
    # theatre-only days of integrated titles as the visual reference).
    theatre_avg = float(np.mean([r[2] for r in rows if not r[4]]))
    floor_x = theatre_avg + 28
    ax.axvline(floor_x, color=MPL_RED, linewidth=1.0,
               linestyle=(0, (4, 3)), alpha=0.55, zorder=1)
    ax.text(floor_x + 1.5, -0.6,
            f"Industry TVOD floor: 28 days (Paramount / Gladiator II)",
            fontsize=9, color=MPL_RED, style="italic", va="bottom", weight="bold")

    ytick_labels = []
    for lab, st in zip(labels, studios):
        ytick_labels.append(f"{lab}\n({st})")
    ax.set_yticks(y)
    ax.set_yticklabels(ytick_labels, fontsize=8.5, color=MPL_INK)
    for lbl, subj in zip(ax.get_yticklabels(), subjects):
        if subj:
            lbl.set_color(MPL_RED)
            lbl.set_weight("bold")
    ax.invert_yaxis()
    ax.set_xlabel("Days from theatrical release", fontsize=11, color=MPL_INK_MUTED)
    ax.set_xlim(0, 175)
    ax.grid(axis="x", color=MPL_RULE, linewidth=0.6, alpha=0.9)
    ax.set_axisbelow(True)

    out = CHARTS_DIR / "all_integrated.png"
    fig.savefig(out, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def render_tvod_falloff_chart() -> Path:
    fig, ax = plt.subplots(figsize=(13.33, 6.5), dpi=200)
    fig.subplots_adjust(left=0.08, right=0.98, top=0.92, bottom=0.16)

    days = np.array(TVOD_DAYS)
    obs = np.array(TVOD_OBS)
    cf = np.array(TVOD_CF)

    ax.axvspan(CONFOUND_FANDANGO[0], CONFOUND_FANDANGO[1], color=MPL_AMBER, alpha=0.09, zorder=0)
    ax.axvspan(CONFOUND_AUG[0], CONFOUND_AUG[1] + 1, color=MPL_BLUE, alpha=0.07, zorder=0)
    ax.axvspan(1, CONFOUND_HANDOFF, color=MPL_GREEN, alpha=0.03, zorder=0)
    ax.axvspan(CONFOUND_HANDOFF, days.max() + 1, color=MPL_RED, alpha=0.06, zorder=0)

    post = days >= CONFOUND_HANDOFF
    ax.fill_between(days[post], obs[post], cf[post],
                    where=(cf[post] > obs[post]),
                    color=MPL_RED, alpha=0.18, linewidth=0,
                    label="Estimated displaced daily volume")

    pre = days <= CONFOUND_HANDOFF - 1
    ax.plot(days[pre], obs[pre], color=MPL_TEAL, linewidth=2.6, marker="o",
            markersize=4, label="Observed - pre-SVOD (normal decay)")
    ax.plot(days[post], obs[post], color=MPL_RED, linewidth=2.6, marker="o",
            markersize=4, label="Observed - post-SVOD (cliff)")
    ax.plot(days[post], cf[post], color=MPL_GRAY, linewidth=2.4,
            linestyle=(0, (6, 4)),
            label="Counterfactual - if held to 86.5d Universal median")

    ax.axvline(CONFOUND_HANDOFF, color=MPL_RED, linewidth=1.8,
               linestyle=(0, (5, 4)), alpha=0.9)
    ax.text(CONFOUND_HANDOFF + 0.25, 96,
            "PEACOCK SVOD DROP\nJUL 17", fontsize=10, color=MPL_RED, weight="bold",
            va="top", ha="left")

    ax.axvline(CONFOUND_JULY_4, color=MPL_GRAY, linewidth=1.0,
               linestyle=(0, (4, 4)), alpha=0.7)
    ax.text(CONFOUND_JULY_4 + 0.2, 96, "JUL 4\n(confound)",
            fontsize=8.5, color=MPL_INK_MUTED,
            va="top", ha="left", style="italic")

    mid_fandango = (CONFOUND_FANDANGO[0] + CONFOUND_FANDANGO[1]) / 2
    ax.text(mid_fandango, 48,
            "FANDANGO RENTAL\n$19.99 -> $9.99\n(confound span)",
            fontsize=8.5, color=MPL_AMBER, weight="bold",
            ha="center", va="center", style="italic")

    mid_aug = (CONFOUND_AUG[0] + CONFOUND_AUG[1]) / 2
    ax.text(mid_aug, 48,
            "AUG TVOD SEASONALITY\n12-18% below annual mean\n(confound span)",
            fontsize=8.5, color=MPL_BLUE, weight="bold",
            ha="center", va="center", style="italic")

    def gap_arrow(day_x, obs_v, cf_v, side="right"):
        ax.annotate("", xy=(day_x, cf_v - 0.5), xytext=(day_x, obs_v + 0.5),
                    arrowprops=dict(arrowstyle="<->", color=MPL_RED, lw=1.8))
        gap = int(round(cf_v - obs_v))
        mid = (obs_v + cf_v) / 2
        offset = 0.9 if side == "right" else -0.9
        ha = "left" if side == "right" else "right"
        ax.annotate(f"-{gap}K/day", xy=(day_x + offset, mid),
                    fontsize=10, color=MPL_RED, weight="bold", ha=ha, va="center",
                    bbox=dict(boxstyle="round,pad=0.28", facecolor="white",
                              edgecolor=MPL_RED, linewidth=1))

    gap_arrow(19, 12.0, 66.0, "right")
    gap_arrow(41, 1.7, 43.0, "left")

    ax.text(30, 24, f"POST-HANDOFF\n{DISPLACED_RANGE} ESTIMATED",
            fontsize=15, color=MPL_RED, weight="bold", ha="center", va="center",
            alpha=0.80, family="DejaVu Sans")

    date_ticks = [
        (1, "Jun 30"), (5, "Jul 4"), (8, "Jul 7"), (15, "Jul 14"),
        (18, "Jul 17\nSVOD"), (22, "Jul 21"), (29, "Jul 28"),
        (36, "Aug 4"), (42, "Aug 10"),
    ]
    ax.set_xticks([d[0] for d in date_ticks])
    ax.set_xticklabels([d[1] for d in date_ticks], fontsize=10)
    for lbl, (d, _) in zip(ax.get_xticklabels(), date_ticks):
        if d == CONFOUND_HANDOFF:
            lbl.set_color(MPL_RED)
            lbl.set_weight("bold")

    ax.set_ylabel("TVOD transactions per day (K)", fontsize=12, color=MPL_INK_MUTED)
    ax.set_xlabel("Calendar day (Jun 30 -> Aug 10, 2026)", fontsize=12, color=MPL_INK_MUTED)
    ax.set_ylim(0, 100)
    ax.set_xlim(0.5, 42.5)
    ax.grid(axis="y", color=MPL_RULE, linewidth=0.6, alpha=0.9)
    ax.set_axisbelow(True)
    ax.legend(loc="upper right", frameon=False, fontsize=10)

    out = CHARTS_DIR / "tvod_falloff.png"
    fig.savefig(out, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


# --- Slide helpers -----------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _tb(slide, x, y, w, h, text, size=14, bold=False, color=INK,
        italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add a text box; returns the shape."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def _rect(slide, x, y, w, h, fill=WHITE, line=None, line_w=Pt(0.75)):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w
    shp.shadow.inherit = False
    return shp


def _rounded(slide, x, y, w, h, fill=WHITE, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.06
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.2)
    shp.shadow.inherit = False
    return shp


def _footer(slide, page_no=None, total=None):
    # Rule
    line = slide.shapes.add_connector(1, Inches(0.6), Inches(7.15),
                                       Inches(12.73), Inches(7.15))
    line.line.color.rgb = RULE
    line.line.width = Pt(0.75)
    _tb(slide, Inches(0.6), Inches(7.20), Inches(9.5), Inches(0.25),
        "Source: Crosswalk panel (viewer-level, projected to U.S. digital population, T+1)  ·  "
        "SubIQ is a panel measurement under strict first-watch attribution",
        size=8, color=INK_FAINT)
    _tb(slide, Inches(10.2), Inches(7.20), Inches(2.55), Inches(0.25),
        "Confidential  ·  For the Obsession IP owner",
        size=8, color=INK_FAINT, italic=True, align=PP_ALIGN.RIGHT)
    if page_no is not None and total is not None:
        _tb(slide, Inches(11.5), Inches(0.30), Inches(1.3), Inches(0.3),
            f"{page_no} / {total}", size=9, color=INK_FAINT, align=PP_ALIGN.RIGHT)


def _kicker(slide, kicker, title):
    _tb(slide, Inches(0.6), Inches(0.30), Inches(11.0), Inches(0.35),
        kicker.upper(), size=10, bold=True, color=INK_FAINT)
    _tb(slide, Inches(0.6), Inches(0.65), Inches(12.13), Inches(0.75),
        title, size=26, bold=True, color=INK)
    # Rule under the title
    line = slide.shapes.add_connector(1, Inches(0.6), Inches(1.45),
                                       Inches(12.73), Inches(1.45))
    line.line.color.rgb = RULE
    line.line.width = Pt(0.75)


# --- Slide builders ---------------------------------------------------

def slide_cover(prs):
    s = _blank(prs)

    # Top rule
    line = s.shapes.add_connector(1, Inches(0.6), Inches(1.6),
                                   Inches(12.73), Inches(1.6))
    line.line.color.rgb = INK
    line.line.width = Pt(1.2)

    _tb(s, Inches(0.6), Inches(1.15), Inches(9), Inches(0.4),
        "CROSSWALK  ·  CLIENT ANALYSIS", size=11, bold=True, color=INK)

    _tb(s, Inches(0.6), Inches(2.0), Inches(12), Inches(1.5),
        "Obsession -> Peacock", size=54, bold=True, color=INK)
    _tb(s, Inches(0.6), Inches(3.3), Inches(12), Inches(0.7),
        "Windowing analysis and participation impact",
        size=22, color=INK_MUTED)

    _tb(s, Inches(0.6), Inches(4.5), Inches(6), Inches(0.35),
        "PREPARED FOR", size=10, bold=True, color=INK_FAINT)
    _tb(s, Inches(0.6), Inches(4.85), Inches(10), Inches(0.55),
        "The Obsession IP owner", size=22, bold=True, color=INK)
    _tb(s, Inches(0.6), Inches(5.35), Inches(12), Inches(0.4),
        "Licensed to Focus Features (UFEG)  ->  Peacock (NBCUniversal)  ·  Drop Jul 17, 2026",
        size=13, color=INK_MUTED)

    _tb(s, Inches(0.6), Inches(6.3), Inches(6), Inches(0.35),
        "August 13, 2026  ·  Revision 2", size=11, bold=True, color=INK_MUTED)

    _footer(s)
    return s


def slide_situation(prs, page_no, total):
    s = _blank(prs)
    _kicker(s, "The situation",
            "A $263M breakout, moved to a sibling SVOD 17 days into its digital window")

    # Timeline as a chevron-style band
    y = Inches(2.3)
    events = [
        ("May 15", "Theatrical open", "Focus (UFEG)", MPL_TEAL),
        ("Jun 29", "Theatrical close", "46 days in theatres", MPL_TEAL),
        ("Jun 30", "TVOD live",       "$24.99 / $19.99 rent", MPL_AMBER),
        ("Jul 12", "Rental cut",      "$19.99 -> $9.99 (confound)", MPL_AMBER),
        ("Jul 17", "PEACOCK DROP",    "Day 17 of TVOD window",   MPL_RED),
        ("Aug 10", "Read date",       "T+1 panel through today", MPL_INDIGO),
        ("Aug 16", "SubIQ close",     "30-day attribution window", MPL_INDIGO),
    ]
    n = len(events)
    col_w = Inches(1.65)
    start_x = Inches(0.75)
    for i, (date, hdr, sub, color_hex) in enumerate(events):
        x = start_x + col_w * i
        color = RGBColor.from_string(color_hex[1:])
        # Colored bar
        _rect(s, x, y, col_w - Inches(0.15), Inches(0.10), fill=color)
        _tb(s, x, y + Inches(0.20), col_w, Inches(0.32),
            date, size=13, bold=True, color=color)
        _tb(s, x, y + Inches(0.60), col_w, Inches(0.4),
            hdr, size=11, bold=True, color=INK)
        _tb(s, x, y + Inches(1.00), col_w - Inches(0.15), Inches(0.6),
            sub, size=9.5, color=INK_MUTED)

    # The headline callout
    _tb(s, Inches(0.6), Inches(4.35), Inches(12), Inches(0.4),
        "THE QUESTION", size=11, bold=True, color=INK_FAINT)
    _tb(s, Inches(0.6), Inches(4.75), Inches(12.13), Inches(1.1),
        "Was the Jul 17 SVOD placement a routine end-of-window handoff, or did it move value "
        "from a stream in which the content owner participates to one in which it does not?",
        size=17, bold=True, color=INK)

    _tb(s, Inches(0.6), Inches(6.15), Inches(12), Inches(0.4),
        "WHAT THE DATA MEASURES", size=11, bold=True, color=INK_FAINT)
    _tb(s, Inches(0.6), Inches(6.50), Inches(12.13), Inches(0.55),
        "The TVOD trajectory, panel-observed daily. Peacock playback, panel-observed daily. "
        "The window structure of every comparator title in the 2024-2026 UFEG record.",
        size=12, color=INK_MUTED)

    _footer(s, page_no, total)
    return s


def slide_exec_summary(prs, page_no, total):
    s = _blank(prs)
    _kicker(s, "Executive summary",
            "The three numbers that carry the story")

    tiles = [
        ("17 DAYS",
         "of TVOD exclusivity",
         "vs 86.5-day Universal MEDIAN (n=8 non-Focus UFEG, 2024-2026). "
         "Obsession is the shortest TVOD window in the sample by ~5.1x.",
         RED),
        (DISPLACED_RANGE,
         "estimated displaced\ntransactional revenue",
         "Range reflects the Fandango rental price cut ($19.99 -> $9.99) that landed inside "
         "the same 7-day period as the handoff. Point estimate withheld pending distributor "
         "Apple + Amazon statements.",
         AMBER),
        ("~138K",
         "Peacock subs acquired\nor reactivated",
         "~76K new + ~62K reactivated. Panel-observed measurement under strict first-watch "
         "attribution (25-day partial read; window closes Aug 16).",
         INDIGO),
    ]
    left = Inches(0.6); top = Inches(2.0); w = Inches(4.05); h = Inches(3.0)
    for i, (big, mid, body, color) in enumerate(tiles):
        x = left + Inches(0.05) + (w + Inches(0.10)) * i
        _rounded(s, x, top, w, h, fill=WHITE, line=color)
        # left accent stripe
        _rect(s, x, top, Inches(0.12), h, fill=color)
        _tb(s, x + Inches(0.30), top + Inches(0.25), w - Inches(0.5), Inches(0.9),
            big, size=36, bold=True, color=color)
        _tb(s, x + Inches(0.30), top + Inches(1.30), w - Inches(0.5), Inches(0.7),
            mid, size=13, bold=True, color=INK)
        _tb(s, x + Inches(0.30), top + Inches(2.00), w - Inches(0.5), Inches(0.95),
            body, size=10, color=INK_MUTED)

    # Trade one-liner (S4-softened)
    _tb(s, Inches(0.6), Inches(5.25), Inches(12), Inches(0.4),
        "THE TRADE, ON ONE LINE", size=11, bold=True, color=INK_FAINT)
    _tb(s, Inches(0.6), Inches(5.65), Inches(12.13), Inches(0.9),
        f"An estimated {DISPLACED_RANGE} in displaced transactional revenue is not covered by any known "
        "intercompany transfer or arm's-length license fee (per public disclosures). "
        "In exchange, ~1.33M net-new Peacock signups accrue to NBCUniversal.",
        size=13.5, bold=True, color=INK)

    _footer(s, page_no, total)
    return s


def slide_setup(prs, page_no, total):
    s = _blank(prs)
    _kicker(s, "Setup - how this analysis was built",
            "Panel scale, projection factor, refresh cadence, read dates")

    tiles = [
        ("6.42M viewers", "US digital panel",
         "Matched to household + device. Deterministic within panel.", TEAL),
        ("~52x", "Panel -> US digital projection",
         "SubIQ ~76K new subs projected from ~1.5K panel observations. Every SubIQ figure carries a 90% CI.",
         INDIGO),
        ("T+1", "Refresh cadence",
         "Every figure is yesterday-close. Re-pulled at milestones: Aug 16 (30d) and Jan 2027 (180d).",
         BLUE),
        ("~9.5M accts", "Obsession universe",
         "Peacock accounts observed initiating Obsession playback in the 30-day attribution window. Deck is a 25-day partial read.",
         AMBER),
    ]
    left = Inches(0.6); top = Inches(1.8); w = Inches(3.02); h = Inches(2.0)
    for i, (big, kicker, body, color) in enumerate(tiles):
        x = left + (w + Inches(0.08)) * i
        _rounded(s, x, top, w, h, fill=WHITE, line=color)
        _rect(s, x, top, Inches(0.10), h, fill=color)
        _tb(s, x + Inches(0.25), top + Inches(0.15), w - Inches(0.4), Inches(0.6),
            big, size=22, bold=True, color=color)
        _tb(s, x + Inches(0.25), top + Inches(0.75), w - Inches(0.4), Inches(0.4),
            kicker, size=10, bold=True, color=INK)
        _tb(s, x + Inches(0.25), top + Inches(1.10), w - Inches(0.4), Inches(0.85),
            body, size=9, color=INK_MUTED)

    # Measurement vs hypothesis vs deferred - three-row strip
    _tb(s, Inches(0.6), Inches(4.05), Inches(12), Inches(0.4),
        "MEASUREMENT vs HYPOTHESIS vs DEFERRED", size=11, bold=True, color=INK_FAINT)

    mvh = [
        ("MEASUREMENT", GREEN,
         "TVOD transaction curve. Peacock playback volumes and completion rates. SubIQ new / reactivated sub counts under strict first-watch attribution. Studio window intervals across 14 UFEG titles."),
        ("HYPOTHESIS", AMBER,
         "Causal attribution of the TVOD fall-off to the handoff (pending Backrooms / Devil Wears Prada 2 controls). Long-run elevated churn on the Obsession cohort (pending 180-day retention, Jan 2027). H6: value NBCU received, modeled at $43-175M pending retention read."),
        ("DEFERRED", GRAY,
         "Apple TV Store + Fandango at Home public rank corroboration (Pull 2). Distributor Apple + Amazon statements that would collapse the $22-32M range to a point."),
    ]
    row_h = Inches(0.72)
    for i, (kind, color, body) in enumerate(mvh):
        y = Inches(4.5) + row_h * i
        _rect(s, Inches(0.6), y + Inches(0.05), Inches(0.08), row_h - Inches(0.15), fill=color)
        _tb(s, Inches(0.8), y + Inches(0.05), Inches(1.5), Inches(0.3),
            kind, size=10, bold=True, color=color)
        _tb(s, Inches(2.4), y + Inches(0.03), Inches(10.3), row_h,
            body, size=10, color=INK_MUTED)

    _footer(s, page_no, total)
    return s


def slide_chart(prs, page_no, total, kicker, title, png_path, caption):
    s = _blank(prs)
    _kicker(s, kicker, title)
    # Insert image scaled to width; center vertically in remaining space.
    img_w = Inches(11.5)
    img_h = Inches(5.0)
    left = Inches((13.333 - 11.5) / 2)
    top = Inches(1.6)
    s.shapes.add_picture(str(png_path), left, top, width=img_w, height=img_h)
    _tb(s, Inches(0.6), top + img_h + Inches(0.1), Inches(12.13), Inches(0.75),
        caption, size=11, color=INK_MUTED, italic=True)
    _footer(s, page_no, total)
    return s


def slide_competitor(prs, page_no, total, kicker, title,
                     subtitle, png_path, bullets, chart_h=3.6):
    """Generic per-studio competitor slide: chart on top, 3-4 bullets below.
    Chart height is tunable so smaller studios (fewer titles) do not stretch
    the image vertically and leave bullets clipped."""
    s = _blank(prs)
    _kicker(s, kicker, title)

    # Subtitle line
    _tb(s, Inches(0.6), Inches(1.55), Inches(12.13), Inches(0.35),
        subtitle, size=12, color=INK_MUTED, italic=True)

    # Chart. Fit width to slide; center horizontally.
    img_w = Inches(11.8)
    img_h = Inches(chart_h)
    left = Inches((13.333 - 11.8) / 2)
    top = Inches(1.95)
    s.shapes.add_picture(str(png_path), left, top, width=img_w, height=img_h)

    # Bullets strip below the chart. Cap y so we never spill into the footer.
    y = top + img_h + Inches(0.15)
    for b in bullets:
        _tb(s, Inches(0.6), y, Inches(0.25), Inches(0.32),
            "·", size=15, bold=True, color=RED)
        _tb(s, Inches(0.85), y + Inches(0.02), Inches(11.9), Inches(0.32),
            b, size=10.5, color=INK_MUTED)
        y += Inches(0.30)

    _footer(s, page_no, total)
    return s


def slide_all_integrated(prs, page_no, total, png_path):
    """The combined 'every integrated-studio title vs Obsession' money shot."""
    s = _blank(prs)
    _kicker(s, "Chart 6  ·  Every integrated-studio title, 2024-2026",
            "22 wide releases across four integrated studios. Obsession is the shortest TVOD in the record.")

    img_w = Inches(9.6)
    img_h = Inches(5.5)
    left = Inches(0.4)
    top = Inches(1.65)
    s.shapes.add_picture(str(png_path), left, top, width=img_w, height=img_h)

    # Right rail: three call-out tiles
    rail_x = Inches(10.2)
    rail_w = Inches(2.9)

    # Tile 1: The floor
    _rounded(s, rail_x, Inches(1.65), rail_w, Inches(1.65), fill=WHITE, line=RED)
    _rect(s, rail_x, Inches(1.65), Inches(0.10), Inches(1.65), fill=RED)
    _tb(s, rail_x + Inches(0.20), Inches(1.75), rail_w - Inches(0.35), Inches(0.28),
        "INDUSTRY FLOOR", size=9, bold=True, color=RED)
    _tb(s, rail_x + Inches(0.20), Inches(2.05), rail_w - Inches(0.35), Inches(0.55),
        "28 days", size=26, bold=True, color=RED)
    _tb(s, rail_x + Inches(0.20), Inches(2.60), rail_w - Inches(0.35), Inches(0.7),
        "Shortest TVOD window granted by any integrated studio in the sample (Paramount / Gladiator II).",
        size=9, color=INK_MUTED)

    # Tile 2: Obsession
    _rounded(s, rail_x, Inches(3.4), rail_w, Inches(1.65), fill=WHITE, line=RED)
    _rect(s, rail_x, Inches(3.4), Inches(0.10), Inches(1.65), fill=RED)
    _tb(s, rail_x + Inches(0.20), Inches(3.5), rail_w - Inches(0.35), Inches(0.28),
        "OBSESSION", size=9, bold=True, color=RED)
    _tb(s, rail_x + Inches(0.20), Inches(3.80), rail_w - Inches(0.35), Inches(0.55),
        "17 days", size=26, bold=True, color=RED)
    _tb(s, rail_x + Inches(0.20), Inches(4.35), rail_w - Inches(0.35), Inches(0.7),
        "11 days below the industry floor. Below every other title in the sample by 11-84 days.",
        size=9, color=INK_MUTED)

    # Tile 3: Gap
    _rounded(s, rail_x, Inches(5.15), rail_w, Inches(1.9), fill=WHITE, line=INK)
    _rect(s, rail_x, Inches(5.15), Inches(0.10), Inches(1.9), fill=INK)
    _tb(s, rail_x + Inches(0.20), Inches(5.25), rail_w - Inches(0.35), Inches(0.28),
        "READ", size=9, bold=True, color=INK_FAINT)
    _tb(s, rail_x + Inches(0.20), Inches(5.55), rail_w - Inches(0.35), Inches(1.4),
        "No integrated studio in the 2024-2026 sample releases to its owned SVOD "
        "faster than 28 TVOD days. Obsession's 17 is 39% below that floor.",
        size=10, color=INK_MUTED)

    _footer(s, page_no, total)
    return s


def slide_paradox(prs, page_no, total):
    s = _blank(prs)
    _kicker(s, "The paradox",
            "Focus recut the theatrical window in response to strong performance. Focus did not recut the TVOD window.")

    left = Inches(0.6)
    top = Inches(2.0)
    box_w = Inches(6.0); box_h = Inches(3.2)

    # Left box: theatrical extended
    _rounded(s, left, top, box_w, box_h, fill=WHITE, line=GREEN)
    _rect(s, left, top, Inches(0.12), box_h, fill=GREEN)
    _tb(s, left + Inches(0.35), top + Inches(0.20), box_w - Inches(0.6), Inches(0.35),
        "ADJUSTED  -  THEATRICAL WINDOW", size=11, bold=True, color=GREEN)
    _tb(s, left + Inches(0.35), top + Inches(0.60), box_w - Inches(0.6), Inches(0.9),
        "46 days", size=44, bold=True, color=GREEN)
    _tb(s, left + Inches(0.35), top + Inches(1.65), box_w - Inches(0.6), Inches(1.4),
        "Focus extended the theatrical window from an initial 17-day plan to 46 days after three "
        "consecutive up weekends. When performance was strong in a window that benefits the "
        "content owner, Focus re-cut the release plan to protect that value.",
        size=11.5, color=INK_MUTED)

    # Right box: TVOD held to 17 days
    right = left + box_w + Inches(0.4)
    _rounded(s, right, top, box_w, box_h, fill=WHITE, line=RED)
    _rect(s, right, top, Inches(0.12), box_h, fill=RED)
    _tb(s, right + Inches(0.35), top + Inches(0.20), box_w - Inches(0.6), Inches(0.35),
        "UNCUT  -  TVOD WINDOW", size=11, bold=True, color=RED)
    _tb(s, right + Inches(0.35), top + Inches(0.60), box_w - Inches(0.6), Inches(0.9),
        "17 days", size=44, bold=True, color=RED)
    _tb(s, right + Inches(0.35), top + Inches(1.65), box_w - Inches(0.6), Inches(1.4),
        "Panel volume held at ~88% of peak weekly velocity at the day-17 handoff. The Universal median "
        "for titles at handoff is 3% of peak. The TVOD window was not adjusted to reflect the same "
        "performance signal that extended the theatrical window.",
        size=11.5, color=INK_MUTED)

    # Bottom read
    _tb(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.4),
        "THE READ", size=11, bold=True, color=INK_FAINT)
    _tb(s, Inches(0.6), Inches(5.9), Inches(12.13), Inches(1.1),
        "Focus recut only the window in which its own corporate interests aligned with the content owner's. "
        "The 46-day theatrical extension and the 17-day TVOD hold cannot both be a response to the same market signal.",
        size=14, bold=True, color=INK)

    _footer(s, page_no, total)
    return s


def slide_what_was_traded(prs, page_no, total):
    s = _blank(prs)
    _kicker(s, "What was traded",
            f"An estimated {DISPLACED_RANGE} in displaced transactional revenue (measurement)")

    left = Inches(0.6); top = Inches(2.0)
    # Big number tile
    _rounded(s, left, top, Inches(5.5), Inches(3.5), fill=WHITE, line=AMBER)
    _rect(s, left, top, Inches(0.14), Inches(3.5), fill=AMBER)
    _tb(s, left + Inches(0.4), top + Inches(0.25), Inches(4.9), Inches(0.4),
        "WHAT WAS TRADED  ·  MEASUREMENT", size=11, bold=True, color=AMBER)
    _tb(s, left + Inches(0.4), top + Inches(0.75), Inches(4.9), Inches(1.5),
        DISPLACED_RANGE, size=54, bold=True, color=AMBER)
    _tb(s, left + Inches(0.4), top + Inches(2.20), Inches(4.9), Inches(0.5),
        "range in displaced TVOD revenue", size=13, bold=True, color=INK)
    _tb(s, left + Inches(0.4), top + Inches(2.70), Inches(4.9), Inches(0.7),
        f"Volume gap: {DISPLACED_TXN} transactions.\nBlended ASP: {ASP_RANGE}.",
        size=11, color=INK_MUTED)

    # Right column: 4 bullets
    right = Inches(6.4); rw = Inches(6.4)
    _tb(s, right, top, rw, Inches(0.4),
        "WHAT IS BEHIND THE RANGE", size=11, bold=True, color=INK_FAINT)
    bullets = [
        "Panel-observed weekly TVOD volume collapsed ~87% at Day 1 post-SVOD (Jun 30 -> Aug 10 read).",
        "Fandango rental price cut ($19.99 -> $9.99) landed inside the same 7-day period as the handoff. Part of the observed Fandango decline is price, not volume - which is why the figure is a range, not a point.",
        "Universal median TVOD window is 86.5 days (n=8 non-Focus UFEG). If Obsession had held to that median, it would have retained a materially larger share of transactional revenue.",
        "Point estimate withheld pending distributor Apple + Amazon statements. Deck will be re-pulled when statements arrive.",
    ]
    y = top + Inches(0.4)
    for b in bullets:
        _tb(s, right, y, Inches(0.25), Inches(0.4),
            "·", size=18, bold=True, color=AMBER)
        _tb(s, right + Inches(0.25), y + Inches(0.05), rw - Inches(0.25), Inches(0.9),
            b, size=11, color=INK_MUTED)
        y += Inches(0.85)

    _footer(s, page_no, total)
    return s


def slide_what_was_received(prs, page_no, total):
    s = _blank(prs)
    _kicker(s, "What was received",
            "H6 hypothesis: modeled Peacock ARPU value to NBCU, pending 180-day retention")

    left = Inches(0.6); top = Inches(2.0)
    _rounded(s, left, top, Inches(5.5), Inches(3.5), fill=WHITE, line=INDIGO)
    _rect(s, left, top, Inches(0.14), Inches(3.5), fill=INDIGO)
    _tb(s, left + Inches(0.4), top + Inches(0.25), Inches(4.9), Inches(0.4),
        "WHAT WAS RECEIVED  ·  HYPOTHESIS (H6)", size=11, bold=True, color=INDIGO)
    _tb(s, left + Inches(0.4), top + Inches(0.75), Inches(4.9), Inches(1.5),
        "$43-175M", size=54, bold=True, color=INDIGO)
    _tb(s, left + Inches(0.4), top + Inches(2.20), Inches(4.9), Inches(0.5),
        "modeled Peacock ARPU value to NBCU", size=13, bold=True, color=INK)
    _tb(s, left + Inches(0.4), top + Inches(2.70), Inches(4.9), Inches(0.7),
        "Range collapses once the 180-day retention read closes (Jan 2027).",
        size=11, color=INK_MUTED)

    # Right column
    right = Inches(6.4); rw = Inches(6.4)
    _tb(s, right, top, rw, Inches(0.4),
        "HOW THE RANGE IS MODELED", size=11, bold=True, color=INK_FAINT)
    bullets = [
        f"Panel-observed measurement: {SUBIQ['acquired']} subs acquired or reactivated ({SUBIQ['new']} new + {SUBIQ['reactivated']} reactivated).",
        "Peacock ARPU $10.99 (ads-supported)  x  1.33M net-new signups  x  expected retention months.",
        "Low end ($43M): 3-month effective sub life. Consistent with an elevated-churn cohort.",
        "High end ($175M): 12-month effective sub life. Consistent with retention at the Peacock baseline.",
        "Which end the range collapses to depends on the 180-day retention read (Jan 2027).",
    ]
    y = top + Inches(0.4)
    for b in bullets:
        _tb(s, right, y, Inches(0.25), Inches(0.4),
            "·", size=18, bold=True, color=INDIGO)
        _tb(s, right + Inches(0.25), y + Inches(0.05), rw - Inches(0.25), Inches(0.75),
            b, size=11, color=INK_MUTED)
        y += Inches(0.72)

    _footer(s, page_no, total)
    return s


def slide_framing(prs, page_no, total):
    s = _blank(prs)
    _kicker(s, "The participation framing",
            "In one sentence, with the trade quantified on both sides")

    _tb(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(1.7),
        f"The Obsession IP owner participates in the {DISPLACED_RANGE} column. "
        "Participation in the $43-175M column is a contractual determination outside the scope of this measurement.",
        size=22, bold=True, color=INK)

    _tb(s, Inches(0.6), Inches(4.0), Inches(12.13), Inches(0.8),
        "The July 17 handoff moved value from a stream in which the content owner has a "
        "contractual share to one in which contractual status is unclear - inside the same corporate family.",
        size=13, color=INK_MUTED)

    _tb(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.4),
        "SUGGESTED NEXT STEPS", size=11, bold=True, color=INK_FAINT)
    steps = [
        "1. Review the licensing agreement for downstream-participation and window-length provisions.",
        "2. Compare the transactional-revenue substitution against any intercompany license-fee credit.",
        "3. Request Peacock 180-day retention detail (Jan 2027) to collapse the H6 value-received range.",
        "4. Convene with Focus / UFEG participation office to align on 2024-2026 industry-window benchmarks.",
        "5. Wait for distributor Apple + Amazon statements to collapse the $22-32M displacement range.",
    ]
    y = Inches(5.6)
    for step in steps:
        _tb(s, Inches(0.6), y, Inches(12.13), Inches(0.32),
            step, size=11.5, color=INK)
        y += Inches(0.3)

    _footer(s, page_no, total)
    return s


def slide_appendix(prs, page_no, total):
    s = _blank(prs)
    _kicker(s, "Appendix",
            "What is deferred, and when it lands")

    _tb(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.4),
        "PULL PLAN AND READ DATES", size=11, bold=True, color=INK_FAINT)

    header_row_y = Inches(2.2)
    col_x = [Inches(0.6), Inches(3.6), Inches(5.4), Inches(7.6), Inches(11.5)]
    col_w = [Inches(3.0), Inches(1.8), Inches(2.2), Inches(3.9), Inches(1.2)]
    headers = ["Pull", "Status", "Read date", "Deliverable", "Section"]
    for x, w, h in zip(col_x, col_w, headers):
        _tb(s, x, header_row_y, w, Inches(0.3),
            h.upper(), size=9.5, bold=True, color=INK_FAINT)
    # Rule under header
    line = s.shapes.add_connector(1, Inches(0.6), Inches(2.55),
                                   Inches(12.73), Inches(2.55))
    line.line.color.rgb = RULE
    line.line.width = Pt(0.75)

    rows = [
        ("Pull 1 - TVOD curve",     "LANDED",  "Aug 11, 2026",       "Panel-observed TVOD trajectory + displacement range",       "Slide 8"),
        ("Pull 2 - Demand + Rank",  "PARTIAL", "Aug 25, 2026",       "Apple TV Store + Fandango at Home public rank corroboration", "Appendix N2"),
        ("Pull 3 - Substitution",   "LANDED",  "Aug 11, 2026",       "TVOD -> Peacock playback substitution cohort (618K)",       "Body"),
        ("Pull 4 - Acq + Durab.",   "PARTIAL", "Aug 16 / Jan 2027",  "SubIQ 30-day close, then 180-day retention (H1 close)",     "Slide 3, 10"),
        ("Pull 5 - Composition",    "LANDED",  "Aug 11, 2026",       "Canonical demos + playback completion",                     "Body"),
    ]
    y = Inches(2.65)
    for pull, status, date, deliv, section in rows:
        color = GREEN if status == "LANDED" else AMBER
        _tb(s, col_x[0], y, col_w[0], Inches(0.4),
            pull, size=11, bold=True, color=INK)
        _tb(s, col_x[1], y, col_w[1], Inches(0.4),
            status, size=11, bold=True, color=color)
        _tb(s, col_x[2], y, col_w[2], Inches(0.4),
            date, size=11, color=INK_MUTED)
        _tb(s, col_x[3], y, col_w[3], Inches(0.4),
            deliv, size=10.5, color=INK_MUTED)
        _tb(s, col_x[4], y, col_w[4], Inches(0.4),
            section, size=10, color=INK_FAINT, italic=True)
        y += Inches(0.42)

    # Bottom callout: what will NOT resolve even after Pull 2 closes
    _tb(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.4),
        "WHAT REMAINS DEFERRED EVEN AFTER PULL 2", size=11, bold=True, color=INK_FAINT)
    _tb(s, Inches(0.6), Inches(5.9), Inches(12.13), Inches(1.3),
        "  ·  Distributor Apple + Amazon statements that would collapse the $22-32M displacement range to a point (90-120 day cadence).\n"
        "  ·  Peacock subscriber ARPU + retention detail beyond panel-observed proxies (needed to collapse the H6 $43-175M range).\n"
        "  ·  First-party confirmation from Peacock / NBCUniversal. Every figure in this deck is Crosswalk panel-observed and projected.",
        size=11, color=INK_MUTED)

    _footer(s, page_no, total)
    return s


# --- Build the deck ----------------------------------------------------

def main():
    print("[1/4] Rendering chart PNGs...")
    studio_png = render_studio_averages_chart()
    titles_png = render_universal_titles_chart()
    disney_png = render_disney_titles_chart()
    wb_png = render_wb_titles_chart()
    paramount_png = render_paramount_titles_chart()
    all_int_png = render_all_integrated_chart()
    tvod_png = render_tvod_falloff_chart()
    print(f"     studio_windows.png      {studio_png.stat().st_size:>7,} bytes")
    print(f"     universal_titles.png    {titles_png.stat().st_size:>7,} bytes")
    print(f"     disney_titles.png       {disney_png.stat().st_size:>7,} bytes")
    print(f"     wb_titles.png           {wb_png.stat().st_size:>7,} bytes")
    print(f"     paramount_titles.png    {paramount_png.stat().st_size:>7,} bytes")
    print(f"     all_integrated.png      {all_int_png.stat().st_size:>7,} bytes")
    print(f"     tvod_falloff.png        {tvod_png.stat().st_size:>7,} bytes")

    print("[2/4] Building PPTX (16:9)...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 15  # non-appendix count
    slide_cover(prs)                                                        # 1
    slide_situation(prs, 2, total)                                          # 2
    slide_exec_summary(prs, 3, total)                                       # 3
    slide_setup(prs, 4, total)                                              # 4
    slide_chart(prs, 5, total,
                "Chart 1  ·  Window structure by studio (averages)",
                "Universal median (n=8) at 86 TVOD days. Paramount at n=1. Obsession at 17.",
                studio_png,
                "Days in Theatre + Days in TVOD exclusive + Days in SVOD (30-day attribution). "
                "Universal is shown at the 86.5-day median; Paramount at n=1 (Gladiator II only, not a slate posture).")   # 5
    slide_chart(prs, 6, total,
                "Chart 2  ·  Every Universal-family title, 2024-2026",
                "Same stacked structure, title-by-title. Obsession is the shortest TVOD in the record.",
                titles_png,
                "Every non-Focus UFEG title received 80-101 TVOD days (median 86.5). Focus non-Obsession slate "
                "lands at 17-38 days but is either awards-driven (Conclave, $32.6M) or below $100M dom BO. "
                "Obsession is the only $250M+ title in the record with a 17-day TVOD window.")                              # 6

    # ---- Competitive proof points (Disney / WB / Paramount / combined) ----
    slide_competitor(
        prs, 7, total,
        "Chart 3  ·  Disney -> Disney+ (n=5, 2024-2026)",
        "Every Disney wide release with a Disney+ handoff: 36-43 TVOD days, 89-105 total.",
        "Disney holds 53-67 days of theatrical exclusivity before opening TVOD, then reserves 36-43 days of TVOD-exclusive protection before Disney+ takes the title free-with-sub.",
        disney_png,
        [
            "Median Disney TVOD window: 40 days. Minimum in the sample: 36 (Inside Out 2, Hoppers).",
            "Even Disney's LOWEST TVOD window (36 days) is 2.1x longer than Obsession's 17.",
            "Disney has the LONGEST theatrical windows in the sample (53-67 days) and STILL reserves 36+ TVOD days before the platform handoff.",
            "None of the five Disney titles were placed on Disney+ inside the first 30 days of TVOD availability.",
        ],
        chart_h=3.6,
    )                                                                       # 7

    slide_competitor(
        prs, 8, total,
        "Chart 4  ·  Warner Bros. -> HBO Max (n=4, 2024-2026)",
        "Warner slate: 31-45 TVOD days, 70-77 days total to HBO Max. Fastest integrated-studio churn other than Paramount.",
        "Warner runs the tightest theatrical and TVOD windows of the major integrated studios, but even its shortest TVOD window is 31 days.",
        wb_png,
        [
            "Median WB TVOD window: 37 days. Minimum: 31 (Sinners).",
            "Superman was expedited ~7 days to co-launch with Peacemaker S2's finale (the CLOSEST industry analog to a platform-scheduled handoff) and STILL held 35 TVOD days.",
            "Even that scheduled cross-property push does not approach Obsession's 17-day window.",
            "WB's most aggressive TVOD compression (Sinners at 31 days) is 82% longer than Obsession's 17.",
        ],
        chart_h=3.4,
    )                                                                       # 8

    slide_competitor(
        prs, 9, total,
        "Chart 5  ·  Paramount -> Paramount+ (n=1, the industry floor)",
        "Gladiator II - the shortest integrated-studio TVOD window in the panel. 28 days.",
        "Paramount compressed its theatrical-to-transactional window from ~90 days in 2022 to ~30 by 2024. Even at that pace, Gladiator II got 28 TVOD days.",
        paramount_png,
        [
            "This is the FLOOR of integrated-studio industry practice. No 2024-2026 wide release in the panel goes lower.",
            "Paramount is presented at n=1 because it has not released enough 2024-2026 wide titles with an owned-SVOD handoff to compute a stable slate average. Read as an outer bound, not a studio norm.",
            "The industry-practice defense for Obsession's 17-day window would need to argue that 17 is a normal compression from Paramount's 28 - a 39% cut against the most aggressive operator observed.",
            "Paramount's compression story documents that integrated studios DO compress TVOD windows. It does not document any studio going below 28 days.",
        ],
        chart_h=2.4,
    )                                                                       # 9

    slide_all_integrated(prs, 10, total, all_int_png)                       # 10

    slide_paradox(prs, 11, total)                                           # 11
    slide_chart(prs, 12, total,
                "Chart 7  ·  TVOD daily fall-off, with and without SVOD",
                f"Panel-observed cliff at Jul 17 vs the counterfactual (86.5d median). {DISPLACED_RANGE} estimated.",
                tvod_png,
                "The OBSERVED trajectory is measurement. The CAUSAL attribution of the entire post-Jul-17 collapse "
                "to the Peacock handoff (vs a natural end-of-window fall-off) is a hypothesis pending Backrooms + "
                "Devil Wears Prada 2 control curves. Four confounds are annotated on the chart.")                            # 12
    slide_what_was_traded(prs, 13, total)                                   # 13
    slide_what_was_received(prs, 14, total)                                 # 14
    slide_framing(prs, 15, total)                                           # 15
    slide_appendix(prs, 16, total + 1)                                      # 16 (appendix)

    print("[3/4] Saving PPTX...")
    out = Path.home() / "Downloads" / "Obsession_Peacock_Client_Deck.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"     wrote {out} ({out.stat().st_size:,} bytes)")

    print("[4/4] Done.")
    return out


if __name__ == "__main__":
    main()

"""Prometheus deck renderer (2026-08-20).

Renders a slide-plan JSON (from prometheus_analysis.DECK_PLAN_SYSTEM_PROMPT)
into a PPTX per the Crosswalk Deck System v2.0 (August 2026):
Graphite Teal / Neutral Off-White alternation, Inter, fixed chrome
(eyebrow dot + caps eyebrow, logo top right, footer, page number),
stat blocks, rounded index tracks, thin-numeral recs, dark cover and close.
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Palette (Deck System v2.0)
GRAPHITE = RGBColor(0x0C, 0x16, 0x18)
SLATE = RGBColor(0x15, 0x25, 0x2A)
OFFWHITE = RGBColor(0xE9, 0xE8, 0xE1)
SIGNAL_GREEN = RGBColor(0xC7, 0xF2, 0x3E)
SIGNAL_OLIVE = RGBColor(0x5E, 0x7E, 0x12)
PAVEMENT = RGBColor(0x3B, 0x3D, 0x38)
TRACK_LIGHT = RGBColor(0xC9, 0xC6, 0xBA)

TITLE_DARK = OFFWHITE          # title on graphite
BODY_DARK = RGBColor(0x9A, 0xA0, 0x9B)
MUTED_DARK = RGBColor(0x5C, 0x64, 0x66)
TITLE_LIGHT = GRAPHITE         # title on off-white
BODY_LIGHT = RGBColor(0x5C, 0x65, 0x60)
MUTED_LIGHT = RGBColor(0x88, 0x8C, 0x89)
FOOTER_GREY = RGBColor(0x5C, 0x64, 0x66)

FONT = 'Inter'
MARGIN = 0.920
BAND_RIGHT = 12.413
BAND_W = BAND_RIGHT - MARGIN
CONTENT_TOP = 2.330
SOURCE_Y = 6.620
FOOTER_Y = 7.080


def _tb(slide, x, y, w, h, text, *, size=11.5, bold=False, color=None,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for para_text in str(text).split('\n'):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = para_text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        if color is not None:
            r.font.color.rgb = color
    return box


def _rect(slide, x, y, w, h, fill, *, rounded=False, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    if rounded and radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def _circle(slide, x, y, d, fill):
    sp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _slashes(slide, x, y, h, color):
    """The three-slash mark, drawn as skewed parallelograms."""
    w = h * 0.20
    gap = h * 0.30
    for i in range(3):
        sp = slide.shapes.add_shape(
            MSO_SHAPE.PARALLELOGRAM, Inches(x + i * gap), Inches(y),
            Inches(w * 1.6), Inches(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
        sp.line.fill.background()
        sp.shadow.inherit = False
        try:
            sp.adjustments[0] = 0.55
        except Exception:
            pass


def _chrome(slide, dark, page_num, eyebrow, logo_white, logo_black,
            footer_text='CROSSWALK  ·  PROFILE IQ'):
    accent_body = BODY_DARK if dark else BODY_LIGHT
    # Eyebrow dot + text
    _circle(slide, MARGIN, 0.573, 0.085, SIGNAL_GREEN)
    _tb(slide, 1.129, 0.540, 5.5, 0.3, str(eyebrow or 'PROFILE IQ').upper(),
        size=10.5, bold=True, color=accent_body)
    # Logo top right (white on dark, black on light)
    logo = logo_white if dark else logo_black
    if logo and os.path.exists(logo):
        try:
            slide.shapes.add_picture(
                logo, Inches(11.083), Inches(0.558), width=Inches(1.330))
        except Exception:
            pass
    else:
        _slashes(slide, 11.95, 0.545, 0.22,
                 OFFWHITE if dark else GRAPHITE)
    # Footer + page number
    _tb(slide, MARGIN, FOOTER_Y, 6.0, 0.25, footer_text,
        size=8, bold=True, color=FOOTER_GREY)
    if page_num is not None:
        _tb(slide, BAND_RIGHT - 1.0, FOOTER_Y, 1.0, 0.25, str(page_num),
            size=8, bold=False, color=FOOTER_GREY, align=PP_ALIGN.RIGHT)


def _bg(slide, prs, color):
    _rect(slide, 0, 0, prs.slide_width / 914400, prs.slide_height / 914400,
          color)


def _title_block(slide, dark, title, subhead=None):
    _tb(slide, MARGIN, 1.050, BAND_W, 1.0, title, size=34, bold=True,
        color=TITLE_DARK if dark else TITLE_LIGHT, line_spacing=1.15)
    if subhead:
        _tb(slide, MARGIN, 1.900, BAND_W, 0.4, subhead, size=16,
            color=BODY_DARK if dark else BODY_LIGHT)


def _slide_cover(slide, prs, plan_slide, logos):
    _bg(slide, prs, GRAPHITE)
    _chrome(slide, True, None, plan_slide.get('eyebrow') or 'PROFILE IQ',
            *logos)
    _slashes(slide, MARGIN, 2.28, 0.55, SIGNAL_GREEN)
    _tb(slide, MARGIN, 3.05, BAND_W - 1.0, 1.8,
        plan_slide.get('title') or 'Audience profile.',
        size=40, bold=True, color=TITLE_DARK, line_spacing=1.12)
    meta = plan_slide.get('meta') or ''
    if meta:
        _tb(slide, MARGIN, 5.35, BAND_W - 2.0, 0.8, meta,
            size=14, color=BODY_DARK, line_spacing=1.3)


def _slide_stats(slide, prs, plan_slide, dark, page_num, logos):
    _bg(slide, prs, GRAPHITE if dark else OFFWHITE)
    _chrome(slide, dark, page_num, plan_slide.get('eyebrow'), *logos)
    _title_block(slide, dark, plan_slide.get('title') or '')
    stats = (plan_slide.get('stats') or [])[:4]
    if stats:
        col_w = (BAND_W - 0.280 * (len(stats) - 1)) / len(stats)
        for i, st in enumerate(stats):
            x = MARGIN + i * (col_w + 0.280)
            fig_color = ((SIGNAL_GREEN if dark else SIGNAL_OLIVE)
                         if i == 0 else (TITLE_DARK if dark else TITLE_LIGHT))
            _tb(slide, x, CONTENT_TOP + 0.35, col_w, 0.75,
                str(st.get('big') or ''), size=30, bold=True, color=fig_color)
            _tb(slide, x, CONTENT_TOP + 1.12, col_w, 0.85,
                str(st.get('label') or ''), size=11.5,
                color=BODY_DARK if dark else BODY_LIGHT, line_spacing=1.3)
    read = plan_slide.get('read')
    if read:
        _tb(slide, MARGIN, SOURCE_Y - 0.75, BAND_W, 0.6, read, size=11.5,
            color=BODY_DARK if dark else BODY_LIGHT, line_spacing=1.35)
    _tb(slide, MARGIN, SOURCE_Y, BAND_W, 0.25,
        'Crosswalk Profile IQ, trailing 12 months.', size=9,
        color=MUTED_DARK if dark else MUTED_LIGHT)


def _slide_chart(slide, prs, plan_slide, dark, page_num, logos):
    _bg(slide, prs, GRAPHITE if dark else OFFWHITE)
    _chrome(slide, dark, page_num, plan_slide.get('eyebrow'), *logos)
    _title_block(slide, dark, plan_slide.get('title') or '')
    rows = (plan_slide.get('rows') or [])[:6]
    unit = plan_slide.get('unit') or '% pen'
    label_w = 2.30
    val_w = 1.55
    track_x = MARGIN + label_w + 0.15
    track_w = BAND_W - label_w - val_w - 0.30
    track_h = 0.17
    avail = (SOURCE_Y - 0.85) - CONTENT_TOP
    row_gap = min(0.50, max(0.30, avail / max(len(rows), 1) - track_h))
    maxv = max([float(r.get('value') or 0) for r in rows] + [0.0001])
    track_c = PAVEMENT if dark else TRACK_LIGHT
    bar_c = SIGNAL_GREEN if dark else SIGNAL_OLIVE
    y = CONTENT_TOP + 0.15
    for r in rows:
        v = float(r.get('value') or 0)
        _tb(slide, MARGIN, y - 0.035, label_w, 0.3,
            str(r.get('label') or ''), size=11, bold=True,
            color=TITLE_DARK if dark else TITLE_LIGHT)
        _rect(slide, track_x, y, track_w, track_h, track_c,
              rounded=True, radius=0.5)
        bw = max(track_w * (v / maxv), track_h)
        _rect(slide, track_x, y, bw, track_h, bar_c,
              rounded=True, radius=0.5)
        note = str(r.get('note') or '').strip()
        val_txt = f"{v:.1f}" + (f"  {note}" if note else '')
        _tb(slide, track_x + track_w + 0.12, y - 0.035, val_w, 0.3,
            val_txt, size=11, bold=True,
            color=TITLE_DARK if dark else TITLE_LIGHT)
        y += track_h + row_gap
    read = plan_slide.get('read')
    if read:
        _tb(slide, MARGIN, SOURCE_Y - 0.62, BAND_W, 0.5, read, size=11.5,
            color=BODY_DARK if dark else BODY_LIGHT, line_spacing=1.35)
    _tb(slide, MARGIN, SOURCE_Y, BAND_W, 0.25,
        f'Crosswalk Profile IQ, trailing 12 months. Values in {unit}.',
        size=9, color=MUTED_DARK if dark else MUTED_LIGHT)


def _slide_recs(slide, prs, plan_slide, dark, page_num, logos):
    _bg(slide, prs, GRAPHITE if dark else OFFWHITE)
    _chrome(slide, dark, page_num, plan_slide.get('eyebrow'), *logos)
    _title_block(slide, dark, plan_slide.get('title') or '')
    recs = (plan_slide.get('recs') or [])[:4]
    if not recs:
        return
    col_gap = 0.280
    col_w = (BAND_W - col_gap * (len(recs) - 1)) / len(recs)
    accent = SIGNAL_GREEN if dark else SIGNAL_OLIVE
    for i, rec in enumerate(recs):
        x = MARGIN + i * (col_w + col_gap)
        _tb(slide, x, CONTENT_TOP + 0.1, col_w, 0.55, f"0{i + 1}",
            size=30, bold=False, color=accent)
        _tb(slide, x, CONTENT_TOP + 0.78, col_w, 0.55,
            str(rec.get('head') or ''), size=11.5, bold=True,
            color=TITLE_DARK if dark else TITLE_LIGHT, line_spacing=1.25)
        _tb(slide, x, CONTENT_TOP + 1.45, col_w, 2.4,
            str(rec.get('body') or ''), size=11.5,
            color=BODY_DARK if dark else BODY_LIGHT, line_spacing=1.35)


def _slide_quadrant(slide, prs, plan_slide, dark, page_num, logos):
    """2x2 map: points plotted on two metrics (typically x = penetration
    for scale, y = index for efficiency) with median split lines and
    corner labels. Used for target maps / prioritization reads."""
    _bg(slide, prs, GRAPHITE if dark else OFFWHITE)
    _chrome(slide, dark, page_num, plan_slide.get('eyebrow'), *logos)
    _title_block(slide, dark, plan_slide.get('title') or '')
    pts = [p for p in (plan_slide.get('points') or [])
           if p.get('x') is not None and p.get('y') is not None][:10]
    if not pts:
        return
    px, py = MARGIN + 0.55, CONTENT_TOP + 0.30
    pw, ph = BAND_W - 0.75, (SOURCE_Y - 0.95) - py
    xs = [float(p['x']) for p in pts]
    ys = [float(p['y']) for p in pts]

    def _rng(vals, split):
        lo, hi = min(vals + ([split] if split is not None else [])), \
            max(vals + ([split] if split is not None else []))
        pad = (hi - lo) * 0.15 or 1.0
        return lo - pad, hi + pad

    sx = plan_slide.get('x_split')
    sy = plan_slide.get('y_split')
    sx = float(sx) if sx is not None else sorted(xs)[len(xs) // 2]
    sy = float(sy) if sy is not None else sorted(ys)[len(ys) // 2]
    x0, x1 = _rng(xs, sx)
    y0, y1 = _rng(ys, sy)
    track_c = PAVEMENT if dark else TRACK_LIGHT
    accent = SIGNAL_GREEN if dark else SIGNAL_OLIVE
    body_c = BODY_DARK if dark else BODY_LIGHT
    # Plot frame + median split lines
    _rect(slide, px, py, pw, 0.014, track_c)
    _rect(slide, px, py + ph, pw, 0.014, track_c)
    _rect(slide, px, py, 0.014, ph, track_c)
    _rect(slide, px + pw, py, 0.014, ph, track_c)
    fx = px + pw * (sx - x0) / (x1 - x0)
    fy = py + ph * (1 - (sy - y0) / (y1 - y0))
    _rect(slide, fx, py, 0.014, ph, track_c)
    _rect(slide, px, fy, pw, 0.014, track_c)
    # Quadrant corner labels
    q = plan_slide.get('q_labels') or {}
    muted = MUTED_DARK if dark else MUTED_LIGHT
    for key, (qx, qy, al) in {
            'tl': (px + 0.10, py + 0.06, PP_ALIGN.LEFT),
            'tr': (px + pw - 2.60, py + 0.06, PP_ALIGN.RIGHT),
            'bl': (px + 0.10, py + ph - 0.34, PP_ALIGN.LEFT),
            'br': (px + pw - 2.60, py + ph - 0.34, PP_ALIGN.RIGHT)}.items():
        if q.get(key):
            _tb(slide, qx, qy, 2.5, 0.28, str(q[key]).upper(), size=9,
                bold=True, color=muted, align=al)
    # Points
    d = 0.16
    for p in pts:
        cx = px + pw * (float(p['x']) - x0) / (x1 - x0)
        cy = py + ph * (1 - (float(p['y']) - y0) / (y1 - y0))
        _circle(slide, cx - d / 2, cy - d / 2, d, accent)
        _tb(slide, cx + d / 2 + 0.04, cy - 0.11, 2.2, 0.26,
            str(p.get('label') or ''), size=9.5, bold=True,
            color=TITLE_DARK if dark else TITLE_LIGHT)
    # Axis labels
    _tb(slide, px, py + ph + 0.10, pw, 0.25,
        str(plan_slide.get('x_label') or ''), size=9.5, bold=True,
        color=body_c, align=PP_ALIGN.CENTER)
    ybox = _tb(slide, px - 1.62, py + ph / 2 - 0.14, 3.0, 0.28,
               str(plan_slide.get('y_label') or ''), size=9.5, bold=True,
               color=body_c, align=PP_ALIGN.CENTER)
    ybox.rotation = -90
    read = plan_slide.get('read')
    if read:
        _tb(slide, MARGIN, SOURCE_Y - 0.42, BAND_W, 0.4, read, size=11,
            color=body_c, line_spacing=1.3)
    _tb(slide, MARGIN, SOURCE_Y + 0.1, BAND_W, 0.25,
        'Crosswalk Profile IQ, trailing 12 months.', size=9, color=muted)


def _slide_personas(slide, prs, plan_slide, dark, page_num, logos):
    """2-3 persona cards: name, sized share, identity line, stat
    receipts, message hook. Cards must carve the audience MECE."""
    _bg(slide, prs, GRAPHITE if dark else OFFWHITE)
    _chrome(slide, dark, page_num, plan_slide.get('eyebrow'), *logos)
    _title_block(slide, dark, plan_slide.get('title') or '')
    cards = (plan_slide.get('cards') or [])[:3]
    if not cards:
        return
    gap = 0.28
    card_w = (BAND_W - gap * (len(cards) - 1)) / len(cards)
    card_y = CONTENT_TOP + 0.15
    card_h = SOURCE_Y - card_y - 0.30
    card_fill = SLATE if dark else RGBColor(0xF7, 0xF6, 0xF0)
    accent = SIGNAL_GREEN if dark else SIGNAL_OLIVE
    title_c = TITLE_DARK if dark else TITLE_LIGHT
    body_c = BODY_DARK if dark else BODY_LIGHT
    for i, c in enumerate(cards):
        x = MARGIN + i * (card_w + gap)
        _rect(slide, x, card_y, card_w, card_h, card_fill,
              rounded=True, radius=0.055)
        pad = 0.26
        cx, cw = x + pad, card_w - 2 * pad
        _tb(slide, cx, card_y + 0.26, cw, 0.62,
            str(c.get('name') or ''), size=16, bold=True, color=title_c,
            line_spacing=1.1)
        _tb(slide, cx, card_y + 0.88, cw, 0.28,
            str(c.get('share') or '').upper(), size=9.5, bold=True,
            color=accent)
        _tb(slide, cx, card_y + 1.22, cw, 0.85,
            str(c.get('identity') or ''), size=10.5, color=body_c,
            line_spacing=1.3)
        sy = card_y + 2.15
        for st in (c.get('stats') or [])[:4]:
            _circle(slide, cx, sy + 0.055, 0.07, accent)
            _tb(slide, cx + 0.17, sy, cw - 0.17, 0.3, str(st),
                size=10, bold=True, color=title_c)
            sy += 0.34
        hook = str(c.get('hook') or '').strip()
        if hook:
            _tb(slide, cx, card_y + card_h - 0.85, cw, 0.72,
                f'\u201c{hook}\u201d', size=10.5, color=body_c,
                line_spacing=1.25)
    _tb(slide, MARGIN, SOURCE_Y, BAND_W, 0.25,
        'Crosswalk Profile IQ, trailing 12 months.', size=9,
        color=MUTED_DARK if dark else MUTED_LIGHT)


def _slide_benchmark(slide, prs, plan_slide, dark, page_num, logos):
    """Paired bars: this audience vs US gen pop per row. Used when the
    contrast against the average American IS the story."""
    _bg(slide, prs, GRAPHITE if dark else OFFWHITE)
    _chrome(slide, dark, page_num, plan_slide.get('eyebrow'), *logos)
    _title_block(slide, dark, plan_slide.get('title') or '')
    rows = (plan_slide.get('rows') or [])[:5]
    if not rows:
        return
    unit = plan_slide.get('unit') or '% pen'
    label_w, val_w = 2.30, 1.85
    track_x = MARGIN + label_w + 0.15
    track_w = BAND_W - label_w - val_w - 0.30
    bar_h = 0.135
    pair_h = bar_h * 2 + 0.05
    avail = (SOURCE_Y - 0.95) - CONTENT_TOP
    row_gap = min(0.45, max(0.22, avail / max(len(rows), 1) - pair_h))
    maxv = max([float(r.get('aud') or 0) for r in rows]
               + [float(r.get('gp') or 0) for r in rows] + [0.0001])
    track_c = PAVEMENT if dark else TRACK_LIGHT
    accent = SIGNAL_GREEN if dark else SIGNAL_OLIVE
    title_c = TITLE_DARK if dark else TITLE_LIGHT
    body_c = BODY_DARK if dark else BODY_LIGHT
    y = CONTENT_TOP + 0.35
    for r in rows:
        aud = float(r.get('aud') or 0)
        gp = float(r.get('gp') or 0)
        _tb(slide, MARGIN, y + 0.02, label_w, 0.3,
            str(r.get('label') or ''), size=11, bold=True, color=title_c)
        _rect(slide, track_x, y, max(track_w * aud / maxv, bar_h), bar_h,
              accent, rounded=True, radius=0.5)
        _rect(slide, track_x, y + bar_h + 0.05,
              max(track_w * gp / maxv, bar_h), bar_h, track_c,
              rounded=True, radius=0.5)
        idx_txt = f"  idx {round(aud / gp * 100)}" if gp >= 0.01 else ''
        _tb(slide, track_x + track_w + 0.12, y - 0.01, val_w, 0.3,
            f"{aud:.1f} vs {gp:.1f}{idx_txt}", size=10.5, bold=True,
            color=title_c)
        y += pair_h + row_gap
    legend_y = CONTENT_TOP - 0.02
    _rect(slide, BAND_RIGHT - 3.55, legend_y + 0.05, 0.28, 0.10, accent,
          rounded=True, radius=0.5)
    _tb(slide, BAND_RIGHT - 3.20, legend_y - 0.04, 1.35, 0.25,
        'This audience', size=9, bold=True, color=body_c)
    _rect(slide, BAND_RIGHT - 1.85, legend_y + 0.05, 0.28, 0.10, track_c,
          rounded=True, radius=0.5)
    _tb(slide, BAND_RIGHT - 1.50, legend_y - 0.04, 1.30, 0.25,
        'US gen pop', size=9, bold=True, color=body_c)
    read = plan_slide.get('read')
    if read:
        _tb(slide, MARGIN, SOURCE_Y - 0.62, BAND_W, 0.5, read, size=11.5,
            color=body_c, line_spacing=1.35)
    _tb(slide, MARGIN, SOURCE_Y, BAND_W, 0.25,
        f'Crosswalk Profile IQ, trailing 12 months. Values in {unit}.',
        size=9, color=MUTED_DARK if dark else MUTED_LIGHT)


def _slide_close(slide, prs, plan_slide, page_num, logos):
    _bg(slide, prs, GRAPHITE)
    _chrome(slide, True, page_num, plan_slide.get('eyebrow') or 'PROFILE IQ',
            *logos)
    big = str(plan_slide.get('big') or '').strip()
    line = str(plan_slide.get('line') or '').strip()
    if big:
        _tb(slide, MARGIN, 2.28, BAND_W, 1.9, big, size=82, bold=True,
            color=SIGNAL_GREEN)
        _tb(slide, MARGIN, 4.45, BAND_W - 2.0, 1.2, line, size=16,
            color=BODY_DARK, line_spacing=1.35)
    else:
        _tb(slide, MARGIN, 2.60, BAND_W - 1.5, 2.4, line or 'Thank you.',
            size=48, bold=True, color=TITLE_DARK, line_spacing=1.15)


def render_deck(plan, out_path, static_dir=None):
    """Render a slide-plan dict to a PPTX at out_path."""
    logo_white = logo_black = None
    if static_dir:
        for cand in ('crosswalk-logo-brand-white.png',
                     'crosswalk-logo-white-hires.png'):
            p = os.path.join(static_dir, cand)
            if os.path.exists(p):
                logo_white = p
                break
        for cand in ('crosswalk-logo-brand-black.png',
                     'crosswalk-logo.png'):
            p = os.path.join(static_dir, cand)
            if os.path.exists(p):
                logo_black = p
                break
    logos = (logo_white, logo_black)

    prs = Presentation()
    prs.slide_width = Emu(int(13.333 * 914400))
    prs.slide_height = Emu(int(7.5 * 914400))
    blank = prs.slide_layouts[6]

    slides = plan.get('slides') or []
    # Middle slides alternate grounds starting light; cover and close
    # are always dark (open dark, close dark, data prefers off-white).
    mid_idx = 0
    for i, ps in enumerate(slides):
        stype = str(ps.get('type') or '').lower()
        slide = prs.slides.add_slide(blank)
        page_num = i + 1
        if stype == 'cover':
            _slide_cover(slide, prs, ps, logos)
            continue
        if stype == 'close':
            _slide_close(slide, prs, ps, page_num, logos)
            continue
        dark = (mid_idx % 2 == 1)
        mid_idx += 1
        if stype == 'stats':
            _slide_stats(slide, prs, ps, dark, page_num, logos)
        elif stype == 'chart':
            _slide_chart(slide, prs, ps, dark, page_num, logos)
        elif stype == 'recs':
            _slide_recs(slide, prs, ps, dark, page_num, logos)
        elif stype == 'quadrant':
            _slide_quadrant(slide, prs, ps, dark, page_num, logos)
        elif stype == 'personas':
            _slide_personas(slide, prs, ps, dark, page_num, logos)
        elif stype == 'benchmark':
            _slide_benchmark(slide, prs, ps, dark, page_num, logos)
        else:
            _slide_stats(slide, prs, ps, dark, page_num, logos)

    prs.save(out_path)
    return out_path

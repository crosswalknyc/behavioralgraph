"""Crosswalk insights-deck renderer for Prometheus (2026-08-26).

Generalizes the reference talent-value deck implementation (the Paige
Bueckers audience-value deck) into a reusable renderer: a slide-plan
JSON in, a finished client-ready PPTX out. Every surface follows the
Crosswalk deck system exactly: Graphite Teal / Neutral Off-White
grounds, Inter everywhere (Geist Mono only on raw clickstream rows),
one accent per surface, 0.920in margins, eyebrow dot + logo chrome,
sentence-case headlines that end in a full stop.

Slide-type vocabulary (the planner picks per slide):
  cover            dark opener: headline, intro, three proof stats
  argument         2x2 numbered cards (dark slate or light card)
  tiles_facts      3-4 stat tiles + up to 6 fact rows below
  bars             ranked bar rows, optional PEN./INDEX columns, read line
  split_stats_bars two stat cards left + labeled bar list right + read
  tiles_row        three tall stat tiles with body copy + read line
  hero             full-bleed Orchid (or dark) single-figure moment
  table            column table with header hairline + read card(s)
  hero_proof       dark: 82pt figure left, three proof cards right
  paths            dark: Geist Mono clickstream rows, lit rows on slate
  close            dark 2x2 numbered next-step cards

The renderer is tolerant: unknown fields are ignored, list lengths are
clamped, missing optionals degrade to a clean layout rather than an
error. Text content is expected to be pre-scrubbed by the caller
(prometheus_analysis.enforce_insights_plan).
"""
from __future__ import annotations

import os

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---- Crosswalk palette (deck-system.md) ------------------------------------
GRAPHITE = RGBColor(0x0C, 0x16, 0x18)
SLATE = RGBColor(0x15, 0x25, 0x2A)
OFFWHITE = RGBColor(0xE9, 0xE8, 0xE1)
CARD = RGBColor(0xE1, 0xE0, 0xD7)
SIGNAL = RGBColor(0xC7, 0xF2, 0x3E)
OLIVE = RGBColor(0x5E, 0x7E, 0x12)
ORCHID = RGBColor(0xE6, 0x82, 0xFF)
AMETHYST = RGBColor(0x8E, 0x3F, 0xA8)
WHITE = RGBColor(0xE9, 0xE8, 0xE1)
BODY_DK = RGBColor(0x9A, 0xA0, 0x9B)
MUTED_DK = RGBColor(0x5C, 0x64, 0x66)
BODY_LT = RGBColor(0x5C, 0x65, 0x60)
MUTED_LT = RGBColor(0x88, 0x8C, 0x89)
FOOTER = RGBColor(0x5C, 0x64, 0x66)
TRACK = RGBColor(0xC9, 0xC6, 0xBA)
TRACK_DK = RGBColor(0x3B, 0x3D, 0x38)
ORCHID_BODY = RGBColor(0xF4, 0xE4, 0xFA)
ORCHID_MUTED = RGBColor(0x5C, 0x2A, 0x6E)
ORCHID_FOOT = RGBColor(0x4A, 0x24, 0x58)

SW, SH = Inches(13.333), Inches(7.500)
M, BAND, GUT = Inches(0.920), Inches(11.493), Inches(0.280)
A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _face(run, name):
    rPr = run._r.get_or_add_rPr()
    for tag in (A_NS + "latin", A_NS + "ea", A_NS + "cs"):
        el = rPr.find(tag)
        if el is None:
            el = etree.SubElement(rPr, tag)
        el.set("typeface", name)


def _spc(run, hundredths):
    run._r.get_or_add_rPr().set("spc", str(int(hundredths)))


def txt(s, x, y, w, h, text, *, size=12, bold=False, color=GRAPHITE,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Inter",
        spc=None, light=False):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    face = "Inter Light" if light else font
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = face
        _face(r, face)
        if spc is not None:
            _spc(r, spc)
    return box


def rect(s, x, y, w, h, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def rrect(s, x, y, w, h, fill, radius=0.14):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    shorter = min(int(w), int(h))
    half = max(shorter / 2, 1)
    sh.adjustments[0] = min(0.5, Inches(radius) / half)
    sh.shadow.inherit = False
    return sh


def dot(s, x, y, fill=SIGNAL):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y,
                            Inches(0.092), Inches(0.092))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


class _Deck:
    """Holds presentation-wide state: logo paths and page counter."""

    def __init__(self, prs, logo_white, logo_black):
        self.prs = prs
        self.logo_white = logo_white
        self.logo_black = logo_black
        self.page = 0

    def logo(self, s, white=True):
        path = self.logo_white if white else self.logo_black
        if path and os.path.exists(path):
            s.shapes.add_picture(str(path), Inches(11.083), Inches(0.558),
                                 width=Inches(1.330))

    def new_slide(self, dark=False, fill=None):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        rect(s, 0, 0, SW, SH, fill or (GRAPHITE if dark else OFFWHITE))
        self.page += 1
        return s

    def chrome(self, s, eyebrow, title, *, dark=False, sub=None,
               source=None, orchid=False):
        ink = WHITE if (dark or orchid) else GRAPHITE
        body = BODY_DK if dark else BODY_LT
        muted = MUTED_DK if dark else MUTED_LT
        if orchid:
            body = ORCHID_BODY
            muted = ORCHID_MUTED
            ink = WHITE
        page = self.page
        eb = f"{page:02d}  {eyebrow}" if eyebrow else f"{page:02d}"
        dot(s, M, Inches(0.540))
        txt(s, Inches(1.129), Inches(0.540), Inches(9.4), Inches(0.22),
            str(eb).upper(), size=10.5, bold=True, color=body, spc=260)
        self.logo(s, white=dark or orchid)
        txt(s, M, Inches(1.050), BAND, Inches(0.72), title,
            size=34, bold=True, color=ink)
        if sub:
            txt(s, M, Inches(1.846), BAND, Inches(0.40), sub,
                size=16, color=body)
        if source:
            txt(s, M, Inches(6.620), BAND, Inches(0.28), source,
                size=9, color=muted)
        foot = ("CONFIDENTIAL" if page == 1
                else "CROSSWALK / BEHAVIORAL INTELLIGENCE")
        foot_c = FOOTER if not orchid else ORCHID_FOOT
        txt(s, M, Inches(7.080), Inches(8.2), Inches(0.22),
            foot, size=8, color=foot_c, spc=160 if page == 1 else 200)
        txt(s, Inches(10.183), Inches(7.080), Inches(2.230), Inches(0.22),
            str(page), size=8, color=foot_c, align=PP_ALIGN.RIGHT)


def _bar(s, x, y, track_w, frac, *, accent=False, h=0.17, dark=False):
    track = TRACK_DK if dark else TRACK
    fill = ((ORCHID if accent else SIGNAL) if dark
            else (AMETHYST if accent else OLIVE))
    rrect(s, x, y, track_w, Inches(h), track, radius=0.08)
    frac = min(max(_num(frac, 0.01), 0.01), 1.0)
    bw = max(int(int(track_w) * frac), 8)
    rrect(s, x, y, Emu(bw), Inches(h), fill, radius=0.08)


def _num(v, default=0.0):
    try:
        return float(v)
    except (TypeError, ValueError):
        return float(default)


def _fmt_val(v, suffix):
    """Bar value display: keep the model's precision, add the suffix."""
    n = _num(v)
    if n == int(n) and abs(n) >= 10 and suffix != "x":
        disp = f"{int(n):,}"
    else:
        disp = f"{n:g}"
    return f"{disp}{suffix}" if suffix else disp


def _items(sl, key, cap):
    out = [i for i in (sl.get(key) or []) if isinstance(i, dict)]
    return out[:cap]


# ---- slide renderers --------------------------------------------------------

def _sl_cover(d, sl):
    s = d.new_slide(dark=True)
    dot(s, M, Inches(0.540))
    txt(s, Inches(1.129), Inches(0.540), Inches(9.4), Inches(0.22),
        str(sl.get("eyebrow") or "CROSSWALK").upper(),
        size=10.5, bold=True, color=BODY_DK, spc=260)
    d.logo(s, True)
    txt(s, M, Inches(1.050), BAND, Inches(1.10),
        sl.get("title") or "", size=34, bold=True, color=WHITE)
    if sl.get("intro"):
        txt(s, M, Inches(2.28), Inches(10.8), Inches(0.70),
            sl["intro"], size=16, color=BODY_DK)
    stats = _items(sl, "stats", 3)
    acc = sl.get("accent_index")
    x, w = M, Inches(3.644)
    for i, st in enumerate(stats):
        color = SIGNAL if (isinstance(acc, int) and i == acc) else WHITE
        txt(s, x, Inches(5.90), w, Inches(0.46),
            st.get("big") or "", size=22, bold=True, color=color)
        txt(s, x, Inches(6.40), w, Inches(0.55),
            st.get("label") or "", size=11.5, color=BODY_DK)
        x += w + GUT
    txt(s, M, Inches(7.080), Inches(8.2), Inches(0.22),
        "CONFIDENTIAL", size=8, color=FOOTER, spc=160)
    txt(s, Inches(10.183), Inches(7.080), Inches(2.230), Inches(0.22),
        str(d.page), size=8, color=FOOTER, align=PP_ALIGN.RIGHT)


def _cards_2x2(d, s, cards, dark):
    cw, ch = Inches(5.607), Inches(1.82)
    for i, c in enumerate(cards[:4]):
        x = M + (i % 2) * (cw + GUT)
        y = Inches(2.330) + (i // 2) * (ch + Inches(0.22))
        rrect(s, x, y, cw, ch, SLATE if dark else CARD)
        txt(s, x + Inches(0.26), y + Inches(0.22), Inches(0.7), Inches(0.42),
            f"{i + 1:02d}", size=30, color=SIGNAL if dark else OLIVE,
            light=True)
        txt(s, x + Inches(1.05), y + Inches(0.28), Inches(4.25), Inches(0.36),
            c.get("head") or "", size=14, bold=True,
            color=WHITE if dark else GRAPHITE)
        txt(s, x + Inches(1.05), y + Inches(0.68), Inches(4.25), Inches(0.96),
            c.get("body") or "", size=11.5,
            color=BODY_DK if dark else BODY_LT)


def _sl_argument(d, sl):
    dark = str(sl.get("ground") or "dark").lower() != "light"
    s = d.new_slide(dark=dark)
    d.chrome(s, sl.get("eyebrow") or "Argument", sl.get("title") or "",
             dark=dark, sub=sl.get("sub"), source=sl.get("source"))
    _cards_2x2(d, s, _items(sl, "cards", 4), dark)


def _sl_tiles_facts(d, sl):
    s = d.new_slide()
    d.chrome(s, sl.get("eyebrow") or "Universe", sl.get("title") or "",
             sub=sl.get("sub"), source=sl.get("source"))
    tiles = _items(sl, "tiles", 4)
    acc = sl.get("accent_index")
    n = max(len(tiles), 1)
    w = Emu(int((int(BAND) - int(GUT) * (n - 1)) / n))
    x = M
    for i, t in enumerate(tiles):
        rrect(s, x, Inches(2.330), w, Inches(1.58), CARD)
        color = OLIVE if (isinstance(acc, int) and i == acc) else GRAPHITE
        txt(s, x + Inches(0.18), Inches(2.48), w - Inches(0.32),
            Inches(0.70), t.get("big") or "", size=30, bold=True,
            color=color)
        txt(s, x + Inches(0.18), Inches(3.22), w - Inches(0.32),
            Inches(0.52), t.get("label") or "", size=11.5, color=BODY_LT)
        x += w + GUT
    y = Inches(4.12)
    for f in _items(sl, "facts", 6):
        txt(s, M, y, Inches(3.40), Inches(0.32),
            f.get("label") or "", size=11.5)
        txt(s, M + Inches(3.50), y, Inches(1.80), Inches(0.32),
            f.get("fig") or "", size=11.5, bold=True, color=OLIVE)
        txt(s, M + Inches(5.50), y, Inches(5.99), Inches(0.32),
            f.get("note") or "", size=11.5, color=BODY_LT)
        y += Inches(0.32)
    if sl.get("read"):
        txt(s, M, Inches(6.30), BAND, Inches(0.52), sl["read"],
            size=11.5, color=BODY_LT)


def _sl_bars(d, sl):
    dark = str(sl.get("ground") or "light").lower() == "dark"
    s = d.new_slide(dark=dark)
    d.chrome(s, sl.get("eyebrow") or "Read", sl.get("title") or "",
             dark=dark, sub=sl.get("sub"), source=sl.get("source"))
    rows = _items(sl, "rows", 9)
    suffix = str(sl.get("value_suffix") if sl.get("value_suffix")
                 is not None else "%")
    show_index = bool(sl.get("show_index")) and any(
        r.get("index") not in (None, "") for r in rows)
    ink = WHITE if dark else GRAPHITE
    body = BODY_DK if dark else BODY_LT
    muted = MUTED_DK if dark else MUTED_LT
    val_accent = ORCHID if dark else AMETHYST
    y = Inches(2.330)
    tw = Inches(6.2) if show_index else Inches(7.4)
    if show_index:
        txt(s, M, y, Inches(2.4), Inches(0.24), "BRAND", size=10.5,
            bold=True, color=muted, spc=100)
        txt(s, M + Inches(9.05), y, Inches(1.1), Inches(0.24), "PEN.",
            size=10.5, bold=True, color=muted, spc=100,
            align=PP_ALIGN.RIGHT)
        txt(s, M + Inches(10.30), y, Inches(1.19), Inches(0.24), "INDEX",
            size=10.5, bold=True, color=muted, spc=100,
            align=PP_ALIGN.RIGHT)
        y = Inches(2.66)
    vmax = max([_num(r.get("value"), 1) for r in rows] or [1.0])
    n = max(len(rows), 1)
    avail = (6.02 if sl.get("read") else 6.60) - float(y) / 914400.0
    step = min(0.52, max(0.36, avail / n))
    for r in rows:
        accent = bool(r.get("accent"))
        txt(s, M, y, Inches(2.70) if not show_index else Inches(2.55),
            Inches(0.28), r.get("label") or "", size=11.5, bold=accent,
            color=ink)
        bx = M + (Inches(2.80) if not show_index else Inches(2.65))
        _bar(s, bx, y + Inches(0.05), tw,
             _num(r.get("value")) / vmax, accent=accent, dark=dark)
        val_c = val_accent if accent else ink
        if show_index:
            txt(s, M + Inches(8.95), y, Inches(1.20), Inches(0.28),
                _fmt_val(r.get("value"), suffix), size=11, bold=True,
                color=val_c, align=PP_ALIGN.RIGHT)
            txt(s, M + Inches(10.30), y, Inches(1.19), Inches(0.28),
                str(r.get("index") or ""), size=11, color=body,
                align=PP_ALIGN.RIGHT)
        else:
            txt(s, bx + tw + Inches(0.12), y, Inches(1.00), Inches(0.28),
                _fmt_val(r.get("value"), suffix), size=11, bold=True,
                color=val_c)
        y += Inches(step)
    if sl.get("read"):
        txt(s, M, Inches(6.02), BAND, Inches(0.62), sl["read"],
            size=11.5, color=body)


def _stat_card(s, x, y, card, *, big_color=GRAPHITE):
    rrect(s, x, y, Inches(3.644), Inches(1.82), CARD)
    txt(s, x + Inches(0.24), y + Inches(0.19), Inches(3.16), Inches(0.24),
        str(card.get("kicker") or "").upper(), size=10.5, bold=True,
        color=MUTED_LT, spc=100)
    txt(s, x + Inches(0.24), y + Inches(0.49), Inches(3.16), Inches(0.7),
        card.get("big") or "", size=34, bold=True,
        color=OLIVE if card.get("accent") else big_color)
    txt(s, x + Inches(0.24), y + Inches(1.23), Inches(3.16), Inches(0.56),
        card.get("label") or "", size=11.5, color=BODY_LT)


def _sl_split_stats_bars(d, sl):
    s = d.new_slide()
    d.chrome(s, sl.get("eyebrow") or "Read", sl.get("title") or "",
             sub=sl.get("sub"), source=sl.get("source"))
    cards = _items(sl, "stat_cards", 2)
    if cards:
        _stat_card(s, M, Inches(2.330), cards[0])
    if len(cards) > 1:
        _stat_card(s, M, Inches(4.40), cards[1])
    rx = M + Inches(3.644) + GUT
    rw = Inches(7.569)
    if sl.get("bars_title"):
        txt(s, rx, Inches(2.330), rw, Inches(0.24),
            str(sl["bars_title"]).upper(), size=10.5, bold=True,
            color=MUTED_LT, spc=100)
    rows = _items(sl, "rows", 8)
    suffix = str(sl.get("value_suffix") if sl.get("value_suffix")
                 is not None else "%")
    vmax = max([_num(r.get("value"), 1) for r in rows] or [1.0])
    y = Inches(2.70)
    tw = Inches(4.6)
    n = max(len(rows), 1)
    step = min(0.50, max(0.40, 3.1 / n))
    for r in rows:
        accent = bool(r.get("accent"))
        txt(s, rx, y, Inches(2.05), Inches(0.28), r.get("label") or "",
            size=11.5, bold=accent)
        _bar(s, rx + Inches(2.15), y + Inches(0.05), tw,
             _num(r.get("value")) / vmax, accent=accent)
        txt(s, rx + Inches(2.15) + tw + Inches(0.10), y, Inches(0.80),
            Inches(0.28), _fmt_val(r.get("value"), suffix), size=11,
            bold=True, color=AMETHYST if accent else GRAPHITE)
        y += Inches(step)
    if sl.get("read"):
        txt(s, rx, Inches(5.90), rw, Inches(0.66), sl["read"],
            size=11.5, color=BODY_LT)


def _sl_tiles_row(d, sl):
    s = d.new_slide()
    d.chrome(s, sl.get("eyebrow") or "Read", sl.get("title") or "",
             sub=sl.get("sub"), source=sl.get("source"))
    tiles = _items(sl, "tiles", 3)
    acc = sl.get("accent_index")
    x, w = M, Inches(3.644)
    for i, t in enumerate(tiles):
        rrect(s, x, Inches(2.330), w, Inches(3.30), CARD)
        color = OLIVE if (isinstance(acc, int) and i == acc) else GRAPHITE
        txt(s, x + Inches(0.24), Inches(2.54), w - Inches(0.48),
            Inches(0.70), t.get("big") or "", size=34, bold=True,
            color=color)
        txt(s, x + Inches(0.24), Inches(3.30), w - Inches(0.48),
            Inches(0.72), t.get("label") or "", size=11.5, bold=True)
        txt(s, x + Inches(0.24), Inches(4.10), w - Inches(0.48),
            Inches(1.30), t.get("body") or "", size=11.5, color=BODY_LT)
        x += w + GUT
    if sl.get("read"):
        txt(s, M, Inches(6.10), BAND, Inches(0.40), sl["read"],
            size=11.5, color=BODY_LT)


def _sl_hero(d, sl):
    orchid = str(sl.get("ground") or "accent").lower() != "dark"
    s = d.new_slide(dark=not orchid, fill=ORCHID if orchid else None)
    d.chrome(s, sl.get("eyebrow") or "Signal", sl.get("title") or "",
             dark=not orchid, orchid=orchid, sub=sl.get("sub"),
             source=sl.get("source"))
    ink = WHITE if orchid else SIGNAL
    body = WHITE if orchid else BODY_DK
    support_c = ORCHID_BODY if orchid else BODY_DK
    txt(s, M, Inches(2.50), BAND, Inches(1.10),
        sl.get("big") or "", size=82, bold=True, color=ink)
    if sl.get("line"):
        txt(s, M, Inches(3.72), BAND, Inches(0.70), sl["line"],
            size=16, color=body)
    if sl.get("support"):
        txt(s, M, Inches(4.70), BAND, Inches(1.20), sl["support"],
            size=16, color=support_c)


def _sl_table(d, sl):
    s = d.new_slide()
    d.chrome(s, sl.get("eyebrow") or "Table", sl.get("title") or "",
             sub=sl.get("sub"), source=sl.get("source"))
    cols = [str(c) for c in (sl.get("columns") or [])][:6]
    rows = [r for r in (sl.get("rows") or []) if isinstance(r, list)][:5]
    if not cols:
        return
    n = len(cols)
    first_w = 3.60 if n >= 5 else 4.20
    rest_w = (11.493 - first_w) / max(n - 1, 1)
    xs = [float(M) / 914400.0]
    for i in range(1, n):
        xs.append(xs[0] + first_w + (i - 1) * rest_w)
    y = Inches(2.330)
    for i, c in enumerate(cols):
        txt(s, Inches(xs[i]), y, Inches(first_w if i == 0 else rest_w),
            Inches(0.24), c.upper(), size=10.5, bold=True, color=MUTED_LT,
            spc=100, align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT)
    rect(s, M, Inches(2.58), BAND, Pt(0.75), TRACK)
    acc_row = sl.get("accent_row")
    acc_col = sl.get("accent_col")
    y = Inches(2.72)
    for ri, row in enumerate(rows):
        row_acc = isinstance(acc_row, int) and ri == acc_row
        for ci in range(n):
            cell = str(row[ci]) if ci < len(row) else ""
            cell_acc = row_acc and isinstance(acc_col, int) and ci == acc_col
            txt(s, Inches(xs[ci]), y,
                Inches(first_w if ci == 0 else rest_w), Inches(0.48),
                cell, size=14, bold=(ci == 0 and row_acc) or cell_acc,
                color=OLIVE if cell_acc else GRAPHITE,
                align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT)
        y += Inches(0.52)
    reads = [t for t in [sl.get("read"), sl.get("read2")] if t]
    if len(reads) == 2:
        rrect(s, M, Inches(5.50), Inches(5.607), Inches(1.10), CARD)
        txt(s, M + Inches(0.24), Inches(5.66), Inches(5.12), Inches(0.86),
            reads[0], size=13, color=BODY_LT)
        rrect(s, M + Inches(5.607) + GUT, Inches(5.50), Inches(5.607),
              Inches(1.10), CARD)
        txt(s, M + Inches(5.887), Inches(5.66), Inches(5.12), Inches(0.86),
            reads[1], size=13, color=BODY_LT)
    elif reads:
        rrect(s, M, Inches(5.50), BAND, Inches(0.92), CARD)
        txt(s, M + Inches(0.24), Inches(5.66), BAND - Inches(0.48),
            Inches(0.64), reads[0], size=14, color=BODY_LT)


def _sl_hero_proof(d, sl):
    s = d.new_slide(dark=True)
    d.chrome(s, sl.get("eyebrow") or "Proof", sl.get("title") or "",
             dark=True, sub=sl.get("sub"), source=sl.get("source"))
    txt(s, M, Inches(2.28), Inches(5.607), Inches(1.20),
        sl.get("big") or "", size=82, bold=True, color=SIGNAL)
    if sl.get("line"):
        txt(s, M, Inches(3.58), Inches(5.607), Inches(0.80), sl["line"],
            size=16, color=BODY_DK)
    rx = M + Inches(5.607) + GUT
    y = Inches(2.330)
    for p in _items(sl, "proofs", 3):
        rrect(s, rx, y, Inches(5.607), Inches(1.18), SLATE)
        txt(s, rx + Inches(0.24), y + Inches(0.16), Inches(5.12),
            Inches(0.44), p.get("fig") or "", size=22, bold=True,
            color=WHITE)
        txt(s, rx + Inches(0.24), y + Inches(0.64), Inches(5.12),
            Inches(0.40), p.get("label") or "", size=11.5, color=BODY_DK)
        y += Inches(1.32)


def _sl_paths(d, sl):
    s = d.new_slide(dark=True)
    d.chrome(s, sl.get("eyebrow") or "Paths", sl.get("title") or "",
             dark=True, sub=sl.get("sub"), source=sl.get("source"))
    y = Inches(2.330)
    for r in _items(sl, "rows", 13):
        lit = bool(r.get("lit"))
        if lit:
            rrect(s, M, y, BAND, Inches(0.245), SLATE, radius=0.06)
        c = SIGNAL if lit else MUTED_DK
        txt(s, M + Inches(0.16), y, Inches(1.1), Inches(0.245),
            str(r.get("kind") or "").upper(), size=8.5, color=c,
            font="Geist Mono", anchor=MSO_ANCHOR.MIDDLE, spc=80)
        txt(s, M + Inches(1.36), y, BAND - Inches(1.52), Inches(0.245),
            r.get("url") or "", size=8.5, color=c, font="Geist Mono",
            anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.28)


def _sl_close(d, sl):
    s = d.new_slide(dark=True)
    d.chrome(s, sl.get("eyebrow") or "Close", sl.get("title") or "",
             dark=True, sub=sl.get("sub"))
    _cards_2x2(d, s, _items(sl, "cards", 4), dark=True)


_RENDERERS = {
    "cover": _sl_cover,
    "argument": _sl_argument,
    "tiles_facts": _sl_tiles_facts,
    "bars": _sl_bars,
    "split_stats_bars": _sl_split_stats_bars,
    "tiles_row": _sl_tiles_row,
    "hero": _sl_hero,
    "table": _sl_table,
    "hero_proof": _sl_hero_proof,
    "paths": _sl_paths,
    "close": _sl_close,
}

SLIDE_TYPES = tuple(_RENDERERS.keys())


def _resolve_logos(static_dir):
    """Prefer the brand lockups shipped in static/; fall back to the
    skill assets when running from the repo root (local tests)."""
    cands_w, cands_k = [], []
    if static_dir:
        cands_w.append(os.path.join(static_dir,
                                    "crosswalk-logo-brand-white.png"))
        cands_k.append(os.path.join(static_dir,
                                    "crosswalk-logo-brand-black.png"))
    here = os.path.dirname(os.path.abspath(__file__))
    skill = os.path.join(os.path.dirname(here), ".cursor", "skills",
                         "crosswalk-brand-standards", "assets")
    cands_w += [os.path.join(here, "static",
                             "crosswalk-logo-brand-white.png"),
                os.path.join(skill, "crosswalk-logo-white.png")]
    cands_k += [os.path.join(here, "static",
                             "crosswalk-logo-brand-black.png"),
                os.path.join(skill, "crosswalk-logo-black.png")]
    lw = next((p for p in cands_w if os.path.exists(p)), None)
    lk = next((p for p in cands_k if os.path.exists(p)), None)
    return lw, lk


def render_insights_deck(plan, out_path, static_dir=None):
    """Render a slide-plan dict to a finished PPTX at out_path.
    Returns the number of slides rendered."""
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    lw, lk = _resolve_logos(static_dir)
    d = _Deck(prs, lw, lk)
    rendered = 0
    for sl in (plan.get("slides") or []):
        if not isinstance(sl, dict):
            continue
        fn = _RENDERERS.get(str(sl.get("type") or "").strip().lower())
        if fn is None:
            continue
        fn(d, sl)
        rendered += 1
    prs.save(out_path)
    return rendered
